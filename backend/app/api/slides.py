import asyncio
import copy
import hashlib
import io
import json
import logging
import os
import re
import shutil
import time
from concurrent.futures import Future, ThreadPoolExecutor
from typing import List, Optional
from urllib.parse import quote
from fastapi import APIRouter, BackgroundTasks, Body, Depends, HTTPException, UploadFile, File, Form
from fastapi.responses import FileResponse, StreamingResponse
from PIL import Image as PILImage
from pptx import Presentation
from pydantic import BaseModel
from sqlalchemy.orm import Session
from sqlalchemy.orm.attributes import flag_modified

from app.models.base import get_db, SessionLocal
from app.models.models import Project, Slide, ReferenceImage, SlideVersion
from app.schemas.project import SlideResponse
from app.core.llm_client import get_llm_client
from app.core.provider_credentials import (
    get_minimax_llm_model,
    get_raw_provider_credentials,
    provider_credentials_context,
    store_current_provider_credentials,
)
from app.utils.project_docs import load_project_documents, load_project_source_packs
from app.utils.text_cleaning import normalize_markdown_content, normalize_markdown_emphasis
from app.utils.reference_image import (
    default_visual_asset_process_mode,
    normalize_visual_asset_kind,
    reference_process_mode_instruction,
)

# 全局运行中任务跟踪（project_id -> asyncio.Task）
_running_tasks: dict = {}
_running_tasks_lock = asyncio.Lock()


def _resolve_file_path(file_path: str) -> str:
    """安全解析文件路径，兼容从不同工作目录启动的情况。"""
    if not file_path:
        return file_path
    if os.path.exists(file_path):
        return file_path
    backend_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    candidate = os.path.join(backend_dir, file_path)
    if os.path.exists(candidate):
        return candidate
    abs_path = os.path.abspath(file_path)
    if os.path.exists(abs_path):
        return abs_path
    return file_path


from app.core.config import settings
from app.services.content_plan import (
    PAGE_MAP_DOCUMENT_LIMIT,
    generate_content_plan,
    infer_effective_content_intent_contract,
    infer_page_count_from_single_ppt,
    content_plan_topic_with_chat_context,
    resolve_content_plan_page_target,
    resolve_requested_content_plan_page_count,
    _auto_reclassify_page_type,
)
from app.services.source_context import (
    DEFAULT_SOURCE_CONTEXT_TOKEN_BUDGET,
    SourceScopeRequired,
    build_source_context,
)
from app.services.source_intent import (
    infer_intent_contract,
    source_diagnostics_from_documents,
)
from app.services.logo_policy import (
    DEFAULT_LOGO_ANCHOR,
    LOGO_ANCHORS,
    logo_anchor_from_ref,
    logo_policy_for_page,
    is_logo_confirmed,
    logo_review_status,
    normalize_logo_anchor,
    should_show_logo,
    should_use_logo_as_scene_asset,
)
from app.services.logo_assets import prepare_logo_lockup_image, prepare_logo_overlay_image, prepare_logo_symbol_image
from app.services.logo_overlay_layout import resolve_logo_render_policy
from app.services.overlay_layers import (
    exact_overlay_asset_ids,
    merge_overlay_layers_into_visual_json,
    normalize_overlay_layers,
    remove_asset_from_overlay_layers,
)
from app.services.visual_plan import (
    VisualPlanGenerationError,
    generate_visual_plan,
    _current_visual_requirements,
    _current_visual_suggestion,
    _guidance_mentions_absent_visible_copy,
    _page_visible_copy_text,
    _recall_visual_assets_for_page,
)
from app.services.prompt_engine import (
    _content_visual_contract,
    _strip_markdown,
    generate_prompt_for_page,
    generate_prompts_for_all_pages,
)
from app.services.pptx_assembler import assemble_pptx
from app.services.style_pack import derive_style_pack_from_content, style_pack_from_selected_style
from app.services.visual_strategy import detect_logo_tone_from_image
from app.services.visual_block_renderer import (
    blocks_to_markdown,
    is_visual_block,
    normalize_content_blocks,
    normalize_route_mode,
    render_visual_block,
    route_to_asset_route_mode,
    stable_block_hash,
)
from app.services.visual_directives import extract_visual_directives
from app.services.image_analyzer import analyze_reference_image, analyze_visual_asset, describe_context_image
from app.services.artifact_versions import artifact_meta, dependency_signature, has_stale_flags, with_artifact_meta, with_stale_flags
from app.tasks import generate_slides_task, redis_client
from app.services.celery_runtime import ensure_celery_worker
from app.services.run_state import (
    apply_project_rollback,
    cancel_active_run,
    create_project_run,
    finish_run,
    get_active_run,
    get_latest_run,
    image_generation_queued_message,
    image_generation_run_stage,
    is_run_active,
    mark_run_running,
    reconcile_project_state,
    serialize_run,
    serialize_workflow_status,
    set_run_task,
    stale_active_run,
    stale_inactive_run_if_needed,
    target_counts,
    update_run_progress,
    generation_progress,
)
from app.services.project_quality_report import build_project_quality_report
from app.celery_app import celery_app
from celery.result import AsyncResult

router = APIRouter(prefix="/projects", tags=["slides"])
logger = logging.getLogger(__name__)


def ensure_generation_worker_ready():
    if not ensure_celery_worker(queue=settings.CELERY_IMAGE_QUEUE):
        raise HTTPException(status_code=503, detail="图片生成服务未启动，任务没有开始。请启动 worker 后重试。")


def _enqueue_generation_task(
    db: Session,
    project_id: str,
    page_nums: Optional[List[int]],
    *,
    run,
    prototype: bool = False,
):
    """Queue image generation and leave a recoverable run state if Redis/Celery fails."""
    try:
        credential_id = store_current_provider_credentials(redis_client)
        task = generate_slides_task.delay(
            project_id,
            page_nums,
            prototype=prototype,
            run_id=run.id,
            credential_id=credential_id,
        )
        set_run_task(db, run.id, task.id)
        db.commit()
    except Exception as exc:
        logger.exception("Failed to enqueue generation task for project %s", project_id)
        message = "后台生成队列暂时不可用。请确认 Docker/Redis 正常运行且磁盘空间充足，然后重试生成。"
        finish_run(db, run.id, status="stale", message=message, error_msg=str(exc)[:500])
        db.commit()
        raise HTTPException(status_code=503, detail=message) from exc
    try:
        redis_client.set(f"project:{project_id}:task_id", task.id, ex=3600)
        redis_client.set(f"project:{project_id}:task_started_at", str(time.time()), ex=3600)
    except Exception as exc:
        logger.warning("Redis 写入生成任务跟踪信息失败，不影响已分发任务: project=%s task=%s error=%s", project_id, task.id, exc)
    return task


def _project_output_dir(project_id: str) -> str:
    return os.path.join(settings.OUTPUT_DIR or "./outputs", project_id)


def _final_pptx_path(project_id: str) -> str:
    return os.path.join(_project_output_dir(project_id), "presentation.pptx")


def _pptx_slide_count(path: str) -> int | None:
    try:
        return len(Presentation(path).slides)
    except Exception as exc:
        logger.warning("Unable to inspect PPTX slide count: path=%s error=%s", path, exc)
        return None


def _completed_slide_image_count(slides: list[Slide]) -> int:
    count = 0
    for slide in slides:
        image_path = _existing_output_path(slide.image_path)
        if slide.status == "completed" and image_path and os.path.exists(image_path):
            count += 1
    return count


def _final_pptx_matches_completed_slides(pptx_path: str, slides: list[Slide]) -> bool:
    if not os.path.exists(pptx_path):
        return False
    if not slides:
        return False
    completed_count = _completed_slide_image_count(slides)
    if completed_count != len(slides):
        return False
    return _pptx_slide_count(pptx_path) == completed_count


def _has_completed_final_pptx(project: Project, slides: list[Slide]) -> tuple[bool, str]:
    pptx_path = _final_pptx_path(project.id)
    has_pptx = (
        project.status == "completed"
        and bool(slides)
        and _final_pptx_matches_completed_slides(pptx_path, slides)
        and not _project_has_stale_artifacts(slides)
    )
    return has_pptx, pptx_path


# 上传限制常量
MAX_UPLOAD_SIZE = 20 * 1024 * 1024  # 20MB
MAX_REFERENCE_IMAGES_PER_PAGE = 10
ALLOWED_UPLOAD_ROLES = {"style_ref", "logo", "template", "visual_asset", "content_ref", "chart_ref", "finetune_ref", "chat_ref"}
ALLOWED_IMAGE_TYPES = {
    "image/jpeg",
    "image/png",
    "image/gif",
    "image/webp",
    "image/svg+xml",
    "image/heic",
    "image/heif",
    "image/bmp",
    "image/tiff",
}

# 全局生成进度存储（内存级，项目重启后丢失），定义在 run_state.py，此处引用
REFERENCE_ANALYSIS_ROLES = {"style_ref", "content_ref", "chart_ref"}


def _asset_analysis_worker_count() -> int:
    try:
        return int(os.getenv("PPTGOD_ASSET_ANALYSIS_WORKERS", "4"))
    except ValueError:
        return 4


ASSET_ANALYSIS_WORKERS = max(2, min(8, _asset_analysis_worker_count()))
ASSET_ANALYSIS_POOL = ThreadPoolExecutor(max_workers=ASSET_ANALYSIS_WORKERS, thread_name_prefix="pptgod-asset")

UNTITLED_PROJECT_TITLES = {"未命名项目", "新建项目", "Untitled Project", "Untitled"}


def _should_autoname_project(project: Project | None) -> bool:
    """Only auto-name projects that still use a default placeholder title."""
    if not project:
        return False
    title = str(project.title or "").strip()
    if not title:
        return True
    return title in UNTITLED_PROJECT_TITLES or title.startswith("未命名项目")


def _first_plain_text(value) -> str:
    if isinstance(value, list):
        parts = [str(item).strip() for item in value if str(item).strip()]
        return " ".join(parts)
    if isinstance(value, dict):
        parts = [str(item).strip() for item in value.values() if str(item).strip()]
        return " ".join(parts)
    return str(value or "").strip()


def _clean_project_title(value: str) -> str:
    text = normalize_markdown_emphasis(str(value or ""))
    text = re.sub(r"\[\[PPTGOD_ATTACHMENT:[^\]]+\]\]", " ", text)
    text = re.sub(r"【(?:文件|图片|附件|用户上传材料)[^】]*】", " ", text)
    text = re.sub(r"\b(?:PDF|Word|PPT|PPTX|Markdown|TXT)\b", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"https?://\S+", " ", text)
    text = re.sub(r"[#*_`>\[\](){}<>|]", " ", text)
    text = re.sub(r"\s+", " ", text).strip(" \t\r\n，。；：、,.!?！？—-")
    for _ in range(3):
        cleaned = re.sub(r"^(?:请|帮我|帮忙|做一份|生成一份|制作一份|我要做|我想做|把|用这个文件|参考)\s*", "", text)
        if cleaned == text:
            break
        text = cleaned
    text = re.sub(r"(?:的)?(?:PPT|ppt|演示文稿|幻灯片|汇报|报告)$", "", text).strip(" \t\r\n，。；：、,.!?！？—-")
    if not text:
        return ""

    # Keep the sidebar name scannable; avoid long brief sentences becoming project titles.
    split = re.split(r"[，,。；;！!?？\n\r]", text, maxsplit=1)
    text = split[0].strip(" \t\r\n，。；：、,.!?！？—-") or text
    text = re.sub(r"(?:的)?(?:PPT|ppt|演示文稿|幻灯片|汇报|报告)$", "", text).strip(" \t\r\n，。；：、,.!?！？—-")
    if len(text) > 24:
        text = text[:24].rstrip(" \t\r\n，。；：、,.!?！？—-")
    return text


def _derive_project_title(topic: str, outline: list[dict]) -> str | None:
    candidates: list[str] = []
    if isinstance(outline, list):
        for item in outline[:4]:
            if not isinstance(item, dict):
                continue
            text_content = item.get("text_content") if isinstance(item.get("text_content"), dict) else {}
            candidates.extend(
                [
                    _first_plain_text(text_content.get("headline")),
                    _first_plain_text(text_content.get("subhead")),
                    _first_plain_text(item.get("title")),
                    _first_plain_text(item.get("section_title")),
                ]
            )

    for line in str(topic or "").splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith(("[[PPTGOD_ATTACHMENT:", "【用户上传材料】", "【文件", "【图片", "已上传文档", "已上传图片")):
            continue
        candidates.append(stripped)

    for candidate in candidates:
        title = _clean_project_title(candidate)
        if len(title) >= 2 and title not in UNTITLED_PROJECT_TITLES:
            return title
    return None


def _style_text_from_selected_style(selected_style: dict | str | None) -> str | None:
    return style_pack_from_selected_style(selected_style)


def _style_override_from_text(style_text: str | None) -> dict | None:
    if not style_text:
        return None
    palette = []
    mood = ""
    style_name = "Content-derived style pack"
    for line in style_text.splitlines():
        stripped = line.strip()
        if stripped.startswith("Style:"):
            style_name = stripped.split(":", 1)[1].strip() or style_name
        elif stripped.startswith("Palette:"):
            palette = [p.strip() for p in stripped.split(":", 1)[1].split(",") if p.strip()][:5]
        elif stripped.startswith("Mood:"):
            mood = stripped.split(":", 1)[1].strip()
    return {
        "meta": {
            "theme": style_name,
            "style_name": style_name,
            "palette": palette,
            "mood": mood,
        },
        "body": style_text,
    }


def _analysis_tokens(*values: str | None) -> list[str]:
    tokens: list[str] = []
    for value in values:
        for token in re.split(r"[\s,，。；;:：、/|()（）\[\]{}\"'“”‘’_-]+", str(value or "")):
            token = token.strip()
            if len(token) >= 2 and token not in tokens:
                tokens.append(token)
            if len(tokens) >= 16:
                return tokens
    return tokens


def _visual_asset_analysis_placeholder(
    *,
    asset_name: str | None,
    asset_kind: str | None,
    usage_note: str | None,
    filename: str,
    user_asset_kind_provided: bool,
    explicit_process_mode: bool,
) -> dict:
    stem = os.path.splitext(os.path.basename(filename or ""))[0]
    subject = asset_name or stem or "用户上传可复用素材"
    kind = normalize_visual_asset_kind(asset_kind)
    return {
        "analysis_status": "queued",
        "analysis_type": "visual_asset",
        "detected_kind": kind,
        "subject": subject,
        "description": usage_note or "",
        "identity_elements": [],
        "distinctive_features": [],
        "must_not_change": [],
        "suggested_keywords": _analysis_tokens(asset_name, usage_note, stem),
        "recommended_usage": usage_note or "",
        "fidelity_note": "",
        "selection_tier": "manual",
        "importance_score": 100,
        "_user_asset_kind_provided": bool(user_asset_kind_provided),
        "_explicit_process_mode": bool(explicit_process_mode),
    }


def _reference_analysis_placeholder(role: str, filename: str, usage_note: str | None = None) -> dict:
    return {
        "analysis_status": "queued",
        "analysis_type": "reference_image",
        "role": role,
        "description": usage_note or os.path.splitext(os.path.basename(filename or ""))[0],
        "suggested_keywords": _analysis_tokens(usage_note, filename),
    }


def _analysis_task_done(label: str, future: Future) -> None:
    try:
        future.result()
    except Exception as exc:
        logger.exception("Asset analysis background task crashed: task=%s error=%s", label, exc)


def _submit_asset_analysis_task(label: str, fn, *args) -> None:
    future = ASSET_ANALYSIS_POOL.submit(fn, *args)
    future.add_done_callback(lambda done: _analysis_task_done(label, done))


def _analyze_visual_asset_for_ref(project_id: str, ref_id: str) -> None:
    db = SessionLocal()
    try:
        ref = db.query(ReferenceImage).filter(
            ReferenceImage.id == ref_id,
            ReferenceImage.project_id == project_id,
            ReferenceImage.role == "visual_asset",
        ).first()
        if not ref or not ref.file_path or not os.path.exists(ref.file_path):
            return
        current = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
        ref.asset_analysis = {**current, "analysis_status": "running", "analysis_type": "visual_asset"}
        db.commit()

        result = analyze_visual_asset(
            ref.file_path,
            asset_name=ref.asset_name or "",
            asset_kind=ref.asset_kind or "other",
            usage_note=ref.usage_note or "",
        )
        result = {
            **current,
            **(result if isinstance(result, dict) else {}),
            "analysis_status": "completed",
            "analysis_type": "visual_asset",
            "selection_tier": current.get("selection_tier") or "manual",
            "importance_score": current.get("importance_score") or 100,
        }
        if not current.get("_user_asset_kind_provided"):
            detected_kind = normalize_visual_asset_kind(result.get("detected_kind"))
            ref.asset_kind = detected_kind
            if not current.get("_explicit_process_mode"):
                ref.process_mode = default_visual_asset_process_mode(detected_kind)
        if not ref.asset_name:
            ref.asset_name = result.get("subject") or os.path.splitext(os.path.basename(ref.file_path))[0]
        ref.asset_analysis = result
        db.commit()
        logger.info("Visual asset analysis completed: project=%s ref=%s", project_id, ref_id)
    except Exception as exc:
        db.rollback()
        logger.warning("Visual asset analysis failed: project=%s ref=%s error=%s", project_id, ref_id, exc)
        try:
            ref = db.query(ReferenceImage).filter(ReferenceImage.id == ref_id, ReferenceImage.project_id == project_id).first()
            if ref:
                current = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
                ref.asset_analysis = {
                    **current,
                    "analysis_status": "failed",
                    "analysis_type": "visual_asset",
                    "error": str(exc)[:500],
                }
                db.commit()
        except Exception:
            db.rollback()
    finally:
        db.close()


def _analyze_reference_image_for_ref(project_id: str, ref_id: str) -> None:
    db = SessionLocal()
    try:
        ref = db.query(ReferenceImage).filter(
            ReferenceImage.id == ref_id,
            ReferenceImage.project_id == project_id,
        ).first()
        if not ref or ref.role not in REFERENCE_ANALYSIS_ROLES or not ref.file_path or not os.path.exists(ref.file_path):
            return
        current = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
        ref.asset_analysis = {**current, "analysis_status": "running", "analysis_type": "reference_image"}
        db.commit()

        result = analyze_reference_image(ref.file_path)
        ref.asset_analysis = {
            **current,
            **(result if isinstance(result, dict) else {}),
            "analysis_status": "completed",
            "analysis_type": "reference_image",
            "role": ref.role,
        }
        db.commit()
        logger.info("Reference image analysis completed: project=%s ref=%s role=%s", project_id, ref_id, ref.role)
    except Exception as exc:
        db.rollback()
        logger.warning("Reference image analysis failed: project=%s ref=%s error=%s", project_id, ref_id, exc)
        try:
            ref = db.query(ReferenceImage).filter(ReferenceImage.id == ref_id, ReferenceImage.project_id == project_id).first()
            if ref:
                current = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
                ref.asset_analysis = {
                    **current,
                    "analysis_status": "failed",
                    "analysis_type": "reference_image",
                    "error": str(exc)[:500],
                }
                db.commit()
        except Exception:
            db.rollback()
    finally:
        db.close()


def _cached_reference_analysis(ref: ReferenceImage) -> dict | None:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    if not analysis:
        return None
    if analysis.get("source_document"):
        return None
    if analysis.get("analysis_status") == "completed":
        return analysis
    if not analysis.get("analysis_status") and any(analysis.get(key) for key in ("description", "colors", "composition_style", "mood")):
        return analysis
    return None


def _wait_for_cached_reference_analysis(ref: ReferenceImage, *, timeout_seconds: float = 4.0) -> dict | None:
    if not ref.id:
        return None
    deadline = time.time() + timeout_seconds
    db = SessionLocal()
    try:
        while time.time() < deadline:
            fresh = db.query(ReferenceImage).filter(ReferenceImage.id == ref.id).first()
            if fresh:
                cached = _cached_reference_analysis(fresh)
                if cached:
                    return cached
                status = (fresh.asset_analysis or {}).get("analysis_status") if isinstance(fresh.asset_analysis, dict) else None
                if status in {"failed"}:
                    return None
            time.sleep(0.25)
    finally:
        db.close()
    return None


def _fallback_reference_summary(ref: ReferenceImage) -> str:
    name = ref.asset_name or os.path.splitext(os.path.basename(ref.file_path or ""))[0]
    pieces = [f"用户上传参考图：{name}"]
    if ref.usage_note:
        pieces.append(f"usage={ref.usage_note}")
    if ref.role:
        pieces.append(f"role={ref.role}")
    return "; ".join(pieces)


def _derive_project_style_pack(project: Project, content_plan: list[dict]) -> str:
    selected = _style_text_from_selected_style(project.selected_style)
    if selected:
        return selected
    analyses = []
    for ref in project.reference_images or []:
        if ref.role != "style_ref" or ref.slide_id:
            continue
        try:
            cached = _cached_reference_analysis(ref)
            if not cached and isinstance(ref.asset_analysis, dict) and ref.asset_analysis.get("analysis_status") in {"queued", "running"}:
                cached = _wait_for_cached_reference_analysis(ref, timeout_seconds=4.0)
            if cached:
                analyses.append(cached)
                continue
            if isinstance(ref.asset_analysis, dict) and ref.asset_analysis.get("analysis_status") in {"queued", "running", "failed"}:
                logger.info("StylePack: using content-derived fallback while style_ref analysis is %s", ref.asset_analysis.get("analysis_status"))
                continue
            if not os.path.exists(ref.file_path):
                logger.info("StylePack: skipping missing style_ref file after cache lookup: %s", ref.file_path)
                continue
            analyses.append(analyze_reference_image(ref.file_path))
        except Exception as exc:
            logger.warning(f"StylePack: failed to analyze style_ref {ref.file_path}: {exc}")
    return derive_style_pack_from_content(content_plan, reference_analyses=analyses)


def _reference_upload_url(project_id: str, file_path: str | None) -> str:
    """Build a stable URL for root uploads and nested PPTX asset files."""
    path = (file_path or "").replace("\\", "/")
    marker = f"/{project_id}/"
    if marker in path:
        suffix = path.split(marker, 1)[1]
    else:
        suffix = os.path.basename(path)
    return f"/uploads/{project_id}/{suffix}"


def _asset_source_document(ref: ReferenceImage) -> str:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    return str(analysis.get("source_document") or "").strip()


def _asset_source_page_num(ref: ReferenceImage) -> int | None:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    for key in ("source_page_num", "pptx_source_page_num", "pdf_source_page_num"):
        page_num = analysis.get(key)
        try:
            value = int(page_num) if page_num is not None else None
        except (TypeError, ValueError):
            value = None
        if value is not None and value > 0:
            return value
    return None


def _source_ref_type(source_doc: str, source_type: str | None = None) -> str:
    kind = str(source_type or "").strip().lower()
    if kind:
        return kind
    lower_doc = str(source_doc or "").strip().lower()
    if lower_doc.endswith(".pdf"):
        return "pdf"
    if lower_doc.endswith((".pptx", ".ppt")):
        return "pptx_slide"
    return ""


def _reference_image_sort_key(ref: ReferenceImage) -> tuple:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    bounds = analysis.get("shape_bounds") if isinstance(analysis.get("shape_bounds"), dict) else {}
    bbox = analysis.get("bbox")
    try:
        group_index = int(analysis.get("asset_group_index") or 9999)
    except (TypeError, ValueError):
        group_index = 9999
    if isinstance(bbox, (list, tuple)) and len(bbox) == 4:
        try:
            left = float(bbox[0])
            top = float(bbox[1])
        except (TypeError, ValueError):
            left = 0.0
            top = 0.0
    else:
        try:
            top = float(bounds.get("top") or 0)
        except (TypeError, ValueError):
            top = 0.0
        try:
            left = float(bounds.get("left") or 0)
        except (TypeError, ValueError):
            left = 0.0
    role_order = {
        "content_ref": 0,
        "chart_ref": 1,
        "visual_asset": 2,
        "style_ref": 3,
        "template": 4,
        "logo": 5,
    }.get(ref.role or "", 9)
    return (
        role_order,
        _asset_source_document(ref),
        _asset_source_page_num(ref) or 10_000,
        str(analysis.get("asset_group_key") or ""),
        group_index,
        top,
        left,
        ref.asset_name or os.path.basename(ref.file_path or ""),
        ref.id,
    )


def _manual_asset_ids(visual_json: dict | None) -> list[str]:
    visual = visual_json if isinstance(visual_json, dict) else {}
    raw = visual.get("manual_visual_asset_ids") or []
    if not isinstance(raw, list):
        return []
    ids: list[str] = []
    for asset_id in raw:
        value = str(asset_id)
        if value and value not in ids:
            ids.append(value)
    return ids


def _manual_asset_usage(visual_json: dict | None) -> dict[str, str]:
    visual = visual_json if isinstance(visual_json, dict) else {}
    raw = visual.get("manual_visual_asset_usage") or {}
    if not isinstance(raw, dict):
        return {}
    return {str(k): str(v) for k, v in raw.items() if k and v}


def _manual_asset_route_modes(visual_json: dict | None, manual_ids: list[str] | None = None) -> dict[str, str]:
    visual = visual_json if isinstance(visual_json, dict) else {}
    raw = visual.get("asset_route_modes") or {}
    if not isinstance(raw, dict):
        return {}
    allowed_ids = set(manual_ids if manual_ids is not None else _manual_asset_ids(visual))
    allowed_modes = {"blend", "double_blend", "overlay"}
    return {
        str(asset_id): str(mode).lower()
        for asset_id, mode in raw.items()
        if str(asset_id) in allowed_ids and str(mode).lower() in allowed_modes
    }


def _visual_asset_ids(visual_json: dict | None) -> list[str]:
    visual = visual_json if isinstance(visual_json, dict) else {}
    raw = visual.get("visual_asset_ids") or []
    if not isinstance(raw, list):
        return []
    ids: list[str] = []
    for asset_id in raw:
        value = str(asset_id)
        if value and value not in ids:
            ids.append(value)
    return ids


def _excluded_asset_ids(visual_json: dict | None) -> list[str]:
    visual = visual_json if isinstance(visual_json, dict) else {}
    raw = visual.get("excluded_visual_asset_ids") or []
    if not isinstance(raw, list):
        return []
    ids: list[str] = []
    for asset_id in raw:
        value = str(asset_id)
        if value and value not in ids:
            ids.append(value)
    return ids


def _without_asset_ids(values: list[str], blocked_ids: set[str]) -> list[str]:
    return [asset_id for asset_id in values if asset_id not in blocked_ids]


def _merge_asset_ids(manual_ids: list[str], auto_ids: list[str] | None) -> list[str]:
    merged: list[str] = []
    for asset_id in [*manual_ids, *(auto_ids or [])]:
        value = str(asset_id)
        if value and value not in merged:
            merged.append(value)
    return merged


def _merge_manual_pins_into_visual_json(visual_json: dict | None, existing_visual_json: dict | None) -> dict:
    visual = copy.deepcopy(visual_json) if isinstance(visual_json, dict) else {}
    existing = existing_visual_json if isinstance(existing_visual_json, dict) else {}
    manual_ids = _manual_asset_ids(existing)
    manual_usage = _manual_asset_usage(existing)
    manual_route_modes = _manual_asset_route_modes(existing, manual_ids)
    excluded_ids = [
        asset_id for asset_id in _excluded_asset_ids(existing)
        if asset_id not in set(manual_ids)
    ]
    excluded_set = set(excluded_ids)
    visual = merge_overlay_layers_into_visual_json(visual, existing)
    visual["visual_asset_ids"] = _without_asset_ids(_visual_asset_ids(visual), excluded_set)
    usage = visual.get("visual_asset_usage") if isinstance(visual.get("visual_asset_usage"), dict) else {}
    visual["visual_asset_usage"] = {
        str(k): str(v)
        for k, v in usage.items()
        if k and v and str(k) not in excluded_set
    }
    route_modes = visual.get("asset_route_modes") if isinstance(visual.get("asset_route_modes"), dict) else {}
    visual["asset_route_modes"] = {
        str(k): str(v)
        for k, v in route_modes.items()
        if k and v and str(k) in set(visual["visual_asset_ids"])
    }
    if excluded_ids:
        visual["excluded_visual_asset_ids"] = excluded_ids
    else:
        visual.pop("excluded_visual_asset_ids", None)
    if not manual_ids:
        return visual
    visual["manual_visual_asset_ids"] = manual_ids
    visual["manual_visual_asset_usage"] = manual_usage
    visual["visual_asset_ids"] = _merge_asset_ids(manual_ids, visual.get("visual_asset_ids") or [])
    visual["visual_asset_usage"] = {**manual_usage, **visual["visual_asset_usage"]}
    route_modes = visual.get("asset_route_modes") if isinstance(visual.get("asset_route_modes"), dict) else {}
    if manual_route_modes or route_modes:
        visual["asset_route_modes"] = {
            **{str(k): str(v) for k, v in route_modes.items() if k and v and str(k) in set(visual["visual_asset_ids"])},
            **manual_route_modes,
        }
    return visual


def _merge_excluded_assets_after_visual_edit(previous_visual: dict | None, next_visual: dict | None) -> list[str]:
    previous = previous_visual if isinstance(previous_visual, dict) else {}
    current = next_visual if isinstance(next_visual, dict) else {}
    previous_ids = _visual_asset_ids(previous)
    next_ids = _visual_asset_ids(current)
    previous_manual = set(_manual_asset_ids(previous))
    current_manual = set(_manual_asset_ids(current))
    current_selected = set(next_ids) | current_manual
    excluded: list[str] = [
        asset_id
        for asset_id in _excluded_asset_ids(current)
        if asset_id not in current_selected
    ]

    for asset_id in previous_ids:
        if asset_id in previous_manual:
            continue
        if asset_id not in current_selected and asset_id not in excluded:
            excluded.append(asset_id)
    return excluded


EXACT_PAGE_REFERENCE_ROLES = {"content_ref", "chart_ref"}
EXACT_PAGE_REFERENCE_EXCLUDED_KINDS = {"source_page_image"}
EXACT_PAGE_REFERENCE_QR_TERMS = ("二维码", "扫码", "小程序码", "条码", "qr", "QR")
EXACT_PAGE_REFERENCE_PRESETS = ("left-card", "right-card", "center-card", "bottom-band", "top-right-small", "bottom-right-small")


def _page_reference_exact_overlay_ref(ref: ReferenceImage) -> bool:
    if ref.role not in EXACT_PAGE_REFERENCE_ROLES:
        return False
    if not ref.slide_id:
        return False
    if str(ref.process_mode or "").lower() != "original":
        return False
    if str(ref.asset_kind or "").lower() in EXACT_PAGE_REFERENCE_EXCLUDED_KINDS:
        return False
    return bool(ref.file_path and os.path.exists(_resolve_file_path(ref.file_path)))


def _page_reference_looks_like_qr(ref: ReferenceImage) -> bool:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    text = " ".join(
        str(value or "")
        for value in (
            ref.asset_name,
            ref.asset_kind,
            ref.usage_note,
            analysis.get("subject"),
            analysis.get("nearby_text"),
            analysis.get("source_slide_text"),
            analysis.get("description"),
        )
    )
    return any(term in text for term in EXACT_PAGE_REFERENCE_QR_TERMS)


def _next_page_reference_overlay_preset(ref: ReferenceImage, used_presets: set[str]) -> str:
    if _page_reference_looks_like_qr(ref) and "bottom-right-small" not in used_presets:
        return "bottom-right-small"
    for preset in EXACT_PAGE_REFERENCE_PRESETS:
        if preset not in used_presets:
            return preset
    return "right-card"


def _mark_reference_as_exact_overlay(ref: ReferenceImage) -> None:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    ref.asset_analysis = {
        **analysis,
        "exact_overlay": True,
        "exact_overlay_reason": "用户选择原样保留，最终页面会保留素材比例和细节。",
    }


def _clear_reference_exact_overlay_mark(ref: ReferenceImage) -> None:
    if not isinstance(ref.asset_analysis, dict):
        return
    analysis = dict(ref.asset_analysis)
    analysis.pop("exact_overlay", None)
    analysis.pop("exact_overlay_reason", None)
    ref.asset_analysis = analysis


def _sync_exact_page_reference_overlay_layers(slide: Slide, refs: list[ReferenceImage] | None = None) -> bool:
    """Keep page-level original refs on the deterministic exact-paste path."""
    page_refs = sorted(refs or list(slide.reference_images or []), key=_reference_image_sort_key)
    exact_refs = [ref for ref in page_refs if _page_reference_exact_overlay_ref(ref)]
    visual = copy.deepcopy(slide.visual_json) if isinstance(slide.visual_json, dict) else {}
    existing_layers = normalize_overlay_layers(visual.get("overlay_layers"), valid_asset_ids=None, strict_assets=False)
    by_asset_id = {str(layer.get("asset_id")): layer for layer in existing_layers if layer.get("asset_id")}
    used_presets = {str(layer.get("preset")) for layer in existing_layers if layer.get("preset")}
    changed = False

    for ref in exact_refs:
        ref_id = str(ref.id)
        _mark_reference_as_exact_overlay(ref)
        if ref_id in by_asset_id:
            continue
        preset = _next_page_reference_overlay_preset(ref, used_presets)
        used_presets.add(preset)
        existing_layers.append({
            "id": f"ov_{ref_id}",
            "asset_id": ref_id,
            "enabled": True,
            "preset": preset,
            "fit": "contain",
            "mode": "exact_cutout",
            "usage_note": f"{ref.asset_name or '本页素材'}：原图保留",
        })
        changed = True

    normalized = normalize_overlay_layers(existing_layers, valid_asset_ids=None, strict_assets=False)
    if changed or ("overlay_layers" in visual and visual.get("overlay_layers") != normalized):
        visual["overlay_layers"] = normalized
        slide.visual_json = visual
        flag_modified(slide, "visual_json")
        return True
    return False


def _load_page_reference_images_for_slides(db: Session, slides: list[Slide]) -> None:
    slide_ids = [slide.id for slide in slides if slide.id]
    if not slide_ids:
        return
    refs = db.query(ReferenceImage).filter(
        ReferenceImage.slide_id.in_(slide_ids),
        ReferenceImage.role != "finetune_ref",
    ).all()
    refs_by_slide: dict[str, list[ReferenceImage]] = {}
    for ref in refs:
        refs_by_slide.setdefault(ref.slide_id, []).append(ref)
    for slide in slides:
        slide.reference_images = sorted(refs_by_slide.get(slide.id, []), key=_reference_image_sort_key)


def _sync_exact_page_reference_overlays_for_generation(db: Session, slides: list[Slide]) -> list[int]:
    _load_page_reference_images_for_slides(db, slides)
    synced_pages: list[int] = []
    for slide in slides:
        if _sync_exact_page_reference_overlay_layers(slide):
            _mark_slide_artifacts_stale(slide, visual=True)
            synced_pages.append(int(slide.page_num))
    return synced_pages


def _serialize_reference_image(
    img: ReferenceImage,
    project_id: str,
    *,
    pinned_slide_ids: list[str] | None = None,
    is_pinned: bool = False,
    relevance: dict | None = None,
) -> dict:
    analysis = img.asset_analysis if isinstance(img.asset_analysis, dict) else {}
    source_page_num = _asset_source_page_num(img)
    source_document = _asset_source_document(img)
    return {
        "id": img.id,
        "role": img.role,
        "slide_id": img.slide_id,
        "process_mode": img.process_mode or "blend",
        "asset_name": img.asset_name,
        "asset_kind": img.asset_kind,
        "usage_note": img.usage_note,
        "asset_analysis": img.asset_analysis,
        "source_document": source_document or None,
        "source_page_num": source_page_num,
        "tags": analysis.get("asset_tags") or analysis.get("suggested_keywords") or [],
        "pinned_slide_ids": pinned_slide_ids or [],
        "is_pinned": is_pinned,
        "matched_terms": (relevance or {}).get("matched_terms", []),
        "relevance_reason": (relevance or {}).get("reason"),
        "review_status": logo_review_status(img) if img.role == "logo" else analysis.get("review_status"),
        "needs_user_review": logo_review_status(img) == "needs_review" if img.role == "logo" else bool(analysis.get("needs_user_review")),
        "confidence_score": analysis.get("confidence_score"),
        "review_reason": analysis.get("review_reason"),
        "detected_names": analysis.get("detected_names") or [],
        "logo_anchor": img.logo_anchor or (DEFAULT_LOGO_ANCHOR if img.role == "logo" else None),
        "file_exists": os.path.exists(img.file_path),
        "url": _reference_upload_url(project_id, img.file_path),
        "overlay_url": _logo_overlay_url(img, project_id) if img.role == "logo" and is_logo_confirmed(img) else None,
        "symbol_overlay_url": _logo_symbol_overlay_url(img, project_id) if img.role == "logo" and is_logo_confirmed(img) else None,
    }


def _template_key_for_page_intent(page_intent: dict | None) -> str:
    page_type = str((page_intent or {}).get("type") or "content").lower()
    layout = str((page_intent or {}).get("layout") or "").lower()
    if page_type == "cover":
        return "cover"
    if page_type == "toc" or layout == "toc":
        return "toc"
    if page_type == "section" or layout == "section":
        return "section"
    if page_type in {"data", "chart", "table"} or layout == "data":
        return "data"
    if page_type in {"hero", "quote"} or layout == "hero":
        return "quote"
    if page_type == "ending":
        return "ending"
    return "content"


def _project_template_refs_for_prompt(project: Project, page_intent: dict | None = None) -> list[dict]:
    """Return at most one template layout hint for the target page."""
    raw_recommendations = getattr(project, "selected_template_recommendations", None)
    recommendations = raw_recommendations if isinstance(raw_recommendations, dict) else {}
    rec = recommendations.get(_template_key_for_page_intent(page_intent))
    if not isinstance(rec, dict):
        fallback = next(
            (
                img for img in (project.reference_images or [])
                if img.role == "template" and img.file_path and os.path.exists(img.file_path)
            ),
            None,
        )
        if not fallback:
            return []
        return [{
            "id": fallback.id,
            "role": "template",
            "process_mode": fallback.process_mode or "blend",
            "description": "Template layout DNA only.",
        }]
    path = rec.get("layout_file_path") or rec.get("file_path")
    if not path or not os.path.exists(path):
        return []
    return [{
        "id": f"template-{rec.get('page_num') or _template_key_for_page_intent(page_intent)}",
        "role": "template",
        "process_mode": "blend",
        "description": "Template layout DNA only.",
        "category": rec.get("category"),
        "application_strength": rec.get("application_strength") or "standard",
    }]


def _project_logo_refs(project: Project) -> list[ReferenceImage]:
    return [
        img for img in project.reference_images or []
        if (
            img.role == "logo"
            and not img.slide_id
            and is_logo_confirmed(img)
            and img.file_path
            and os.path.exists(img.file_path)
        )
    ]


def _project_logo_refs_for_prompt(project: Project, page_intent: dict | None = None) -> list[dict]:
    logo_refs = _project_logo_refs(project)
    if not logo_refs or not any(should_use_logo_as_scene_asset(page_intent or {}, img) for img in logo_refs):
        return []

    logo_path = prepare_logo_lockup_image([img.file_path for img in logo_refs]) if logo_refs else None
    if not logo_path:
        logo_path = logo_refs[0].file_path
    return [{
        "id": ",".join(str(img.id) for img in logo_refs),
        "role": "logo",
        "process_mode": "blend",
        "description": logo_path,
        "asset_name": "联合标识" if len(logo_refs) > 1 else logo_refs[0].asset_name,
        "logo_anchor": logo_anchor_from_ref(logo_refs[0]),
    }]


def _project_logo_ref(project: Project) -> ReferenceImage | None:
    refs = _project_logo_refs(project)
    return refs[0] if refs else None


def _logo_overlay_url(ref: ReferenceImage, project_id: str) -> str | None:
    if ref.role != "logo" or not ref.file_path or not os.path.exists(ref.file_path):
        return None
    overlay_path = prepare_logo_overlay_image(ref.file_path)
    if not overlay_path or not os.path.exists(overlay_path):
        return None
    return _reference_upload_url(project_id, overlay_path)


def _logo_symbol_overlay_url(ref: ReferenceImage, project_id: str) -> str | None:
    if ref.role != "logo" or not ref.file_path or not os.path.exists(ref.file_path):
        return None
    symbol_path = prepare_logo_symbol_image(ref.file_path)
    if not symbol_path or not os.path.exists(symbol_path):
        return None
    return _reference_upload_url(project_id, symbol_path)


def _with_project_logo_policy(page_intent: dict | None, project: Project) -> dict | None:
    if not isinstance(page_intent, dict):
        return page_intent
    logo_ref = _project_logo_ref(project)
    intent = copy.deepcopy(page_intent)
    policy = intent.get("logo_policy") if isinstance(intent.get("logo_policy"), dict) else {}
    if not logo_ref:
        policy = dict(policy)
        policy["show_logo"] = False
        policy["use_as_scene_asset"] = False
        policy.setdefault("placement", DEFAULT_LOGO_ANCHOR)
        policy.setdefault("scale", "small")
        policy.pop("resolved_overlay_box", None)
        intent["logo_policy"] = policy
        return intent
    page_type = str(intent.get("type") or "").lower()
    if page_type == "ending":
        policy["placement"] = logo_anchor_from_ref(logo_ref)
        policy["scale"] = "small"
        policy.pop("resolved_overlay_box", None)
    elif page_type != "cover":
        policy["placement"] = logo_anchor_from_ref(logo_ref)
        policy.pop("resolved_overlay_box", None)
    normalized_policy = logo_policy_for_page({**intent, "logo_policy": policy})
    if "use_as_scene_asset" in policy:
        normalized_policy["use_as_scene_asset"] = bool(policy.get("use_as_scene_asset"))
    intent["logo_policy"] = normalized_policy
    return intent


def _valid_overlay_asset_ids_for_slide(project: Project, slide: Slide | None) -> set[str]:
    if slide is None:
        return set()
    valid: set[str] = set()
    seen: set[str] = set()
    refs = [*(project.reference_images or []), *(slide.reference_images or [])]
    for ref in refs:
        ref_id = str(getattr(ref, "id", "") or "")
        if not ref_id or ref_id in seen:
            continue
        seen.add(ref_id)
        if not getattr(ref, "file_path", None) or not os.path.exists(ref.file_path):
            continue
        if ref.role == "visual_asset" and not ref.slide_id:
            valid.add(ref_id)
        elif ref.role in {"content_ref", "chart_ref"} and (not ref.slide_id or ref.slide_id == slide.id):
            valid.add(ref_id)
    return valid


def _with_prompt_asset_policies(page_intent: dict | None, project: Project, slide: Slide | None = None) -> dict | None:
    intent = _with_project_logo_policy(page_intent, project)
    if not isinstance(intent, dict) or slide is None:
        return intent
    valid_overlay_ids = _valid_overlay_asset_ids_for_slide(project, slide)
    if isinstance(intent.get("overlay_layers"), list):
        intent["overlay_layers"] = normalize_overlay_layers(
            intent.get("overlay_layers"),
            valid_asset_ids=valid_overlay_ids,
            strict_assets=True,
        )
    intent["available_overlay_asset_ids"] = sorted(valid_overlay_ids)
    return intent


def _existing_output_path(path: str | None) -> str | None:
    if not path:
        return None
    if os.path.exists(path):
        return path
    output_dir = settings.OUTPUT_DIR or "./outputs"
    if path.startswith("./outputs"):
        candidate = os.path.join(os.path.dirname(os.path.abspath(output_dir)), path[2:])
        if os.path.exists(candidate):
            return candidate
    return path


def _with_resolved_logo_overlay_box(page_intent: dict | None, slide: Slide, project: Project) -> dict | None:
    intent = _with_project_logo_policy(page_intent, project)
    if not isinstance(intent, dict) or not slide.image_path or not should_show_logo(intent):
        return intent
    policy = intent.get("logo_policy") if isinstance(intent.get("logo_policy"), dict) else {}
    normalized_policy = logo_policy_for_page(intent)
    if (
        isinstance(policy.get("resolved_overlay_box"), dict)
        and str(policy.get("render_variant") or "").strip().lower() not in {"", "omit"}
    ):
        return intent
    render_policy_input = {**policy, **normalized_policy}
    if "render_variant" not in normalized_policy:
        render_policy_input.pop("render_variant", None)
    logo_refs = _project_logo_refs(project)
    if not logo_refs:
        return intent
    logo_path = prepare_logo_lockup_image([ref.file_path for ref in logo_refs])
    if not logo_path:
        return intent
    symbol_path = prepare_logo_symbol_image(logo_refs[0].file_path) if len(logo_refs) == 1 else None
    render_policy = resolve_logo_render_policy(
        _existing_output_path(slide.image_path),
        logo_path,
        symbol_path,
        str(intent.get("type") or slide.type or "content").lower(),
        normalized_policy.get("placement") or logo_anchor_from_ref(logo_refs[0]),
        normalized_policy.get("scale") or "small",
        render_policy_input,
    )
    next_intent = copy.deepcopy(intent)
    next_policy = dict(normalized_policy)
    next_policy.update({k: v for k, v in render_policy.items() if v is not None})
    next_intent["logo_policy"] = next_policy
    return next_intent


def _repair_project_logo_policies(project: Project, slides: list[Slide]) -> bool:
    """Normalize stale logo policies so old projects do not keep hiding required logos."""
    if not _project_logo_refs(project):
        return False
    changed = False
    for slide in slides:
        current = copy.deepcopy(slide.visual_json) if isinstance(slide.visual_json, dict) else {}
        repaired = _with_resolved_logo_overlay_box(current, slide, project)
        if not isinstance(repaired, dict):
            continue
        if repaired.get("logo_policy") == current.get("logo_policy"):
            continue
        slide.visual_json = repaired
        flag_modified(slide, "visual_json")
        changed = True
    return changed


def _visual_asset_summary(ref: ReferenceImage) -> str:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    parts = []
    if ref.asset_name:
        parts.append(f"name={ref.asset_name}")
    if ref.asset_kind:
        parts.append(f"kind={ref.asset_kind}")
    if ref.usage_note:
        parts.append(f"user_note={ref.usage_note}")
    subject = analysis.get("subject")
    if subject and subject != ref.asset_name:
        parts.append(f"subject={subject}")
    if analysis.get("description"):
        parts.append(f"description={analysis['description']}")
    identity_elements = analysis.get("identity_elements")
    if isinstance(identity_elements, list) and identity_elements:
        parts.append("identity_elements=" + "、".join(str(x) for x in identity_elements[:6]))
    features = analysis.get("distinctive_features")
    if isinstance(features, list) and features:
        parts.append("features=" + "、".join(str(x) for x in features[:5]))
    must_not_change = analysis.get("must_not_change")
    if isinstance(must_not_change, list) and must_not_change:
        parts.append("must_not_change=" + "、".join(str(x) for x in must_not_change[:6]))
    keywords = analysis.get("suggested_keywords")
    if isinstance(keywords, list) and keywords:
        parts.append("keywords=" + "、".join(str(x) for x in keywords[:16]))
    source_page = _asset_source_page_num(ref)
    source_doc = analysis.get("source_document")
    if source_doc and source_page:
        source_type = str(analysis.get("source_type") or "").upper()
        source_label = "PDF" if source_type == "PDF" else "PPT"
        parts.append(f"source=来自上传 {source_label}「{source_doc}」第{source_page}页")
    source_text = analysis.get("source_slide_text") or analysis.get("nearby_text")
    if isinstance(source_text, str) and source_text.strip():
        parts.append(f"source_slide_text={source_text[:240]}")
    tags = analysis.get("asset_tags")
    if isinstance(tags, list) and tags:
        parts.append("tags=" + "、".join(str(x) for x in tags[:16]))
    if analysis.get("recommended_usage"):
        parts.append(f"recommended_usage={analysis['recommended_usage']}")
    if analysis.get("fidelity_note"):
        parts.append(f"fidelity={analysis['fidelity_note']}")
    return "; ".join(parts)


PPTX_LEGACY_IDENTITY_ASSET_TERMS = (
    "产品", "包装", "瓶", "sku", "SKU", "货架", "终端", "陈列", "样品", "实物",
    "设备", "取件机", "快递柜", "人物", "模特", "代言人", "团队", "肖像",
    "物料", "主视觉", "KV", "品牌物料", "产品物料", "海报", "立牌", "展架",
    "吊旗", "标识", "导视",
)
PPTX_LEGACY_LOW_VALUE_TERMS = (
    "背景", "氛围", "氛围感", "风景", "插画", "装饰", "纹理", "光效", "底图", "配图",
)
PPTX_LEGACY_NON_REUSABLE_TERMS = (
    "二维码", "扫码", "身份码", "取件码", "小程序码", "条码", "链接码",
    "手机边框", "手机框", "手机壳", "手机外框", "手机界面", "界面外框", "屏幕框",
    "app界面", "APP界面", "小程序界面", "小程序", "mockup", "Mockup",
)


def _float_metadata(value, default: float = 0.0) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _visual_asset_planning_score(ref: ReferenceImage) -> float:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    if not analysis.get("source_document"):
        return 100.0
    score = _float_metadata(analysis.get("importance_score"), 0.0)
    if score:
        return score
    kind = str(ref.asset_kind or analysis.get("detected_kind") or "other").lower()
    if kind in {"product", "material", "person"}:
        return 55.0
    return 0.0


def _is_planning_visual_asset(ref: ReferenceImage) -> bool:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    if not analysis.get("source_document"):
        return True
    source_text = " ".join([
        str(ref.asset_name or ""),
        str(analysis.get("source_slide_text") or ""),
        " ".join(str(tag) for tag in (analysis.get("asset_tags") or []) if tag),
    ])
    source_lower = source_text.lower()
    if any(term.lower() in source_lower for term in PPTX_LEGACY_NON_REUSABLE_TERMS):
        return False
    reason_lower = str(analysis.get("selection_reason") or "").lower()
    if "qr/identity code" in reason_lower or "phone/ui container" in reason_lower or "layout chrome" in reason_lower:
        return False
    tier = str(analysis.get("selection_tier") or "").lower()
    if tier:
        return tier == "core_global"

    # Legacy PPT assets created before the selection-tier metadata existed were
    # over-inclusive. Keep only brand/product identity assets; charts,
    # screenshots, maps, and diagrams stay page-level evidence.
    kind = str(ref.asset_kind or analysis.get("detected_kind") or "other").lower()
    area_ratio = _float_metadata(analysis.get("area_ratio"), 0.0)
    has_identity_term = any(term.lower() in source_lower for term in PPTX_LEGACY_IDENTITY_ASSET_TERMS)
    has_low_value_term = any(term.lower() in source_lower for term in PPTX_LEGACY_LOW_VALUE_TERMS)
    if kind in {"product", "person"}:
        return (not area_ratio or area_ratio >= 0.006) and has_identity_term
    if kind == "material":
        if area_ratio and area_ratio < 0.006:
            return False
        if has_low_value_term and not has_identity_term:
            return False
        return has_identity_term
    if kind == "scene":
        return False
    if kind == "other":
        return False
    return False


def _project_visual_assets_for_planning(project: Project) -> list[dict]:
    assets = []
    for ref in project.reference_images or []:
        if ref.role != "visual_asset" or ref.slide_id:
            continue
        if not ref.file_path or not os.path.exists(ref.file_path):
            continue
        if not _is_planning_visual_asset(ref):
            continue
        importance_score = _visual_asset_planning_score(ref)
        analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
        assets.append({
            "id": ref.id,
            "name": ref.asset_name or os.path.splitext(os.path.basename(ref.file_path or ""))[0],
            "kind": ref.asset_kind or "other",
            "process_mode": ref.process_mode or default_visual_asset_process_mode(ref.asset_kind),
            "usage_note": ref.usage_note or "",
            "selection_tier": analysis.get("selection_tier") or ("manual" if not analysis.get("source_document") else "legacy_candidate"),
            "importance_score": importance_score,
            "source_document": _asset_source_document(ref),
            "source_page_num": _asset_source_page_num(ref),
            "analysis_summary": _visual_asset_summary(ref),
        })
    assets.sort(key=lambda item: (-float(item.get("importance_score") or 0), str(item.get("name") or "")))
    return assets


def _project_refs_for_prompt(
    project: Project,
    visual_asset_ids: list[str] | None = None,
    page_intent: dict | None = None,
) -> list[dict]:
    refs = []
    wanted = set(visual_asset_ids or [])
    overlay_wanted = exact_overlay_asset_ids(page_intent or {})
    route_modes = (page_intent or {}).get("asset_route_modes") or {}
    if not isinstance(route_modes, dict):
        route_modes = {}
    route_modes = {str(k): str(v).lower() for k, v in route_modes.items() if k and v}
    wanted = {
        asset_id for asset_id in wanted
        if asset_id not in overlay_wanted and route_modes.get(str(asset_id)) != "overlay"
    }
    if wanted:
        for img in project.reference_images or []:
            if img.role != "visual_asset" or img.id not in wanted:
                continue
            if not img.file_path or not os.path.exists(img.file_path):
                continue
            route_mode = str(route_modes.get(str(img.id)) or "").lower()
            if not route_mode:
                route_mode = "double_blend" if str(img.asset_kind or "").lower() in {"product", "material"} else "blend"
            effective_process_mode = "crop" if route_mode == "double_blend" else (
                "original" if route_mode == "overlay" else "blend"
            )
            refs.append({
                "id": img.id,
                "role": img.role,
                "asset_route_mode": route_mode,
                "process_mode": effective_process_mode,
                "description": _visual_asset_summary(img) or img.file_path,
                "asset_name": img.asset_name,
                "asset_kind": img.asset_kind,
                "usage_note": img.usage_note,
            })
    refs.extend(_project_logo_refs_for_prompt(project, page_intent))
    refs.extend(_project_template_refs_for_prompt(project, page_intent))
    return refs


def _slide_refs_for_prompt(slide: Slide, visual_json: dict | None = None) -> list[dict]:
    visual = visual_json if isinstance(visual_json, dict) else {}
    overlay_asset_ids = exact_overlay_asset_ids(visual)
    refs: list[dict] = []
    for ref in sorted(slide.reference_images or [], key=_reference_image_sort_key):
        if str(ref.id) in overlay_asset_ids:
            continue
        analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
        is_content_block = analysis.get("source") == "content_block"
        description = _format_pptx_reference_analysis(ref) or os.path.basename(ref.file_path or "")
        if is_content_block:
            description = (
                f"Content block visual asset: {ref.asset_name or '画面素材'}. "
                f"Type={analysis.get('content_block_kind') or analysis.get('visual_type') or 'diagram'}. "
                "Keep the core nodes, arrows, table cells, and relationship structure from the rendered reference."
            )
        refs.append({
            "id": ref.id,
            "role": ref.role,
            "description": description,
            "process_mode": ref.process_mode or "blend",
            "asset_name": ref.asset_name,
            "asset_kind": ref.asset_kind,
            "usage_note": ref.usage_note,
            "asset_route_mode": analysis.get("asset_route_mode") if is_content_block else None,
            "asset_analysis": analysis,
        })
    return refs


def _invalidate_visual_asset_dependent_outputs(project: Project):
    """
    Core visual assets affect matching, prompts, and generated images, but not the
    selected style. Keep the style choice intact and move downstream outputs back
    to visual planning.
    """
    if project.selected_style:
        project.content_plan_confirmed = True
    for slide in project.slides or []:
        _mark_slide_artifacts_stale(slide, content=True)
        slide.error_msg = None
    if project.slides:
        if project.content_plan_confirmed and project.status in {"draft", "planning"}:
            project.status = "visual_ready"
        elif not project.content_plan_confirmed:
            project.status = "planning"


def _clear_slide_generation_outputs(slide: Slide, *, clear_visual: bool):
    if clear_visual:
        slide.visual_json = {}
    slide.prompt_text = None
    slide.image_path = None
    slide.error_msg = None


def _mark_slide_artifacts_stale(slide: Slide, **flags: bool):
    slide.visual_json = with_stale_flags(slide.visual_json if isinstance(slide.visual_json, dict) else {}, **flags)
    flag_modified(slide, "visual_json")


def _project_has_stale_artifacts(slides: list[Slide]) -> bool:
    return any(has_stale_flags(slide.visual_json) for slide in slides)


_PROMPT_VISIBLE_LABELS = {
    "headline": "headline",
    "subhead": "subhead",
    "body": "body",
    "info block body": "body",
    "linked caption body": "body",
}


def _content_page_for_generation_check(slide: Slide) -> dict:
    content = copy.deepcopy(slide.content_json) if isinstance(slide.content_json, dict) else {}
    content.setdefault("page_num", slide.page_num)
    content.setdefault("type", slide.type or "content")
    if slide.type:
        content["type"] = slide.type
    return content


def _normalize_visible_copy_for_generation(value) -> str:
    cleaned = _strip_markdown(str(value or ""))
    cleaned = cleaned.replace("\\", "")
    cleaned = re.sub(r"^[\s•·\-*]+", "", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned)
    cleaned = re.sub(r"^(?:\[表格\]\s*){2,}", "[表格] ", cleaned)
    return cleaned.strip()


def _content_body_lines_for_generation(body) -> list[str]:
    if isinstance(body, str):
        return [line.strip() for line in body.splitlines() if line.strip()]
    if isinstance(body, list):
        lines = []
        for item in body:
            value = item.get("content") if isinstance(item, dict) else item
            if str(value or "").strip():
                lines.append(str(value).strip())
        return lines
    return []


def _expected_prompt_visible_slots(slide: Slide) -> dict[str, list[str]]:
    content = slide.content_json if isinstance(slide.content_json, dict) else {}
    raw_text = content.get("text_content") if isinstance(content.get("text_content"), dict) else {}
    visual = slide.visual_json if isinstance(slide.visual_json, dict) else {}
    page_type = str(visual.get("type") or slide.type or content.get("type") or "").strip().lower()
    layout = str(visual.get("layout") or "").strip().lower()
    is_punchline = page_type in {"hero", "quote"} or layout == "hero"
    content_text, _visual_intents, _diagram_labels = _content_visual_contract(raw_text or {})

    expected: dict[str, list[str]] = {"headline": [], "subhead": [], "body": []}
    for slot in ("headline", "subhead"):
        normalized = _normalize_visible_copy_for_generation(content_text.get(slot))
        if normalized:
            expected[slot].append(normalized)
    if not is_punchline:
        for line in _content_body_lines_for_generation(content_text.get("body")):
            normalized = _normalize_visible_copy_for_generation(line)
            if normalized:
                expected["body"].append(normalized)
    return expected


def _prompt_visible_slots(prompt_text: str | None) -> dict[str, list[str]]:
    prompt = str(prompt_text or "")
    match = re.search(r"Visible Text:\s*\n(?P<section>.*?)(?:\n\s*\n|$)", prompt, flags=re.S)
    if not match:
        return {}
    slots: dict[str, list[str]] = {"headline": [], "subhead": [], "body": []}
    section = match.group("section")
    directive_re = re.compile(
        r'^(Headline|Subhead|Body|Info block body|Linked caption body):\s*"(?P<value>.*?)"'
        r'(?=\n(?:Headline|Subhead|Body|Info block body|Linked caption body|Visible text rule:)|\s*$)',
        flags=re.M | re.S,
    )
    for item in directive_re.finditer(section):
        label = item.group(1)
        slot = _PROMPT_VISIBLE_LABELS.get(label.strip().lower())
        if not slot:
            continue
        value = item.group("value").strip()
        normalized = _normalize_visible_copy_for_generation(value)
        if normalized:
            slots.setdefault(slot, []).append(normalized)
    return {slot: values for slot, values in slots.items() if values}


def _slide_prompt_visible_text_is_current(slide: Slide) -> bool:
    expected = _expected_prompt_visible_slots(slide)
    expected_all = {
        value
        for values in expected.values()
        for value in values
    }
    prompt_slots = _prompt_visible_slots(slide.prompt_text)
    if expected_all and not prompt_slots:
        return False
    for slot, values in expected.items():
        for value in values:
            if value not in prompt_slots.get(slot, []):
                return False
    for slot, values in prompt_slots.items():
        if slot not in {"headline", "subhead", "body"}:
            continue
        for value in values:
            if value not in expected_all:
                return False
    return True


def _slide_generation_guidance_matches_current_content(slide: Slide) -> bool:
    page = _content_page_for_generation_check(slide)
    visible_text = _page_visible_copy_text(page)
    content = slide.content_json if isinstance(slide.content_json, dict) else {}
    guidance_sources = [
        content.get("visual_suggestion"),
        content.get("visual_requirements"),
    ]
    return not any(
        _guidance_mentions_absent_visible_copy(source, visible_text)
        for source in guidance_sources
        if source
    )


def _slide_prompt_dependencies_match_generation_sources(slide: Slide, current_deps: dict[str, str]) -> bool:
    visual = slide.visual_json if isinstance(slide.visual_json, dict) else {}
    meta = artifact_meta(visual)
    prompt_deps = meta.get("prompt_dependencies") if isinstance(meta.get("prompt_dependencies"), dict) else None
    if not prompt_deps:
        return True
    for key in ("content", "selected_style", "style_assets", "visual_assets"):
        if prompt_deps.get(key) != current_deps.get(key):
            return False
    return True


def _slide_prompt_text_matches_current_generation_context(slide: Slide) -> bool:
    return _slide_prompt_visible_text_is_current(slide)


def _raise_if_generation_inputs_not_current(project: Project, slides: list[Slide], target_slides: list[Slide]) -> None:
    stale_plan_pages = [
        s.page_num
        for s in target_slides
        if has_stale_flags(s.visual_json, "content", "visual")
    ]
    if stale_plan_pages:
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, stale_plan_pages))} 页有未应用的内容或画面修改，请先更新画面方案。"
        )

    stale_guidance_pages = [
        s.page_num
        for s in target_slides
        if not _slide_generation_guidance_matches_current_content(s)
    ]
    if stale_guidance_pages:
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, stale_guidance_pages))} 页的画面方案与当前内容不一致，请先更新画面方案。"
        )

    current_deps = dependency_signature(project, slides)
    dependency_stale_pages = [
        s.page_num
        for s in target_slides
        if s.prompt_text and not _slide_prompt_dependencies_match_generation_sources(s, current_deps)
    ]
    if dependency_stale_pages:
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, dependency_stale_pages))} 页的生图 Prompt 不是基于最新风格或素材，请先更新画面方案。"
        )

    prompt_stale_pages = [
        s.page_num
        for s in target_slides
        if s.prompt_text and not _slide_prompt_text_matches_current_generation_context(s)
    ]
    if prompt_stale_pages:
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, prompt_stale_pages))} 页的生图 Prompt 不是基于最新内容，请先更新画面方案。"
        )


def _invalidate_content_dependent_outputs(project: Project):
    """Content changes make downstream page artifacts stale without deleting them."""
    slides = list(project.slides or [])
    if not slides:
        project.status = "draft"
        project.content_plan_confirmed = False
        project.style_proposal = None
        project.selected_style = None
        return
    for slide in slides:
        _mark_slide_artifacts_stale(slide, content=True)
        slide.error_msg = None
    if project.content_plan_confirmed or project.selected_style:
        project.content_plan_confirmed = True
        if project.status in {"draft", "planning"}:
            project.status = "visual_ready"
    else:
        project.status = "planning"


def _invalidate_style_dependent_outputs(project: Project):
    """Style-source changes keep confirmed content and current images as stale references."""
    slides = list(project.slides or [])
    if slides:
        project.content_plan_confirmed = True
    for slide in slides:
        _mark_slide_artifacts_stale(slide, content=True)
        slide.error_msg = None
    if project.status in {"draft", "planning"} and slides:
        project.status = "visual_ready"


def _invalidate_visual_plan_dependent_outputs(project: Project, slides: list[Slide]):
    """Visual edits/page refs invalidate prompts and images, not confirmed content."""
    for slide in slides:
        _mark_slide_artifacts_stale(slide, visual=True)
        slide.error_msg = None


def _reassemble_generated_pptx_if_possible(project: Project) -> bool:
    """
    Logo review/anchor changes only affect the PPTX overlay layer once images exist.
    Rebuild the export in place instead of sending users back through image generation.
    """
    if project.status not in {"prototype_ready", "completed"}:
        return False
    slides = sorted(project.slides or [], key=lambda slide: slide.page_num or 0)
    if not slides:
        return False

    slide_images = []
    for slide in slides:
        image_path = _existing_output_path(slide.image_path)
        if slide.status != "completed" or not image_path or not os.path.exists(image_path):
            return False
        visual_json = _with_resolved_logo_overlay_box(slide.visual_json, slide, project)
        slide_images.append({
            "page_num": slide.page_num,
            "type": slide.type,
            "image_path": image_path,
            "speaker_notes": (slide.content_json or {}).get("speaker_notes") if isinstance(slide.content_json, dict) else "",
            "visual_json": visual_json if isinstance(visual_json, dict) else {},
        })

    logo_refs = _project_logo_refs(project)
    logo_config = (
        {
            "file_paths": [ref.file_path for ref in logo_refs],
            "anchor": logo_anchor_from_ref(logo_refs[0]),
        }
        if logo_refs else None
    )
    overlay_assets = {
        str(ref.id): {
            "file_path": ref.file_path,
            "asset_name": ref.asset_name,
            "asset_kind": ref.asset_kind,
        }
        for ref in project.reference_images or []
        if (
            (
                ref.role == "visual_asset" and not ref.slide_id
            )
            or (
                ref.role in {"content_ref", "chart_ref"} and ref.slide_id
            )
        )
        and ref.file_path
        and os.path.exists(ref.file_path)
    }
    filename = "presentation.pptx" if project.status == "completed" else "prototype.pptx"
    output_path = os.path.join(settings.OUTPUT_DIR or "./outputs", project.id, filename)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    assemble_pptx(
        slide_images=slide_images,
        output_path=output_path,
        logo_config=logo_config,
        overlay_assets=overlay_assets,
    )
    return True


def _project_export_overlay_assets(project: Project) -> dict:
    return {
        str(ref.id): {
            "file_path": ref.file_path,
            "asset_name": ref.asset_name,
            "asset_kind": ref.asset_kind,
        }
        for ref in project.reference_images or []
        if (
            (
                ref.role == "visual_asset" and not ref.slide_id
            )
            or (
                ref.role in {"content_ref", "chart_ref"} and ref.slide_id
            )
        )
        and ref.file_path
        and os.path.exists(ref.file_path)
    }


def _project_export_logo_config(project: Project) -> dict | None:
    logo_refs = _project_logo_refs(project)
    if not logo_refs:
        return None
    return {
        "file_paths": [ref.file_path for ref in logo_refs],
        "anchor": logo_anchor_from_ref(logo_refs[0]),
    }


def _assemble_partial_project_pptx(project: Project, slides: list[Slide], output_path: str) -> int:
    """
    Build an export with one PPTX page per current slide card.

    Slides without a generated image are kept as blank pages so users can export
    partial work without losing the planned deck structure.
    """
    if not slides:
        return 0

    slide_images = []
    generated_count = 0
    for slide in sorted(slides, key=lambda item: item.page_num or 0):
        image_path = _existing_output_path(slide.image_path)
        if image_path and os.path.exists(image_path):
            generated_count += 1
        else:
            image_path = None
        visual_json = (
            _with_resolved_logo_overlay_box(slide.visual_json, slide, project)
            if image_path
            else (slide.visual_json if isinstance(slide.visual_json, dict) else {})
        )
        slide_images.append({
            "page_num": slide.page_num,
            "type": slide.type,
            "image_path": image_path,
            "speaker_notes": (slide.content_json or {}).get("speaker_notes") if isinstance(slide.content_json, dict) else "",
            "visual_json": visual_json if isinstance(visual_json, dict) else {},
        })

    if generated_count == 0:
        return 0

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    assemble_pptx(
        slide_images=slide_images,
        output_path=output_path,
        logo_config=_project_export_logo_config(project),
        overlay_assets=_project_export_overlay_assets(project),
    )
    return generated_count


def _ensure_completed_final_pptx(project: Project, slides: list[Slide]) -> tuple[bool, str]:
    pptx_path = _final_pptx_path(project.id)
    if (
        project.status != "completed"
        or not slides
        or _project_has_stale_artifacts(slides)
        or any(s.status != "completed" or not s.image_path for s in slides)
    ):
        return False, pptx_path
    if _final_pptx_matches_completed_slides(pptx_path, slides):
        return True, pptx_path

    generated_count = _assemble_partial_project_pptx(project, slides, pptx_path)
    if generated_count != len(slides):
        return False, pptx_path
    return _final_pptx_matches_completed_slides(pptx_path, slides), pptx_path


def _resolve_generation_page_nums(slides: list[Slide], requested_page_nums: list[int] | None, prototype: bool) -> list[int] | None:
    """Resolve generation targets. Prototype without explicit pages samples the first 3 slides."""
    if prototype and not requested_page_nums:
        return [s.page_num for s in slides[:3]]
    return requested_page_nums


def _normalize_content_json_markdown(content_json: dict) -> dict:
    """Normalize Markdown fields at API boundaries while preserving structure."""
    if not isinstance(content_json, dict):
        return content_json
    content = copy.deepcopy(content_json)
    text_content = content.get("text_content")
    if isinstance(text_content, dict):
        for key in ("headline", "subhead", "body"):
            value = text_content.get(key)
            if isinstance(value, str):
                text_content[key] = normalize_markdown_content(value)
            elif isinstance(value, list):
                text_content[key] = [
                    normalize_markdown_content(item) if isinstance(item, str) else item
                    for item in value
                ]
    if isinstance(content.get("speaker_notes"), str):
        content["speaker_notes"] = normalize_markdown_content(content["speaker_notes"])
    content["content_blocks"] = normalize_content_blocks(content)
    body_from_blocks = blocks_to_markdown(content["content_blocks"])
    if body_from_blocks:
        text_content = content.get("text_content")
        if not isinstance(text_content, dict):
            text_content = {}
        text_content["body"] = normalize_markdown_content(body_from_blocks)
        content["text_content"] = text_content
    return content


def _json_equivalent(left, right) -> bool:
    try:
        return json.dumps(left, ensure_ascii=False, sort_keys=True, default=str) == json.dumps(
            right,
            ensure_ascii=False,
            sort_keys=True,
            default=str,
        )
    except TypeError:
        return left == right


def _content_json_for_generation(slide: Slide) -> dict:
    content = copy.deepcopy(slide.content_json) if isinstance(slide.content_json, dict) else {}
    changed = False
    if content.get("visual_suggestion") != "":
        content["visual_suggestion"] = ""
        changed = True
    if content.get("visual_requirements") != []:
        content["visual_requirements"] = []
        changed = True
    if changed:
        slide.content_json = content
        flag_modified(slide, "content_json")
    return content


def _content_block_refs(db: Session, project_id: str, slide_id: str) -> list[ReferenceImage]:
    refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project_id,
        ReferenceImage.slide_id == slide_id,
        ReferenceImage.role == "chart_ref",
    ).all()
    return [
        ref for ref in refs
        if isinstance(ref.asset_analysis, dict)
        and ref.asset_analysis.get("source") == "content_block"
    ]


def _content_block_asset_dir(project_id: str, slide_id: str) -> str:
    base_dir = os.path.join(settings.UPLOAD_DIR, project_id, "content_blocks", slide_id)
    if not os.path.isabs(base_dir):
        base_dir = os.path.abspath(base_dir)
    os.makedirs(base_dir, exist_ok=True)
    return base_dir


def _content_block_ref_payload(block: dict, file_path: str) -> dict:
    route_mode = normalize_route_mode(block.get("route_mode"), block.get("kind"))
    return {
        "source": "content_block",
        "analysis_type": "content_block",
        "analysis_status": "completed",
        "content_block_id": block.get("id"),
        "content_block_kind": block.get("kind"),
        "visual_type": block.get("visual_type") or block.get("kind"),
        "source_hash": block.get("source_hash") or stable_block_hash(block),
        "asset_route_mode": route_to_asset_route_mode(route_mode),
        "rendered_file": os.path.basename(file_path),
    }


def _sync_content_block_assets(db: Session, project: Project, slide: Slide, content: dict) -> dict:
    """Render visual content blocks into page-level chart refs and keep overlay state aligned."""
    blocks = normalize_content_blocks(content)
    content["content_blocks"] = blocks
    text_content = content.get("text_content") if isinstance(content.get("text_content"), dict) else {}
    content["text_content"] = {
        **text_content,
        "body": normalize_markdown_content(blocks_to_markdown(blocks)),
    }

    visual_blocks = [block for block in blocks if is_visual_block(block)]
    visual_by_id = {str(block.get("id")): block for block in visual_blocks}
    existing_refs = _content_block_refs(db, project.id, slide.id)
    refs_by_block_id = {
        str(ref.asset_analysis.get("content_block_id")): ref
        for ref in existing_refs
        if isinstance(ref.asset_analysis, dict) and ref.asset_analysis.get("content_block_id")
    }

    visual = copy.deepcopy(slide.visual_json) if isinstance(slide.visual_json, dict) else {}
    overlay_layers = normalize_overlay_layers(visual.get("overlay_layers"), valid_asset_ids=None, strict_assets=False)

    for ref in existing_refs:
        block_id = str((ref.asset_analysis or {}).get("content_block_id") or "")
        if block_id and block_id not in visual_by_id:
            overlay_layers = [layer for layer in overlay_layers if str(layer.get("asset_id")) != str(ref.id)]
            db.delete(ref)

    output_dir = _content_block_asset_dir(project.id, slide.id)
    for block in visual_blocks:
        route_mode = normalize_route_mode(block.get("route_mode"), block.get("kind"))
        block["route_mode"] = route_mode
        block["source_hash"] = stable_block_hash(block)
        file_path = render_visual_block(block, output_dir)
        block_id = str(block.get("id"))
        ref = refs_by_block_id.get(block_id)
        payload = _content_block_ref_payload(block, file_path)
        if ref is None:
            ref = ReferenceImage(
                project_id=project.id,
                slide_id=slide.id,
                file_path=file_path,
                role="chart_ref",
                process_mode=route_mode,
                asset_name=str(block.get("title") or "画面素材"),
                usage_note=f"由正文中的「{block.get('title') or '画面素材'}」生成",
                asset_analysis=payload,
            )
            db.add(ref)
            db.flush()
        else:
            ref.file_path = file_path
            ref.process_mode = route_mode
            ref.asset_name = str(block.get("title") or ref.asset_name or "画面素材")
            ref.usage_note = f"由正文中的「{block.get('title') or '画面素材'}」生成"
            ref.asset_analysis = {
                **(ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}),
                **payload,
            }
        block["rendered_asset_id"] = ref.id

        overlay_layers = [layer for layer in overlay_layers if str(layer.get("asset_id")) != str(ref.id)]
        if route_mode == "original":
            overlay_layers.append({
                "id": f"ov_{ref.id}",
                "asset_id": ref.id,
                "enabled": True,
                "preset": "center-card",
                "fit": "contain",
                "mode": "exact_card",
                "usage_note": f"{ref.asset_name or '画面素材'}：原样保留",
            })

    visual["overlay_layers"] = normalize_overlay_layers(overlay_layers, valid_asset_ids=None, strict_assets=False)
    slide.visual_json = visual
    content["content_blocks"] = blocks
    return content


def _plain_markdown_value(value) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value.replace("\r\n", "\n").replace("\r", "\n").strip()
    if isinstance(value, list):
        parts = []
        for item in value:
            if isinstance(item, dict):
                parts.append(str(item.get("content") or "").strip())
            else:
                parts.append(str(item).strip())
        return "\n\n".join(part for part in parts if part)
    return str(value).strip()


def _markdown_section(label: str, value: str) -> str:
    return f"### {label}\n\n{value.strip()}\n"


def _safe_export_filename(title: str) -> str:
    name = re.sub(r'[\\/:*?"<>|]+', "-", str(title or "").strip()).strip(" .")
    return (name[:80] or "内容规划") + ".md"


def _build_content_plan_markdown(project: Project, slides: list[Slide]) -> str:
    lines: list[str] = [
        f"# {project.title} - 内容规划导出",
        "",
        "<!--",
        "PPTGOD_EXPORT_KIND: content_plan_markdown",
        f"project_id: {project.id}",
        f"project_status: {project.status}",
        f"content_plan_confirmed: {1 if project.content_plan_confirmed else 0}",
        "说明：可以直接修改每页的标题、副标题、正文、备注。请保留 PPTGOD_PAGE_START / PPTGOD_PAGE_END 注释，方便后续识别页边界。",
        "-->",
        "",
    ]

    for slide in slides:
        content = slide.content_json if isinstance(slide.content_json, dict) else {}
        text_content = content.get("text_content") if isinstance(content.get("text_content"), dict) else {}
        page_num = int(slide.page_num or content.get("page_num") or 0)
        page_type = str(slide.type or content.get("type") or "content")
        section_title = _plain_markdown_value(content.get("section_title"))
        headline = _plain_markdown_value(text_content.get("headline"))
        subhead = _plain_markdown_value(text_content.get("subhead"))
        body = _plain_markdown_value(text_content.get("body"))
        speaker_notes = _plain_markdown_value(content.get("speaker_notes"))

        lines.extend([
            "---",
            f"<!-- PPTGOD_PAGE_START page_num={page_num} type={page_type} section_title={json.dumps(section_title, ensure_ascii=False)} -->",
            f"## P{page_num} · {page_type}" + (f" · {section_title}" if section_title else ""),
            "",
            _markdown_section("标题", headline),
            _markdown_section("副标题", subhead),
            _markdown_section("正文", body),
            _markdown_section("备注", speaker_notes),
            f"<!-- PPTGOD_PAGE_END page_num={page_num} -->",
            "",
        ])

    return "\n".join(lines).rstrip() + "\n"


def _update_progress(project_id: str, data: dict, run_id: str | None = None):
    """更新指定项目的生成进度。"""
    generation_progress[project_id] = data
    if run_id:
        db = SessionLocal()
        try:
            update_run_progress(
                db,
                run_id,
                stage=data.get("stage"),
                message=data.get("message"),
                completed_count=data.get("current_page"),
                total_count=data.get("total_pages") or data.get("target_count"),
            )
            db.commit()
        except Exception as exc:
            db.rollback()
            logger.warning(f"Failed to persist run progress for {run_id}: {exc}")
        finally:
            db.close()


def _mark_generation_idle(project: Project | None, db: Session, reason: str):
    """Return a stale generating project to a recoverable state."""
    if not project:
        return
    slides = db.query(Slide).filter(
        Slide.project_id == project.id,
        Slide.status == "generating",
    ).all()
    for slide in slides:
        slide.status = "prompt_ready"
        slide.error_msg = reason
        try:
            redis_client.delete(f"project:{project.id}:slide:{slide.page_num}:generating")
        except Exception as exc:
            logger.warning(
                "Failed to release stale page lock project=%s page=%s: %s",
                project.id, slide.page_num, exc,
            )
    if project.status not in {"draft", "planning", "visual_ready", "prompt_ready", "prototype_ready", "completed", "failed"}:
        project.status = "prompt_ready"
    db.commit()


def _cleanup_stale_generating_slides_for_project(db: Session, project_id: str) -> int:
    """Release orphaned Redis locks and reset generating slides to prompt_ready.

    Called before starting a new generation run to ensure stale locks from
    crashed or interrupted tasks do not block the new run.
    """
    slides = db.query(Slide).filter(
        Slide.project_id == project_id,
        Slide.status == "generating",
    ).all()
    cleaned = 0
    for slide in slides:
        lock_key = f"project:{project_id}:slide:{slide.page_num}:generating"
        try:
            # If lock is gone the task clearly died; if lock exists we also
            # delete it because the caller has already verified no active run.
            redis_client.delete(lock_key)
        except Exception as exc:
            logger.warning(
                "Pre-clean failed to delete lock project=%s page=%s: %s",
                project_id, slide.page_num, exc,
            )
        slide.status = "prompt_ready"
        slide.error_msg = None
        cleaned += 1
    if cleaned:
        db.commit()
        logger.info("Pre-cleaned %s stale generating slides for project %s", cleaned, project_id)
    return cleaned


def _celery_task_id_from_payload(payload: dict | None) -> str | None:
    if not isinstance(payload, dict):
        return None
    task_id = payload.get("id")
    if task_id:
        return str(task_id)
    request = payload.get("request")
    if isinstance(request, dict) and request.get("id"):
        return str(request["id"])
    return None


def _celery_task_present_in_worker(task_id: str) -> bool | None:
    """Return whether a task is visible to the currently connected workers."""
    try:
        from app.celery_app import celery_app

        inspector = celery_app.control.inspect(timeout=1.0)
        snapshots = [inspector.active(), inspector.reserved(), inspector.scheduled()]
    except Exception as exc:
        logger.warning(f"Failed to inspect Celery workers for task {task_id}: {exc}")
        return None

    saw_worker = False
    saw_snapshot = False
    for snapshot in snapshots:
        if snapshot is None:
            continue
        if not isinstance(snapshot, dict):
            continue
        saw_snapshot = True
        if snapshot:
            saw_worker = True
        for tasks in snapshot.values():
            if not isinstance(tasks, list):
                continue
            for item in tasks:
                if _celery_task_id_from_payload(item) == task_id:
                    return True
    if not saw_snapshot or not saw_worker:
        return None
    return False


def _celery_workers_online() -> bool | None:
    try:
        from app.celery_app import celery_app

        replies = celery_app.control.inspect(timeout=1.0).ping()
    except Exception as exc:
        logger.warning(f"Failed to ping Celery workers: {exc}")
        return None
    if replies is None:
        return False
    if not isinstance(replies, dict):
        return None
    return bool(replies)


def _celery_result_state(task_id: str) -> str:
    try:
        from app.celery_app import celery_app

        return str(AsyncResult(task_id, app=celery_app).state)
    except Exception as exc:
        logger.warning(f"Failed to read Celery result state for task {task_id}: {exc}")
        return "UNKNOWN"


def _run_seconds_since(value) -> float:
    if not value:
        return 0
    try:
        return max(0, time.time() - float(value.timestamp()))
    except Exception:
        return 0


def _task_age_seconds(project_id: str, active_run) -> float:
    started_raw = redis_client.get(f"project:{project_id}:task_started_at")
    try:
        started_at = float(started_raw.decode() if isinstance(started_raw, bytes) else started_raw)
    except (TypeError, ValueError):
        started_at = 0
    if started_at > 0:
        return max(0, time.time() - started_at)
    return _run_seconds_since(getattr(active_run, "started_at", None))


def _queued_celery_timeout_seconds() -> int:
    try:
        return max(120, int(settings.CELERY_QUEUE_WAIT_TIMEOUT_SECONDS or 0))
    except (TypeError, ValueError):
        return 3600


def _stale_missing_celery_task_if_needed(project: Project | None, db: Session, active_run) -> bool:
    if not project or not active_run or not active_run.task_id:
        return False

    try:
        timeout = int(settings.GENERATION_PENDING_TIMEOUT_SECONDS or 0)
    except (TypeError, ValueError):
        timeout = 0
    timeout = max(30, timeout)

    task_id = str(active_run.task_id)
    state = _celery_result_state(task_id)
    task_age = _task_age_seconds(project.id, active_run)
    inactive_for = _run_seconds_since(active_run.updated_at or active_run.started_at)
    reason = None

    if state == "PENDING" and task_age > timeout:
        workers_online = _celery_workers_online()
        if workers_online is False:
            reason = "后台生成服务未在线，本次生成已停止。请启动 worker 后重试未完成页面。"
        elif active_run.status == "running":
            present = _celery_task_present_in_worker(task_id)
            if present is False:
                reason = "后台生成服务中断，本次生成已停止。请重试未完成页面。"
        elif task_age > _queued_celery_timeout_seconds():
            present = _celery_task_present_in_worker(task_id)
            if present is not True:
                reason = "后台生成服务长时间没有接收本次任务，已停止本次生成。请重试未完成页面。"
    elif state in {"STARTED", "RETRY"} and inactive_for > timeout:
        present = _celery_task_present_in_worker(task_id)
        if present is False:
            reason = "后台生成服务中断，本次生成已停止。请重试未完成页面。"

    if not reason:
        return False

    logger.warning(
        "Marking lost Celery task stale: project=%s run=%s task=%s state=%s",
        project.id,
        active_run.id,
        task_id,
        state,
    )
    _mark_generation_idle(project, db, reason)
    stale_active_run(db, project.id, reason)
    generation_progress.pop(project.id, None)
    redis_client.delete(f"project:{project.id}:task_id")
    redis_client.delete(f"project:{project.id}:task_started_at")
    db.commit()
    _auto_restart_generation_if_needed(project, db, active_run)
    return True


def _active_run_for_project_action(project: Project | None, db: Session):
    """Return the current active run after clearing executor state that is known lost."""
    if not project:
        return None
    active_run = get_active_run(db, project.id)
    lost_celery_task = _stale_missing_celery_task_if_needed(project, db, active_run)
    if lost_celery_task:
        return get_active_run(db, project.id)
    stale_run = stale_inactive_run_if_needed(db, project.id)
    if stale_run and stale_run.status == "stale":
        db.commit()
    return get_active_run(db, project.id)


def _auto_restart_generation_if_needed(project: Project, db: Session, previous_run) -> bool:
    """After a Celery task is lost, automatically restart generation so the user does not need to manually retry.

    Conservative policy: only auto-restart if the previous run had already made some
    progress (completed_count > 0). This avoids wasting tokens when the root cause is
    a permanent failure (e.g. worker misconfiguration) rather than a transient crash.
    """
    if previous_run.kind not in {"batch_generation", "page_generation", "prototype_generation", "retry_failed"}:
        return False

    # If the run never completed a single page, the failure is likely permanent
    # (bad config, bad deployment, etc.). Don't auto-restart — let the user investigate.
    if (previous_run.completed_count or 0) == 0:
        return False

    restart_key = f"project:{project.id}:auto_restart_run_id"
    try:
        already_restarted = redis_client.get(restart_key)
        if already_restarted and already_restarted.decode() == previous_run.id:
            logger.info("Auto-restart already attempted for run %s, skipping", previous_run.id)
            return False
    except Exception:
        pass

    slides = (
        db.query(Slide)
        .filter(
            Slide.project_id == project.id,
            Slide.status.in_(["prompt_ready", "failed"]),
            Slide.prompt_text.isnot(None),
        )
        .order_by(Slide.page_num)
        .all()
    )
    slides = [s for s in slides if str(s.prompt_text or "").strip()]
    if not slides:
        return False

    page_nums = [s.page_num for s in slides]
    prototype = previous_run.kind == "prototype_generation"
    run_kind = previous_run.kind if previous_run.kind == "prototype_generation" else "batch_generation"
    run_stage = image_generation_run_stage(kind=run_kind, prototype=prototype, page_nums=page_nums)

    run = create_project_run(
        db,
        project.id,
        kind=run_kind,
        stage=run_stage,
        target_page_nums=page_nums,
        total_count=len(page_nums),
        message="后台任务中断，已自动恢复生成",
    )
    run.status = "queued"
    db.commit()

    try:
        redis_client.set(restart_key, previous_run.id, ex=3600)
    except Exception:
        pass

    try:
        _enqueue_generation_task(db, project.id, page_nums, run=run, prototype=prototype)
        logger.info(
            "Auto-restarted generation for project=%s previous_run=%s new_run=%s pages=%s",
            project.id, previous_run.id, run.id, page_nums,
        )
        return True
    except Exception as exc:
        logger.exception("Auto-restart failed for project=%s", project.id)
        return False


def _format_reference_analysis(analysis: dict) -> str:
    parts = []
    if analysis.get("description"):
        parts.append(str(analysis["description"]))
    if analysis.get("composition_style"):
        parts.append(f"composition={analysis['composition_style']}")
    if analysis.get("mood"):
        parts.append(f"mood={analysis['mood']}")
    colors = analysis.get("colors")
    if isinstance(colors, dict):
        color_text = ", ".join(f"{k}:{v}" for k, v in colors.items() if v)
        if color_text:
            parts.append(f"colors={color_text}")
    if analysis.get("font_suggestion"):
        parts.append(f"font={analysis['font_suggestion']}")
    return "; ".join(parts)


def _format_pptx_reference_analysis(ref: ReferenceImage) -> str | None:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    if not analysis.get("source_document"):
        return None
    parts = []
    if ref.asset_name:
        parts.append(f"asset={ref.asset_name}")
    source_doc = analysis.get("source_document")
    source_page = _asset_source_page_num(ref)
    if source_doc and source_page:
        source_type = str(analysis.get("source_type") or "").upper()
        source_label = "PDF" if source_type == "PDF" else "PPT"
        parts.append(f"source=来自上传{source_label}「{source_doc}」第{source_page}页")
    if analysis.get("classification"):
        parts.append(f"classification={analysis['classification']}")
    if analysis.get("asset_group_role") == "parallel_page_reference_set":
        group_index = analysis.get("asset_group_index")
        group_size = analysis.get("asset_group_size")
        if group_index and group_size:
            parts.append(f"group=同页并列图片组 {group_index}/{group_size}")
    if analysis.get("area_ratio") is not None:
        parts.append(f"area_ratio={analysis['area_ratio']}")
    tags = analysis.get("asset_tags") or analysis.get("suggested_keywords")
    if isinstance(tags, list) and tags:
        parts.append("tags=" + "、".join(str(x) for x in tags[:12]))
    source_text = analysis.get("source_slide_text") or analysis.get("nearby_text")
    if isinstance(source_text, str) and source_text.strip():
        parts.append(f"source_slide_text={source_text[:240]}")
    if ref.usage_note:
        parts.append(f"usage={ref.usage_note}")
    return "; ".join(parts) if parts else None


def _pending_pptx_asset_extractions(project_id: str) -> int:
    docs_dir = os.path.join(settings.UPLOAD_DIR, project_id, "docs")
    if not os.path.exists(docs_dir):
        return 0
    pending = 0
    for filename in os.listdir(docs_dir):
        if not filename.endswith(".assets.json"):
            continue
        try:
            with open(os.path.join(docs_dir, filename), "r", encoding="utf-8") as f:
                data = json.load(f)
            if isinstance(data, dict) and data.get("status") in {"queued", "running"}:
                pending += 1
        except (OSError, json.JSONDecodeError):
            continue
    return pending


def _soft_wait_for_pptx_asset_extractions(project_id: str, *, timeout_seconds: float = 4.0) -> int:
    """Use a small slice of background time so freshly created slides can link refs."""
    deadline = time.time() + timeout_seconds
    pending = _pending_pptx_asset_extractions(project_id)
    while pending and time.time() < deadline:
        time.sleep(0.4)
        pending = _pending_pptx_asset_extractions(project_id)
    return pending


def _build_slide_reference_contexts(
    slides: list[Slide],
) -> tuple[dict[int, list[str]], dict[int, str]]:
    """Analyze page-level reference images for visual/prompt regeneration.

    Returns:
        contexts: verbose lines for LLM (existing behavior)
        user_hints: page_num -> short Chinese copy for humans (must appear in 画面描述)
    """
    contexts: dict[int, list[str]] = {}
    user_hints: dict[int, str] = {}
    for slide in slides:
        page_contexts = []
        hint_parts: list[str] = []
        ref_count = len(slide.reference_images or [])
        logger.info(f"RefContext: slide {slide.page_num} has {ref_count} reference_images")
        _sync_exact_page_reference_overlay_layers(slide)
        overlay_asset_ids = exact_overlay_asset_ids(
            slide.visual_json if isinstance(slide.visual_json, dict) else {}
        )
        for idx, ref in enumerate(sorted(slide.reference_images or [], key=_reference_image_sort_key), start=1):
            if str(ref.id) in overlay_asset_ids:
                logger.info(f"RefContext: slide {slide.page_num} ref {idx} skipped (overlay)")
                continue
            resolved_file = _resolve_file_path(ref.file_path)
            file_exists = os.path.exists(resolved_file)
            logger.info(f"RefContext: slide {slide.page_num} ref {idx} file={ref.file_path} exists={file_exists} role={ref.role}")
            if not file_exists:
                continue
            summary = _format_pptx_reference_analysis(ref)
            if not summary:
                cached = _cached_reference_analysis(ref)
                status = (ref.asset_analysis or {}).get("analysis_status") if isinstance(ref.asset_analysis, dict) else None
                if not cached and status in {"queued", "running"}:
                    cached = _wait_for_cached_reference_analysis(ref, timeout_seconds=4.0)
                if cached:
                    summary = _format_reference_analysis(cached)
                elif status in {"queued", "running", "failed"}:
                    summary = _fallback_reference_summary(ref)
                else:
                    analysis = analyze_reference_image(ref.file_path)
                    summary = _format_reference_analysis(analysis)
            page_contexts.append(
                f"Reference Image {idx}: role={ref.role}; process_mode={ref.process_mode or 'blend'}; "
                f"intent={reference_process_mode_instruction(ref.process_mode)};"
                f"file={os.path.basename(ref.file_path)}; "
                f"actual_input=uploaded_to_image_model_as_reference_{idx}; analysis={summary}"
            )
            hint_parts.append(f"参考图{idx}（AI 参考）：{summary}")

        if page_contexts:
            contexts[slide.page_num] = page_contexts
        if hint_parts:
            user_hints[slide.page_num] = "\n".join(hint_parts)

    logger.info(f"RefContext: built contexts for pages={list(contexts.keys())}, hints for pages={list(user_hints.keys())}")
    return contexts, user_hints


def _link_pending_pptx_page_refs(project_id: str, db: Session) -> int:
    """Attach source-extracted page refs uploaded before slides existed."""
    slides = db.query(Slide).filter(Slide.project_id == project_id).all()
    slide_by_page = {slide.page_num: slide for slide in slides}
    slide_by_source: dict[tuple[str, int], Slide] = {}
    slide_by_figure: dict[tuple[str, int, str], list[Slide]] = {}
    for slide in slides:
        content = slide.content_json if isinstance(slide.content_json, dict) else {}
        source_refs = content.get("source_refs") if isinstance(content.get("source_refs"), list) else []
        figure_refs = content.get("figure_refs") if isinstance(content.get("figure_refs"), list) else []
        for item in [*source_refs, *figure_refs]:
            if not isinstance(item, dict):
                continue
            source_doc = str(item.get("source_document") or "").strip()
            try:
                source_page = int(item.get("source_page_num") or 0)
            except (TypeError, ValueError):
                source_page = 0
            if source_doc and source_page > 0:
                source_kind = _source_ref_type(source_doc, str(item.get("source_type") or ""))
                figure_id = str(item.get("figure_id") or item.get("id") or "").strip()
                if figure_id:
                    slide_by_figure.setdefault((source_doc, source_page, figure_id), []).append(slide)
                if source_kind != "pdf":
                    slide_by_source.setdefault((source_doc, source_page), slide)
    if not slide_by_page:
        return 0

    pending_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project_id,
        ReferenceImage.slide_id.is_(None),
        ReferenceImage.role == "content_ref",
    ).all()
    existing_linked_keys: set[tuple[str, int, str, str]] = set()
    linked_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project_id,
        ReferenceImage.slide_id.isnot(None),
        ReferenceImage.role == "content_ref",
    ).all()
    for linked_ref in linked_refs:
        analysis = linked_ref.asset_analysis if isinstance(linked_ref.asset_analysis, dict) else {}
        source_doc = str(analysis.get("source_document") or "").strip()
        page_num = _asset_source_page_num(linked_ref)
        figure_id = str(analysis.get("id") or analysis.get("figure_id") or "").strip()
        if source_doc and isinstance(page_num, int) and page_num > 0 and figure_id and linked_ref.slide_id:
            existing_linked_keys.add((source_doc, page_num, figure_id, linked_ref.slide_id))
    linked = 0
    for ref in pending_refs:
        analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
        page_num = _asset_source_page_num(ref)
        if not isinstance(page_num, int):
            continue
        source_doc = str(analysis.get("source_document") or "").strip()
        asset_source_type = _source_ref_type(source_doc, str(analysis.get("source_type") or ""))
        figure_id = str(analysis.get("id") or analysis.get("figure_id") or "").strip()
        figure_slides = (
            slide_by_figure.get((source_doc, page_num, figure_id), [])
            if source_doc and figure_id
            else []
        )
        if figure_slides:
            first_linked = False
            for slide in figure_slides:
                key = (source_doc, page_num, figure_id, slide.id)
                if key in existing_linked_keys:
                    continue
                if not first_linked and ref.slide_id is None:
                    ref.slide_id = slide.id
                    first_linked = True
                else:
                    db.add(ReferenceImage(
                        project_id=ref.project_id,
                        slide_id=slide.id,
                        file_path=ref.file_path,
                        role=ref.role,
                        process_mode=ref.process_mode,
                        asset_name=ref.asset_name,
                        asset_kind=ref.asset_kind,
                        usage_note=ref.usage_note,
                        asset_analysis=copy.deepcopy(analysis),
                        logo_anchor=ref.logo_anchor,
                    ))
                existing_linked_keys.add(key)
                linked += 1
            continue
        slide = None
        if not slide and asset_source_type != "pdf":
            slide = slide_by_source.get((source_doc, page_num)) if source_doc else None
        if not slide and asset_source_type in {"", "pptx", "pptx_slide"}:
            slide = slide_by_page.get(page_num)
        if not slide:
            continue
        ref.slide_id = slide.id
        linked += 1
    if linked:
        logger.info(f"Linked {linked} PPTX-extracted page refs to slides for project={project_id}")
    return linked


class PageNumsRequest(BaseModel):
    page_nums: Optional[List[int]] = None
    prototype: bool = False
    stage_context: Optional[str] = None


class ContentPlanRequest(BaseModel):
    topic: Optional[str] = None
    page_count: Optional[int] = None
    attachment_ids: Optional[List[str]] = None
    chat_context: Optional[str] = None


class CreateSlideRequest(BaseModel):
    page_num: int
    content_json: dict


class UpdateContentRequest(BaseModel):
    page_num: int
    slide_id: Optional[str] = None
    content_json: dict


class UpdateVisualRequest(BaseModel):
    page_num: int
    slide_id: Optional[str] = None
    visual_json: dict


class UpdateTypeRequest(BaseModel):
    page_num: int
    slide_id: Optional[str] = None
    type: str


class AssetPinsRequest(BaseModel):
    asset_ids: List[str]
    usage: dict[str, str] = {}


class OverlayLayerRequest(BaseModel):
    id: Optional[str] = None
    asset_id: str
    enabled: bool = True
    preset: str = "right-card"
    fit: str = "contain"
    mode: str = "exact_card"
    usage_note: Optional[str] = None
    z_index: Optional[int] = None


class OverlayLayersRequest(BaseModel):
    layers: List[OverlayLayerRequest] = []


class ReorderRequest(BaseModel):
    page_nums: List[int]


def _content_plan_attachment_context(db: Session, project_id: str, attachment_ids: list[str] | None) -> str:
    ids = [str(item) for item in (attachment_ids or []) if str(item).strip()]
    if not ids:
        return ""
    refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project_id,
        ReferenceImage.id.in_(ids),
        ReferenceImage.role.in_(["chat_ref", "content_ref", "chart_ref", "visual_asset"]),
    ).all()
    by_id = {str(ref.id): ref for ref in refs if ref.file_path and os.path.exists(ref.file_path)}
    parts: list[str] = []
    role_label = {
        "chat_ref": "本轮对话图片",
        "content_ref": "内容素材图",
        "chart_ref": "图表参考图",
        "visual_asset": "项目素材图",
    }
    for idx, ref_id in enumerate(ids, start=1):
        ref = by_id.get(ref_id)
        if not ref:
            continue
        label = ref.asset_name or os.path.basename(ref.file_path)
        description = describe_context_image(ref.file_path, label, role_label.get(ref.role, "参考图"), "内容规划生成")
        if not description:
            analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
            description = analysis.get("description") or analysis.get("subject") or "图片已上传，但当前没有可用解读。"
        parts.append(f"--- 本轮附件 {idx}: {label} ---\n{description}")
    return "\n\n".join(parts).strip()


def _content_plan_run_is_active_for_writeback(db: Session, run_id: str | None, logger, project_id: str, phase: str) -> bool:
    db.expire_all()
    if not run_id:
        return True
    if is_run_active(db, run_id):
        latest_active = get_active_run(db, project_id)
        if latest_active and str(latest_active.id) == str(run_id):
            return True
        logger.info(
            "[ContentPlan BG] Run %s was superseded during %s; skipping writeback for project=%s",
            run_id,
            phase,
            project_id,
        )
        return False
    logger.info(
        "[ContentPlan BG] Run %s is no longer active during %s; skipping writeback for project=%s",
        run_id,
        phase,
        project_id,
    )
    return False


def _generate_content_plan_bg(
    project_id: str,
    topic: str,
    page_count: int | None = None,
    run_id: str | None = None,
    attachment_ids: list[str] | None = None,
    chat_context: str | None = None,
):
    """后台任务：异步生成 Content Plan。"""
    import logging
    logger = logging.getLogger(__name__)
    from app.models.base import SessionLocal
    db = SessionLocal()
    try:
        logger.info(f"[ContentPlan BG] Starting for project={project_id}, topic={topic[:30]}...")
        mark_run_running(db, run_id, stage="document_parse", message="正在读取上传材料...")
        db.commit()
        _update_progress(project_id, {"stage": "document_parse", "message": "正在读取上传材料...", "current_page": 0, "total_pages": page_count or 10}, run_id)

        def report_document_progress(payload: dict):
            current_page = payload.get("current_page")
            total_pages = payload.get("total_pages")
            message = str(payload.get("message") or "正在读取上传材料...")
            _update_progress(project_id, {
                "stage": "document_parse",
                "message": message,
                "current_page": current_page or 0,
                "total_pages": total_pages or page_count or 10,
            }, run_id)

        try:
            document_wait_seconds = max(0.0, float(os.getenv("PPTGOD_CONTENT_PLAN_DOCUMENT_WAIT_SECONDS", "900")))
        except ValueError:
            document_wait_seconds = 900.0

        source_packs = load_project_source_packs(
            project_id,
            parse_missing=True,
            running_wait_seconds=document_wait_seconds,
            progress_callback=report_document_progress,
        )
        try:
            source_context = build_source_context(
                topic,
                source_packs,
                token_budget=DEFAULT_SOURCE_CONTEXT_TOKEN_BUDGET,
            )
            documents = source_context.text
        except SourceScopeRequired as exc:
            payload = exc.payload
            message = payload.get("reason") or "上传材料过长，需要先选择范围后再生成。"
            update_run_progress(
                db,
                run_id,
                stage="needs_scope",
                message=message,
                total_count=0,
                completed_count=0,
            )
            finish_run(
                db,
                run_id,
                status="failed",
                message=message,
                error_msg=json.dumps(payload, ensure_ascii=False)[:500],
            )
            db.commit()
            _update_progress(project_id, {
                "stage": "needs_scope",
                "status": "needs_scope",
                "message": message,
                "reason": message,
                "suggested_scopes": payload.get("suggested_scopes") or [],
                "source_stats": payload.get("source_stats") or {},
            }, run_id)
            return
        if not documents:
            documents = load_project_documents(
                project_id,
                parse_missing=True,
                text_limit=PAGE_MAP_DOCUMENT_LIMIT,
                preserve_ppt_sources=True,
                running_wait_seconds=document_wait_seconds,
                progress_callback=report_document_progress,
            )
        explicit_attachment_context = _content_plan_attachment_context(db, project_id, attachment_ids)
        if explicit_attachment_context:
            documents = "\n\n".join(
                part for part in [
                    documents,
                    "【本轮 Agent 指令附件】\n这些素材来自触发本次生成的同一条用户消息，优先级高于历史项目素材。\n"
                    + explicit_attachment_context,
                ] if part
            )
        explicit_attachment_ids = {str(item) for item in (attachment_ids or []) if str(item).strip()}
        image_context_parts: list[str] = []
        content_refs = db.query(ReferenceImage).filter(
            ReferenceImage.project_id == project_id,
            ReferenceImage.role == "content_ref",
            ReferenceImage.slide_id.is_(None),
        ).all()
        for idx, ref in enumerate(content_refs[:8], start=1):
            if str(ref.id) in explicit_attachment_ids:
                continue
            analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
            if _asset_source_page_num(ref):
                continue
            if not ref.file_path or not os.path.exists(ref.file_path):
                continue
            label = ref.asset_name or os.path.basename(ref.file_path)
            description = describe_context_image(ref.file_path, label, "内容素材图", "内容规划生成")
            if description:
                image_context_parts.append(f"--- 图片素材 {idx}: {label} ---\n{description}")
        if image_context_parts:
            image_context = "\n\n".join(image_context_parts)
            documents = "\n\n".join(part for part in [documents, image_context] if part)
        intent_topic = content_plan_topic_with_chat_context(topic, chat_context)
        intent_contract = None
        project_for_intent = db.query(Project).filter(Project.id == project_id).first()
        if project_for_intent:
            legacy_contract = infer_intent_contract(
                intent_topic,
                source_diagnostics=source_diagnostics_from_documents(documents),
            )
            intent_contract = infer_effective_content_intent_contract(
                intent_topic,
                documents,
                legacy_contract,
            )
            project_for_intent.intent_contract = intent_contract
            flag_modified(project_for_intent, "intent_contract")
            db.commit()
        if documents:
            logger.info(f"[ContentPlan BG] Loaded documents, length={len(documents)}")
            inferred_page_count = infer_page_count_from_single_ppt(documents, intent_topic, intent_contract=intent_contract)
            if page_count is None and inferred_page_count:
                page_count = inferred_page_count
                update_run_progress(db, run_id, total_count=page_count)
                db.commit()
                logger.info(f"[ContentPlan BG] Inferred page_count={page_count} from single uploaded PPT")

        target_page_count, _, _ = resolve_content_plan_page_target(
            intent_topic,
            page_count,
            documents,
            intent_contract=intent_contract,
        )
        update_run_progress(db, run_id, total_count=target_page_count)
        db.commit()

        mark_run_running(db, run_id, stage="content_plan", message="开始生成 Content Plan...")
        db.commit()
        _update_progress(project_id, {"stage": "content_plan", "message": "开始生成 Content Plan...", "current_page": 0, "total_pages": target_page_count}, run_id)

        def report_progress(data: dict):
            _update_progress(project_id, data, run_id)

        outline = generate_content_plan(
            topic=topic,
            page_count=page_count,
            documents=documents,
            on_progress=report_progress,
            intent_contract=intent_contract,
            chat_context=chat_context,
        )
        if not _content_plan_run_is_active_for_writeback(db, run_id, logger, project_id, "model_return"):
            return
        logger.info(f"[ContentPlan BG] Generated {len(outline)} pages")
        update_run_progress(
            db,
            run_id,
            stage="saving",
            message="正在保存结果...",
            total_count=len(outline),
            completed_count=len(outline),
        )
        db.commit()
        _update_progress(project_id, {"stage": "saving", "message": "正在保存结果...", "current_page": len(outline), "total_pages": len(outline)}, run_id)
        if not _content_plan_run_is_active_for_writeback(db, run_id, logger, project_id, "before_slide_replace"):
            return
        # PPT 上传早于/晚于内容规划都要能重新挂回新生成的页面。
        # 生成新 slides 前，先把从 PPT 拆出的页面参考图解绑，避免旧 slide 被删除后素材丢失在孤儿关系里。
        pptx_refs = db.query(ReferenceImage).filter(
            ReferenceImage.project_id == project_id,
            ReferenceImage.role == "content_ref",
        ).all()
        for ref in pptx_refs:
            analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
            if _asset_source_page_num(ref):
                ref.slide_id = None
        db.flush()
        # 删除旧的 slides（如果存在）
        db.query(Slide).filter(Slide.project_id == project_id).delete()
        # 保存新的 slides
        for item in outline:
            item = copy.deepcopy(item)
            item.setdefault("page_num", item.get("page_num", 0))
            slide = Slide(
                project_id=project_id,
                page_num=item["page_num"],
                type=item.get("type", "content"),
                content_json=item,
            )
            db.add(slide)
        db.flush()
        pending_pptx_assets = _pending_pptx_asset_extractions(project_id)
        if pending_pptx_assets:
            _update_progress(project_id, {
                "stage": "asset_linking",
                "message": "正在挂接已解析的 PPT 图片素材...",
                "current_page": len(outline),
                "total_pages": len(outline),
            }, run_id)
            pending_pptx_assets = _soft_wait_for_pptx_asset_extractions(project_id, timeout_seconds=4.0)
            if pending_pptx_assets:
                logger.info(
                    "[ContentPlan BG] %s PPTX asset extraction task(s) still running; continuing without hard gate",
                    pending_pptx_assets,
                )
        linked_refs = _link_pending_pptx_page_refs(project_id, db)
        project = db.query(Project).filter(Project.id == project_id).first()
        if project:
            project.status = "planning"
            project.content_plan_confirmed = False
            project.style_proposal = None
            project.selected_style = None
            if _should_autoname_project(project):
                derived_title = _derive_project_title(topic, outline)
                if derived_title:
                    project.title = derived_title
        if not _content_plan_run_is_active_for_writeback(db, run_id, logger, project_id, "before_final_commit"):
            db.rollback()
            return
        finish_run(
            db,
            run_id,
            status="succeeded",
            message=f"内容规划已生成，共 {len(outline)} 页" + (f"，已绑定 {linked_refs} 张原 PPT 页面素材" if linked_refs else ""),
            completed_count=len(outline),
        )
        # 内容规划生成完成，设置未读通知
        if project:
            project.has_unread_notification = True
            project.unread_notification_message = "内容规划已生成"
        db.commit()
        logger.info(f"[ContentPlan BG] Completed for project={project_id}")
    except Exception as e:
        db.rollback()
        logger.error(f"[ContentPlan BG] Failed for project={project_id}: {e}")
        if not _content_plan_run_is_active_for_writeback(db, run_id, logger, project_id, "failure"):
            return
        _update_progress(project_id, {"stage": "error", "message": f"生成失败：{str(e)[:100]}"}, run_id)
        # 标记项目状态为失败，让前端可以检测
        try:
            project = db.query(Project).filter(Project.id == project_id).first()
            if project:
                project.status = "draft"
            finish_run(db, run_id, status="failed", message="内容规划生成失败", error_msg=str(e)[:500])
            db.commit()
        except Exception:
            pass
    finally:
        # 生成结束后清理内存进度，避免过时进度残留误导前端
        generation_progress.pop(project_id, None)
        db.close()


@router.post("/{project_id}/content-plan")
def create_content_plan(
    project_id: str,
    background_tasks: BackgroundTasks,
    body: ContentPlanRequest = ContentPlanRequest(),
    db: Session = Depends(get_db),
):
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    if project.status not in {"draft", "planning", "content_plan_ready"} or project.content_plan_confirmed:
        raise HTTPException(status_code=409, detail="当前阶段不能重新生成内容规划，请先回退到内容规划。")
    if _active_run_for_project_action(project, db):
        raise HTTPException(status_code=409, detail="当前项目已有任务正在运行，请等待完成后再开始下一步")

    # 优先使用用户传入的 topic，否则用项目标题
    topic = body.topic.strip() if body.topic else project.title
    chat_context = (body.chat_context or "").strip() or None
    intent_topic = content_plan_topic_with_chat_context(topic, chat_context)
    page_count = resolve_requested_content_plan_page_count(intent_topic, body.page_count)
    target_page_count, _, _ = resolve_content_plan_page_target(intent_topic, page_count)

    try:
        run = create_project_run(
            db,
            project_id,
            kind="content_plan",
            stage="content_plan",
            total_count=target_page_count,
            message="内容规划生成已排队",
        )
        db.commit()
    except RuntimeError as exc:
        raise HTTPException(status_code=409, detail=str(exc))

    background_tasks.add_task(
        _generate_content_plan_bg,
        project_id,
        topic,
        page_count,
        run.id,
        body.attachment_ids or [],
        chat_context,
    )
    return {"message": "Content plan generation started", "status": project.status, "run": serialize_run(run)}


@router.get("/{project_id}/slides", response_model=List[SlideResponse])
def list_slides(project_id: str, db: Session = Depends(get_db)):
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    if _repair_project_logo_policies(project, slides):
        db.commit()
    # 手动加载所有参考图，避免 joinedload 在 SQLAlchemy 2.0 下的兼容问题
    slide_ids = [s.id for s in slides]
    refs = db.query(ReferenceImage).filter(
        ReferenceImage.slide_id.in_(slide_ids),
        ReferenceImage.role != "finetune_ref",
    ).all() if slide_ids else []
    refs_by_slide = {}
    for ref in refs:
        refs_by_slide.setdefault(ref.slide_id, []).append(ref)

    return [
        {
            "id": s.id,
            "project_id": s.project_id,
            "page_num": s.page_num,
            "type": s.type,
            "status": s.status,
            "content_json": s.content_json,
            "visual_json": _with_resolved_logo_overlay_box(s.visual_json, s, project),
            "prompt_text": s.prompt_text,
            "image_path": s.image_path,
            "error_msg": s.error_msg,
            "reference_images": [
                {
                    "id": ref.id,
                    "role": ref.role,
                    "process_mode": ref.process_mode or "blend",
                    "asset_name": ref.asset_name,
                    "asset_kind": ref.asset_kind,
                    "usage_note": ref.usage_note,
                    "asset_analysis": ref.asset_analysis,
                    "logo_anchor": ref.logo_anchor or (DEFAULT_LOGO_ANCHOR if ref.role == "logo" else None),
                    "url": _reference_upload_url(project_id, ref.file_path),
                    "overlay_url": _logo_overlay_url(ref, project_id) if ref.role == "logo" and is_logo_confirmed(ref) else None,
                    "symbol_overlay_url": _logo_symbol_overlay_url(ref, project_id) if ref.role == "logo" and is_logo_confirmed(ref) else None,
                }
                for ref in sorted(refs_by_slide.get(s.id, []), key=_reference_image_sort_key)
            ],
        }
        for s in slides
    ]


@router.get("/{project_id}/slides/export-markdown")
def export_slides_markdown(project_id: str, db: Session = Depends(get_db)):
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    if not slides:
        raise HTTPException(status_code=400, detail="当前项目还没有可导出的内容页")

    markdown = _build_content_plan_markdown(project, slides)
    filename = _safe_export_filename(f"{project.title}-内容规划")
    quoted_filename = quote(filename)
    return StreamingResponse(
        io.BytesIO(markdown.encode("utf-8")),
        media_type="text/markdown; charset=utf-8",
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8''{quoted_filename}",
        },
    )


@router.post("/{project_id}/visual-plan")
def create_visual_plan(
    project_id: str,
    body: PageNumsRequest = PageNumsRequest(),
    db: Session = Depends(get_db),
):
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    if not slides:
        raise HTTPException(status_code=400, detail="No content plan found. Generate content plan first.")

    # 手动加载参考图，避免 joinedload 在 SQLAlchemy 2.0 下的兼容问题
    slide_ids = [s.id for s in slides]
    refs = db.query(ReferenceImage).filter(
        ReferenceImage.slide_id.in_(slide_ids),
        ReferenceImage.role != "finetune_ref",
    ).all() if slide_ids else []
    refs_by_slide = {}
    for ref in refs:
        refs_by_slide.setdefault(ref.slide_id, []).append(ref)
    for s in slides:
        s.reference_images = sorted(refs_by_slide.get(s.id, []), key=_reference_image_sort_key)

    ref_contexts, ref_user_hints = _build_slide_reference_contexts(slides)
    stage_context = (body.stage_context or "").strip()[:4000]
    content_plan = []
    for s in slides:
        item = _content_json_for_generation(s)
        item["page_num"] = s.page_num
        item["manual_visual_asset_ids"] = _manual_asset_ids(s.visual_json)
        item["manual_visual_asset_usage"] = _manual_asset_usage(s.visual_json)
        item["asset_route_modes"] = _manual_asset_route_modes(s.visual_json, item["manual_visual_asset_ids"])
        item["excluded_visual_asset_ids"] = _excluded_asset_ids(s.visual_json)
        visual = s.visual_json if isinstance(s.visual_json, dict) else {}
        if visual.get("overlay_layers"):
            item["overlay_layers"] = visual["overlay_layers"]
        if stage_context:
            item["global_user_requirements"] = stage_context
        if ref_contexts.get(s.page_num):
            item["reference_context"] = "\n".join(ref_contexts[s.page_num])
        if ref_user_hints.get(s.page_num):
            item["reference_user_hint"] = ref_user_hints[s.page_num]
        content_plan.append(item)

    # 打样：只处理选中的页
    if body.page_nums:
        content_plan = [p for p in content_plan if p["page_num"] in body.page_nums]

    # 获取参考图 ID（如果有）
    ref_images = project.reference_images
    ref_ids = [img.id for img in ref_images] if ref_images else None
    global_visual_assets = _project_visual_assets_for_planning(project)
    style_text_override = _derive_project_style_pack(project, content_plan)
    style_override = _style_override_from_text(style_text_override)

    try:
        visual_plan = generate_visual_plan(
            content_plan=content_plan,
            style_id=project.style_id or "default",
            reference_image_ids=ref_ids,
            style_override=style_override,
            global_visual_assets=global_visual_assets,
            has_project_logo=bool(_project_logo_refs(project)),
        )
    except VisualPlanGenerationError as exc:
        raise HTTPException(status_code=502, detail=f"视觉方案生成失败：{str(exc)[:300]}") from exc
    artifact_deps = dependency_signature(project, slides)

    # 保存 visual_plan 到每页 slide（只更新选中的页，或全部）
    visual_by_page = {v["page_num"]: v for v in visual_plan}
    target_slides = slides
    if body.page_nums:
        target_slides = [s for s in slides if s.page_num in body.page_nums]

    for slide in target_slides:
        existing_visual = slide.visual_json if isinstance(slide.visual_json, dict) else {}
        next_visual = _merge_manual_pins_into_visual_json(
            with_artifact_meta(
                visual_by_page.get(slide.page_num, {}),
                kind="visual_plan",
                dependencies=artifact_deps,
            ),
            existing_visual,
        )
        slide.visual_json = next_visual
        slide.prompt_text = None
        slide.image_path = None
        slide.error_msg = None
        slide.status = "visual_ready"

    # 如果全部都有 visual，项目状态推进
    if all(s.visual_json for s in slides):
        project.status = "visual_ready"
    db.commit()

    return {"message": "Visual plan generated", "slides_count": len(visual_plan), "prototype": bool(body.page_nums)}


@router.post("/{project_id}/prompts")
def create_prompts(
    project_id: str,
    body: PageNumsRequest = PageNumsRequest(),
    db: Session = Depends(get_db),
):
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    if not slides:
        raise HTTPException(status_code=400, detail="No slides found")

    # 手动加载参考图，避免 joinedload 在 SQLAlchemy 2.0 下的兼容问题
    slide_ids = [s.id for s in slides]
    refs = db.query(ReferenceImage).filter(
        ReferenceImage.slide_id.in_(slide_ids),
        ReferenceImage.role != "finetune_ref",
    ).all() if slide_ids else []
    refs_by_slide = {}
    for ref in refs:
        refs_by_slide.setdefault(ref.slide_id, []).append(ref)
    for s in slides:
        s.reference_images = sorted(refs_by_slide.get(s.id, []), key=_reference_image_sort_key)

    # 检查是否有 visual_plan
    if not any(s.visual_json for s in slides):
        raise HTTPException(status_code=400, detail="No visual plan found. Generate visual plan first.")

    # 过滤选中的页
    target_slides = slides
    if body.page_nums:
        target_slides = [s for s in slides if s.page_num in body.page_nums]

    ref_contexts, ref_user_hints = _build_slide_reference_contexts(target_slides)
    stage_context = (body.stage_context or "").strip()[:4000]
    content_plan = []
    for s in target_slides:
        item = _content_json_for_generation(s)
        item["page_num"] = s.page_num
        item["manual_visual_asset_ids"] = _manual_asset_ids(s.visual_json)
        item["manual_visual_asset_usage"] = _manual_asset_usage(s.visual_json)
        item["asset_route_modes"] = _manual_asset_route_modes(s.visual_json, item["manual_visual_asset_ids"])
        item["excluded_visual_asset_ids"] = _excluded_asset_ids(s.visual_json)
        visual = s.visual_json if isinstance(s.visual_json, dict) else {}
        if visual.get("overlay_layers"):
            item["overlay_layers"] = visual["overlay_layers"]
        if stage_context:
            item["global_user_requirements"] = stage_context
        if ref_contexts.get(s.page_num):
            item["reference_context"] = "\n".join(ref_contexts[s.page_num])
        if ref_user_hints.get(s.page_num):
            item["reference_user_hint"] = ref_user_hints[s.page_num]
        content_plan.append(item)
    visual_plan = []
    for s in target_slides:
        if not s.visual_json:
            continue
        intent = _with_prompt_asset_policies(s.visual_json, project, s) or {}
        if isinstance(intent, dict) and not intent.get("page_num"):
            intent = {**intent, "page_num": s.page_num}
        visual_plan.append(intent)
    visual_intent_by_page = {
        int(intent.get("page_num") or 0): intent
        for intent in visual_plan
        if isinstance(intent, dict)
    }

    # style_ref 已被转成 style_text；Logo 按页面级 policy 融合，visual_asset 按页选择。
    ref_images_by_page = {
        s.page_num: (
            _project_refs_for_prompt(
                project,
                (s.visual_json or {}).get("visual_asset_ids") if isinstance(s.visual_json, dict) else [],
                _with_prompt_asset_policies(s.visual_json, project, s) if isinstance(s.visual_json, dict) else None,
            )
            + _slide_refs_for_prompt(s, s.visual_json if isinstance(s.visual_json, dict) else None)
        )
        for s in target_slides
    }
    style_text_override = _style_text_from_selected_style(project.selected_style) or _derive_project_style_pack(project, content_plan)
    artifact_deps = dependency_signature(project, slides)

    prompts = generate_prompts_for_all_pages(
        visual_plan=visual_plan,
        content_plan=content_plan,
        style_id="default",
        reference_images_by_page=ref_images_by_page,
        style_text_override=style_text_override,
    )

    # 保存 prompt 到每页 slide
    prompt_by_page = {p["page_num"]: p["prompt"] for p in prompts}
    for slide in target_slides:
        slide.prompt_text = prompt_by_page.get(slide.page_num)
        if slide.prompt_text:
            visual_intent = visual_intent_by_page.get(slide.page_num, slide.visual_json)
            slide.visual_json = with_artifact_meta(
                visual_intent,
                kind="visual_plan",
                dependencies=artifact_deps,
                prompt_dependencies=artifact_deps,
            )
            slide.image_path = None
            slide.error_msg = None
            slide.visual_json = with_stale_flags(slide.visual_json, content=False, visual=False, image=False)
            flag_modified(slide, "visual_json")
            slide.status = "prompt_ready"

    # 如果全部都有 prompt，项目状态推进
    if all(s.prompt_text for s in slides):
        project.status = "prompt_ready"
    db.commit()

    return {"message": "Prompts generated", "slides_count": len(prompts), "prototype": bool(body.page_nums)}


@router.post("/{project_id}/visual-prompts")
async def create_visual_and_prompts(
    project_id: str,
    body: PageNumsRequest = PageNumsRequest(),
    db: Session = Depends(get_db),
):
    """一步生成视觉方案和生图 Prompt（SSE 流式返回真实进度）。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    if not slides:
        raise HTTPException(status_code=400, detail="No content plan found. Generate content plan first.")

    target_slides = slides
    if body.page_nums:
        target_slides = [s for s in slides if s.page_num in body.page_nums]

    content_plan = []
    for s in target_slides:
        item = _content_json_for_generation(s)
        item["page_num"] = s.page_num
        item["manual_visual_asset_ids"] = _manual_asset_ids(s.visual_json)
        item["manual_visual_asset_usage"] = _manual_asset_usage(s.visual_json)
        item["asset_route_modes"] = _manual_asset_route_modes(s.visual_json, item["manual_visual_asset_ids"])
        item["excluded_visual_asset_ids"] = _excluded_asset_ids(s.visual_json)
        visual = s.visual_json if isinstance(s.visual_json, dict) else {}
        if visual.get("overlay_layers"):
            item["overlay_layers"] = visual["overlay_layers"]
        content_plan.append(item)

    ref_images_project = project.reference_images
    ref_ids = [img.id for img in ref_images_project] if ref_images_project else None

    style_text_override = _derive_project_style_pack(project, content_plan)
    style_override = _style_override_from_text(style_text_override)
    target_page_nums = [s.page_num for s in target_slides] if body.page_nums else None
    try:
        run = create_project_run(
            db,
            project_id,
            kind="visual_prompts",
            stage="visual_planning",
            target_page_nums=target_page_nums,
            total_count=len(target_slides),
            message="画面描述和生图 Prompt 生成已排队",
        )
        db.commit()
    except RuntimeError as exc:
        raise HTTPException(status_code=409, detail=str(exc))

    # 启动后台任务，独立于 HTTP 连接运行
    async with _running_tasks_lock:
        existing_task = _running_tasks.get(project_id)
        if existing_task and not existing_task.done():
            existing_task.cancel()
            try:
                await existing_task
            except asyncio.CancelledError:
                pass

        task_credentials = get_raw_provider_credentials()
        task = asyncio.create_task(
            _do_generate_visual_and_prompts(project_id, target_page_nums, run.id, task_credentials, (body.stage_context or "").strip()[:4000])
        )
        _running_tasks[project_id] = task

    return {"status": "started", "message": "视觉方案和生图 Prompt 生成已启动，请稍候。", "run": serialize_run(run)}


async def _do_generate_visual_and_prompts(
    project_id: str,
    page_nums: Optional[List[int]] = None,
    run_id: str | None = None,
    provider_credentials=None,
    stage_context: str | None = None,
):
    """后台任务：生成视觉方案和 Prompt，完成后更新数据库。"""
    if provider_credentials is not None:
        with provider_credentials_context(provider_credentials):
            return await _do_generate_visual_and_prompts(project_id, page_nums, run_id, None, stage_context)

    db = SessionLocal()
    try:
        project = db.query(Project).filter(Project.id == project_id).first()
        if not project:
            logger.warning(f"Project {project_id} not found for background generation")
            return

        mark_run_running(db, run_id, stage="visual_planning", message="正在分析内容结构，为每一页设计视觉方案...")
        db.commit()
        if not is_run_active(db, run_id):
            logger.info(f"VisualPrompts BG: run {run_id} is no longer active before work; skipping")
            return

        slides = (
            db.query(Slide)
            .filter(Slide.project_id == project_id)
            .order_by(Slide.page_num)
            .all()
        )
        if not slides:
            logger.warning(f"No slides found for project {project_id}")
            project.status = "planning"
            db.commit()
            return

        target_slides = slides
        if page_nums:
            target_slides = [s for s in slides if s.page_num in page_nums]

        # 手动加载参考图，避免 joinedload 在 SQLAlchemy 2.0 下的兼容问题
        slide_ids = [s.id for s in target_slides]
        refs = db.query(ReferenceImage).filter(
            ReferenceImage.slide_id.in_(slide_ids),
            ReferenceImage.role != "finetune_ref",
        ).all() if slide_ids else []
        refs_by_slide = {}
        for ref in refs:
            refs_by_slide.setdefault(ref.slide_id, []).append(ref)
        for s in target_slides:
            s.reference_images = sorted(refs_by_slide.get(s.id, []), key=_reference_image_sort_key)

        ref_contexts, ref_user_hints = _build_slide_reference_contexts(target_slides)
        logger.info(f"VisualPrompts BG: project={project_id}, ref_contexts pages={list(ref_contexts.keys())}, ref_user_hints pages={list(ref_user_hints.keys())}")
        content_plan = []
        stage_context = (stage_context or "").strip()[:4000]
        for s in target_slides:
            item = _content_json_for_generation(s)
            item["page_num"] = s.page_num
            item["manual_visual_asset_ids"] = _manual_asset_ids(s.visual_json)
            item["manual_visual_asset_usage"] = _manual_asset_usage(s.visual_json)
            item["asset_route_modes"] = _manual_asset_route_modes(s.visual_json, item["manual_visual_asset_ids"])
            item["excluded_visual_asset_ids"] = _excluded_asset_ids(s.visual_json)
            visual = s.visual_json if isinstance(s.visual_json, dict) else {}
            if visual.get("overlay_layers"):
                item["overlay_layers"] = visual["overlay_layers"]
            if stage_context:
                item["global_user_requirements"] = stage_context
            if ref_contexts.get(s.page_num):
                item["reference_context"] = "\n".join(ref_contexts[s.page_num])
            if ref_user_hints.get(s.page_num):
                item["reference_user_hint"] = ref_user_hints[s.page_num]
            content_plan.append(item)

        ref_images_project = project.reference_images
        ref_ids = [img.id for img in ref_images_project] if ref_images_project else None
        global_visual_assets = _project_visual_assets_for_planning(project)

        style_text_override = _derive_project_style_pack(project, content_plan)
        style_override = _style_override_from_text(style_text_override)
        # 注意：页面级参考图已通过 content_plan 的 reference_context 传递；visual_asset 按页选择。

        # Step 1: 生成 Visual Plan
        _update_progress(project_id, {
            "stage": "visual_planning",
            "message": "正在分析内容结构，为每一页设计视觉方案...",
            "current_page": 0,
            "total_pages": len(target_slides),
        }, run_id)
        visual_plan = await asyncio.to_thread(
            generate_visual_plan,
            content_plan=content_plan,
            style_id=project.style_id or "default",
            reference_image_ids=ref_ids,
            style_override=style_override,
            global_visual_assets=global_visual_assets,
            has_project_logo=bool(_project_logo_refs(project)),
            progress_callback=lambda progress: _update_progress(project_id, {
                "stage": progress.get("stage", "visual_planning") if isinstance(progress, dict) else "visual_planning",
                "message": progress.get("message", "正在生成视觉方案") if isinstance(progress, dict) else str(progress),
                "current_page": progress.get("current_page", 0) if isinstance(progress, dict) else 0,
                "total_pages": progress.get("total_pages", len(target_slides)) if isinstance(progress, dict) else len(target_slides),
            }, run_id),
        )
        db.expire_all()
        if not is_run_active(db, run_id):
            logger.info(f"VisualPrompts BG: run {run_id} is no longer active after visual plan; skipping stale writeback")
            return
        project = db.query(Project).filter(Project.id == project_id).first()
        if not project:
            return

        # 更新数据库：visual plan（先提交，避免后续 prompt 失败导致 visual plan 也被回滚）
        visual_by_page = {v["page_num"]: v for v in visual_plan}
        artifact_deps = dependency_signature(project, slides)
        for slide in target_slides:
            existing_visual = slide.visual_json if isinstance(slide.visual_json, dict) else {}
            next_visual = _merge_manual_pins_into_visual_json(
                with_artifact_meta(
                    visual_by_page.get(slide.page_num, {}),
                    kind="visual_plan",
                    dependencies=artifact_deps,
                ),
                existing_visual,
            )
            slide.visual_json = next_visual
            slide.prompt_text = None
            slide.image_path = None
            slide.error_msg = None
            flag_modified(slide, "visual_json")
            slide.status = "visual_ready"
        db.commit()

        # Step 2: 并发生成 Prompts（最多 5 个并发，避免 API 限流）
        visual_plan_for_prompts = [
            _with_prompt_asset_policies(s.visual_json, project, s)
            for s in target_slides
            if s.visual_json
        ]
        content_plan_for_prompts = content_plan
        content_by_page = {}
        for item in content_plan_for_prompts:
            content_text = dict(item.get("text_content", {}) or {})
            if item.get("global_user_requirements"):
                content_text["global_user_requirements"] = item["global_user_requirements"]
            content_by_page[item["page_num"]] = content_text

        total_prompt_pages = len(visual_plan_for_prompts)
        completed_count = 0
        semaphore = asyncio.Semaphore(5)
        progress_lock = asyncio.Lock()
        slide_by_page_for_refs = {s.page_num: s for s in target_slides}

        async def _gen_one(intent: dict) -> dict:
            nonlocal completed_count
            page_num = intent["page_num"]
            try:
                async with semaphore:
                    content_text = content_by_page.get(page_num, {})
                    prompt_refs = _project_refs_for_prompt(project, intent.get("visual_asset_ids") or [], intent)
                    prompt_refs.extend(_slide_refs_for_prompt(slide_by_page_for_refs.get(page_num), intent) if slide_by_page_for_refs.get(page_num) else [])
                    prompt = await asyncio.to_thread(
                        generate_prompt_for_page,
                        page_intent=intent,
                        content_text=content_text,
                        style_id=project.style_id or "default",
                        reference_images=prompt_refs or None,
                        style_text_override=style_text_override,
                    )
                async with progress_lock:
                    completed_count += 1
                    _update_progress(project_id, {
                        "stage": "prompt_writing",
                        "message": "正在撰写生图 Prompt",
                        "current_page": completed_count,
                        "total_pages": total_prompt_pages,
                    }, run_id)
                return {"page_num": page_num, "prompt": prompt}
            except Exception as e:
                logger.error(f"PromptEngine: 第 {page_num} 页 Prompt 生成失败: {e}")
                async with progress_lock:
                    completed_count += 1
                    _update_progress(project_id, {
                        "stage": "prompt_writing",
                        "message": "部分页面 Prompt 生成失败，继续处理剩余页面",
                        "current_page": completed_count,
                        "total_pages": total_prompt_pages,
                    }, run_id)
                return {"page_num": page_num, "prompt": "", "error": str(e)}

        tasks = [_gen_one(intent) for intent in visual_plan_for_prompts]
        prompts = await asyncio.gather(*tasks, return_exceptions=False)
        db.expire_all()
        if not is_run_active(db, run_id):
            logger.info(f"VisualPrompts BG: run {run_id} is no longer active after prompts; skipping stale writeback")
            return
        project = db.query(Project).filter(Project.id == project_id).first()
        if not project:
            return
        _update_progress(project_id, {
            "stage": "saving",
            "message": "正在保存结果...",
            "current_page": total_prompt_pages,
            "total_pages": total_prompt_pages,
        }, run_id)

        # 更新数据库：prompts
        prompt_by_page = {
            p["page_num"]: str(p.get("prompt") or "").strip()
            for p in prompts
            if isinstance(p, dict)
        }
        prompt_failures = {
            int(p.get("page_num") or 0): str(p.get("error") or "生图 Prompt 为空").strip()
            for p in prompts
            if isinstance(p, dict) and not str(p.get("prompt") or "").strip()
        }
        prompt_failures = {page_num: reason for page_num, reason in prompt_failures.items() if page_num > 0}
        for slide in target_slides:
            slide.prompt_text = prompt_by_page.get(slide.page_num) or None
            if slide.prompt_text:
                slide.visual_json = with_artifact_meta(
                    slide.visual_json,
                    kind="visual_plan",
                    dependencies=artifact_deps,
                    prompt_dependencies=artifact_deps,
                )
                slide.image_path = None
                slide.error_msg = None
                slide.visual_json = with_stale_flags(slide.visual_json, content=False, visual=False, image=False)
                flag_modified(slide, "visual_json")
                slide.status = "prompt_ready"

        # 状态更新：优先按目标页判断（支持部分生成），再回退到全局判断
        target_all_visual = all(s.visual_json for s in target_slides)
        target_all_prompt = all(s.prompt_text for s in target_slides)
        if target_all_prompt:
            project.status = "prompt_ready"
        elif target_all_visual:
            project.status = "visual_ready"
        elif all(s.prompt_text for s in slides):
            project.status = "prompt_ready"
        elif all(s.visual_json for s in slides):
            project.status = "visual_ready"

        if prompt_failures:
            failed_pages = sorted(prompt_failures)
            failed_pages_text = "、".join(str(page_num) for page_num in failed_pages[:12])
            if len(failed_pages) > 12:
                failed_pages_text += " 等"
            finish_run(
                db,
                run_id,
                status="failed",
                message=f"第 {failed_pages_text} 页生图 Prompt 生成失败，请重试缺失页。",
                completed_count=max(0, total_prompt_pages - len(prompt_failures)),
                failed_count=len(prompt_failures),
                error_msg="; ".join(
                    f"第 {page_num} 页: {reason[:120]}"
                    for page_num, reason in sorted(prompt_failures.items())
                )[:500],
            )
        else:
            finish_run(
                db,
                run_id,
                status="succeeded",
                message=f"画面描述和生图 Prompt 已生成，共 {total_prompt_pages} 页",
                completed_count=total_prompt_pages,
            )
        db.commit()
        logger.info(f"Visual plan and prompts generated for project {project_id}")
    except asyncio.CancelledError:
        logger.info(f"VisualPrompts BG: project={project_id} task was cancelled")
        # 任务被取消时只标记 run，不再根据旧 slides 推断并写回项目阶段。
        # 回退接口已经决定了目标阶段，取消分支晚返回时不能覆盖它。
        try:
            db.rollback()
            finish_run(db, run_id, status="cancelled", message="任务被取消", error_msg="任务被取消")
            db.commit()
        except Exception as status_err:
            logger.warning(f"Failed to reset project status after cancellation: {status_err}")
        _update_progress(project_id, {
            "stage": "error",
            "message": "任务被取消",
            "current_page": 0,
            "total_pages": len(target_slides) if page_nums else 0,
        }, run_id)
        raise  # 重新抛出，让 asyncio 框架识别取消
    except Exception as e:
        db.rollback()
        logger.exception(f"Failed to generate visual plan for project {project_id}: {e}")
        # 重置项目状态，避免永远卡在 generating
        try:
            project = db.query(Project).filter(Project.id == project_id).first()
            if project:
                if all(s.visual_json for s in project.slides):
                    project.status = "visual_ready"
                elif all(s.content_json for s in project.slides):
                    project.status = "planning"
                else:
                    project.status = "draft"
                finish_run(db, run_id, status="failed", message="画面方案生成失败", error_msg=str(e)[:500])
                db.commit()
        except Exception as status_err:
            logger.warning(f"Failed to reset project status after error: {status_err}")
        _update_progress(project_id, {
            "stage": "error",
            "message": f"生成失败：{str(e)[:100]}",
            "current_page": 0,
            "total_pages": len(target_slides) if page_nums else 0,
        }, run_id)
    finally:
        db.close()
        current_task = asyncio.current_task()
        async with _running_tasks_lock:
            owns_registry_entry = _running_tasks.get(project_id) is current_task
            if owns_registry_entry:
                _running_tasks.pop(project_id, None)
        if owns_registry_entry or project_id not in _running_tasks:
            generation_progress.pop(project_id, None)


@router.get("/{project_id}/generation-status")
async def get_generation_status(project_id: str, db: Session = Depends(get_db)):
    """查询正在运行的生成任务状态（同时检查 asyncio 后台任务和 Celery 任务）。"""
    task = _running_tasks.get(project_id)
    project = db.query(Project).filter(Project.id == project_id).first()
    active_run = get_active_run(db, project_id)
    lost_celery_task = _stale_missing_celery_task_if_needed(project, db, active_run)
    if lost_celery_task:
        active_run = None
    if project and not active_run:
        slides = db.query(Slide).filter(Slide.project_id == project_id).order_by(Slide.page_num).all()
        before_status = project.status
        reconcile_project_state(project, slides, active_run)
        if project.status != before_status or db.dirty:
            db.commit()

    # 检查 asyncio 后台任务
    if active_run and task and not task.done():
        return {"generation_status": "running", "project_status": project.status if project else None, "active_run": serialize_run(active_run)}

    # 检查 Celery 任务（通过 Redis 中的 task_id）
    try:
        task_id = active_run.task_id if active_run and active_run.task_id else redis_client.get(f"project:{project_id}:task_id")
        if task_id:
            from app.celery_app import celery_app
            celery_task = AsyncResult(task_id.decode() if isinstance(task_id, bytes) else task_id, app=celery_app)
            if celery_task.state in ("PENDING", "STARTED", "RETRY"):
                return {"generation_status": "running", "project_status": project.status if project else None, "active_run": serialize_run(active_run)}
    except Exception as e:
        logger.warning(f"Failed to check Celery status for {project_id}: {e}")

    if active_run:
        return {"generation_status": "running", "project_status": project.status if project else None, "active_run": serialize_run(active_run)}

    return {"generation_status": "idle", "project_status": project.status if project else None, "active_run": None}


# ==================== P3: Generation & Download ====================

@router.post("/{project_id}/generate")
def start_generation(
    project_id: str,
    body: PageNumsRequest = PageNumsRequest(),
    db: Session = Depends(get_db),
):
    """启动生成流水线（Celery 异步执行）。支持打样：只生成选中的页。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    if _active_run_for_project_action(project, db):
        raise HTTPException(status_code=409, detail="当前项目已有任务正在运行，请等待完成后再开始下一步")

    if project.status == "completed" and not body.page_nums:
        raise HTTPException(status_code=400, detail="已全部完成，如需重新生成请指定页码")

    # 检查项目是否有 slides
    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    if not slides:
        raise HTTPException(status_code=400, detail="No slides found. Generate content plan first.")

    page_nums = _resolve_generation_page_nums(slides, body.page_nums, bool(body.prototype))

    # 记录本次生成的目标页码，供前端进度显示用
    target_slides = [s for s in slides if s.page_num in page_nums] if page_nums else slides
    synced_exact_pages = _sync_exact_page_reference_overlays_for_generation(db, target_slides)
    if synced_exact_pages:
        db.commit()
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, synced_exact_pages))} 页有原样保留素材，请先重新生成画面描述以预留位置。"
        )
    _raise_if_generation_inputs_not_current(project, slides, target_slides)
    missing_prompt_pages = [
        s.page_num
        for s in target_slides
        if not s.prompt_text or not str(s.prompt_text).strip()
    ]
    if missing_prompt_pages:
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, missing_prompt_pages))} 页缺少生图 Prompt，请先生成画面描述。"
        )
    _cleanup_stale_generating_slides_for_project(db, project_id)
    ensure_generation_worker_ready()
    run_kind = "prototype_generation" if body.prototype else ("page_generation" if page_nums else "batch_generation")
    run_stage = image_generation_run_stage(kind=run_kind)
    run = create_project_run(
        db,
        project_id,
        kind=run_kind,
        stage=run_stage,
        target_page_nums=[s.page_num for s in target_slides],
        total_count=len(target_slides),
        message=image_generation_queued_message(run_kind),
    )
    db.commit()
    _update_progress(project_id, {
        "target_page_nums": page_nums,
        "target_count": len(target_slides),
    }, run.id)

    # 使用 Celery 异步任务（Celery worker 内部有 Redis 锁兜底防重）
    task = _enqueue_generation_task(db, project_id, page_nums, run=run, prototype=body.prototype)

    return {
        "message": "Generation started",
        "project_id": project_id,
        "prototype": body.prototype or bool(page_nums),
        "page_nums": page_nums,
        "task_id": task.id,
        "run": serialize_run(run),
    }


@router.post("/{project_id}/confirm-prototype")
def confirm_prototype(
    project_id: str,
    db: Session = Depends(get_db),
):
    """确认打样结果，启动全量生成。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    if _active_run_for_project_action(project, db):
        raise HTTPException(status_code=409, detail="当前项目已有任务正在运行，请等待完成后再开始下一步")

    if project.status != "prototype_ready":
        raise HTTPException(
            status_code=400,
            detail=f"当前状态为 {project.status}，不支持确认打样。请先完成打样生成。"
        )

    # 找出所有未完成或失败的页，打样成功的页面不重复生成。
    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )

    pending_slides = [
        s for s in slides
        if s.status != "failed" and (s.status != "completed" or has_stale_flags(s.visual_json))
    ]
    failed_slides = [s for s in slides if s.status == "failed"]

    # 如果有失败的，先重试；如果有未生成的，一起生成
    target_page_nums = [s.page_num for s in pending_slides + failed_slides]

    if not target_page_nums:
        # 所有页都已完成，直接标记 completed
        project.status = "completed"
        db.commit()
        return {"message": "All slides already completed", "status": "completed"}

    target_slides = pending_slides + failed_slides
    missing_prompt_slides = [
        s
        for s in target_slides
        if not s.prompt_text or not str(s.prompt_text).strip()
    ]
    stale_prompt_slides = [
        s
        for s in target_slides
        if s.prompt_text and str(s.prompt_text).strip() and not _slide_prompt_text_matches_current_generation_context(s)
    ]
    prompt_refresh_slides = []
    seen_prompt_refresh = set()
    for slide in [*missing_prompt_slides, *stale_prompt_slides]:
        if slide.id in seen_prompt_refresh:
            continue
        seen_prompt_refresh.add(slide.id)
        prompt_refresh_slides.append(slide)
    missing_visual_pages = [
        s.page_num
        for s in prompt_refresh_slides
        if not isinstance(s.visual_json, dict) or not str(s.visual_json.get("visual_description") or "").strip()
    ]
    if missing_visual_pages:
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, missing_visual_pages))} 页缺少画面描述，请先生成画面方案。"
        )

    if prompt_refresh_slides:
        logger.info(
            "ConfirmPrototype: project=%s 自动刷新缺失或过期 Prompt pages=%s",
            project_id,
            [s.page_num for s in prompt_refresh_slides],
        )
        for slide in prompt_refresh_slides:
            prompt = _regenerate_slide_prompt(slide, project, db)
            if prompt:
                slide.status = "prompt_ready"
        refreshed_slides = (
            db.query(Slide)
            .filter(Slide.project_id == project_id)
            .order_by(Slide.page_num)
            .all()
        )
        refreshed_deps = dependency_signature(project, refreshed_slides)
        refresh_ids = {slide.id for slide in prompt_refresh_slides}
        for refreshed_slide in refreshed_slides:
            if refreshed_slide.id not in refresh_ids or not refreshed_slide.prompt_text:
                continue
            refreshed_slide.visual_json = with_artifact_meta(
                refreshed_slide.visual_json,
                kind="visual_plan",
                dependencies=refreshed_deps,
                prompt_dependencies=refreshed_deps,
            )
            flag_modified(refreshed_slide, "visual_json")
        slides = refreshed_slides
        target_slides = [s for s in refreshed_slides if s.page_num in target_page_nums]
        db.flush()

    missing_prompt_pages = [
        s.page_num
        for s in target_slides
        if not s.prompt_text or not str(s.prompt_text).strip()
    ]
    if missing_prompt_pages:
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, missing_prompt_pages))} 页缺少生图 Prompt，请先生成画面描述。"
        )
    _raise_if_generation_inputs_not_current(project, slides, target_slides)

    ensure_generation_worker_ready()
    run = create_project_run(
        db,
        project_id,
        kind="batch_generation",
        stage=image_generation_run_stage(kind="batch_generation"),
        target_page_nums=target_page_nums,
        total_count=len(target_page_nums),
        message=image_generation_queued_message("batch_generation"),
    )
    db.commit()

    task = _enqueue_generation_task(db, project_id, target_page_nums, run=run)

    return {
        "message": "Full generation started",
        "project_id": project_id,
        "page_nums": target_page_nums,
        "task_id": task.id,
        "run": serialize_run(run),
    }


@router.post("/{project_id}/stop-generation")
def stop_generation(
    project_id: str,
    db: Session = Depends(get_db),
):
    """停止当前生成任务，重置项目和 slide 状态。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    active_run = get_active_run(db, project_id)
    if not active_run:
        return {"message": "No generation in progress", "status": project.status}

    # 重置所有 generating 状态的 slide
    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .all()
    )
    for slide in slides:
        if slide.status == "generating":
            slide.status = "prompt_ready"
            slide.error_msg = "用户手动停止"
    cancel_active_run(db, project_id, "用户手动停止")
    reconcile_project_state(project, slides, active_run)

    db.commit()

    # 撤销 Celery 任务，真正阻止 worker 继续生图
    task_id = redis_client.get(f"project:{project_id}:task_id")
    if task_id:
        try:
            AsyncResult(task_id.decode() if isinstance(task_id, bytes) else task_id, app=celery_app).revoke(terminate=True)
            logger.info(f"Revoked Celery task {task_id} for project {project_id}")
        except Exception as e:
            logger.warning(f"Failed to revoke task {task_id}: {e}")
        redis_client.delete(f"project:{project_id}:task_id")
        redis_client.delete(f"project:{project_id}:task_started_at")

    lock_page_nums = set()
    if active_run and active_run.target_page_nums:
        lock_page_nums.update(int(p) for p in active_run.target_page_nums if str(p).isdigit())
    lock_page_nums.update(s.page_num for s in slides if s.status == "generating")
    for page_num in sorted(lock_page_nums):
        redis_client.delete(f"project:{project_id}:slide:{page_num}:generating")

    # 清除进度缓存
    if project_id in generation_progress:
        del generation_progress[project_id]

    return {"message": "Generation stopped", "status": project.status}


@router.get("/{project_id}/status")
def get_project_status(project_id: str, db: Session = Depends(get_db)):
    """获取项目生成状态和每页进度。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )

    active_run = get_active_run(db, project_id)
    lost_celery_task = _stale_missing_celery_task_if_needed(project, db, active_run)
    stale_run = None if lost_celery_task else stale_inactive_run_if_needed(db, project_id)
    if stale_run and stale_run.status == "stale":
        db.commit()
    active_run = get_active_run(db, project_id)

    slide_status = []
    completed_count = 0
    for s in slides:
        # 只有真正生成完成的页面才算 completed
        if s.status == "completed":
            completed_count += 1
        slide_status.append({
            "page_num": s.page_num,
            "status": s.status,
            "error_msg": s.error_msg,
        })

    before_status = project.status
    reconcile_project_state(project, slides, active_run)
    logo_policy_changed = _repair_project_logo_policies(project, slides)
    if project.status != before_status or logo_policy_changed or db.dirty:
        db.commit()

    pptx_filename = "prototype.pptx" if project.status == "prototype_ready" else "presentation.pptx"
    pptx_path = os.path.join(_project_output_dir(project_id), pptx_filename)
    if project.status == "completed":
        has_pptx, pptx_path = _has_completed_final_pptx(project, slides)
    else:
        has_pptx = os.path.exists(pptx_path) and not _project_has_stale_artifacts(slides)
    target_count, target_completed, target_failed = target_counts(active_run, slides)
    progress = generation_progress.get(project_id, {})
    return {
        "project_id": project_id,
        "project_status": project.status,
        "total_slides": len(slides),
        "completed_slides": target_completed,  # backward compatible: now target scoped
        "total_completed_slides": completed_count,
        "target_completed_slides": target_completed,
        "target_failed_slides": target_failed,
        "target_count": target_count or progress.get("target_count") or len(slides),
        "target_page_nums": (active_run.target_page_nums if active_run else progress.get("target_page_nums")),
        "active_run": serialize_run(active_run, slides),
        "has_pptx": has_pptx,
        "pptx_path": pptx_path if has_pptx else None,
        "slides": slide_status,
    }


@router.get("/{project_id}/workflow-status")
def get_project_workflow_status(project_id: str, db: Session = Depends(get_db)):
    """获取统一的项目阶段、任务进度和页面状态。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )

    active_run = get_active_run(db, project_id)
    lost_celery_task = _stale_missing_celery_task_if_needed(project, db, active_run)
    stale_run = None if lost_celery_task else stale_inactive_run_if_needed(db, project_id)
    if stale_run and stale_run.status == "stale":
        db.commit()
    active_run = get_active_run(db, project_id)
    latest_run = get_latest_run(db, project_id)
    before_status = project.status
    reconcile_project_state(project, slides, active_run)
    logo_policy_changed = _repair_project_logo_policies(project, slides)
    if project.status != before_status or logo_policy_changed or db.dirty:
        db.commit()

    pptx_filename = "prototype.pptx" if project.status == "prototype_ready" else "presentation.pptx"
    pptx_path = os.path.join(_project_output_dir(project_id), pptx_filename)
    if project.status == "completed":
        has_pptx, pptx_path = _has_completed_final_pptx(project, slides)
    else:
        has_pptx = os.path.exists(pptx_path) and not _project_has_stale_artifacts(slides)
    quality_report = (
        build_project_quality_report(project, slides, has_pptx=has_pptx, pptx_path=pptx_path)
        if not active_run
        else None
    )

    return serialize_workflow_status(
        project,
        slides,
        active_run=active_run,
        latest_run=latest_run,
        has_pptx=has_pptx,
        pptx_path=pptx_path,
        quality_report=quality_report,
    )


@router.get("/{project_id}/generation-progress")
def get_generation_progress(project_id: str, db: Session = Depends(get_db)):
    """获取 Content Plan 后台生成的实时进度。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    active_run = get_active_run(db, project_id)
    lost_celery_task = _stale_missing_celery_task_if_needed(project, db, active_run)
    stale_run = None if lost_celery_task else stale_inactive_run_if_needed(db, project_id)
    if stale_run and stale_run.status == "stale":
        db.commit()
    active_run = get_active_run(db, project_id)
    if not active_run:
        generation_progress.pop(project_id, None)
        slides = db.query(Slide).filter(Slide.project_id == project_id).order_by(Slide.page_num).all()
        before_status = project.status
        reconcile_project_state(project, slides, active_run)
        if project.status != before_status or db.dirty:
            db.commit()
        target_count = len(slides)
        completed_count = sum(1 for s in slides if s.status == "completed")
        return {
            "project_id": project_id,
            "project_status": project.status,
            "active_run": None,
            "stage": None,
            "message": None,
            "current_page": completed_count,
            "total_pages": target_count,
            "think": None,
        }

    progress = generation_progress.get(project_id, {})
    slides = db.query(Slide).filter(Slide.project_id == project_id).order_by(Slide.page_num).all()
    run_data = serialize_run(active_run, slides)
    return {
        "project_id": project_id,
        "project_status": project.status,
        "active_run": run_data,
        "stage": progress.get("stage") or active_run.stage,
        "message": progress.get("message") or active_run.message,
        "current_page": run_data["completed_count"],
        "total_pages": run_data["total_count"],
        "think": progress.get("think"),
    }


@router.get("/{project_id}/download")
def download_pptx(project_id: str, prototype: bool = False, db: Session = Depends(get_db)):
    """下载生成的 PPTX 文件。未生成的页面会保留为空白页。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    slides = db.query(Slide).filter(Slide.project_id == project_id).order_by(Slide.page_num).all()

    output_dir = os.path.join(
        settings.OUTPUT_DIR or "./outputs",
        project_id,
    )
    if not prototype and project.status == "completed":
        has_final_pptx, final_pptx_path = _ensure_completed_final_pptx(project, slides)
        if has_final_pptx:
            display_name = f"{project.title}.pptx"
            return FileResponse(
                path=final_pptx_path,
                filename=display_name,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

    filename = "prototype.pptx" if prototype else "partial_presentation.pptx"
    pptx_path = os.path.join(output_dir, filename)
    generated_count = _assemble_partial_project_pptx(project, slides, pptx_path)
    if generated_count == 0:
        raise HTTPException(status_code=404, detail="还没有生成过页面，暂时不能导出 PPTX。")

    display_name = f"{project.title}_prototype.pptx" if prototype else f"{project.title}.pptx"
    return FileResponse(
        path=pptx_path,
        filename=display_name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@router.post("/{project_id}/upload")
def upload_file(
    project_id: str,
    file: UploadFile = File(...),
    role: str = Form("style_ref"),
    slide_id: Optional[str] = Form(None),
    process_mode: Optional[str] = Form(None),
    asset_name: Optional[str] = Form(None),
    asset_kind: Optional[str] = Form(None),
    usage_note: Optional[str] = Form(None),
    logo_anchor: Optional[str] = Form(None),
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
):
    """上传参考图或 Logo 到项目目录。支持按页上传（传 slide_id）。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    if role not in ALLOWED_UPLOAD_ROLES:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid role '{role}'. Allowed: {', '.join(sorted(ALLOWED_UPLOAD_ROLES))}",
        )

    explicit_process_mode = process_mode is not None
    if process_mode is not None and process_mode not in {"blend", "crop", "original"}:
        raise HTTPException(status_code=400, detail="Invalid process_mode. Allowed: blend, crop, original")
    normalized_logo_anchor = normalize_logo_anchor(logo_anchor)

    normalized_asset_kind = None
    if role == "visual_asset":
        if slide_id:
            raise HTTPException(status_code=400, detail="visual_asset must be uploaded as a project-level asset")
        normalized_asset_kind = normalize_visual_asset_kind(asset_kind)
        if process_mode is None:
            process_mode = default_visual_asset_process_mode(normalized_asset_kind)
    elif role == "logo" and process_mode is None:
        process_mode = "original"
    elif process_mode is None:
        process_mode = "blend"

    # 如果传了 slide_id，校验该 slide 存在
    if slide_id:
        slide = db.query(Slide).filter(Slide.id == slide_id, Slide.project_id == project_id).first()
        if not slide:
            raise HTTPException(status_code=404, detail="Slide not found")

    # 数量限制：页面级长期参考图最多 10 张；项目级素材库不设硬上限。
    if slide_id and role not in {"visual_asset", "finetune_ref", "chat_ref"}:
        existing_count = db.query(ReferenceImage).filter(
            ReferenceImage.project_id == project_id,
            ReferenceImage.slide_id == slide_id,
            ReferenceImage.role != "finetune_ref",
        ).count()
        if existing_count >= MAX_REFERENCE_IMAGES_PER_PAGE:
            raise HTTPException(
                status_code=400,
                detail=f"该页面已有 {existing_count} 张参考图，上限 {MAX_REFERENCE_IMAGES_PER_PAGE} 张",
            )

    # 读取并校验文件
    file_bytes = file.file.read()
    if len(file_bytes) == 0:
        raise HTTPException(status_code=400, detail="上传的文件为空")
    if len(file_bytes) > MAX_UPLOAD_SIZE:
        raise HTTPException(
            status_code=400,
            detail=f"文件大小超过限制（最大 {MAX_UPLOAD_SIZE // 1024 // 1024}MB）",
        )
    if file.content_type and file.content_type not in ALLOWED_IMAGE_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件类型 '{file.content_type}'。允许: JPEG, PNG, GIF, WebP, SVG",
        )

    project_upload_dir = os.path.join(settings.UPLOAD_DIR, project_id)
    if not os.path.isabs(project_upload_dir):
        project_upload_dir = os.path.abspath(project_upload_dir)
    os.makedirs(project_upload_dir, exist_ok=True)

    # 安全检查：拒绝路径遍历攻击
    safe_name = file.filename.replace("\\", "/").split("/")[-1]
    if ".." in safe_name or not safe_name or safe_name.startswith("."):
        raise HTTPException(status_code=400, detail="非法文件名")

    # 尝试用 PIL 打开并统一转 PNG，兼容 JPEG/PNG/GIF/WebP/HEIC/BMP/TIFF 等
    try:
        img = PILImage.open(io.BytesIO(file_bytes))
        if img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGBA")
        else:
            img = img.convert("RGB")
        safe_name = os.path.splitext(safe_name)[0] + ".png"
        prefix = f"slide_{slide_id}_" if slide_id else ""
        if role == "finetune_ref":
            safe_name = f"{int(time.time() * 1000)}_{safe_name}"
        filename = f"{prefix}{role}_{safe_name}"
        file_path = os.path.join(project_upload_dir, filename)
        img.save(file_path, "PNG")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"图片格式无法处理: {e}")

    asset_analysis = None
    cleaned_asset_name = (asset_name or "").strip() or None
    cleaned_usage_note = (usage_note or "").strip() or None
    if role == "visual_asset":
        if background_tasks is not None:
            asset_analysis = _visual_asset_analysis_placeholder(
                asset_name=cleaned_asset_name,
                asset_kind=normalized_asset_kind or "other",
                usage_note=cleaned_usage_note,
                filename=safe_name,
                user_asset_kind_provided=bool(asset_kind and str(asset_kind).strip()),
                explicit_process_mode=explicit_process_mode,
            )
            if not cleaned_asset_name:
                cleaned_asset_name = asset_analysis.get("subject")
        else:
            try:
                asset_analysis = analyze_visual_asset(
                    file_path,
                    asset_name=cleaned_asset_name or os.path.splitext(os.path.basename(safe_name))[0],
                    asset_kind=normalized_asset_kind or "other",
                    usage_note=cleaned_usage_note or "",
                )
                if isinstance(asset_analysis, dict):
                    asset_analysis = {
                        **asset_analysis,
                        "analysis_status": "completed",
                        "analysis_type": "visual_asset",
                        "selection_tier": "manual",
                        "importance_score": 100,
                    }
            except Exception as e:
                logger.warning(f"Visual asset analysis failed for {file_path}: {e}")
                asset_analysis = _visual_asset_analysis_placeholder(
                    asset_name=cleaned_asset_name,
                    asset_kind=normalized_asset_kind or "other",
                    usage_note=cleaned_usage_note,
                    filename=safe_name,
                    user_asset_kind_provided=bool(asset_kind and str(asset_kind).strip()),
                    explicit_process_mode=explicit_process_mode,
                )
                asset_analysis["analysis_status"] = "failed"
                asset_analysis["error"] = str(e)[:500]
            if not asset_kind and isinstance(asset_analysis, dict):
                normalized_asset_kind = normalize_visual_asset_kind(asset_analysis.get("detected_kind"))
                if not explicit_process_mode:
                    process_mode = default_visual_asset_process_mode(normalized_asset_kind)
            if not cleaned_asset_name:
                analyzed_subject = asset_analysis.get("subject") if isinstance(asset_analysis, dict) else None
                cleaned_asset_name = analyzed_subject or os.path.splitext(os.path.basename(safe_name))[0]
    elif role in REFERENCE_ANALYSIS_ROLES:
        asset_analysis = _reference_analysis_placeholder(role, safe_name, cleaned_usage_note)
    elif role == "logo":
        asset_analysis = {
            "analysis_status": "tone_detected",
            "analysis_type": "logo",
            "review_status": "auto_confirmed",
            "review_reason": "用户手动上传的品牌 Logo",
            "confidence_score": 1.0,
            "needs_user_review": False,
            **detect_logo_tone_from_image(file_path),
        }

    ref_image = ReferenceImage(
        project_id=project_id,
        slide_id=slide_id,
        file_path=file_path,
        role=role,
        process_mode=process_mode,
        asset_name=cleaned_asset_name if role == "visual_asset" or role in {"content_ref", "chart_ref"} else None,
        asset_kind=normalized_asset_kind if role == "visual_asset" else None,
        usage_note=cleaned_usage_note if role == "visual_asset" or role in {"content_ref", "chart_ref"} else None,
        asset_analysis=asset_analysis if role == "visual_asset" or role in REFERENCE_ANALYSIS_ROLES or role == "logo" else None,
        logo_anchor=normalized_logo_anchor if role == "logo" else None,
    )
    db.add(ref_image)
    if not slide_id and role in {"style_ref", "logo", "template"}:
        _invalidate_style_dependent_outputs(project)
    if role == "visual_asset":
        _invalidate_visual_asset_dependent_outputs(project)
    if slide_id and role != "finetune_ref" and slide:
        _invalidate_visual_plan_dependent_outputs(project, [slide])
    db.commit()
    db.refresh(ref_image)

    if background_tasks is not None:
        if role == "visual_asset":
            _submit_asset_analysis_task("visual_asset", _analyze_visual_asset_for_ref, project_id, ref_image.id)
        elif role in REFERENCE_ANALYSIS_ROLES:
            _submit_asset_analysis_task("reference_image", _analyze_reference_image_for_ref, project_id, ref_image.id)

    return {
        "id": ref_image.id,
        "file_path": file_path,
        "role": role,
        "slide_id": slide_id,
        "process_mode": process_mode,
        "asset_name": ref_image.asset_name,
        "asset_kind": ref_image.asset_kind,
        "usage_note": ref_image.usage_note,
        "asset_analysis": ref_image.asset_analysis,
        "logo_anchor": ref_image.logo_anchor,
        "url": _reference_upload_url(project_id, file_path),
        "overlay_url": _logo_overlay_url(ref_image, project_id) if role == "logo" and is_logo_confirmed(ref_image) else None,
        "symbol_overlay_url": _logo_symbol_overlay_url(ref_image, project_id) if role == "logo" and is_logo_confirmed(ref_image) else None,
    }


@router.get("/{project_id}/reference-images")
def list_reference_images(
    project_id: str,
    slide_id: Optional[str] = None,
    q: Optional[str] = None,
    asset_kind: Optional[str] = None,
    source_document: Optional[str] = None,
    source_page_num: Optional[int] = None,
    process_mode: Optional[str] = None,
    pinned_slide_id: Optional[str] = None,
    recommend_slide_id: Optional[str] = None,
    limit: Optional[int] = None,
    offset: int = 0,
    db: Session = Depends(get_db),
):
    """列出项目参考图；带筛选参数时作为项目素材库查询。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    has_library_params = any([
        q,
        asset_kind,
        source_document,
        source_page_num is not None,
        process_mode,
        pinned_slide_id,
        recommend_slide_id,
        limit is not None,
        offset,
    ])

    query = db.query(ReferenceImage).filter(ReferenceImage.project_id == project_id)
    if slide_id:
        query = query.filter(ReferenceImage.slide_id == slide_id)
    else:
        # 不传 slide_id 时，只返回项目级（slide_id 为 null）的参考图
        query = query.filter(ReferenceImage.slide_id.is_(None))

    query = query.filter(ReferenceImage.role != "finetune_ref")
    images = query.all()
    if not slide_id and not has_library_params:
        images = [
            img for img in images
            if img.role != "visual_asset" or _is_planning_visual_asset(img)
        ]

    slides = db.query(Slide).filter(Slide.project_id == project_id).all()
    pinned_by_asset: dict[str, list[str]] = {}
    for slide in slides:
        for asset_id in _manual_asset_ids(slide.visual_json):
            pinned_by_asset.setdefault(asset_id, []).append(slide.id)

    recommend_lookup: dict[str, dict] = {}
    if recommend_slide_id:
        rec_slide = next((s for s in slides if s.id == recommend_slide_id), None)
        if rec_slide:
            rec_page = copy.deepcopy(rec_slide.content_json) or {}
            rec_page["page_num"] = rec_slide.page_num
            rec_page["type"] = rec_slide.type or rec_page.get("type", "content")
            global_assets = _project_visual_assets_for_planning(project)
            for item in _recall_visual_assets_for_page(rec_page, global_assets):
                if item.get("id"):
                    recommend_lookup[str(item["id"])] = item

    def _matches(img: ReferenceImage) -> bool:
        analysis = img.asset_analysis if isinstance(img.asset_analysis, dict) else {}
        haystack_parts = [
            img.asset_name or "",
            img.asset_kind or "",
            img.usage_note or "",
            os.path.basename(img.file_path or ""),
            _asset_source_document(img),
            str(_asset_source_page_num(img) or ""),
            str(analysis.get("source_slide_text") or ""),
            " ".join(str(x) for x in (analysis.get("asset_tags") or analysis.get("suggested_keywords") or [])),
        ]
        if q and str(q).strip().lower() not in "\n".join(haystack_parts).lower():
            return False
        if asset_kind and (img.asset_kind or "other") != asset_kind:
            return False
        if source_document and _asset_source_document(img) != source_document:
            return False
        if source_page_num is not None and _asset_source_page_num(img) != source_page_num:
            return False
        if process_mode and (img.process_mode or "blend") != process_mode:
            return False
        if pinned_slide_id and pinned_slide_id not in pinned_by_asset.get(img.id, []):
            return False
        return True

    filtered = [img for img in images if _matches(img)]
    filtered.sort(
        key=lambda img: (
            0 if img.id in pinned_by_asset else 1,
            0 if img.id in recommend_lookup else 1,
            _asset_source_document(img),
            _asset_source_page_num(img) or 10_000,
            img.asset_name or os.path.basename(img.file_path or ""),
        )
    )

    facets = {
        "source_documents": sorted({doc for doc in (_asset_source_document(img) for img in images) if doc}),
        "asset_kinds": sorted({(img.asset_kind or "other") for img in images if img.role == "visual_asset"}),
        "process_modes": sorted({(img.process_mode or "blend") for img in images}),
        "source_page_nums": sorted({page for page in (_asset_source_page_num(img) for img in images) if page is not None}),
    }

    total = len(filtered)
    start = max(0, int(offset or 0))
    if limit is None:
        paged = filtered[start:]
    else:
        paged = filtered[start:start + max(1, min(int(limit), 200))]
    items = [
        _serialize_reference_image(
            img,
            project_id,
            pinned_slide_ids=pinned_by_asset.get(img.id, []),
            is_pinned=bool(pinned_slide_id and pinned_slide_id in pinned_by_asset.get(img.id, [])),
            relevance=recommend_lookup.get(img.id),
        )
        for img in paged
    ]

    if has_library_params:
        return {"items": items, "total": total, "facets": facets}
    return items


@router.post("/{project_id}/suggest-reference-images")
def suggest_reference_images(
    project_id: str,
    db: Session = Depends(get_db),
):
    """内容总监：根据内容大纲推荐参考图片。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    if not slides:
        return {"suggestions": []}

    outline = []
    for s in slides:
        content = s.content_json or {}
        tc = content.get("text_content", {})
        outline.append({
            "page_num": s.page_num,
            "type": s.type or "content",
            "headline": tc.get("headline", ""),
            "subhead": tc.get("subhead", ""),
            "body": tc.get("body", "")[:300],
        })

    system_prompt = (
        "你是资深 PPT 内容总监。根据用户提供的 PPT 大纲，为每一页分析是否需要参考图片，"
        "并给出具体建议。必须且只能输出合法 JSON 数组，严禁添加任何额外说明文本。"
    )
    user_prompt = f"""请根据以下 PPT 大纲，推荐需要参考图片的页面。

要求：
1. 只建议"有明确视觉主体"的页面（产品、人物、场景、数据图表、品牌展示等）
2. 纯文字过渡页、目录页不要推荐
3. 输出 JSON 数组，每个元素包含：page_num(int), type(str), reason(str), recommended_mode(str)
4. recommended_mode 是内部生成策略：
   - 人像/产品/场景作为画面参考 → "blend"
   - Logo/多图并排/图标作为身份参考 → "crop"
   - 证书/严格比例/不允许改动的参考 → "original"

大纲：
{json.dumps(outline, ensure_ascii=False, indent=2)}

只输出 JSON 数组，不要任何其他文字。"""

    try:
        client = get_llm_client()
        resp = client.chat.completions.create(
            model=get_minimax_llm_model(),
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.5,
        )
        content = resp.choices[0].message.content or ""
        # 清理可能的 markdown 代码块
        content = content.strip()
        if content.startswith("```"):
            content = content.split("\n", 1)[1] if "\n" in content else ""
            if content.endswith("```"):
                content = content.rsplit("\n", 1)[0]
            content = content.strip()
        suggestions = json.loads(content)
        if not isinstance(suggestions, list):
            suggestions = []
    except Exception as e:
        logger.warning(f"LLM suggest_reference_images failed: {e}")
        suggestions = []

    return {"suggestions": suggestions}


@router.delete("/{project_id}/reference-images/{ref_id}")
def delete_reference_image(
    project_id: str,
    ref_id: str,
    db: Session = Depends(get_db),
):
    """删除指定参考图/Logo/模板，同时清理文件。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    ref = db.query(ReferenceImage).filter(
        ReferenceImage.id == ref_id,
        ReferenceImage.project_id == project_id,
    ).first()
    if not ref:
        raise HTTPException(status_code=404, detail="Reference image not found")

    # 如果删除的是模板，需要删除该项目的所有模板页记录和文件
    if ref.role == "template":
        all_template_refs = db.query(ReferenceImage).filter(
            ReferenceImage.project_id == project_id,
            ReferenceImage.role == "template",
        ).all()
        for t_ref in all_template_refs:
            db.delete(t_ref)
        if project.selected_template_recommendations:
            project.selected_template_recommendations = None
            logger.info(f"Cleared template recommendations for project {project_id}")
        _invalidate_style_dependent_outputs(project)
        db.commit()
        # 数据库提交成功后再删物理文件，文件删失败不影响事务
        for t_ref in all_template_refs:
            try:
                analysis = t_ref.asset_analysis if isinstance(t_ref.asset_analysis, dict) else {}
                paths = {t_ref.file_path, analysis.get("layout_file_path"), analysis.get("preview_file_path")}
                for path in [p for p in paths if p]:
                    if os.path.exists(path):
                        os.remove(path)
                        logger.info(f"Deleted template file: {path}")
            except Exception as e:
                logger.warning(f"Failed to delete template file {t_ref.file_path}: {e}")
        return {"message": "Deleted all template pages", "count": len(all_template_refs)}

    # 普通参考图/Logo 删除
    if ref.role in {"style_ref", "logo"} and not ref.slide_id:
        _invalidate_style_dependent_outputs(project)
    if ref.slide_id and ref.role != "finetune_ref":
        slide = db.query(Slide).filter(Slide.id == ref.slide_id, Slide.project_id == project_id).first()
        if slide:
            if ref.role in {"content_ref", "chart_ref"}:
                slide.visual_json = remove_asset_from_overlay_layers(slide.visual_json, ref_id)
            _invalidate_visual_plan_dependent_outputs(project, [slide])
    if ref.role == "visual_asset":
        slides = db.query(Slide).filter(Slide.project_id == project_id).all()
        for slide in slides:
            visual = copy.deepcopy(slide.visual_json) or {}
            ids = visual.get("visual_asset_ids")
            if isinstance(ids, list) and ref_id in ids:
                visual["visual_asset_ids"] = [x for x in ids if x != ref_id]
                usage = visual.get("visual_asset_usage")
                if isinstance(usage, dict):
                    usage.pop(ref_id, None)
                    visual["visual_asset_usage"] = usage
            manual_ids = visual.get("manual_visual_asset_ids")
            if isinstance(manual_ids, list) and ref_id in manual_ids:
                visual["manual_visual_asset_ids"] = [x for x in manual_ids if x != ref_id]
                manual_usage = visual.get("manual_visual_asset_usage")
                if isinstance(manual_usage, dict):
                    manual_usage.pop(ref_id, None)
                    visual["manual_visual_asset_usage"] = manual_usage
            excluded_ids = visual.get("excluded_visual_asset_ids")
            if isinstance(excluded_ids, list) and ref_id in excluded_ids:
                visual["excluded_visual_asset_ids"] = [x for x in excluded_ids if x != ref_id]
            visual = remove_asset_from_overlay_layers(visual, ref_id)
            slide.visual_json = visual
            flag_modified(slide, "visual_json")
        _invalidate_visual_asset_dependent_outputs(project)
    db.delete(ref)
    db.commit()
    try:
        if os.path.exists(ref.file_path):
            os.remove(ref.file_path)
            logger.info(f"Deleted file: {ref.file_path}")
    except Exception as e:
        logger.warning(f"Failed to delete file {ref.file_path}: {e}")
    return {"message": "Deleted", "id": ref_id}


@router.patch("/{project_id}/reference-images/{ref_id}")
def update_reference_image(
    project_id: str,
    ref_id: str,
    payload: dict = Body(...),
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
):
    """更新参考图的处理模式（blend/crop/original）和视觉资产元信息。"""
    process_mode = payload.get("process_mode")
    if process_mode is not None and process_mode not in {"blend", "crop", "original"}:
        raise HTTPException(status_code=400, detail="Invalid process_mode. Allowed: blend, crop, original")
    logo_anchor = payload.get("logo_anchor")
    if logo_anchor is not None and str(logo_anchor).strip().lower().replace("_", "-") not in LOGO_ANCHORS:
        raise HTTPException(status_code=400, detail="Invalid logo_anchor. Allowed: top-left, top-right, bottom-left, bottom-right")
    review_status = payload.get("review_status")
    if review_status is not None:
        review_status = str(review_status).strip().lower()
        if review_status not in {"auto_confirmed", "user_confirmed", "needs_review", "dismissed", "not_logo"}:
            raise HTTPException(status_code=400, detail="Invalid review_status")

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    ref = db.query(ReferenceImage).filter(
        ReferenceImage.id == ref_id,
        ReferenceImage.project_id == project_id,
    ).first()
    if not ref:
        raise HTTPException(status_code=404, detail="Reference image not found")

    changed_visual_asset = False
    changed_logo = False
    reanalyze_visual_asset = False
    if process_mode is not None:
        ref.process_mode = process_mode
        changed_visual_asset = ref.role == "visual_asset"
        changed_logo = ref.role == "logo"
        if ref.slide_id and ref.role in EXACT_PAGE_REFERENCE_ROLES:
            slide = db.query(Slide).filter(Slide.id == ref.slide_id, Slide.project_id == project_id).first()
            if slide:
                if process_mode == "original":
                    _sync_exact_page_reference_overlay_layers(slide, [ref])
                else:
                    next_visual = remove_asset_from_overlay_layers(slide.visual_json, ref.id)
                    if next_visual != (slide.visual_json if isinstance(slide.visual_json, dict) else {}):
                        slide.visual_json = next_visual
                        flag_modified(slide, "visual_json")
                    _clear_reference_exact_overlay_mark(ref)

    if ref.role == "logo" and logo_anchor is not None:
        ref.logo_anchor = normalize_logo_anchor(logo_anchor)
        changed_logo = True

    if ref.role == "logo" and review_status is not None:
        current = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
        ref.asset_analysis = {
            **current,
            "review_status": review_status,
            "needs_user_review": review_status == "needs_review",
            "review_reason": payload.get("review_reason") or current.get("review_reason"),
        }
        changed_logo = True
        if review_status == "not_logo":
            ref.role = "visual_asset"
            ref.process_mode = "crop"
            ref.asset_kind = normalize_visual_asset_kind(payload.get("asset_kind") or "material")
            ref.asset_name = (payload.get("asset_name") or ref.asset_name or "PPT 提取素材").strip()
            ref.usage_note = (payload.get("usage_note") or ref.usage_note or "用户从疑似 Logo 改为普通可复用素材").strip()
            ref.logo_anchor = None
            ref.asset_analysis = {
                **ref.asset_analysis,
                "analysis_type": "visual_asset",
                "detected_kind": ref.asset_kind,
                "subject": ref.asset_name,
                "selection_tier": "manual_review",
                "recommended_usage": ref.usage_note,
            }
            changed_visual_asset = True

    if ref.role == "visual_asset":
        if "asset_name" in payload:
            ref.asset_name = (payload.get("asset_name") or "").strip() or None
            changed_visual_asset = True
        if "asset_kind" in payload:
            ref.asset_kind = normalize_visual_asset_kind(payload.get("asset_kind"))
            changed_visual_asset = True
        if "usage_note" in payload:
            ref.usage_note = (payload.get("usage_note") or "").strip() or None
            changed_visual_asset = True
        if payload.get("reanalyze"):
            current = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
            ref.asset_analysis = {
                **current,
                "analysis_status": "queued",
                "analysis_type": "visual_asset",
                "_user_asset_kind_provided": True,
                "_explicit_process_mode": True,
            }
            reanalyze_visual_asset = True
            changed_visual_asset = True

    if not process_mode and not changed_visual_asset and not changed_logo:
        raise HTTPException(status_code=400, detail="No supported fields to update")

    if changed_logo:
        reassembled = _reassemble_generated_pptx_if_possible(project)
        if not reassembled:
            _invalidate_style_dependent_outputs(project)
    if changed_visual_asset and ref.role == "visual_asset":
        _invalidate_visual_asset_dependent_outputs(project)
    elif ref.slide_id and ref.role != "finetune_ref":
        slide = db.query(Slide).filter(Slide.id == ref.slide_id, Slide.project_id == project_id).first()
        if slide:
            _invalidate_visual_plan_dependent_outputs(project, [slide])
    db.commit()
    db.refresh(ref)
    if reanalyze_visual_asset and background_tasks is not None:
        _submit_asset_analysis_task("visual_asset_reanalysis", _analyze_visual_asset_for_ref, project_id, ref.id)
    return {
        "id": ref.id,
        "role": ref.role,
        "process_mode": ref.process_mode,
        "asset_name": ref.asset_name,
        "asset_kind": ref.asset_kind,
        "usage_note": ref.usage_note,
        "asset_analysis": ref.asset_analysis,
        "review_status": logo_review_status(ref) if ref.role == "logo" else (ref.asset_analysis or {}).get("review_status"),
        "logo_anchor": ref.logo_anchor or (DEFAULT_LOGO_ANCHOR if ref.role == "logo" else None),
        "url": _reference_upload_url(project_id, ref.file_path),
        "overlay_url": _logo_overlay_url(ref, project_id) if ref.role == "logo" and is_logo_confirmed(ref) else None,
        "symbol_overlay_url": _logo_symbol_overlay_url(ref, project_id) if ref.role == "logo" and is_logo_confirmed(ref) else None,
    }


@router.patch("/{project_id}/slides/{slide_id}/asset-pins")
def update_slide_asset_pins(
    project_id: str,
    slide_id: str,
    payload: AssetPinsRequest,
    db: Session = Depends(get_db),
):
    """Replace one slide's manually pinned project visual assets."""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    slide = db.query(Slide).filter(Slide.id == slide_id, Slide.project_id == project_id).first()
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    requested_ids: list[str] = []
    for asset_id in payload.asset_ids or []:
        value = str(asset_id)
        if value and value not in requested_ids:
            requested_ids.append(value)

    if requested_ids:
        valid_assets = db.query(ReferenceImage).filter(
            ReferenceImage.project_id == project_id,
            ReferenceImage.slide_id.is_(None),
            ReferenceImage.role == "visual_asset",
            ReferenceImage.id.in_(requested_ids),
        ).all()
        valid_ids = {asset.id for asset in valid_assets}
        missing = [asset_id for asset_id in requested_ids if asset_id not in valid_ids]
        if missing:
            raise HTTPException(status_code=400, detail=f"Invalid visual_asset ids: {', '.join(missing[:5])}")

    visual = copy.deepcopy(slide.visual_json) if isinstance(slide.visual_json, dict) else {}
    manual_usage = {
        str(k): str(v)
        for k, v in (payload.usage or {}).items()
        if str(k) in requested_ids and v
    }
    existing_usage = visual.get("visual_asset_usage") if isinstance(visual.get("visual_asset_usage"), dict) else {}
    auto_ids = [
        str(asset_id)
        for asset_id in (visual.get("visual_asset_ids") or [])
        if str(asset_id) not in requested_ids
    ]
    excluded_ids = [
        asset_id
        for asset_id in _excluded_asset_ids(visual)
        if asset_id not in set(requested_ids)
    ]
    excluded_set = set(excluded_ids)
    visual["manual_visual_asset_ids"] = requested_ids
    visual["manual_visual_asset_usage"] = manual_usage
    visual["visual_asset_ids"] = _merge_asset_ids(
        requested_ids,
        [asset_id for asset_id in auto_ids if asset_id not in excluded_set],
    )
    if excluded_ids:
        visual["excluded_visual_asset_ids"] = excluded_ids
    else:
        visual.pop("excluded_visual_asset_ids", None)
    if isinstance(visual.get("overlay_layers"), list):
        requested_set = set(requested_ids)
        global_asset_ids = {
            asset.id for asset in db.query(ReferenceImage).filter(
                ReferenceImage.project_id == project_id,
                ReferenceImage.slide_id.is_(None),
                ReferenceImage.role == "visual_asset",
            ).all()
        }
        visual["overlay_layers"] = [
            layer for layer in normalize_overlay_layers(visual.get("overlay_layers"), valid_asset_ids=None, strict_assets=False)
            if str(layer.get("asset_id")) not in global_asset_ids or str(layer.get("asset_id")) in requested_set
        ]
    visual["visual_asset_usage"] = {
        **{str(k): str(v) for k, v in existing_usage.items() if k and v and str(k) not in manual_usage and str(k) not in excluded_set},
        **manual_usage,
    }
    slide.visual_json = visual
    _invalidate_visual_plan_dependent_outputs(project, [slide])
    db.commit()
    db.refresh(slide)
    return {
        "slide_id": slide.id,
        "page_num": slide.page_num,
        "visual_json": slide.visual_json or {},
        "manual_visual_asset_ids": _manual_asset_ids(slide.visual_json),
    }


@router.patch("/{project_id}/slides/{slide_id}/overlay-layers")
def update_slide_overlay_layers(
    project_id: str,
    slide_id: str,
    payload: OverlayLayersRequest,
    db: Session = Depends(get_db),
):
    """Replace one slide's exact overlay layers."""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    slide = db.query(Slide).filter(Slide.id == slide_id, Slide.project_id == project_id).first()
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    raw_layers = [
        layer.model_dump() if hasattr(layer, "model_dump") else layer.dict()
        for layer in (payload.layers or [])
    ]
    requested_asset_ids = {
        str(layer.get("asset_id"))
        for layer in raw_layers
        if layer.get("asset_id")
    }
    valid_assets = []
    if requested_asset_ids:
        candidates = db.query(ReferenceImage).filter(
            ReferenceImage.project_id == project_id,
            ReferenceImage.id.in_(requested_asset_ids),
        ).all()
        valid_assets = [
            asset for asset in candidates
            if asset.file_path and os.path.exists(_resolve_file_path(asset.file_path))
        ]
    valid_ids = {str(asset.id) for asset in valid_assets}
    missing = [asset_id for asset_id in requested_asset_ids if asset_id not in valid_ids]
    if missing:
        logger.warning(
            "update_slide_overlay_layers: requested asset_ids not found or missing file: %s",
            missing,
        )
        raise HTTPException(status_code=400, detail=f"Invalid overlay asset ids: {', '.join(missing[:5])}")

    normalized_layers = normalize_overlay_layers(raw_layers, valid_asset_ids=valid_ids, strict_assets=True)
    overlay_asset_ids = {str(layer.get("asset_id")) for layer in normalized_layers if layer.get("asset_id")}
    for asset in valid_assets:
        if str(asset.id) in overlay_asset_ids:
            asset.process_mode = "original"
            analysis = asset.asset_analysis if isinstance(asset.asset_analysis, dict) else {}
            asset.asset_analysis = {
                **analysis,
                "exact_overlay": True,
                "exact_overlay_reason": "用户选择原样保留，最终页面会保留素材比例和细节。",
            }
    visual = copy.deepcopy(slide.visual_json) if isinstance(slide.visual_json, dict) else {}
    visual["overlay_layers"] = normalized_layers
    slide.visual_json = visual
    flag_modified(slide, "visual_json")
    _invalidate_visual_plan_dependent_outputs(project, [slide])
    db.commit()
    db.refresh(slide)
    return {
        "slide_id": slide.id,
        "page_num": slide.page_num,
        "visual_json": slide.visual_json or {},
        "overlay_layers": normalized_layers,
    }


@router.post("/{project_id}/retry-failed")
def retry_failed_slides(
    project_id: str,
    db: Session = Depends(get_db),
):
    """批量重试 project 下所有失败的 slides。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    if _active_run_for_project_action(project, db):
        raise HTTPException(status_code=409, detail="当前项目已有任务正在运行，请等待完成后再开始下一步")

    failed_slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id, Slide.status == "failed")
        .order_by(Slide.page_num)
        .all()
    )

    if not failed_slides:
        raise HTTPException(status_code=400, detail="没有失败的页面需要重试")

    synced_exact_pages = _sync_exact_page_reference_overlays_for_generation(db, failed_slides)
    if synced_exact_pages:
        db.commit()
        raise HTTPException(
            status_code=400,
            detail=f"第 {', '.join(map(str, synced_exact_pages))} 页有原样保留素材，请先重新生成画面描述以预留位置。"
        )

    project_slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    _raise_if_generation_inputs_not_current(project, project_slides, failed_slides)

    # 过滤掉正在生成中的 slide，避免重复触发
    page_nums = []
    for slide in failed_slides:
        if slide.status == "generating":
            continue
        slide.status = "generating"
        slide.error_msg = None
        page_nums.append(slide.page_num)

    if not page_nums:
        raise HTTPException(status_code=400, detail="所有失败页面正在重试中，请稍候")

    _cleanup_stale_generating_slides_for_project(db, project_id)
    ensure_generation_worker_ready()
    run = create_project_run(
        db,
        project_id,
        kind="retry_failed",
        stage=image_generation_run_stage(kind="retry_failed"),
        target_page_nums=page_nums,
        total_count=len(page_nums),
        message=image_generation_queued_message("retry_failed"),
    )
    db.commit()

    # 记录本次重试的目标页码和数量，供前端进度显示用
    _update_progress(project_id, {
        "target_page_nums": page_nums,
        "target_count": len(page_nums),
    }, run.id)

    task = _enqueue_generation_task(
        db,
        project_id,
        page_nums,
        run=run,
        prototype=project.status == "prototype_ready",
    )

    return {
        "message": "Retry started",
        "page_nums": page_nums,
        "count": len(page_nums),
        "run": serialize_run(run),
    }


class RetryRequest(BaseModel):
    regenerate_prompt: bool = False
    user_feedback: Optional[str] = None


def _regenerate_slide_prompt(slide: Slide, project: Project, db: Session, user_feedback: str | None = None) -> str:
    """为单页重新生成 prompt_text，基于最新的 content_json 和 visual_json。"""
    slide.reference_images = sorted(
        db.query(ReferenceImage).filter(
            ReferenceImage.slide_id == slide.id,
            ReferenceImage.role != "finetune_ref",
        ).all(),
        key=_reference_image_sort_key,
    )
    _sync_exact_page_reference_overlay_layers(slide)
    content_json = _content_json_for_generation(slide)
    visual_json = slide.visual_json or {}

    page_intent = {
        "page_num": slide.page_num,
        "type": visual_json.get("type") or content_json.get("type", "content"),
        "layout": visual_json.get("layout"),
        "visual_summary": visual_json.get("visual_summary", ""),
        "visual_evidence": visual_json.get("visual_evidence", ""),
        "visual_description": visual_json.get("visual_description", ""),
        "design_notes": visual_json.get("design_notes", ""),
        "visual_asset_ids": visual_json.get("visual_asset_ids", []),
        "visual_asset_usage": visual_json.get("visual_asset_usage", {}),
        "overlay_layers": visual_json.get("overlay_layers", []),
        "logo_policy": visual_json.get("logo_policy"),
    }
    page_intent = _with_prompt_asset_policies(page_intent, project, slide) or page_intent

    content_text = content_json.get("text_content", {})
    if isinstance(content_text, dict):
        content_text = {
            "headline": content_text.get("headline", ""),
            "subhead": content_text.get("subhead", ""),
            "body": content_text.get("body", ""),
        }
        if content_json.get("visual_requirements"):
            content_text["visual_requirements"] = content_json.get("visual_requirements")
    else:
        content_text = {"headline": "", "subhead": "", "body": ""}

    # 收集参考图描述：项目级可复用素材优先，页面级参考图随后补充。
    reference_images = _project_refs_for_prompt(project, visual_json.get("visual_asset_ids") or [], page_intent)
    reference_images.extend(_slide_refs_for_prompt(slide, visual_json))

    # 获取项目级风格覆盖；没有确认风格时从内容推导紧凑 style pack。
    style_text_override = _style_text_from_selected_style(project.selected_style) or _derive_project_style_pack(
        project,
        [{"type": slide.type or "content", "text_content": content_text}],
    )

    prompt = generate_prompt_for_page(
        page_intent=page_intent,
        content_text=content_text,
        style_id=project.style_id or "default",
        reference_images=reference_images or None,
        style_text_override=style_text_override,
        user_feedback=user_feedback,
    )

    slide.prompt_text = prompt
    project_slides = db.query(Slide).filter(Slide.project_id == project.id).order_by(Slide.page_num).all()
    artifact_deps = dependency_signature(project, project_slides)
    slide.visual_json = with_artifact_meta(
        slide.visual_json,
        kind="visual_plan",
        dependencies=artifact_deps,
        prompt_dependencies=artifact_deps,
    )
    slide.image_path = None
    slide.error_msg = None
    slide.visual_json = with_stale_flags(slide.visual_json, content=False, visual=False, image=False)
    flag_modified(slide, "visual_json")
    db.commit()
    logger.info(f"RetrySlide: 已重新生成第 {slide.page_num} 页 prompt")
    return prompt


class FinetuneRegionBox(BaseModel):
    x: float
    y: float
    width: float
    height: float


class FinetuneRegionRequest(BaseModel):
    id: Optional[str] = None
    label: Optional[str] = None
    bbox: FinetuneRegionBox


class FinetuneRequest(BaseModel):
    # Backward compatible: older frontend versions send a fully composed prompt.
    new_prompt: Optional[str] = None
    # Preferred path: send the user's plain edit request and let the image model
    # edit the current slide image directly.
    instruction: Optional[str] = None
    attachment_ids: Optional[List[str]] = None
    regions: Optional[List[FinetuneRegionRequest]] = None


def _finetune_should_attach_project_product_assets(instruction: str) -> bool:
    text = (instruction or "").lower()
    action_terms = ("换成", "替换", "换上", "改成", "改为", "replace", "use the uploaded", "上传")
    asset_terms = (
        "核心资产", "客户", "产品", "瓶", "包装",
        "visual asset", "product", "bottle", "packaging",
    )
    return any(term in text for term in action_terms) and any(term in text for term in asset_terms)


def _project_visual_asset_ids_for_finetune(project: Project, slide: Slide, instruction: str) -> List[str]:
    if not _finetune_should_attach_project_product_assets(instruction):
        return []

    selected_ids: list[str] = []
    if slide.visual_json and isinstance(slide.visual_json, dict):
        raw_ids = slide.visual_json.get("visual_asset_ids") or []
        if isinstance(raw_ids, list):
            selected_ids = [str(x) for x in raw_ids]

    candidate_refs = [
        ref for ref in (project.reference_images or [])
        if ref.role == "visual_asset"
        and not ref.slide_id
        and str(ref.asset_kind or "").lower() in {"product", "material"}
        and os.path.exists(ref.file_path)
    ]
    candidate_refs.sort(
        key=lambda ref: (
            selected_ids.index(str(ref.id)) if str(ref.id) in selected_ids else 999,
            ref.asset_name or "",
        )
    )
    return [str(ref.id) for ref in candidate_refs[:3]]


def _finetune_region_to_dict(region) -> dict:
    if isinstance(region, dict):
        return region
    if hasattr(region, "model_dump"):
        return region.model_dump()
    if hasattr(region, "dict"):
        return region.dict()
    return {}


def _normalize_finetune_regions(regions: Optional[List[object]]) -> List[dict]:
    normalized: List[dict] = []
    for index, raw_region in enumerate(regions or []):
        region = _finetune_region_to_dict(raw_region)
        raw_bbox = region.get("bbox") if isinstance(region, dict) else None
        bbox = _finetune_region_to_dict(raw_bbox)
        try:
            x = max(0.0, min(1.0, float(bbox.get("x", 0))))
            y = max(0.0, min(1.0, float(bbox.get("y", 0))))
            width = max(0.0, min(1.0 - x, float(bbox.get("width", 0))))
            height = max(0.0, min(1.0 - y, float(bbox.get("height", 0))))
        except (TypeError, ValueError):
            continue
        if width <= 0 or height <= 0:
            continue
        normalized.append({
            "id": str(region.get("id") or f"region-{index + 1}"),
            "label": str(region.get("label") or f"Region {index + 1}"),
            "bbox": {
                "x": round(x, 4),
                "y": round(y, 4),
                "width": round(width, 4),
                "height": round(height, 4),
            },
        })
        if len(normalized) >= 8:
            break
    return normalized


def _format_finetune_regions_note(regions: Optional[List[object]]) -> str:
    normalized_regions = _normalize_finetune_regions(regions)
    if not normalized_regions:
        return ""
    lines = ["\nSelected edit regions (normalized to the full 16:9 slide image):"]
    for index, region in enumerate(normalized_regions, start=1):
        bbox = region["bbox"]
        label = region.get("label") or f"Region {index}"
        lines.append(
            "- Region "
            f"{index} ({label}): "
            f"x={bbox['x'] * 100:.1f}%, "
            f"y={bbox['y'] * 100:.1f}%, "
            f"width={bbox['width'] * 100:.1f}%, "
            f"height={bbox['height'] * 100:.1f}%."
        )
    lines.append(
        "These selected region(s) are supplied to the image edit request as an edit mask. "
        "When available, an additional marked copy of the slide shows the same region(s) with red outlines; "
        "use those red outlines only as a location guide and do not draw them in the result. "
        "Only the selected/masked region(s) may change. Keep everything outside the selected region(s) unchanged. "
        "Keep every same-looking or same-word item outside the selected region(s) unchanged unless the user explicitly asks otherwise."
    )
    return "\n".join(lines)


def _build_direct_finetune_prompt(
    slide: Slide,
    instruction: str,
    attachment_count: int = 0,
    project_visual_asset_count: int = 0,
    regions: Optional[List[object]] = None,
) -> str:
    """Compose a direct image-edit prompt for single-slide finetuning."""
    content = slide.content_json or {}
    text_content = content.get("text_content") if isinstance(content, dict) else {}
    headline = ""
    subhead = ""
    body = ""
    if isinstance(text_content, dict):
        headline = str(text_content.get("headline") or "")
        subhead = str(text_content.get("subhead") or "")
        raw_body = text_content.get("body") or ""
        if isinstance(raw_body, list):
            body = "\n".join(str(item.get("content") if isinstance(item, dict) else item) for item in raw_body[:5])
        else:
            body = str(raw_body)

    ref_count = attachment_count + project_visual_asset_count
    ref_note = (
        f"\nThere are {ref_count} additional reference image(s) after the current slide. "
        "These additional images are the user's current-turn references. "
        "When the user says 'this image', 'the image', 'this person', or similar wording, resolve it to the most prominent subject in these current-turn reference images, not to anything already printed on the slide, product packaging, labels, icons, or earlier page assets. "
        "Use them only when they help satisfy the user's request. Infer natural placement, scale, and cropping from the slide layout."
        if ref_count
        else ""
    )
    project_asset_note = (
        "\nProject visual asset rule: the last "
        f"{project_visual_asset_count} additional reference image(s) are protected project product/material assets. "
        "If the request asks to replace a product, bottle, packaging, or client asset, use these protected project asset images as the authoritative source. "
        "Do not keep, copy, or regenerate any conflicting brand/product already visible in the current slide image."
        if project_visual_asset_count
        else ""
    )
    region_note = _format_finetune_regions_note(regions)

    return f"""Edit the FIRST supplied image, which is the current PPT slide.

User edit request:
{instruction.strip()}
{region_note}

Rules:
- Apply the request immediately. Do not ask clarifying questions.
- Preserve all unmentioned text, layout, colors, typography, icons, charts, image areas, and background exactly as much as possible.
- If selected edit regions are provided, apply the request inside those regions first and keep the rest of the slide stable.
- If the user supplied extra reference images, they are the current-turn references. Integrate or borrow from them only where the request implies it.
- Do not copy people or illustrations printed on product packaging unless the user explicitly asks for the packaging artwork.
- Keep the final result as a polished 16:9 presentation slide, not a mockup.
- Keep text readable and do not invent new copy unless the user explicitly asked for it.
{ref_note}
{project_asset_note}

Slide text context for preservation:
Headline: {headline}
Subhead: {subhead}
Body: {body}
"""


@router.post("/{project_id}/slides/{slide_id}/finetune")
def finetune_slide(
    project_id: str,
    slide_id: str,
    body: FinetuneRequest,
    db: Session = Depends(get_db),
):
    """单页微调：存档当前图片 → 更新 prompt → 触发重新生成。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    if _active_run_for_project_action(project, db):
        raise HTTPException(status_code=409, detail="当前项目已有任务正在运行，请等待完成后再开始下一步")

    slide = db.query(Slide).filter(Slide.id == slide_id, Slide.project_id == project_id).first()
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    if slide.status == "generating":
        raise HTTPException(status_code=400, detail="该页面正在生成中，请等待完成后再微调")

    instruction = (body.instruction or "").strip()
    new_prompt = (body.new_prompt or "").strip()
    attachment_ids = body.attachment_ids or []
    finetune_regions = _normalize_finetune_regions(body.regions or [])
    if not instruction and not new_prompt:
        raise HTTPException(status_code=400, detail="instruction 不能为空")

    if instruction and (not slide.image_path or not os.path.exists(slide.image_path)):
        raise HTTPException(status_code=400, detail="当前页面没有可用于微调的图片")

    # 1. 存档当前图片为历史版本
    archived_version = _archive_current_image(slide, db)

    # 2. 更新 prompt。instruction 模式下，把当前页图片作为第一张参考图走 image edit。
    if instruction:
        if not archived_version or not os.path.exists(archived_version.image_path):
            raise HTTPException(status_code=400, detail="当前页面图片存档失败，无法微调")
        valid_attachment_ids = []
        if attachment_ids:
            refs = db.query(ReferenceImage).filter(
                ReferenceImage.project_id == project_id,
                ReferenceImage.slide_id == slide_id,
                ReferenceImage.id.in_(attachment_ids),
                ReferenceImage.role == "finetune_ref",
            ).all()
            found_ids = {str(ref.id) for ref in refs}
            valid_attachment_ids = [ref_id for ref_id in attachment_ids if ref_id in found_ids]

        project_visual_asset_ids = _project_visual_asset_ids_for_finetune(project, slide, instruction)
        slide.prompt_text = _build_direct_finetune_prompt(
            slide,
            instruction,
            len(valid_attachment_ids),
            len(project_visual_asset_ids),
            regions=finetune_regions,
        )
        visual_json = copy.deepcopy(slide.visual_json) or {}
        visual_json["finetune_base_image_path"] = archived_version.image_path
        visual_json["finetune_instruction"] = instruction
        visual_json["finetune_attachment_ids"] = valid_attachment_ids
        visual_json["finetune_visual_asset_ids"] = project_visual_asset_ids
        visual_json["finetune_regions"] = finetune_regions
        slide.visual_json = visual_json
        flag_modified(slide, "visual_json")
    else:
        slide.prompt_text = new_prompt
    slide.status = "generating"
    slide.error_msg = None
    ensure_generation_worker_ready()
    run = create_project_run(
        db,
        project_id,
        kind="finetune",
        stage=image_generation_run_stage(kind="finetune"),
        target_page_nums=[slide.page_num],
        total_count=1,
        message=image_generation_queued_message("finetune"),
    )
    db.commit()

    # 3. 触发异步生图
    _update_progress(project_id, {
        "target_page_nums": [slide.page_num],
        "target_count": 1,
    }, run.id)

    task = _enqueue_generation_task(db, project_id, [slide.page_num], run=run)

    return {
        "message": "Finetune started",
        "slide_id": slide_id,
        "page_num": slide.page_num,
        "mode": "direct_edit" if instruction else "prompt",
        "run": serialize_run(run),
    }


@router.post("/{project_id}/slides/{slide_id}/retry")
def retry_slide(
    project_id: str,
    slide_id: str,
    body: RetryRequest = Body(default_factory=RetryRequest),
    db: Session = Depends(get_db),
):
    """重试单页生成。支持重新生成 prompt 后再生图。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    if _active_run_for_project_action(project, db):
        raise HTTPException(status_code=409, detail="当前项目已有任务正在运行，请等待完成后再开始下一步")

    slide = db.query(Slide).filter(Slide.id == slide_id, Slide.project_id == project_id).first()
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    if slide.status == "generating":
        raise HTTPException(status_code=400, detail="该页面正在生成中，请勿重复重试")

    # 如果要求重新生成 prompt，基于最新 content/visual 重新生成
    if body.regenerate_prompt:
        _regenerate_slide_prompt(slide, project, db, user_feedback=(body.user_feedback or "").strip() or None)
    else:
        synced_exact_pages = _sync_exact_page_reference_overlays_for_generation(db, [slide])
        if synced_exact_pages:
            db.commit()
            raise HTTPException(
                status_code=400,
                detail="该页面有原样保留素材，请先重新生成画面描述以预留位置。",
            )

    project_slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )
    _raise_if_generation_inputs_not_current(project, project_slides, [slide])

    if not slide.prompt_text:
        raise HTTPException(status_code=400, detail="Slide has no prompt. Generate prompts first.")

    slide.status = "generating"
    slide.error_msg = None
    ensure_generation_worker_ready()
    run = create_project_run(
        db,
        project_id,
        kind="page_generation",
        stage=image_generation_run_stage(kind="page_generation"),
        target_page_nums=[slide.page_num],
        total_count=1,
        message=image_generation_queued_message("page_generation"),
    )
    db.commit()

    # 记录本次重试的目标页码和数量，供前端进度显示用
    _update_progress(project_id, {
        "target_page_nums": [slide.page_num],
        "target_count": 1,
    }, run.id)

    _enqueue_generation_task(
        db,
        project_id,
        [slide.page_num],
        run=run,
        prototype=project.status == "prototype_ready",
    )

    return {"message": "Retry started", "slide_id": slide_id, "page_num": slide.page_num, "run": serialize_run(run)}


@router.patch("/{project_id}/slides/content")
def update_slide_content(
    project_id: str,
    body: UpdateContentRequest,
    db: Session = Depends(get_db),
):
    """更新指定页码的 slide content_json。安全 merge：只更新正文相关字段。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    if body.slide_id:
        slide = db.query(Slide).filter(Slide.project_id == project_id, Slide.id == body.slide_id).first()
    else:
        slide = (
            db.query(Slide)
            .filter(Slide.project_id == project_id, Slide.page_num == body.page_num)
            .first()
        )
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    # 必须用深拷贝，否则 SQLAlchemy 检测不到 JSON 字段的变更
    existing = copy.deepcopy(slide.content_json) or {}
    new_content = _normalize_content_json_markdown(body.content_json)
    original_text = existing.get("text_content") if isinstance(existing.get("text_content"), dict) else {}
    original_text = copy.deepcopy(original_text)
    original_visual_suggestion = existing.get("visual_suggestion")
    original_visual_requirements = copy.deepcopy(existing.get("visual_requirements"))

    # 安全 merge：只替换 text_content 和 speaker_notes，保留其他字段
    if "text_content" in new_content:
        existing_text = existing.get("text_content") or {}
        if not isinstance(existing_text, dict):
            existing_text = {}
        existing_text.update(new_content["text_content"])
        existing["text_content"] = existing_text
    merged_text = existing.get("text_content") if isinstance(existing.get("text_content"), dict) else {}
    text_content_changed = any(
        not _json_equivalent(original_text.get(key), merged_text.get(key))
        for key in ("headline", "subhead", "body")
    )
    if "speaker_notes" in new_content:
        existing["speaker_notes"] = new_content["speaker_notes"]
    if "visual_suggestion" in new_content:
        if text_content_changed and _json_equivalent(new_content["visual_suggestion"], original_visual_suggestion):
            existing["visual_suggestion"] = ""
        else:
            existing["visual_suggestion"] = new_content["visual_suggestion"]
    elif text_content_changed:
        existing["visual_suggestion"] = ""
    if "visual_requirements" in new_content:
        if text_content_changed and _json_equivalent(new_content["visual_requirements"], original_visual_requirements):
            existing["visual_requirements"] = []
        else:
            existing["visual_requirements"] = new_content["visual_requirements"]
    elif text_content_changed:
        existing["visual_requirements"] = []
    if existing.get("visual_suggestion"):
        existing["visual_suggestion"] = _current_visual_suggestion(existing)
    if existing.get("visual_requirements"):
        existing["visual_requirements"] = _current_visual_requirements(existing)
    if "content_blocks" in new_content:
        existing["content_blocks"] = new_content["content_blocks"]

    existing = _sync_content_block_assets(db, project, slide, existing)
    text_content = existing.get("text_content") if isinstance(existing.get("text_content"), dict) else {}
    directive_extraction = extract_visual_directives(str(text_content.get("body") or ""))
    slide.content_json = existing
    _mark_slide_artifacts_stale(slide, content=True)
    db.commit()

    # 自动重分类：若未锁定，根据内容特征判断是否需要切换类型
    type_changed = False
    new_type = None
    if not slide.type_locked:
        detected = _auto_reclassify_page_type(slide.content_json, slide.type)
        if detected and detected != slide.type:
            slide.type = detected
            db.commit()
            type_changed = True
            new_type = detected

    return {
        "message": "Slide content updated",
        "page_num": slide.page_num,
        "slide_id": slide.id,
        "type_changed": type_changed,
        "type": slide.type,
        "new_type": new_type,
        "visual_directive_suggestions": directive_extraction.get("suggestions", []),
    }


@router.patch("/{project_id}/slides/visual")
def update_slide_visual(
    project_id: str,
    body: UpdateVisualRequest,
    db: Session = Depends(get_db),
):
    """更新指定页码的 slide visual_json。安全 merge：只更新传入的字段。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    if body.slide_id:
        slide = db.query(Slide).filter(Slide.project_id == project_id, Slide.id == body.slide_id).first()
    else:
        slide = (
            db.query(Slide)
            .filter(Slide.project_id == project_id, Slide.page_num == body.page_num)
            .first()
        )
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    existing = copy.deepcopy(slide.visual_json) or {}
    previous_visual = copy.deepcopy(existing)
    new_visual = body.visual_json

    # 安全 merge：替换传入的字段，保留其他字段
    allowed_fields = {
        "visual_evidence",
        "visual_description",
        "visual_summary",
        "design_notes",
        "layout",
        "visual_asset_ids",
        "visual_asset_usage",
        "asset_route_modes",
        "excluded_visual_asset_ids",
        "image_slots",
    }
    for key in allowed_fields:
        if key in new_visual:
            existing[key] = new_visual[key]
    if "visual_asset_ids" in new_visual:
        excluded_ids = _merge_excluded_assets_after_visual_edit(previous_visual, existing)
        if excluded_ids:
            existing["excluded_visual_asset_ids"] = excluded_ids
        else:
            existing.pop("excluded_visual_asset_ids", None)

    if "overlay_layers" in new_visual:
        raw_layers = new_visual.get("overlay_layers") or []
        if not isinstance(raw_layers, list):
            raise HTTPException(status_code=400, detail="overlay_layers must be a list")
        requested_asset_ids = {
            str(layer.get("asset_id"))
            for layer in raw_layers
            if isinstance(layer, dict) and layer.get("asset_id")
        }
        valid_assets = []
        if requested_asset_ids:
            candidates = db.query(ReferenceImage).filter(
                ReferenceImage.project_id == project_id,
                ReferenceImage.id.in_(requested_asset_ids),
            ).all()
            valid_assets = [
                asset for asset in candidates
                if asset.file_path and os.path.exists(_resolve_file_path(asset.file_path))
            ]
        valid_ids = {str(asset.id) for asset in valid_assets}
        missing = [asset_id for asset_id in requested_asset_ids if asset_id not in valid_ids]
        if missing:
            logger.warning(
                "update_slide_visual: requested overlay asset_ids not found or missing file: %s",
                missing,
            )
            raise HTTPException(status_code=400, detail=f"Invalid overlay asset ids: {', '.join(missing[:5])}")

        normalized_layers = normalize_overlay_layers(raw_layers, valid_asset_ids=valid_ids, strict_assets=True)
        overlay_asset_ids = {str(layer.get("asset_id")) for layer in normalized_layers if layer.get("asset_id")}
        for asset in valid_assets:
            if str(asset.id) in overlay_asset_ids:
                asset.process_mode = "original"
                analysis = asset.asset_analysis if isinstance(asset.asset_analysis, dict) else {}
                asset.asset_analysis = {
                    **analysis,
                    "exact_overlay": True,
                    "exact_overlay_reason": "用户选择原样保留，最终页面会保留素材比例和细节。",
                }
        existing["overlay_layers"] = normalized_layers

    project_slides = db.query(Slide).filter(Slide.project_id == project_id).order_by(Slide.page_num).all()
    slide.visual_json = with_artifact_meta(
        existing,
        kind="manual_visual_edit",
        dependencies=dependency_signature(project, project_slides),
    )
    _invalidate_visual_plan_dependent_outputs(project, [slide])
    db.commit()

    return {
        "message": "Slide visual updated",
        "page_num": slide.page_num,
        "slide_id": slide.id,
    }


@router.patch("/{project_id}/slides/type")
def update_slide_type(
    project_id: str,
    body: UpdateTypeRequest,
    db: Session = Depends(get_db),
):
    """手动更新 slide 类型并锁定，后续自动重分类不再覆盖。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    if body.slide_id:
        slide = db.query(Slide).filter(Slide.project_id == project_id, Slide.id == body.slide_id).first()
    else:
        slide = (
            db.query(Slide)
            .filter(Slide.project_id == project_id, Slide.page_num == body.page_num)
            .first()
        )
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    slide.type = body.type
    slide.type_locked = True
    _mark_slide_artifacts_stale(slide, content=True)
    db.commit()

    return {
        "message": "Slide type updated and locked",
        "page_num": slide.page_num,
        "slide_id": slide.id,
        "type": slide.type,
        "type_locked": slide.type_locked,
    }


@router.delete("/{project_id}/slides/{slide_id}")
def delete_slide(
    project_id: str,
    slide_id: str,
    db: Session = Depends(get_db),
):
    """删除指定 slide，并自动压缩后续页码。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slide = db.query(Slide).filter(Slide.id == slide_id, Slide.project_id == project_id).first()
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    deleted_page_num = slide.page_num
    db.delete(slide)

    # 后续页面页码减 1
    later_slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id, Slide.page_num > deleted_page_num)
        .order_by(Slide.page_num)
        .all()
    )
    for s in later_slides:
        s.page_num -= 1
        # 同步更新 content_json 中的 page_num（必须用深拷贝，否则 SQLAlchemy 检测不到变更）
        if s.content_json and isinstance(s.content_json, dict):
            updated = copy.deepcopy(s.content_json)
            updated["page_num"] = s.page_num
            s.content_json = updated
        # 同步更新 visual_json 中的 page_num
        visual_json = getattr(s, "visual_json", None)
        if visual_json and isinstance(visual_json, dict):
            updated = copy.deepcopy(visual_json)
            updated["page_num"] = s.page_num
            s.visual_json = updated

    _invalidate_content_dependent_outputs(project)
    db.commit()
    return {"message": "Slide deleted", "slide_id": slide_id, "deleted_page_num": deleted_page_num}


@router.post("/{project_id}/slides")
def create_slide(
    project_id: str,
    body: CreateSlideRequest,
    db: Session = Depends(get_db),
):
    """在指定页码位置插入新 slide，后续页面页码自动后移。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )

    insert_page_num = body.page_num

    # 如果插入位置超出末尾，直接追加到最后
    if not slides:
        insert_page_num = 1
    elif insert_page_num > len(slides) + 1:
        insert_page_num = len(slides) + 1

    # 后续页面页码加 1（从后往前遍历，避免重复更新）
    for s in reversed(slides):
        if s.page_num >= insert_page_num:
            s.page_num += 1
            if s.content_json and isinstance(s.content_json, dict):
                updated = copy.deepcopy(s.content_json)
                updated["page_num"] = s.page_num
                s.content_json = updated
            visual_json = getattr(s, "visual_json", None)
            if visual_json and isinstance(visual_json, dict):
                updated = copy.deepcopy(visual_json)
                updated["page_num"] = s.page_num
                s.visual_json = updated

    new_content = _normalize_content_json_markdown(copy.deepcopy(body.content_json))
    new_content.setdefault("page_num", insert_page_num)
    new_slide = Slide(
        project_id=project_id,
        page_num=insert_page_num,
        type=new_content.get("type", "content"),
        content_json=new_content,
    )
    db.add(new_slide)
    db.flush()
    new_slide.content_json = _sync_content_block_assets(db, project, new_slide, new_content)
    _invalidate_content_dependent_outputs(project)
    db.commit()
    db.refresh(new_slide)

    return {
        "message": "Slide created",
        "slide_id": new_slide.id,
        "page_num": new_slide.page_num,
    }


@router.post("/{project_id}/reorder")
def reorder_slides(
    project_id: str,
    body: ReorderRequest,
    db: Session = Depends(get_db),
):
    """根据传入的 page_nums 顺序重新排序 slides。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    slides = (
        db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )

    # 建立 page_num -> slide 映射
    slide_by_page = {s.page_num: s for s in slides}

    if len(body.page_nums) != len(slides):
        raise HTTPException(status_code=400, detail="页码数量与项目不符")
    if len(set(body.page_nums)) != len(body.page_nums):
        raise HTTPException(status_code=400, detail="页码列表中存在重复")

    for new_idx, old_page_num in enumerate(body.page_nums, start=1):
        slide = slide_by_page.get(old_page_num)
        if not slide:
            raise HTTPException(status_code=400, detail=f"页码 {old_page_num} 不存在")
        slide.page_num = new_idx
        if slide.content_json and isinstance(slide.content_json, dict):
            updated = copy.deepcopy(slide.content_json)
            updated["page_num"] = new_idx
            slide.content_json = updated

    _invalidate_content_dependent_outputs(project)
    db.commit()
    return {"message": "Slides reordered", "new_order": body.page_nums}


TEMPLATE_RECOMMENDATION_KEYS = ("cover", "toc", "section", "content", "data", "quote", "ending")


def _file_sha1(path: str | None) -> str | None:
    if not path or not os.path.exists(path):
        return None
    sha1 = hashlib.sha1()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            sha1.update(chunk)
    return sha1.hexdigest()


def _template_page_analysis(page: dict, source_filename: str, total_pages: int) -> dict:
    return {
        "analysis_type": "template_page",
        "analysis_status": "completed",
        "source_document": source_filename,
        "template_page_num": page.get("page_num"),
        "template_total_pages": total_pages,
        "template_category": page.get("category") or "content",
        "category_confidence": page.get("category_confidence") or 0.6,
        "source_kind": page.get("source_kind") or "template",
        "template_application_strength": "standard",
        "preview_file_path": page.get("file_path"),
        "layout_file_path": page.get("layout_file_path") or page.get("file_path"),
        "text_density": page.get("text_density") or 0,
        "logo_removed": bool(page.get("logo_removed")),
    }


def _template_recommendation_payload(page: dict | None) -> dict | None:
    if not isinstance(page, dict):
        return None
    return {
        "page_num": page.get("page_num"),
        "file_path": page.get("layout_file_path") or page.get("file_path"),
        "preview_file_path": page.get("file_path"),
        "layout_file_path": page.get("layout_file_path") or page.get("file_path"),
        "category": page.get("category") or "content",
        "category_confidence": page.get("category_confidence") or 0.6,
        "source_kind": page.get("source_kind") or "template",
        "application_strength": "standard",
        "logo_removed": bool(page.get("logo_removed")),
    }


def _template_page_num_from_ref(ref: ReferenceImage) -> int:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    try:
        page_num = int(analysis.get("template_page_num") or 0)
    except (TypeError, ValueError):
        page_num = 0
    if page_num > 0:
        return page_num
    match = re.search(r"page_(\d+)", os.path.basename(ref.file_path or ""))
    return int(match.group(1)) if match else 0


def _template_page_payload_from_ref(project_id: str, ref: ReferenceImage, total: int | None = None) -> dict:
    analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
    page_num = _template_page_num_from_ref(ref)
    layout_file_path = analysis.get("layout_file_path") or ref.file_path
    layout_name = os.path.basename(layout_file_path) if layout_file_path else os.path.basename(ref.file_path)
    preview_name = os.path.basename(ref.file_path)
    return {
        "id": ref.id,
        "page_num": page_num,
        "url": f"/uploads/{project_id}/templates/{preview_name}",
        "layout_url": f"/uploads/{project_id}/templates/{layout_name}",
        "category": analysis.get("template_category") or _infer_template_page_category_safe(page_num, total or 0),
        "category_confidence": analysis.get("category_confidence") or 0.6,
        "source_kind": analysis.get("source_kind") or "template",
        "application_strength": analysis.get("template_application_strength") or "standard",
        "logo_removed": bool(analysis.get("logo_removed")),
        "file_exists": os.path.exists(ref.file_path),
    }


def _infer_template_page_category_safe(page_num: int, total: int) -> str:
    from app.services.template_extractor import _infer_page_category

    return _infer_page_category(page_num, total)


def _logo_anchor_from_bounds(bounds: dict | None) -> str:
    if not isinstance(bounds, dict):
        return DEFAULT_LOGO_ANCHOR
    try:
        cx = float(bounds.get("left") or 0) + float(bounds.get("width") or 0) / 2
        cy = float(bounds.get("top") or 0) + float(bounds.get("height") or 0) / 2
    except (TypeError, ValueError):
        return DEFAULT_LOGO_ANCHOR
    horizontal = "left" if cx < 0.5 else "right"
    vertical = "top" if cy < 0.5 else "bottom"
    return f"{vertical}-{horizontal}"


def _logo_regions_from_pptx_assets(assets: list) -> dict[int, list[dict]]:
    regions_by_page: dict[int, list[dict]] = {}
    for asset in assets:
        if getattr(asset, "role", None) != "logo":
            continue
        metadata = getattr(asset, "metadata", {}) if isinstance(getattr(asset, "metadata", {}), dict) else {}
        bounds = metadata.get("shape_bounds")
        if not isinstance(bounds, dict) or not bounds:
            continue
        repeated = metadata.get("pptx_repeated_page_nums") or [getattr(asset, "source_page_num", None)]
        for page_num in repeated:
            try:
                normalized_page = int(page_num)
            except (TypeError, ValueError):
                continue
            regions_by_page.setdefault(normalized_page, []).append({
                "bbox_norm": bounds,
                "anchor": _logo_anchor_from_bounds(bounds),
            })
    return regions_by_page


def _project_has_confirmed_manual_logo(project: Project, db: Session | None = None) -> bool:
    refs = (
        db.query(ReferenceImage).filter(ReferenceImage.project_id == project.id, ReferenceImage.role == "logo").all()
        if db is not None
        else (project.reference_images or [])
    )
    for ref in refs:
        if ref.role != "logo" or ref.slide_id:
            continue
        analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
        if analysis.get("source_document") or analysis.get("template_logo_source"):
            continue
        if is_logo_confirmed(ref):
            return True
    return False


def _template_logo_analysis(asset, source_filename: str, *, needs_review: bool) -> dict:
    metadata = getattr(asset, "metadata", {}) if not isinstance(asset, dict) else asset.get("metadata", {})
    metadata = metadata if isinstance(metadata, dict) else {}
    file_path = getattr(asset, "file_path", None) if not isinstance(asset, dict) else asset.get("file_path")
    digest = metadata.get("pptx_image_sha1") or _file_sha1(file_path)
    classification = (
        getattr(asset, "classification", None)
        if not isinstance(asset, dict)
        else asset.get("classification")
    ) or metadata.get("classification") or "logo"
    return {
        **metadata,
        "analysis_status": "completed",
        "analysis_type": "logo",
        "source_document": source_filename,
        "template_logo_source": "layout_template",
        "classification": classification,
        "pptx_image_sha1": digest,
        "review_status": "needs_review" if needs_review else "auto_confirmed",
        "needs_user_review": bool(needs_review),
        "confidence_score": 0.62 if needs_review else 0.9,
        "review_reason": (
            "模板中识别出的 Logo 候选，已有项目 Logo，请确认是否替换或合并使用。"
            if needs_review
            else "从版式模板中识别出的 Logo，已作为项目级品牌标识使用。"
        ),
    }


def _attach_template_logo_assets(project: Project, assets: list, source_filename: str, db: Session) -> int:
    if not assets:
        return 0
    existing_logo_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.slide_id.is_(None),
        ReferenceImage.role == "logo",
    ).all()
    existing_hashes = {
        str((ref.asset_analysis or {}).get("pptx_image_sha1"))
        for ref in existing_logo_refs
        if isinstance(ref.asset_analysis, dict) and ref.asset_analysis.get("pptx_image_sha1")
    }
    has_manual_logo = _project_has_confirmed_manual_logo(project, db)
    auto_confirmed_template_logo_attached = any(
        is_logo_confirmed(ref)
        and isinstance(ref.asset_analysis, dict)
        and ref.asset_analysis.get("template_logo_source") == "layout_template"
        for ref in existing_logo_refs
    )
    attached = 0
    for asset in assets:
        file_path = getattr(asset, "file_path", None) if not isinstance(asset, dict) else asset.get("file_path")
        if not file_path or not os.path.exists(file_path):
            continue
        metadata = getattr(asset, "metadata", {}) if not isinstance(asset, dict) else asset.get("metadata", {})
        metadata = metadata if isinstance(metadata, dict) else {}
        digest = metadata.get("pptx_image_sha1") or _file_sha1(file_path)
        if digest and digest in existing_hashes:
            continue
        needs_review = bool(has_manual_logo or auto_confirmed_template_logo_attached)
        analysis = _template_logo_analysis(asset, source_filename, needs_review=needs_review)
        if needs_review and not has_manual_logo:
            analysis["review_reason"] = "同一模板中识别出的额外 Logo 候选，需用户确认后才会合并为项目 Logo。"
        bounds = metadata.get("shape_bounds") if isinstance(metadata.get("shape_bounds"), dict) else None
        logo_anchor = metadata.get("logo_anchor") or _logo_anchor_from_bounds(bounds)
        db.add(ReferenceImage(
            project_id=project.id,
            slide_id=None,
            file_path=file_path,
            role="logo",
            process_mode="original",
            asset_name=(getattr(asset, "asset_name", None) if not isinstance(asset, dict) else asset.get("asset_name")),
            usage_note=(getattr(asset, "usage_note", None) if not isinstance(asset, dict) else asset.get("usage_note")),
            asset_analysis=analysis,
            logo_anchor=normalize_logo_anchor(logo_anchor),
        ))
        if digest:
            existing_hashes.add(str(digest))
        if not needs_review:
            auto_confirmed_template_logo_attached = True
        attached += 1
    return attached


def _delete_old_template_logo_assets(project_id: str, db: Session) -> int:
    old_logo_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project_id,
        ReferenceImage.slide_id.is_(None),
        ReferenceImage.role == "logo",
    ).all()
    deleted = 0
    for ref in old_logo_refs:
        analysis = ref.asset_analysis if isinstance(ref.asset_analysis, dict) else {}
        if analysis.get("template_logo_source") != "layout_template":
            continue
        db.delete(ref)
        deleted += 1
    if deleted:
        db.flush()
    return deleted


TEMPLATE_EXTRACTION_STATUS_FILE = "template_extract_status.json"


def _template_status_path(project_id: str) -> str:
    project_dir = os.path.join(settings.UPLOAD_DIR, project_id)
    os.makedirs(project_dir, exist_ok=True)
    return os.path.join(project_dir, TEMPLATE_EXTRACTION_STATUS_FILE)


def _write_template_status(project_id: str, status: str, **payload) -> None:
    data = {
        "status": status,
        "updated_at": time.time(),
        **payload,
    }
    try:
        with open(_template_status_path(project_id), "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False)
    except OSError as exc:
        logger.warning("写入模板提取状态失败: project=%s error=%s", project_id, exc)


def _read_template_status(project_id: str) -> dict:
    path = _template_status_path(project_id)
    if not os.path.exists(path):
        return {"status": "not_found"}
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, dict) else {"status": "unknown"}
    except (OSError, json.JSONDecodeError):
        return {"status": "unknown"}


def _is_current_template_job(project_id: str, job_id: str | None) -> bool:
    if not job_id:
        return True
    return _read_template_status(project_id).get("job_id") == job_id


def _extract_template_into_project(
    project: Project,
    file_path: str,
    safe_name: str,
    ext: str,
    db: Session,
    *,
    job_id: str | None = None,
) -> dict:
    pptx_logo_assets = []
    logo_regions_by_page: dict[int, list[dict]] = {}
    if ext == ".pptx":
        try:
            from app.services.pptx_asset_extractor import extract_pptx_image_assets

            with open(file_path, "rb") as f:
                file_bytes = f.read()
            template_assets_dir = os.path.join(settings.UPLOAD_DIR, project.id, "template_assets")
            os.makedirs(template_assets_dir, exist_ok=True)
            extracted_assets = extract_pptx_image_assets(
                file_bytes=file_bytes,
                source_filename=safe_name,
                output_dir=template_assets_dir,
                max_assets_per_slide=1,
                max_total_assets=24,
            )
            pptx_logo_assets = [asset for asset in extracted_assets if asset.role == "logo"]
            logo_regions_by_page = _logo_regions_from_pptx_assets(pptx_logo_assets)
        except Exception as exc:
            logger.warning("模板 PPTX Logo 提取失败，继续提取模板页: project=%s file=%s error=%s", project.id, safe_name, exc)

    from app.services.template_extractor import extract_template_package, promote_template_package

    package = extract_template_package(
        file_path,
        project.id,
        settings.UPLOAD_DIR,
        logo_regions_by_page=logo_regions_by_page,
        source_filename=safe_name,
        finalize=False,
    )

    if not _is_current_template_job(project.id, job_id):
        work_dir = package.get("work_dir") if isinstance(package, dict) else None
        if work_dir and os.path.exists(work_dir):
            shutil.rmtree(work_dir, ignore_errors=True)
        db.rollback()
        return {
            "status": "stale",
            "message": "Template extraction skipped because a newer upload exists",
            "filename": safe_name,
        }

    package = promote_template_package(package)
    pages = package["pages"]

    old_template_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "template",
    ).all()
    for old_ref in old_template_refs:
        db.delete(old_ref)

    recommendations = package.get("recommendations") or {}
    total_pages = len(pages)
    for page in pages:
        ref = ReferenceImage(
            project_id=project.id,
            file_path=page["file_path"],
            role="template",
            process_mode="blend",
            asset_analysis=_template_page_analysis(page, safe_name, total_pages),
        )
        db.add(ref)

    rendered_logo_assets = package.get("rendered_logo_candidates") or []
    _delete_old_template_logo_assets(project.id, db)
    attached_logo_count = _attach_template_logo_assets(
        project,
        [*pptx_logo_assets, *rendered_logo_assets],
        safe_name,
        db,
    )

    project.selected_template_recommendations = {
        key: _template_recommendation_payload(recommendations.get(key))
        for key in TEMPLATE_RECOMMENDATION_KEYS
    }
    _invalidate_style_dependent_outputs(project)
    db.commit()

    return {
        "message": "Template extracted",
        "filename": safe_name,
        "total_pages": len(pages),
        "status": "completed",
        "document_kind": package.get("document_kind") or "template",
        "extracted_logos": attached_logo_count,
        "pages": [
            {
                "page_num": p["page_num"],
                "url": p["url"],
                "layout_url": p.get("layout_url"),
                "category": p.get("category", "content"),
                "category_confidence": p.get("category_confidence", 0.6),
                "source_kind": p.get("source_kind", "template"),
                "logo_removed": bool(p.get("logo_removed")),
            }
            for p in pages
        ],
        "recommendations": {
            key: (
                {
                    "page_num": value["page_num"],
                    "url": value["url"],
                    "layout_url": value.get("layout_url"),
                    "category": value.get("category", "content"),
                    "category_confidence": value.get("category_confidence", 0.6),
                    "source_kind": value.get("source_kind", "template"),
                    "logo_removed": bool(value.get("logo_removed")),
                }
                if isinstance(value := recommendations.get(key), dict)
                else None
            )
            for key in TEMPLATE_RECOMMENDATION_KEYS
        },
    }


def _process_template_upload_bg(project_id: str, file_path: str, safe_name: str, ext: str, job_id: str) -> None:
    _write_template_status(project_id, "running", filename=safe_name, job_id=job_id)
    db = SessionLocal()
    try:
        project = db.query(Project).filter(Project.id == project_id).first()
        if not project:
            raise RuntimeError("Project not found")
        result = _extract_template_into_project(project, file_path, safe_name, ext, db, job_id=job_id)
        if result.get("status") == "stale":
            return
        _write_template_status(
            project_id,
            "completed",
            filename=safe_name,
            job_id=job_id,
            total_pages=result.get("total_pages", 0),
            document_kind=result.get("document_kind") or "template",
            extracted_logos=result.get("extracted_logos", 0),
        )
    except Exception as exc:
        db.rollback()
        logger.warning("模板提取后台任务失败: project=%s file=%s error=%s", project_id, safe_name, exc)
        if _is_current_template_job(project_id, job_id):
            _write_template_status(project_id, "failed", filename=safe_name, job_id=job_id, error=str(exc)[:500])
    finally:
        db.close()


@router.post("/{project_id}/extract-template")
def extract_template(
    project_id: str,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
):
    """上传 PPT/PDF 并提取每页缩略图作为模板参考。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    # 安全检查：拒绝路径遍历攻击
    safe_name = file.filename.replace("\\", "/").split("/")[-1]
    if ".." in safe_name or not safe_name or safe_name.startswith("."):
        raise HTTPException(status_code=400, detail="非法文件名")
    ext = os.path.splitext(safe_name.lower())[1]
    if ext not in {".ppt", ".pptx", ".pdf"}:
        raise HTTPException(status_code=400, detail="版式模板仅支持 PPT、PPTX 或 PDF")

    file_bytes = file.file.read()
    if len(file_bytes) == 0:
        raise HTTPException(status_code=400, detail="上传的文件为空")

    # 保存上传文件
    upload_dir = os.path.join(settings.UPLOAD_DIR, project_id)
    os.makedirs(upload_dir, exist_ok=True)
    file_path = os.path.join(upload_dir, f"template_{safe_name}")
    with open(file_path, "wb") as buffer:
        buffer.write(file_bytes)

    job_id = hashlib.sha1(f"{project_id}:{safe_name}:{len(file_bytes)}:{time.time()}".encode("utf-8")).hexdigest()[:16]
    _write_template_status(project_id, "queued", filename=safe_name, job_id=job_id)
    background_tasks.add_task(_process_template_upload_bg, project_id, file_path, safe_name, ext, job_id)
    return {
        "message": "Template extraction queued",
        "status": "processing",
        "filename": safe_name,
        "job_id": job_id,
    }


@router.get("/{project_id}/template-status")
def get_template_status(
    project_id: str,
    db: Session = Depends(get_db),
):
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    return _read_template_status(project_id)


@router.get("/{project_id}/template-pages")
def list_template_pages(
    project_id: str,
    db: Session = Depends(get_db),
):
    """列出已提取的模板页面。"""
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project_id,
        ReferenceImage.role == "template",
    ).all()
    refs.sort(key=_template_page_num_from_ref)

    total = len(refs)

    return [
        _template_page_payload_from_ref(project_id, ref, total)
        for ref in refs
    ]


# ========== 单页微调：版本管理 ==========

MAX_VERSIONS_PER_SLIDE = 10


def _archive_current_image(slide: Slide, db: Session) -> SlideVersion | None:
    """将 slide 当前图片存档为一个历史版本。若版本数超限则删除最老的。"""
    if not slide.image_path or not os.path.exists(slide.image_path):
        return None

    # 计算下一个版本号
    max_ver = db.query(SlideVersion).filter(
        SlideVersion.slide_id == slide.id
    ).order_by(SlideVersion.version_number.desc()).first()
    next_ver = (max_ver.version_number + 1) if max_ver else 1

    version_dir = os.path.join(settings.OUTPUT_DIR or "./outputs", slide.project_id, "versions")
    os.makedirs(version_dir, exist_ok=True)
    version_path = os.path.join(version_dir, f"slide_{slide.page_num:02d}_v{next_ver}.png")
    shutil.copy2(slide.image_path, version_path)

    version = SlideVersion(
        slide_id=slide.id,
        project_id=slide.project_id,
        image_path=version_path,
        prompt_text=slide.prompt_text,
        version_number=next_ver,
    )
    db.add(version)

    # 清理超限版本
    all_versions = db.query(SlideVersion).filter(
        SlideVersion.slide_id == slide.id
    ).order_by(SlideVersion.version_number.asc()).all()
    if len(all_versions) > MAX_VERSIONS_PER_SLIDE:
        for v in all_versions[:len(all_versions) - MAX_VERSIONS_PER_SLIDE]:
            # Older versions used to point at the live slide image. Never delete
            # the file currently shown by the slide while pruning version rows.
            if v.image_path != slide.image_path and v.image_path != version_path and os.path.exists(v.image_path):
                try:
                    os.remove(v.image_path)
                except Exception:
                    pass
            db.delete(v)

    db.flush()
    return version


def _output_url_for_path(path: str) -> str:
    output_root = os.path.abspath(settings.OUTPUT_DIR or "./outputs")
    abs_path = os.path.abspath(path)
    try:
        rel_path = os.path.relpath(abs_path, output_root)
    except ValueError:
        rel_path = os.path.basename(path)
    return "/outputs/" + rel_path.replace(os.sep, "/")


@router.get("/{project_id}/slides/{slide_id}/versions")
def list_slide_versions(
    project_id: str,
    slide_id: str,
    db: Session = Depends(get_db),
):
    """返回某页的所有历史版本（按版本号倒序）。"""
    slide = db.query(Slide).filter(Slide.id == slide_id, Slide.project_id == project_id).first()
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    versions = db.query(SlideVersion).filter(
        SlideVersion.slide_id == slide_id
    ).order_by(SlideVersion.version_number.desc()).all()

    return [
        {
            "id": v.id,
            "version_number": v.version_number,
            "image_url": _output_url_for_path(v.image_path),
            "prompt_text": v.prompt_text,
            "created_at": v.created_at.isoformat() if v.created_at else None,
        }
        for v in versions
    ]


@router.delete("/{project_id}/slides/{slide_id}/versions/{version_id}")
def delete_slide_version(
    project_id: str,
    slide_id: str,
    version_id: str,
    db: Session = Depends(get_db),
):
    """删除某个历史版本（含物理文件）。"""
    version = db.query(SlideVersion).filter(
        SlideVersion.id == version_id,
        SlideVersion.slide_id == slide_id,
        SlideVersion.project_id == project_id,
    ).first()
    if not version:
        raise HTTPException(status_code=404, detail="Version not found")

    # 删除物理文件
    if os.path.exists(version.image_path):
        try:
            os.remove(version.image_path)
        except Exception as e:
            logger.warning(f"Failed to delete version file {version.image_path}: {e}")

    db.delete(version)
    db.commit()
    return {"message": "Version deleted", "version_id": version_id}


class RestoreVersionBody(BaseModel):
    pass


@router.post("/{project_id}/slides/{slide_id}/versions/{version_id}/restore")
def restore_slide_version(
    project_id: str,
    slide_id: str,
    version_id: str,
    db: Session = Depends(get_db),
):
    """将某历史版本恢复为当前（新旧 image_path 互换）。"""
    slide = db.query(Slide).filter(Slide.id == slide_id, Slide.project_id == project_id).first()
    if not slide:
        raise HTTPException(status_code=404, detail="Slide not found")

    version = db.query(SlideVersion).filter(
        SlideVersion.id == version_id,
        SlideVersion.slide_id == slide_id,
    ).first()
    if not version:
        raise HTTPException(status_code=404, detail="Version not found")

    if not os.path.exists(version.image_path):
        raise HTTPException(status_code=400, detail="Version image file no longer exists")

    # 交换当前图片和版本图片的路径
    current_path = slide.image_path
    version_path = version.image_path

    slide.image_path = version_path
    slide.prompt_text = version.prompt_text
    version.image_path = current_path

    db.commit()
    return {
        "message": "Version restored",
        "slide_id": slide_id,
        "new_image_url": _output_url_for_path(version_path),
    }
