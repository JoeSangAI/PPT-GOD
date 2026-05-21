import logging
import os
from typing import Dict, List

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches

from app.services.logo_assets import prepare_logo_lockup_image, prepare_logo_symbol_image
from app.services.logo_overlay_layout import logo_geometry_from_resolved_box, resolve_logo_overlay_box, resolve_logo_render_policy
from app.services.logo_policy import LOGO_HEIGHT_RATIOS, LOGO_WIDTH_RATIOS, logo_policy_for_page, normalize_logo_placement, should_show_logo
from app.services.overlay_layers import contained_picture_box, enabled_overlay_layers, overlay_box
from app.services.text_region_detector import compute_safe_overlay_box

logger = logging.getLogger(__name__)


def _resolve_file_path(file_path: str) -> str:
    """安全解析文件路径，兼容从不同工作目录启动的情况。"""
    if not file_path:
        return file_path
    if os.path.exists(file_path):
        return file_path
    backend_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    candidate = os.path.join(backend_dir, file_path)
    if os.path.exists(candidate):
        return candidate
    abs_path = os.path.abspath(file_path)
    if os.path.exists(abs_path):
        return abs_path
    return file_path


def _logo_geometry(
    prs: Presentation,
    logo_path: str,
    slide_type: str,
    placement: str,
    scale: str = "small",
    slide_image_path: str | None = None,
    resolved_box: Dict | None = None,
):
    resolved_geometry = logo_geometry_from_resolved_box(resolved_box, prs.slide_width, prs.slide_height)
    if resolved_geometry:
        return resolved_geometry

    resolved_box = resolve_logo_overlay_box(slide_image_path, logo_path, slide_type, placement, scale)
    resolved_geometry = logo_geometry_from_resolved_box(resolved_box, prs.slide_width, prs.slide_height)
    if resolved_geometry:
        return resolved_geometry

    placement = normalize_logo_placement(placement)
    margin = Inches(0.28)
    is_large = scale == "large" or slide_type == "cover"
    size_key = "large" if is_large else "small"
    max_width = int(prs.slide_width * LOGO_WIDTH_RATIOS[size_key])
    max_height = int(prs.slide_height * LOGO_HEIGHT_RATIOS[size_key])

    with Image.open(logo_path) as img:
        w_px, h_px = img.size
    ratio = h_px / max(w_px, 1)
    width = max_width
    height = int(width * ratio)
    if height > max_height:
        height = max_height
        width = int(height / max(ratio, 0.01))

    if placement == "center":
        left = int((prs.slide_width - width) / 2)
        top = int((prs.slide_height - height) / 2)
    elif placement == "lower-center":
        left = int((prs.slide_width - width) / 2)
        top = int(prs.slide_height * 0.68)
    elif placement == "title-block-center":
        left = int(prs.slide_width * 0.68 - width / 2)
        top = int(prs.slide_height * 0.70)
    else:
        left = margin if placement.endswith("left") else prs.slide_width - margin - width
        top = margin if placement.startswith("top") else prs.slide_height - margin - height
    return left, top, width, height


def assemble_pptx(
    slide_images: List[Dict],
    output_path: str,
    logo_config: Dict | None = None,
    overlay_assets: Dict[str, Dict] | None = None,
) -> str:
    """
    将生成的图片组装为 PPTX。

    slide_images: [{"page_num": int, "image_path": str, "speaker_notes": str}]
    output_path: 输出文件完整路径
    """
    prs = Presentation()
    # 16:9 (1792x1024)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 按 page_num 排序
    sorted_slides = sorted(slide_images, key=lambda x: x["page_num"])

    # 安全获取空白布局：优先索引 6，不足时回退到最后一个
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    logo_paths = []
    if logo_config:
        raw_paths = logo_config.get("file_paths")
        if isinstance(raw_paths, list):
            logo_paths.extend(str(path) for path in raw_paths if path)
        elif logo_config.get("file_path"):
            logo_paths.append(str(logo_config["file_path"]))
    logo_path_for_overlay = prepare_logo_lockup_image(logo_paths) if logo_paths else None
    logo_symbol_path_for_overlay = prepare_logo_symbol_image(logo_paths[0]) if len(logo_paths) == 1 else None

    for slide_data in sorted_slides:
        slide = prs.slides.add_slide(blank_layout)

        # 插入背景图（铺满整张幻灯片）
        img_path = slide_data.get("image_path")
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(
                img_path,
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

        overlay_layers = enabled_overlay_layers(slide_data.get("visual_json") or {})
        if overlay_layers:
            logger.info(
                "Assembler: page %s has %s overlay layers, overlay_assets keys=%s",
                slide_data.get("page_num"),
                len(overlay_layers),
                list((overlay_assets or {}).keys()),
            )
        for layer in overlay_layers:
            asset = (overlay_assets or {}).get(str(layer.get("asset_id")))
            asset_path = asset.get("file_path") if isinstance(asset, dict) else None
            resolved_asset_path = _resolve_file_path(asset_path)
            if not resolved_asset_path or not os.path.exists(resolved_asset_path):
                logger.warning(
                    "Assembler: overlay asset missing for page %s asset=%s asset_exists=%s overlay_assets_count=%s path=%s resolved=%s",
                    slide_data.get("page_num"),
                    layer.get("asset_id"),
                    bool(asset),
                    len(overlay_assets or {}),
                    asset_path,
                    resolved_asset_path,
                )
                continue
            left, top, width, height = overlay_box(prs, str(layer.get("preset") or "right-card"))
            # 文字避让：exact_cutout 模式下，如果已预计算文字区域，则调整位置
            if layer.get("mode") == "exact_cutout":
                text_regions = slide_data.get("text_regions")
                if text_regions:
                    left, top, width, height = compute_safe_overlay_box(
                        left, top, width, height,
                        text_regions,
                        float(prs.slide_width),
                        float(prs.slide_height),
                    )
            if layer.get("mode") == "exact_card":
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                card.line.color.rgb = RGBColor(218, 226, 235)
                card.line.width = Inches(0.008)
                inset = int(min(width, height) * 0.055)
                pic_left, pic_top, pic_width, pic_height = contained_picture_box(
                    resolved_asset_path,
                    left + inset,
                    top + inset,
                    max(1, width - inset * 2),
                    max(1, height - inset * 2),
                    valign=str(layer.get("valign") or "bottom"),
                )
            else:
                pic_left, pic_top, pic_width, pic_height = contained_picture_box(
                    resolved_asset_path,
                    left,
                    top,
                    width,
                    height,
                    valign=str(layer.get("valign") or "bottom"),
                )
            slide.shapes.add_picture(
                resolved_asset_path,
                left=pic_left,
                top=pic_top,
                width=pic_width,
                height=pic_height,
            )

        if (
            logo_config
            and logo_path_for_overlay
            and os.path.exists(logo_path_for_overlay)
            and should_show_logo(slide_data)
        ):
            raw_policy = (slide_data.get("visual_json") or {}).get("logo_policy") or {}
            policy = logo_policy_for_page(slide_data)
            render_policy_input = {**raw_policy, **policy}
            if "render_variant" not in policy:
                render_policy_input.pop("render_variant", None)
            render_policy = resolve_logo_render_policy(
                img_path,
                logo_path_for_overlay,
                logo_symbol_path_for_overlay,
                str(slide_data.get("type") or "content").lower(),
                policy.get("placement") or logo_config.get("anchor") or "top-right",
                policy.get("scale") or "small",
                render_policy_input,
            )
            if render_policy.get("show_logo") is False:
                continue
            logo_path_to_render = logo_path_for_overlay
            left, top, width, height = _logo_geometry(
                prs,
                logo_path_to_render,
                str(slide_data.get("type") or "content").lower(),
                render_policy.get("placement") or policy.get("placement") or logo_config.get("anchor") or "top-right",
                render_policy.get("scale") or policy.get("scale") or "small",
                slide_image_path=img_path,
                resolved_box=render_policy.get("resolved_overlay_box") or raw_policy.get("resolved_overlay_box"),
            )
            slide.shapes.add_picture(
                logo_path_to_render,
                left=left,
                top=top,
                width=width,
                height=height,
            )

        # 写入 Speaker Notes
        notes = slide_data.get("speaker_notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = notes

    prs.save(output_path)
    logger.info(f"Assembler: PPTX 已保存至 {output_path}，共 {len(sorted_slides)} 页")
    return output_path
