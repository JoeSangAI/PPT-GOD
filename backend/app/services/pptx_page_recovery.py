from __future__ import annotations

import json
import os
import re
import tempfile
from typing import Any

from app.services.image_analyzer import describe_context_image


def slide_text_is_sparse(lines: list[str], min_chars: int = 20) -> bool:
    text = re.sub(r"\s+", "", "\n".join(str(line or "") for line in lines))
    return len(text) < min_chars


def _iter_picture_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)
            continue
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
            yield shape


def _slide_area(slide) -> int:
    try:
        presentation = slide.part.presentation
        return int(presentation.slide_width) * int(presentation.slide_height)
    except Exception:
        return 0


def extract_largest_slide_picture_to_tempfile(slide, *, min_area_ratio: float = 0.35) -> str:
    pictures = list(_iter_picture_shapes(slide.shapes))
    if not pictures:
        return ""
    largest = max(pictures, key=lambda shape: int(getattr(shape, "width", 0) or 0) * int(getattr(shape, "height", 0) or 0))
    picture_area = int(getattr(largest, "width", 0) or 0) * int(getattr(largest, "height", 0) or 0)
    slide_area = _slide_area(slide)
    if slide_area and picture_area / slide_area < min_area_ratio:
        return ""

    image = largest.image
    ext = str(getattr(image, "ext", "png") or "png").lower().lstrip(".")
    with tempfile.NamedTemporaryFile(prefix="pptgod_ppt_page_", suffix=f".{ext}", delete=False) as tmp:
        tmp.write(image.blob)
        return tmp.name


def parse_page_recovery_description(description: str) -> dict[str, Any]:
    text = (description or "").strip()
    if not text:
        return {}
    clean = re.sub(r"^```(?:json)?\s*|```$", "", text, flags=re.MULTILINE | re.IGNORECASE).strip()
    start = clean.find("{")
    end = clean.rfind("}")
    if start != -1 and end > start:
        try:
            parsed = json.loads(clean[start:end + 1])
            if isinstance(parsed, dict):
                return parsed
        except Exception:
            pass
    ocr_text = extract_ocr_text_from_unstructured_description(clean)
    has_recovery_sections = _has_recovery_analysis_sections(clean)
    return {
        "ocr_text": ocr_text if (ocr_text or has_recovery_sections) else text,
        "page_intent": "",
        "key_facts": [],
        "confidence": 0.6,
    }


_OCR_HEADER_RE = re.compile(
    r"^\s*(?:#{1,6}\s*)?(?:\*\*)?\s*(?:\d+[.)、]\s*)?"
    r"(?:OCR\s*文字|OCR文字|截图识别文字)\s*(?:[:：]\s*(.*?))?(?:\*\*)?\s*$",
    flags=re.IGNORECASE,
)
_NON_OCR_SECTION_RE = re.compile(
    r"^\s*(?:#{1,6}\s*)?(?:\*\*)?\s*(?:\d+[.)、]\s*)?"
    r"(?:图像内容|图片内容|可用于\s*PPT\s*的信息|视觉参考|页面意图|修改建议|关键事实|识别置信度)"
    r"\s*(?:[:：].*)?(?:\*\*)?\s*$",
    flags=re.IGNORECASE,
)
_RECOVERY_NOISE_RE = re.compile(
    r"(?:我是|你好.*我是|作为)\s*PPT\s*Agent\s*的读图助手|以下是|详细解读|详细解析",
    flags=re.IGNORECASE,
)
_OCR_LABEL_RE = re.compile(
    r"^(?:"
    r"大标题(?:（.*?）)?|主标题(?:（.*?）)?|左侧大标题|右侧大标题|卡片标题|右侧卡片标题|"
    r"标题/正文|标题|副标题(?:（.*?）)?|正文内容|正文要点\s*\d*|核心卖点(?:（.*?）)?|"
    r"官方网址|装饰符号|其他标识|品牌名|产品定位|页面文字"
    r")\s*[:：]\s*(.*)$"
)


def _has_recovery_analysis_sections(description: str) -> bool:
    return any(
        _OCR_HEADER_RE.match(line.strip()) or _NON_OCR_SECTION_RE.match(line.strip())
        for line in str(description or "").splitlines()
    )


def _clean_recovered_ocr_line(line: str) -> str:
    value = str(line or "").strip()
    if not value:
        return ""
    value = re.sub(r"^\s*(?:[-*+]|\d+[.)、])\s*", "", value).strip()
    value = value.replace("`", "")
    value = re.sub(r"\*\*([^*]+)\*\*", r"\1", value).strip()
    value = re.sub(r"^#{1,6}\s*", "", value).strip()
    value = value.strip("：:；;，,。")
    if not value or _RECOVERY_NOISE_RE.search(value):
        return ""
    label_match = _OCR_LABEL_RE.match(value)
    if label_match:
        value = label_match.group(1).strip()
    if not value or value in {"无明显文字", "无明显可读文字", "无"}:
        return ""
    if value.startswith("无明显"):
        return ""
    return value


def extract_ocr_text_from_unstructured_description(description: str) -> str:
    """Extract only the OCR section from a non-JSON vision response.

    `describe_context_image` returns a human-readable analysis by design. PPT
    page recovery needs original page text, not the assistant's visual analysis;
    otherwise sections such as "视觉参考" become slide body copy.
    """
    lines = str(description or "").splitlines()
    collected: list[str] = []
    in_ocr = False
    saw_structured_header = False
    for raw in lines:
        line = raw.strip()
        if not line:
            continue
        header = _OCR_HEADER_RE.match(line)
        if header:
            saw_structured_header = True
            in_ocr = True
            inline = _clean_recovered_ocr_line(header.group(1) or "")
            if inline:
                collected.append(inline)
            continue
        if _NON_OCR_SECTION_RE.match(line):
            if in_ocr:
                break
            saw_structured_header = True
            continue
        if not in_ocr:
            continue
        cleaned = _clean_recovered_ocr_line(line)
        if cleaned:
            collected.append(cleaned)

    if collected:
        return "\n".join(collected)
    if saw_structured_header:
        return ""
    simple_lines = [_clean_recovered_ocr_line(line) for line in lines]
    simple_lines = [line for line in simple_lines if line]
    return "\n".join(simple_lines[:20])


def read_ppt_page_image(
    image_path: str,
    *,
    page_num: int,
    source_filename: str,
    timeout_seconds: float | None = None,
) -> dict[str, Any]:
    description = describe_context_image(
        image_path,
        f"{source_filename} 第{page_num}页",
        "原 PPT 页面截图",
        "恢复这页 PPT 中可读文字、界面标签、关键事实和页面意图",
        timeout_seconds=timeout_seconds,
    )
    return parse_page_recovery_description(description)


def render_recovered_page_sections(recovered: dict[str, Any]) -> list[str]:
    if not isinstance(recovered, dict):
        return []
    sections: list[str] = []
    ocr_text = str(recovered.get("ocr_text") or recovered.get("text") or "").strip()
    if ocr_text:
        sections.append("【截图识别文字】\n" + ocr_text)
    return sections


def recover_sparse_slide_text(
    slide,
    *,
    page_num: int,
    source_filename: str,
    existing_lines: list[str],
    timeout_seconds: float | None = None,
) -> list[str]:
    if not slide_text_is_sparse(existing_lines):
        return []
    image_path = extract_largest_slide_picture_to_tempfile(slide)
    if not image_path:
        return []
    try:
        recovered = read_ppt_page_image(
            image_path,
            page_num=page_num,
            source_filename=source_filename,
            timeout_seconds=timeout_seconds,
        )
        return render_recovered_page_sections(recovered)
    except Exception:
        return []
    finally:
        try:
            os.unlink(image_path)
        except OSError:
            pass
