import hashlib
import io
import os
import re
import zipfile
from dataclasses import dataclass, field
from typing import Iterable
from xml.etree import ElementTree as ET

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class PptxImageAsset:
    file_path: str
    source_page_num: int
    repeated_page_nums: list[int]
    classification: str
    role: str
    process_mode: str
    asset_kind: str | None
    asset_name: str
    usage_note: str
    metadata: dict = field(default_factory=dict)


@dataclass
class _ImageOccurrence:
    page_num: int
    blob: bytes
    ext: str
    left: int
    top: int
    width: int
    height: int
    slide_width: int
    slide_height: int
    slide_text: str
    raw_sha1: str | None = None
    crop: tuple[float, float, float, float] = (0.0, 0.0, 0.0, 0.0)

    @property
    def area_ratio(self) -> float:
        slide_area = max(1, self.slide_width * self.slide_height)
        return max(0, self.width) * max(0, self.height) / slide_area


@dataclass
class _LibraryPromotion:
    promote: bool
    score: int
    tier: str
    reason: str
    kind: str | None = None


@dataclass
class _ContentRefCandidate:
    priority: float
    asset: PptxImageAsset


PARALLEL_PAGE_REF_MIN_COUNT = 4
PARALLEL_PAGE_REF_MAX_COUNT = 8


PRODUCT_CONTEXT_TERMS = (
    "产品", "主产品", "包装", "瓶", "瓶身", "sku", "SKU", "货架", "终端", "陈列",
    "样品", "实物", "设备", "取件机", "快递柜", "商品", "礼盒", "门店物料",
)
PERSON_CONTEXT_TERMS = ("人物", "模特", "代言人", "创始人", "讲师", "专家", "团队", "肖像", "合影")
MATERIAL_CONTEXT_TERMS = (
    "物料", "主视觉", "KV", "key visual", "品牌物料", "产品物料", "海报",
    "立牌", "展架", "吊旗", "标识", "导视", "贴片",
)
CONTENT_GRAPHIC_CONTEXT_TERMS = (
    "图表", "数据", "表格", "截图", "界面", "系统", "小程序", "大屏", "看板", "流程",
    "路径", "动线", "地图", "矩阵", "架构", "模型", "漏斗", "时间轴", "二维码",
)
IDENTITY_CODE_CONTEXT_TERMS = (
    "二维码", "扫码", "身份码", "取件码", "小程序码", "条码", "链接码",
)
UI_CONTAINER_CONTEXT_TERMS = (
    "手机边框", "手机框", "手机壳", "手机外框", "手机界面", "界面外框", "屏幕框",
    "app界面", "APP界面", "小程序界面", "小程序", "mockup", "Mockup",
)
LOW_VALUE_CONTEXT_TERMS = (
    "背景", "氛围", "氛围感", "风景", "插画", "装饰", "纹理", "光效", "底图",
    "配图", "意境", "校园风景", "天空", "草地",
)
CORE_ASSET_CONTEXT_TERMS = (
    *PRODUCT_CONTEXT_TERMS,
    *PERSON_CONTEXT_TERMS,
    *MATERIAL_CONTEXT_TERMS,
)


def _safe_stem(filename: str) -> str:
    stem = os.path.splitext(os.path.basename(filename or "pptx"))[0]
    stem = re.sub(r"[^A-Za-z0-9._-]+", "_", stem).strip("._-")
    return stem or "pptx"


def _iter_picture_shapes(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
            yield shape


def _shape_image_ext(shape) -> str:
    try:
        return (shape.image.ext or "png").lower()
    except Exception:
        return "png"


def _shape_crop_tuple(shape) -> tuple[float, float, float, float]:
    try:
        return (
            float(getattr(shape, "crop_left", 0.0) or 0.0),
            float(getattr(shape, "crop_top", 0.0) or 0.0),
            float(getattr(shape, "crop_right", 0.0) or 0.0),
            float(getattr(shape, "crop_bottom", 0.0) or 0.0),
        )
    except Exception:
        return (0.0, 0.0, 0.0, 0.0)


def _visible_picture_blob(blob: bytes, crop: tuple[float, float, float, float]) -> bytes:
    """Return the actual visible image after PowerPoint picture cropping."""
    crop_left, crop_top, crop_right, crop_bottom = crop
    if max(crop_left, crop_top, crop_right, crop_bottom) <= 0.0001:
        return blob

    try:
        with Image.open(io.BytesIO(blob)) as img:
            width, height = img.size
            left = int(width * max(0.0, crop_left))
            top = int(height * max(0.0, crop_top))
            right = int(width * (1.0 - max(0.0, crop_right)))
            bottom = int(height * (1.0 - max(0.0, crop_bottom)))
            left = max(0, min(width - 1, left))
            top = max(0, min(height - 1, top))
            right = max(left + 1, min(width, right))
            bottom = max(top + 1, min(height, bottom))
            cropped = img.crop((left, top, right, bottom))
            out = io.BytesIO()
            mode = "RGBA" if cropped.mode in {"RGBA", "LA", "P"} else "RGB"
            cropped.convert(mode).save(out, "PNG")
            return out.getvalue()
    except Exception:
        return blob


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            parts.append(shape.text.strip())
    return "\n".join(parts)


def _image_stats(blob: bytes) -> tuple[int, int, float, float, bool]:
    with Image.open(io.BytesIO(blob)) as img:
        rgba = img.convert("RGBA")
        width, height = rgba.size
        alpha = rgba.getchannel("A")
        alpha_min, alpha_max = alpha.getextrema()
        has_transparency = alpha_min < 250

        small = rgba.copy()
        small.thumbnail((160, 160))
        rgb = small.convert("RGB")
        colors = rgb.getcolors(maxcolors=4096) or []
        total = max(1, rgb.size[0] * rgb.size[1])
        dominant_share = max((count for count, _color in colors), default=0) / total

        stat = ImageStat.Stat(rgb)
        channel_std = sum(stat.stddev) / max(1, len(stat.stddev))
        return width, height, dominant_share, channel_std, has_transparency or alpha_max < 255


def _is_decorative(occ: _ImageOccurrence, dominant_share: float, channel_std: float) -> bool:
    if occ.width <= 0 or occ.height <= 0:
        return True
    if occ.area_ratio < 0.0015:
        return True
    if min(occ.width, occ.height) < 24:
        return True
    aspect = occ.width / max(1, occ.height)
    if (aspect > 12 or aspect < 1 / 12) and occ.area_ratio < 0.02:
        return True
    if dominant_share > 0.985 and channel_std < 8:
        return True
    return False


def _asset_kind_for_occ(occ: _ImageOccurrence, has_transparency: bool) -> str:
    if has_transparency and occ.area_ratio <= 0.12:
        return "material"
    if occ.area_ratio >= 0.22:
        return "scene"
    return "other"


def _is_cover_logo_candidate(
    occ: _ImageOccurrence,
    slide_count: int,
    dominant_share: float,
    channel_std: float,
    has_transparency: bool,
) -> bool:
    if occ.page_num not in {1, slide_count}:
        return False
    if occ.area_ratio < 0.0015 or occ.area_ratio > 0.05:
        return False
    aspect = occ.width / max(1, occ.height)
    if aspect > 10 or aspect < 0.1:
        return False
    y_mid = (occ.top + occ.height / 2) / max(1, occ.slide_height)
    in_brand_zone = y_mid <= 0.30 or y_mid >= 0.72
    if not in_brand_zone:
        return False
    has_crop = max(occ.crop or (0.0, 0.0, 0.0, 0.0)) > 0.0001
    if has_crop and occ.area_ratio <= 0.015:
        return True
    # Logo marks are usually flat, transparent, or text-like. This rejects many
    # photographic cover thumbnails while still allowing multi-color brand marks.
    if not has_transparency and dominant_share < 0.48 and channel_std > 45:
        return False
    return True


def _text_has_any(text: str, terms: tuple[str, ...]) -> bool:
    lower = (text or "").lower()
    return any(term.lower() in lower for term in terms)


def _is_full_slide_background(occ: _ImageOccurrence) -> bool:
    if occ.area_ratio >= 0.62:
        return True
    slide_w = max(1, occ.slide_width)
    slide_h = max(1, occ.slide_height)
    return (
        occ.left <= slide_w * 0.025
        and occ.top <= slide_h * 0.025
        and occ.width >= slide_w * 0.92
        and occ.height >= slide_h * 0.90
    )


def _is_ui_container_chrome(occ: _ImageOccurrence) -> bool:
    source_text = " ".join(occ.slide_text.split())
    if not _text_has_any(source_text, UI_CONTAINER_CONTEXT_TERMS):
        return False
    aspect = occ.width / max(1, occ.height)
    if 0.26 <= aspect <= 0.72 and occ.area_ratio <= 0.28:
        return True
    return False


def _suggest_promoted_kind(kind: str, source_text: str) -> str:
    if kind and kind != "other":
        return kind
    if _text_has_any(source_text, PRODUCT_CONTEXT_TERMS):
        return "product"
    if _text_has_any(source_text, PERSON_CONTEXT_TERMS):
        return "person"
    if _text_has_any(source_text, MATERIAL_CONTEXT_TERMS):
        return "material"
    return kind or "other"


def _library_promotion_for_asset(
    occurrences: list[_ImageOccurrence],
    kind: str,
    dominant_share: float,
    channel_std: float,
    has_transparency: bool,
    unique_image_count_on_page: int,
) -> _LibraryPromotion:
    first = occurrences[0]
    repeated_pages = {occ.page_num for occ in occurrences}
    source_text = " ".join(first.slide_text.split())
    has_identity_context = _text_has_any(source_text, CORE_ASSET_CONTEXT_TERMS)
    has_low_value_context = _text_has_any(source_text, LOW_VALUE_CONTEXT_TERMS)
    has_content_graphic_context = _text_has_any(source_text, CONTENT_GRAPHIC_CONTEXT_TERMS)
    has_identity_code_context = _text_has_any(source_text, IDENTITY_CODE_CONTEXT_TERMS)
    has_ui_container_context = _text_has_any(source_text, UI_CONTAINER_CONTEXT_TERMS)

    if _is_full_slide_background(first):
        return _LibraryPromotion(False, 0, "page_ref_only", "full-slide/background image", kind)
    if has_ui_container_context and _is_ui_container_chrome(first):
        return _LibraryPromotion(False, 0, "layout_chrome", "phone/UI container is layout chrome, not reusable material", kind)
    if has_identity_code_context:
        return _LibraryPromotion(False, 16, "page_ref_only", "QR/identity code kept as page-specific evidence", kind)
    if len(repeated_pages) >= 2 and first.area_ratio >= 0.18:
        return _LibraryPromotion(False, 0, "page_ref_only", "large repeated template/background image", kind)
    if kind == "scene":
        return _LibraryPromotion(False, 8, "page_ref_only", "large scene kept only as page reference", kind)
    if has_content_graphic_context and not has_identity_context:
        return _LibraryPromotion(False, 18, "page_ref_only", "content graphic kept as page-level evidence", kind)
    if has_low_value_context and not has_identity_context:
        return _LibraryPromotion(False, 0, "low_value", "ambient/decorative context", kind)

    score = 0
    reasons: list[str] = []
    if has_transparency and 0.006 <= first.area_ratio <= 0.16 and has_identity_context:
        score += 45
        reasons.append("transparent identity/material asset with matching slide context")
    if len(repeated_pages) >= 2 and first.area_ratio <= 0.12 and has_identity_context and not has_low_value_context:
        score += 35
        reasons.append("small/medium repeated identity asset")
    if (
        has_identity_context
        and unique_image_count_on_page <= 2
        and 0.035 <= first.area_ratio <= 0.20
        and channel_std >= 12
        and dominant_share <= 0.96
    ):
        score += 35
        reasons.append("primary identity/product image on a relevant slide")
    if dominant_share > 0.94 and channel_std < 12:
        score -= 30
        reasons.append("mostly flat/solid image")
    if has_low_value_context:
        score -= 20
        reasons.append("low-value visual context")

    promoted_kind = _suggest_promoted_kind(kind, source_text)
    if score >= 35:
        return _LibraryPromotion(True, score, "core_global", "; ".join(reasons), promoted_kind)
    return _LibraryPromotion(False, max(0, score), "page_ref_only", "; ".join(reasons) or "not distinctive enough for global library", kind)


def _content_ref_priority(
    occ: _ImageOccurrence,
    kind: str,
    dominant_share: float,
    channel_std: float,
    has_transparency: bool,
) -> float:
    source_text = " ".join(occ.slide_text.split())
    if _is_full_slide_background(occ):
        base = 5.0
    else:
        base = min(40.0, occ.area_ratio * 260)
    has_identity_context = _text_has_any(source_text, CORE_ASSET_CONTEXT_TERMS)
    has_content_graphic_context = _text_has_any(source_text, CONTENT_GRAPHIC_CONTEXT_TERMS)
    if has_identity_context:
        base += 18
    elif has_content_graphic_context:
        base += 14
    if kind in {"material", "product", "person"}:
        base += 16
    if kind == "scene":
        base -= 8
    if has_transparency:
        base += 8
    if _text_has_any(source_text, LOW_VALUE_CONTEXT_TERMS) and not (has_identity_context or has_content_graphic_context):
        base -= 18
    if dominant_share > 0.95 and channel_std < 12:
        base -= 20
    return base


def _classify_asset(
    occurrences: list[_ImageOccurrence],
    dominant_share: float,
    channel_std: float,
    has_transparency: bool,
    slide_count: int,
) -> tuple[str, str, str | None]:
    first = occurrences[0]
    repeated_pages = {occ.page_num for occ in occurrences}
    if _is_decorative(first, dominant_share, channel_std):
        return "decorative", "ignore", None
    if _is_ui_container_chrome(first):
        return "layout_chrome", "ignore", None
    # Be conservative with auto-logo detection. Real PPTs often reuse phone
    # frames, QR-code screenshots, icons, and UI chrome across pages; promoting
    # those to a global logo is worse than leaving brand marks as recallable
    # library assets. Only very small repeated transparent marks become logo.
    if len(repeated_pages) >= 3 and first.area_ratio <= 0.006 and has_transparency:
        return "logo", "logo", None
    if len(repeated_pages) >= 4 and first.area_ratio <= 0.004:
        return "logo", "logo", None
    if _is_cover_logo_candidate(first, slide_count, dominant_share, channel_std, has_transparency):
        return "logo_candidate", "logo", None
    kind = _asset_kind_for_occ(first, has_transparency)
    return "useful", "content_ref", kind


def _slide_xml_blip_targets(file_bytes: bytes) -> dict[int, list[tuple[str, bytes]]]:
    """Find images referenced by slide XML, including backgrounds/fill blips.

    python-pptx exposes ordinary picture shapes well, but background images and
    shape fills may only appear as a:blip relationships. This fallback keeps
    those visual sources in the asset library without adding a new dependency.
    """
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    }
    result: dict[int, list[tuple[str, bytes]]] = {}
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            slide_names = sorted(
                (name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)),
                key=lambda name: int(re.search(r"slide(\d+)\.xml$", name).group(1)),
            )
            for slide_name in slide_names:
                page_num = int(re.search(r"slide(\d+)\.xml$", slide_name).group(1))
                rels_name = slide_name.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
                if rels_name not in zf.namelist():
                    continue
                rels_root = ET.fromstring(zf.read(rels_name))
                rel_targets = {
                    rel.attrib.get("Id"): rel.attrib.get("Target", "")
                    for rel in rels_root.findall("rel:Relationship", ns)
                }
                slide_root = ET.fromstring(zf.read(slide_name))
                for blip in slide_root.findall(".//a:blip", ns):
                    rel_id = blip.attrib.get(f"{{{ns['r']}}}embed") or blip.attrib.get(f"{{{ns['r']}}}link")
                    target = rel_targets.get(rel_id)
                    if not target:
                        continue
                    media_path = os.path.normpath(os.path.join("ppt/slides", target)).replace("\\", "/")
                    if media_path.startswith("ppt/slides/../"):
                        media_path = "ppt/" + media_path[len("ppt/slides/../"):]
                    if media_path not in zf.namelist():
                        continue
                    result.setdefault(page_num, []).append((media_path, zf.read(media_path)))
    except Exception:
        return result
    return result


def _occurrence_key(digest: str, occ: _ImageOccurrence) -> tuple[str, int, int, int, int, int]:
    return (
        digest,
        occ.page_num,
        int(occ.left or 0),
        int(occ.top or 0),
        int(occ.width or 0),
        int(occ.height or 0),
    )


def _candidate_area(candidate: _ContentRefCandidate) -> float:
    try:
        return float(candidate.asset.metadata.get("area_ratio") or 0)
    except (TypeError, ValueError, AttributeError):
        return 0.0


def _candidate_spatial_key(candidate: _ContentRefCandidate) -> tuple[float, float, str]:
    bounds = candidate.asset.metadata.get("shape_bounds") if isinstance(candidate.asset.metadata, dict) else {}
    if not isinstance(bounds, dict):
        bounds = {}
    try:
        top = float(bounds.get("top") or 0)
    except (TypeError, ValueError):
        top = 0.0
    try:
        left = float(bounds.get("left") or 0)
    except (TypeError, ValueError):
        left = 0.0
    return (top, left, candidate.asset.file_path)


def _parallel_page_ref_candidates(candidates: list[_ContentRefCandidate]) -> list[_ContentRefCandidate]:
    """
    Detect same-page picture sets that should stay together.

    A single PPT page often uses 4-8 peer images as a grid, examples, personas,
    channel cards, or side-by-side screenshots. Ranking those images individually
    and keeping only the top few breaks the author's original evidence set, so
    once a page has enough useful non-background images we keep the compact group
    as a unit. The cap protects prompt/image-input latency.
    """
    group = [
        candidate for candidate in candidates
        if 0.003 <= _candidate_area(candidate) <= 0.45
    ]
    if len(group) < PARALLEL_PAGE_REF_MIN_COUNT:
        return []
    group.sort(key=_candidate_spatial_key)
    return group[:PARALLEL_PAGE_REF_MAX_COUNT]


def _keyword_tags(source_filename: str, occ: _ImageOccurrence, repeated_pages: list[int]) -> list[str]:
    tags = [
        "原PPT素材",
        f"第{occ.page_num}页",
        f"ppt_page_{occ.page_num}",
    ]
    if len(repeated_pages) > 1:
        tags.append("跨页重复素材")
        tags.extend(f"第{page_num}页" for page_num in repeated_pages[:8])
    for token in re.split(r"[\s,，。；;:：、/|()（）\[\]{}\"'“”‘’_-]+", occ.slide_text):
        token = token.strip()
        if 2 <= len(token) <= 24 and token not in tags:
            tags.append(token)
        if len(tags) >= 24:
            break
    source_stem = _safe_stem(source_filename)
    if source_stem not in tags:
        tags.append(source_stem)
    return tags


def _base_metadata(
    source_filename: str,
    digest: str,
    classification: str,
    occ: _ImageOccurrence,
    repeated_pages: list[int],
    pixel_w: int,
    pixel_h: int,
    dominant_share: float,
) -> dict:
    source_text = " ".join(occ.slide_text.split())
    return {
        "source_document": source_filename,
        "pptx_source_page_num": occ.page_num,
        "pptx_repeated_page_nums": repeated_pages,
        "pptx_image_sha1": digest,
        "pptx_raw_image_sha1": occ.raw_sha1 or digest,
        "pptx_crop": [round(value, 6) for value in occ.crop],
        "classification": classification,
        "area_ratio": round(occ.area_ratio, 5),
        "shape_bounds": {
            "left": round(occ.left / max(1, occ.slide_width), 5),
            "top": round(occ.top / max(1, occ.slide_height), 5),
            "width": round(occ.width / max(1, occ.slide_width), 5),
            "height": round(occ.height / max(1, occ.slide_height), 5),
        },
        "pixel_size": [pixel_w, pixel_h],
        "dominant_color_share": round(dominant_share, 4),
        "source_slide_text": source_text[:800],
        "asset_tags": _keyword_tags(source_filename, occ, repeated_pages),
        "asset_lock": {
            "scope": "pptx_source_page",
            "source_document": source_filename,
            "page_num": occ.page_num,
            "repeated_page_nums": repeated_pages,
        },
    }


def extract_pptx_image_assets(
    file_bytes: bytes,
    source_filename: str,
    output_dir: str,
    *,
    max_assets_per_slide: int = 3,
    max_total_assets: int | None = None,
) -> list[PptxImageAsset]:
    """Extract useful images from a PPTX and classify them for the deck pipeline.

    This is intentionally conservative: it skips solid decorations and tiny
    icons, promotes repeated small transparent marks to logo, and keeps useful
    slide images as page references with source_page_num metadata.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    os.makedirs(output_dir, exist_ok=True)

    occurrences_by_hash: dict[str, list[_ImageOccurrence]] = {}
    shape_digests_by_page: dict[int, set[str]] = {}
    seen_occurrences: set[tuple[str, int, int, int, int, int]] = set()
    for page_num, slide in enumerate(prs.slides, start=1):
        slide_text = _slide_text(slide)
        for shape in _iter_picture_shapes(slide.shapes):
            raw_blob = shape.image.blob
            raw_digest = hashlib.sha1(raw_blob).hexdigest()
            crop = _shape_crop_tuple(shape)
            blob = _visible_picture_blob(raw_blob, crop)
            digest = hashlib.sha1(blob).hexdigest()
            shape_digests_by_page.setdefault(page_num, set()).update({raw_digest, digest})
            occurrence = _ImageOccurrence(
                page_num=page_num,
                blob=blob,
                ext=_shape_image_ext(shape),
                left=int(shape.left or 0),
                top=int(shape.top or 0),
                width=int(shape.width or 0),
                height=int(shape.height or 0),
                slide_width=int(prs.slide_width),
                slide_height=int(prs.slide_height),
                slide_text=slide_text,
                raw_sha1=raw_digest,
                crop=crop,
            )
            occurrence_key = _occurrence_key(digest, occurrence)
            if occurrence_key in seen_occurrences:
                continue
            seen_occurrences.add(occurrence_key)
            occurrences_by_hash.setdefault(digest, []).append(occurrence)

    xml_blips = _slide_xml_blip_targets(file_bytes)
    slide_text_by_page = {
        page_num: _slide_text(slide)
        for page_num, slide in enumerate(prs.slides, start=1)
    }
    for page_num, blips in xml_blips.items():
        for _media_path, blob in blips:
            digest = hashlib.sha1(blob).hexdigest()
            if digest in shape_digests_by_page.get(page_num, set()):
                continue
            occurrence = _ImageOccurrence(
                page_num=page_num,
                blob=blob,
                ext="png",
                left=0,
                top=0,
                width=int(prs.slide_width),
                height=int(prs.slide_height),
                slide_width=int(prs.slide_width),
                slide_height=int(prs.slide_height),
                slide_text=slide_text_by_page.get(page_num, ""),
            )
            occurrence_key = _occurrence_key(digest, occurrence)
            if occurrence_key in seen_occurrences:
                continue
            seen_occurrences.add(occurrence_key)
            occurrences_by_hash.setdefault(digest, []).append(occurrence)

    assets: list[PptxImageAsset] = []
    content_ref_candidates: list[_ContentRefCandidate] = []
    seen_content_ref_candidates: set[tuple[str, int]] = set()
    per_slide_count: dict[int, int] = {}
    source_stem = _safe_stem(source_filename)
    unique_image_count_by_page: dict[int, int] = {}
    for occurrences in occurrences_by_hash.values():
        for page_num in {occ.page_num for occ in occurrences}:
            unique_image_count_by_page[page_num] = unique_image_count_by_page.get(page_num, 0) + 1

    slide_count = len(prs.slides)
    for digest, occurrences in occurrences_by_hash.items():
        occurrences.sort(key=lambda occ: (occ.page_num, -occ.area_ratio))
        first = occurrences[0]
        try:
            pixel_w, pixel_h, dominant_share, channel_std, has_transparency = _image_stats(first.blob)
        except Exception:
            continue

        classification, role, kind = _classify_asset(occurrences, dominant_share, channel_std, has_transparency, slide_count)
        if role == "ignore":
            continue

        page_nums = sorted({occ.page_num for occ in occurrences})
        safe_ext = "png"
        out_name = f"{source_stem}_p{first.page_num:03d}_{digest[:10]}.{safe_ext}"
        file_path = os.path.join(output_dir, out_name)
        if not os.path.exists(file_path):
            with Image.open(io.BytesIO(first.blob)) as img:
                mode = "RGBA" if img.mode in {"RGBA", "LA", "P"} else "RGB"
                img.convert(mode).save(file_path, "PNG")

        base_metadata = _base_metadata(
            source_filename,
            digest,
            classification,
            first,
            page_nums,
            pixel_w,
            pixel_h,
            dominant_share,
        )

        if role == "logo":
            assets.append(
                PptxImageAsset(
                    file_path=file_path,
                    source_page_num=first.page_num,
                    repeated_page_nums=page_nums,
                    classification=classification,
                    role="logo",
                    process_mode="original",
                    asset_kind=None,
                    asset_name=f"{source_stem} extracted logo",
                    usage_note="从上传 PPT 中识别出的 Logo / 联合标识，作为全局品牌标识使用。",
                    metadata=base_metadata,
                )
            )
            continue

        promotion = _library_promotion_for_asset(
            occurrences,
            kind or "other",
            dominant_share,
            channel_std,
            has_transparency,
            unique_image_count_by_page.get(first.page_num, 1),
        )
        promoted_kind = promotion.kind or kind
        process_mode = "crop" if promoted_kind in {"product", "material", "other"} else "blend"

        if promotion.promote:
            assets.append(
                PptxImageAsset(
                    file_path=file_path,
                    source_page_num=first.page_num,
                    repeated_page_nums=page_nums,
                    classification="library_asset",
                    role="visual_asset",
                    process_mode=process_mode,
                    asset_kind=promoted_kind,
                    asset_name=f"{source_stem} p{first.page_num} core asset",
                    usage_note="从上传 PPT 中筛选出的高价值核心素材；仅在页面内容明确相关时自动召回。",
                    metadata={
                        **base_metadata,
                        "classification": "library_asset",
                        "library_role": "core_global_asset",
                        "selection_tier": promotion.tier,
                        "importance_score": promotion.score,
                        "selection_reason": promotion.reason,
                        "unique_image_count_on_page": unique_image_count_by_page.get(first.page_num, 1),
                    },
                )
            )

        for occ in occurrences:
            content_ref_key = (digest, occ.page_num)
            if content_ref_key in seen_content_ref_candidates:
                continue
            seen_content_ref_candidates.add(content_ref_key)
            metadata = _base_metadata(
                source_filename,
                digest,
                classification,
                occ,
                page_nums,
                pixel_w,
                pixel_h,
                dominant_share,
            )
            priority = _content_ref_priority(occ, promoted_kind or kind or "other", dominant_share, channel_std, has_transparency)
            content_ref_candidates.append(
                _ContentRefCandidate(
                    priority=priority,
                    asset=PptxImageAsset(
                        file_path=file_path,
                        source_page_num=occ.page_num,
                        repeated_page_nums=page_nums,
                        classification=classification,
                        role="content_ref",
                        process_mode="crop" if (promoted_kind or kind) in {"product", "material", "other"} else "blend",
                        asset_kind=promoted_kind or kind,
                        asset_name=f"{source_stem} p{occ.page_num} image",
                        usage_note="从上传 PPT 对应页提取的有用图片，优先作为本页参考图使用。",
                        metadata={
                            **metadata,
                            "selection_tier": "page_ref",
                            "importance_score": round(priority, 2),
                            "selection_reason": promotion.reason,
                        },
                    ),
                )
            )

    candidates_by_page: dict[int, list[_ContentRefCandidate]] = {}
    for candidate in content_ref_candidates:
        candidates_by_page.setdefault(candidate.asset.source_page_num, []).append(candidate)

    for page_num in sorted(candidates_by_page):
        page_candidates = candidates_by_page[page_num]
        page_candidates.sort(
            key=lambda item: (
                -item.priority,
                -_candidate_area(item),
                item.asset.file_path,
            )
        )

        parallel_group = _parallel_page_ref_candidates(page_candidates)
        if parallel_group:
            selected_candidates = parallel_group
            group_key = f"pptx_page_{page_num}_parallel_refs"
            group_size = len(selected_candidates)
            for idx, candidate in enumerate(selected_candidates, start=1):
                candidate.asset.metadata = {
                    **candidate.asset.metadata,
                    "asset_group_key": group_key,
                    "asset_group_index": idx,
                    "asset_group_size": group_size,
                    "asset_group_role": "parallel_page_reference_set",
                    "selection_reason": "same-page parallel image set kept together",
                }
        else:
            selected_candidates = page_candidates[:max_assets_per_slide]

        for candidate in selected_candidates:
            if max_total_assets is not None and len(assets) >= max_total_assets:
                break
            asset = candidate.asset
            if not parallel_group and per_slide_count.get(asset.source_page_num, 0) >= max_assets_per_slide:
                continue
            per_slide_count[asset.source_page_num] = per_slide_count.get(asset.source_page_num, 0) + 1
            assets.append(asset)
        if max_total_assets is not None and len(assets) >= max_total_assets:
            break

    return assets[:max_total_assets] if max_total_assets is not None else assets
