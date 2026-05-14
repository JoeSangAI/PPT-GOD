import logging
import os
import shutil
import subprocess
import tempfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Dict, List, Optional

import fitz  # PyMuPDF
from PIL import Image, ImageChops, ImageDraw, ImageStat

logger = logging.getLogger(__name__)


COMMON_SOFFICE_PATHS = (
    "/usr/bin/soffice",
    "/usr/local/bin/soffice",
    "/opt/homebrew/bin/soffice",
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
)


def find_soffice_binary() -> str | None:
    """Return a usable LibreOffice executable path, if available."""
    candidates = [
        os.getenv("SOFFICE_PATH"),
        shutil.which("soffice"),
        *COMMON_SOFFICE_PATHS,
    ]
    for candidate in candidates:
        if candidate and os.path.exists(candidate) and os.access(candidate, os.X_OK):
            return candidate
    return None


def convert_ppt_to_pdf(ppt_path: str, output_dir: str) -> str:
    """使用 LibreOffice 将 PPT/PPTX 转换为 PDF。"""
    soffice = find_soffice_binary()
    if not soffice:
        raise RuntimeError(
            "LibreOffice 未安装或 soffice 不在 PATH。请在 Docker 镜像中安装 LibreOffice，"
            "或设置 SOFFICE_PATH 指向可执行文件。"
        )

    profile_dir = tempfile.mkdtemp(prefix="pptgod-lo-profile-")
    try:
        cmd = [
            soffice,
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--nolockcheck",
            f"-env:UserInstallation={Path(profile_dir).resolve().as_uri()}",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            ppt_path,
        ]
        logger.info("TemplateExtractor: 转换 PPT -> PDF: %s", ppt_path)
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            logger.error("LibreOffice 转换失败: %s", result.stderr)
            raise RuntimeError(f"PPT 转换失败: {result.stderr}")
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)

    generated_pdf = os.path.join(output_dir, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")
    if not os.path.exists(generated_pdf):
        raise RuntimeError("PDF 文件未生成")
    return generated_pdf


def extract_pdf_thumbnails(pdf_path: str, output_dir: str, dpi: int = 150) -> List[str]:
    """使用 PyMuPDF 提取 PDF 每页为 PNG 缩略图。"""
    doc = fitz.open(pdf_path)
    paths = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        img_path = os.path.join(output_dir, f"page_{page_num + 1:03d}.png")
        pix.save(img_path)
        paths.append(img_path)
        logger.info("TemplateExtractor: 提取第 %s 页缩略图 -> %s", page_num + 1, img_path)
    doc.close()
    return paths


def _extract_pdf_texts(pdf_path: str) -> list[str]:
    texts: list[str] = []
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            texts.append((page.get_text("text") or "").strip())
        doc.close()
    except Exception as exc:
        logger.warning("TemplateExtractor: PDF 文本提取失败: %s", exc)
    return texts


def _extract_pptx_texts(ppt_path: str) -> list[str]:
    try:
        from pptx import Presentation
    except Exception:
        return []
    texts: list[str] = []
    try:
        prs = Presentation(ppt_path)
        for slide in prs.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    parts.append(shape.text.strip())
            texts.append("\n".join(parts))
    except Exception as exc:
        logger.warning("TemplateExtractor: PPTX 文本提取失败: %s", exc)
    return texts


def _text_density(text: str) -> int:
    return len("".join(str(text or "").split()))


def infer_template_source_kind(page_texts: list[str], total_pages: int) -> str:
    """Distinguish a clean layout template from a finished deck, conservatively."""
    if total_pages <= 0:
        return "template"
    densities = [_text_density(text) for text in page_texts]
    avg_density = sum(densities) / max(1, len(densities))
    dense_pages = sum(1 for value in densities if value >= 90)
    placeholder_terms = ("示例", "占位", "placeholder", "sample", "lorem", "标题", "正文")
    placeholder_hits = sum(
        1
        for text in page_texts
        if any(term.lower() in str(text or "").lower() for term in placeholder_terms)
    )
    if total_pages >= 8 and (avg_density >= 70 or dense_pages >= max(4, total_pages // 2)):
        return "finished_ppt"
    if total_pages >= 5 and dense_pages >= max(3, int(total_pages * 0.65)) and placeholder_hits <= 1:
        return "finished_ppt"
    return "template"


def _infer_page_category(page_num: int, total: int, text: str = "") -> str:
    """基于页码、文本线索和位置推断页面类别。"""
    normalized = "".join(str(text or "").split()).lower()
    if page_num == 1:
        return "cover"
    if page_num == total and total > 1:
        return "ending"
    if any(term in normalized for term in ("目录", "agenda", "contents", "大纲")):
        return "toc"
    if any(term in normalized for term in ("数据", "指标", "同比", "环比", "%", "kpi", "roi", "增长率")):
        return "data"
    if any(term in normalized for term in ("金句", "quote", "观点", "核心结论")) and _text_density(text) <= 120:
        return "quote"
    if any(term in normalized for term in ("章节", "章", "chapter", "section", "part")) and _text_density(text) <= 120:
        return "section"
    if page_num == 2 and total >= 3:
        return "toc"
    if page_num == 3 and total >= 5 and _text_density(text) <= 80:
        return "section"
    return "content"


def _category_confidence(category: str, page_num: int, text: str = "") -> float:
    normalized = "".join(str(text or "").split()).lower()
    if category in {"cover", "ending"}:
        return 0.92
    if category == "toc" and any(term in normalized for term in ("目录", "agenda", "contents", "大纲")):
        return 0.9
    if category == "data" and any(term in normalized for term in ("数据", "指标", "%", "kpi", "roi")):
        return 0.82
    if category == "quote":
        return 0.72
    if category == "section":
        return 0.66
    if category == "toc" and page_num == 2:
        return 0.64
    return 0.6


def _border_background_color(img: Image.Image) -> tuple[int, int, int]:
    rgb = img.convert("RGB")
    width, height = rgb.size
    edge = max(1, min(width, height) // 100)
    strips = [
        rgb.crop((0, 0, width, edge)),
        rgb.crop((0, height - edge, width, height)),
        rgb.crop((0, 0, edge, height)),
        rgb.crop((width - edge, 0, width, height)),
    ]
    pixels = []
    for strip in strips:
        small = strip.resize((max(1, strip.width // 8), max(1, strip.height // 8)))
        pixels.extend(list(small.getdata()))
    if not pixels:
        return (255, 255, 255)
    probe = Image.new("RGB", (len(pixels), 1))
    probe.putdata(pixels)
    stat = ImageStat.Stat(probe)
    return tuple(int(value) for value in stat.median[:3])


def _fill_logo_region(img: Image.Image, bbox: tuple[int, int, int, int]) -> None:
    width, height = img.size
    left, top, right, bottom = bbox
    pad = max(6, int(max(width, height) * 0.006))
    left = max(0, left - pad)
    top = max(0, top - pad)
    right = min(width, right + pad)
    bottom = min(height, bottom + pad)
    ImageDraw.Draw(img).rectangle((left, top, right, bottom), fill=_border_background_color(img))


def _bbox_from_normalized(bounds: dict, width: int, height: int) -> tuple[int, int, int, int] | None:
    try:
        left = float(bounds.get("left") or 0) * width
        top = float(bounds.get("top") or 0) * height
        box_width = float(bounds.get("width") or 0) * width
        box_height = float(bounds.get("height") or 0) * height
    except (TypeError, ValueError, AttributeError):
        return None
    if box_width <= 1 or box_height <= 1:
        return None
    return (
        max(0, int(left)),
        max(0, int(top)),
        min(width, int(left + box_width)),
        min(height, int(top + box_height)),
    )


def _layout_reference_path(preview_path: str, output_dir: str, logo_regions: list[dict]) -> str:
    if not logo_regions:
        return preview_path
    with Image.open(preview_path) as source:
        img = source.convert("RGB")
        width, height = img.size
        for region in logo_regions:
            bbox = None
            if isinstance(region.get("bbox"), (list, tuple)) and len(region["bbox"]) == 4:
                bbox = tuple(int(value) for value in region["bbox"])
            elif isinstance(region.get("bbox_norm"), dict):
                bbox = _bbox_from_normalized(region["bbox_norm"], width, height)
            if bbox:
                _fill_logo_region(img, bbox)
        layout_path = os.path.join(output_dir, "layout_" + os.path.basename(preview_path))
        img.save(layout_path, "PNG")
        return layout_path


def _corner_candidates_for_page(img_path: str, page_num: int) -> list[dict]:
    candidates: list[dict] = []
    try:
        with Image.open(img_path) as source:
            img = source.convert("RGB")
            width, height = img.size
            zone_w = int(width * 0.22)
            zone_h = int(height * 0.16)
            zones = {
                "top-left": (0, 0, zone_w, zone_h),
                "top-right": (width - zone_w, 0, width, zone_h),
                "bottom-left": (0, height - zone_h, zone_w, height),
                "bottom-right": (width - zone_w, height - zone_h, width, height),
            }
            for anchor, zone in zones.items():
                crop = img.crop(zone)
                bg = _border_background_color(crop)
                diff = ImageChops.difference(crop, Image.new("RGB", crop.size, bg)).convert("L")
                mask = diff.point(lambda value: 255 if value > 32 else 0)
                bbox = mask.getbbox()
                if not bbox:
                    continue
                left, top, right, bottom = bbox
                box_w = right - left
                box_h = bottom - top
                area_ratio = (box_w * box_h) / max(1, width * height)
                if area_ratio < 0.0007 or area_ratio > 0.035:
                    continue
                aspect = box_w / max(1, box_h)
                if aspect > 8 or aspect < 0.16:
                    continue
                abs_bbox = (zone[0] + left, zone[1] + top, zone[0] + right, zone[1] + bottom)
                logo_crop = img.crop(abs_bbox)
                small = logo_crop.resize((16, 16)).convert("L")
                avg = sum(small.getdata()) / 256
                hash_bits = "".join("1" if value >= avg else "0" for value in small.getdata())
                candidates.append({
                    "page_num": page_num,
                    "anchor": anchor,
                    "bbox": abs_bbox,
                    "hash": hash_bits,
                    "crop": logo_crop,
                })
    except Exception as exc:
        logger.warning("TemplateExtractor: 角标 Logo 检测失败 page=%s error=%s", page_num, exc)
    return candidates


def _hash_distance(left: str, right: str) -> int:
    return sum(1 for a, b in zip(left, right) if a != b) + abs(len(left) - len(right))


def detect_repeated_corner_logos(pages: list[Dict], output_dir: str) -> tuple[list[dict], dict[int, list[dict]]]:
    """Conservative rendered-page fallback for repeated corner logos."""
    if len(pages) < 2:
        return [], {}
    groups: list[dict] = []
    for page in pages:
        for candidate in _corner_candidates_for_page(page["file_path"], int(page["page_num"])):
            matched = None
            for group in groups:
                if group["anchor"] == candidate["anchor"] and _hash_distance(group["hash"], candidate["hash"]) <= 18:
                    matched = group
                    break
            if matched is None:
                groups.append({"anchor": candidate["anchor"], "hash": candidate["hash"], "items": [candidate]})
            else:
                matched["items"].append(candidate)

    min_repeats = 2 if len(pages) <= 3 else 3
    logo_assets: list[dict] = []
    regions_by_page: dict[int, list[dict]] = {}
    for idx, group in enumerate(groups, start=1):
        page_nums = sorted({item["page_num"] for item in group["items"]})
        if len(page_nums) < min_repeats:
            continue
        first = group["items"][0]
        file_path = os.path.join(output_dir, f"detected_logo_{idx:02d}_{group['anchor']}.png")
        first["crop"].save(file_path, "PNG")
        for item in group["items"]:
            regions_by_page.setdefault(int(item["page_num"]), []).append({
                "bbox": item["bbox"],
                "anchor": group["anchor"],
            })
        logo_assets.append({
            "file_path": file_path,
            "source_page_num": first["page_num"],
            "repeated_page_nums": page_nums,
            "classification": "rendered_repeated_logo",
            "role": "logo",
            "process_mode": "original",
            "asset_name": f"模板提取 Logo {idx}",
            "usage_note": "从版式模板重复角标中识别出的 Logo，作为项目级品牌标识候选。",
            "metadata": {
                "pptx_source_page_num": first["page_num"],
                "pptx_repeated_page_nums": page_nums,
                "classification": "rendered_repeated_logo",
                "template_logo_detection": "rendered_corner_repeat",
                "logo_anchor": group["anchor"],
            },
        })
    return logo_assets, regions_by_page


def extract_template_package(
    file_path: str,
    project_id: str,
    upload_dir: str,
    *,
    logo_regions_by_page: Optional[dict[int, list[dict]]] = None,
    source_filename: str | None = None,
    finalize: bool = True,
) -> Dict:
    """
    Extract template pages and stable metadata.
    Returns {pages, recommendations, document_kind, rendered_logo_candidates}.
    """
    ext = os.path.splitext(file_path)[1].lower()
    project_dir = os.path.join(upload_dir, project_id)
    final_dir = os.path.join(project_dir, "templates")
    os.makedirs(project_dir, exist_ok=True)
    output_dir = tempfile.mkdtemp(prefix="templates_work_", dir=project_dir)

    page_texts: list[str] = []
    img_paths: list[str] = []
    pdf_path = file_path
    if ext in (".ppt", ".pptx"):
        with ThreadPoolExecutor(max_workers=2, thread_name_prefix="template-pptx") as pool:
            text_future = pool.submit(_extract_pptx_texts, file_path)
            with tempfile.TemporaryDirectory() as tmpdir:
                pdf_path = convert_ppt_to_pdf(file_path, tmpdir)
                dest_pdf = os.path.join(output_dir, "converted.pdf")
                shutil.copy2(pdf_path, dest_pdf)
                pdf_path = dest_pdf
            page_texts = text_future.result()
    elif ext == ".pdf":
        with ThreadPoolExecutor(max_workers=2, thread_name_prefix="template-pdf") as pool:
            text_future = pool.submit(_extract_pdf_texts, file_path)
            if not os.path.exists(pdf_path):
                raise RuntimeError("PDF 文件不存在")
            img_paths = extract_pdf_thumbnails(pdf_path, output_dir)
            page_texts = text_future.result()

    if not os.path.exists(pdf_path):
        raise RuntimeError("PDF 文件不存在")

    if not img_paths:
        img_paths = extract_pdf_thumbnails(pdf_path, output_dir)
    total = len(img_paths)
    source_kind = infer_template_source_kind(page_texts, total)
    base_pages = [{"page_num": idx + 1, "file_path": path} for idx, path in enumerate(img_paths)]
    rendered_logos, rendered_regions = detect_repeated_corner_logos(base_pages, output_dir)

    combined_regions: dict[int, list[dict]] = {}
    for source in (logo_regions_by_page or {}, rendered_regions):
        for page_num, regions in source.items():
            combined_regions.setdefault(int(page_num), []).extend(regions or [])

    pages: list[Dict] = []
    for idx, img_path in enumerate(img_paths, start=1):
        text = page_texts[idx - 1] if idx - 1 < len(page_texts) else ""
        category = _infer_page_category(idx, total, text)
        regions = combined_regions.get(idx) or []
        layout_path = _layout_reference_path(img_path, output_dir, regions)
        pages.append({
            "page_num": idx,
            "file_path": img_path,
            "layout_file_path": layout_path,
            "url": f"/uploads/{project_id}/templates/{os.path.basename(img_path)}",
            "layout_url": f"/uploads/{project_id}/templates/{os.path.basename(layout_path)}",
            "category": category,
            "category_confidence": _category_confidence(category, idx, text),
            "source_kind": source_kind,
            "text_density": _text_density(text),
            "logo_removed": bool(regions),
            "source_document": source_filename or os.path.basename(file_path),
        })

    package = {
        "pages": pages,
        "recommendations": recommend_template_pages(pages),
        "document_kind": source_kind,
        "rendered_logo_candidates": rendered_logos,
        "work_dir": output_dir,
        "final_dir": final_dir,
    }
    if finalize:
        package = promote_template_package(package)

    logger.info(
        "TemplateExtractor: 共提取 %s 页模板缩略图 source_kind=%s rendered_logos=%s",
        len(package["pages"]),
        source_kind,
        len(rendered_logos),
    )
    return package


def _replace_path_prefix(path: str, source_dir: str, target_dir: str) -> str:
    return path.replace(source_dir, target_dir, 1) if path.startswith(source_dir) else path


def promote_template_package(package: Dict) -> Dict:
    """Move a successfully extracted work directory into the project's active template directory."""
    work_dir = package.get("work_dir")
    final_dir = package.get("final_dir")
    if not work_dir or not final_dir:
        return package
    if not os.path.exists(work_dir):
        raise RuntimeError("模板临时目录不存在")
    if os.path.exists(final_dir):
        shutil.rmtree(final_dir)
    shutil.move(work_dir, final_dir)
    for page in package.get("pages") or []:
        if isinstance(page, dict):
            if page.get("file_path"):
                page["file_path"] = _replace_path_prefix(page["file_path"], work_dir, final_dir)
            if page.get("layout_file_path"):
                page["layout_file_path"] = _replace_path_prefix(page["layout_file_path"], work_dir, final_dir)
    for logo in package.get("rendered_logo_candidates") or []:
        if isinstance(logo, dict) and logo.get("file_path"):
            logo["file_path"] = _replace_path_prefix(logo["file_path"], work_dir, final_dir)
    package["work_dir"] = None
    return package


def extract_template_images(file_path: str, project_id: str, upload_dir: str) -> List[Dict]:
    """
    从上传的 PPT/PDF 中提取每页缩略图。
    返回 [{page_num, file_path, url}, ...]
    """
    return extract_template_package(file_path, project_id, upload_dir)["pages"]


def recommend_template_pages(pages: List[Dict], content_plan: Optional[List[Dict]] = None) -> Dict[str, Optional[Dict]]:
    """
    根据页数和 Content Plan 推荐代表性页面。
    返回 cover/toc/section/content/data/quote/ending 的代表页。
    """
    total = len(pages)
    keys = ("cover", "toc", "section", "content", "data", "quote", "ending")
    if total == 0:
        return {key: None for key in keys}

    for page in pages:
        if not page.get("category"):
            page["category"] = _infer_page_category(int(page.get("page_num") or 0), total)

    def first_by_category(*categories: str) -> Optional[Dict]:
        return next((p for p in pages if p.get("category") in categories), None)

    cover = first_by_category("cover") or pages[0]
    ending = first_by_category("ending") or pages[-1]
    toc = first_by_category("toc") or (pages[1] if total >= 2 else None)
    section = first_by_category("section") or toc
    content_candidates = [p for p in pages if p.get("category") in {"content", "data"}]
    content_page = content_candidates[len(content_candidates) // 2] if content_candidates else pages[total // 2]
    data = first_by_category("data") or content_page
    quote = first_by_category("quote") or content_page

    return {
        "cover": cover,
        "toc": toc,
        "section": section,
        "content": content_page,
        "data": data,
        "quote": quote,
        "ending": ending,
    }
