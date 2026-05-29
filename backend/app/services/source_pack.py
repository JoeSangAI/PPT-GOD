import hashlib
import json
import os
import re
from io import BytesIO
from typing import Callable


LOW_TEXT_CHAR_THRESHOLD = 80
PDF_PAGE_REFERENCE_MAX_PAGES = 120
SLIDE_LIKE_ASPECT_RATIO_MIN = 1.45
SLIDE_LIKE_ASPECT_RATIO_MAX = 1.95
PDF_PAGE_REFERENCE_TARGET_WIDTH = 1920
REPEATED_PAGE_ELEMENT_MIN_PAGES = 3
REPEATED_PAGE_ELEMENT_MAX_HEIGHT = 90
REPEATED_PAGE_ELEMENT_TOP_BAND = 90
PDF_REPEATED_LOGO_MIN_AREA_RATIO = 0.00015
PDF_REPEATED_LOGO_MAX_AREA_RATIO = 0.035
PDF_REPEATED_LOGO_POSITION_TOLERANCE = 0.04
PDF_DOCUMENT_LOGO_MIN_AREA_RATIO = 0.002
PDF_DOCUMENT_LOGO_MAX_AREA_RATIO = 0.04
PDF_DOCUMENT_LOGO_MIN_WIDTH_RATIO = 0.08
PDF_DOCUMENT_LOGO_MAX_WIDTH_RATIO = 0.36
PDF_DOCUMENT_LOGO_MAX_HEIGHT_RATIO = 0.18
PDF_DOCUMENT_LOGO_MIN_ASPECT = 1.55
PDF_DOCUMENT_LOGO_MAX_ASPECT = 8.0
PDF_DOCUMENT_LOGO_EDGE_BAND = 0.18
PDF_DOCUMENT_LOGO_TERMINAL_PAGES = 2
PDF_DOCUMENT_LOGO_KEYWORDS = (
    "logo",
    "brand",
    "capital",
    "financial advisor",
    "financial adviser",
    "advisor",
    "adviser",
    "ventures",
    "partners",
    "品牌",
    "标识",
    "资本",
    "投资",
    "基金",
    "财务顾问",
)
PDF_TEXT_INLINE_LABEL_MAX_CHARS = 14
PDF_TEXT_INLINE_LABEL_GAP_MAX = 80.0
PDF_TEXT_ROW_CENTER_TOLERANCE = 20.0

FRONT_MATTER_ROLE_TERMS = {
    "preface": ("译者序", "前言", "序言", "preface"),
    "intro": ("绪论", "序论", "导论", "引言", "introduction"),
    "guide": ("阅读指南", "reading guide"),
}


def estimate_tokens(text: str) -> int:
    """Conservative planning-budget estimate; exact tokenizer is unnecessary here."""
    cleaned = re.sub(r"\s+", "", text or "")
    if not cleaned:
        return 0
    ascii_words = re.findall(r"[A-Za-z0-9_]+", cleaned)
    non_ascii_chars = len(re.sub(r"[A-Za-z0-9_]", "", cleaned))
    return max(1, int(len(ascii_words) * 1.3 + non_ascii_chars))


def _document_kind(filename: str) -> str:
    ext = os.path.splitext(filename or "")[1].lower()
    if ext == ".pdf":
        return "pdf"
    if ext in {".md", ".markdown"}:
        return "markdown"
    if ext in {".txt", ".json", ".csv", ".html", ".htm"}:
        return "text"
    if ext in {".docx", ".doc"}:
        return "docx"
    if ext in {".pptx", ".ppt"}:
        return "pptx"
    return "text"


def _safe_decode(file_bytes: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


def _safe_asset_filename(filename: str, page_num: int, xref: int, ext: str, payload: bytes) -> str:
    stem = re.sub(r"[^A-Za-z0-9_.-]+", "_", os.path.splitext(filename)[0]).strip("_") or "source"
    digest = hashlib.sha1(payload).hexdigest()[:12]
    safe_ext = re.sub(r"[^A-Za-z0-9]+", "", ext or "png") or "png"
    return f"{stem}_p{page_num}_x{xref}_{digest}.{safe_ext}"


def _safe_page_reference_filename(filename: str, page_num: int, payload: bytes) -> str:
    stem = re.sub(r"[^A-Za-z0-9_.-]+", "_", os.path.splitext(filename)[0]).strip("_") or "source"
    digest = hashlib.sha1(payload).hexdigest()[:12]
    return f"{stem}_p{page_num:03d}_page_{digest}.png"


def _text_preview(text: str, limit: int = 260) -> str:
    cleaned = re.sub(r"\s+", " ", text or "").strip()
    return cleaned[:limit]


def _bbox_metrics(bbox: list[float]) -> tuple[float, float, float]:
    if not isinstance(bbox, list) or len(bbox) < 4:
        return 0.0, 0.0, 0.0
    try:
        width = max(0.0, float(bbox[2]) - float(bbox[0]))
        height = max(0.0, float(bbox[3]) - float(bbox[1]))
    except (TypeError, ValueError):
        return 0.0, 0.0, 0.0
    return width, height, width * height


def _pdf_figure_role(width: int, height: int, bbox: list[float]) -> tuple[str, str]:
    bbox_width, bbox_height, bbox_area = _bbox_metrics(bbox)
    try:
        y0 = float(bbox[1])
    except (TypeError, ValueError, IndexError):
        y0 = 10_000.0
    pixel_area = max(0, int(width or 0)) * max(0, int(height or 0))
    large_pixels = pixel_area >= 40_000 and min(width or 0, height or 0) >= 90
    large_bbox = bbox_area >= 20_000 and min(bbox_width, bbox_height) >= 80
    tiny_pixels = pixel_area > 0 and (pixel_area < 12_000 or min(width or 0, height or 0) < 45)
    tiny_bbox = bbox_area > 0 and (bbox_area < 3_000 or min(bbox_width, bbox_height) < 32)
    top_header_like = y0 <= REPEATED_PAGE_ELEMENT_TOP_BAND and 0 < bbox_height <= REPEATED_PAGE_ELEMENT_MAX_HEIGHT
    if top_header_like:
        return "auxiliary", "low"
    if large_pixels or large_bbox:
        return "content", "high"
    if tiny_pixels or tiny_bbox:
        return "auxiliary", "low"
    return "content", "medium"


def _pdf_page_aspect_ratio(page) -> float:
    try:
        width = float(page.rect.width)
        height = float(page.rect.height)
    except Exception:
        return 0.0
    if width <= 0 or height <= 0:
        return 0.0
    return width / height


def _is_slide_like_pdf_page(page) -> bool:
    ratio = _pdf_page_aspect_ratio(page)
    return SLIDE_LIKE_ASPECT_RATIO_MIN <= ratio <= SLIDE_LIKE_ASPECT_RATIO_MAX


_PDF_TEXT_LINE_LIST_START_RE = re.compile(
    r"^(?:[•·*▪▫◦]\s*|[-–—]\s+|\d+[.)、．]\s*|[一二三四五六七八九十]+[、.)．]\s*)"
)
_PDF_TEXT_HEADING_RE = re.compile(
    r"^(?:目录|前言|序言|绪论|序论|导论|阅读指南|第\s*[0-9一二三四五六七八九十百]+\s*[章节部篇讲课](?:\s+|[:：、.\-]).*)$",
    flags=re.IGNORECASE,
)
def _pdf_text_span_needs_separator(previous: dict, current: dict) -> bool:
    previous_text = re.sub(r"\s+", "", str(previous.get("text") or ""))
    current_text = re.sub(r"\s+", "", str(current.get("text") or ""))
    if not previous_text or not current_text:
        return False
    if len(previous_text) > PDF_TEXT_INLINE_LABEL_MAX_CHARS or len(current_text) < 8:
        return False
    try:
        previous_x0, previous_y0, previous_x1, _previous_y1 = [float(value) for value in previous.get("bbox")[:4]]
        current_x0, current_y0, _current_x1, _current_y1 = [float(value) for value in current.get("bbox")[:4]]
        previous_size = float(previous.get("size") or 0.0)
        current_size = float(current.get("size") or 0.0)
    except (TypeError, ValueError):
        return False
    starts_after_label = current_x0 >= previous_x0 + max(16.0, previous_size * 1.1)
    visually_distinct = abs(current_y0 - previous_y0) >= 3.0 or abs(current_size - previous_size) >= 0.5
    boundary_is_text = bool(re.search(r"[\w\u4e00-\u9fff]$", previous_text) and re.match(r"^[\w\u4e00-\u9fff]", current_text))
    return bool(starts_after_label and visually_distinct and boundary_is_text)


def _pdf_text_line_from_spans(line: dict) -> str:
    parts: list[str] = []
    previous_span: dict | None = None
    for span in line.get("spans") or []:
        text = str(span.get("text") or "")
        if not text:
            continue
        if previous_span and _pdf_text_span_needs_separator(previous_span, span):
            parts.append(" ")
        parts.append(text)
        previous_span = span
    return "".join(parts).strip()


def _pdf_text_is_heading_line(line: str) -> bool:
    text = str(line or "").strip()
    if not text:
        return False
    if _PDF_TEXT_HEADING_RE.match(text):
        return True
    return len(text) <= 12 and bool(re.search(r"(章|节|团队|目录|融资|优势|路线图|策略|计划|介绍)$", text))


def _pdf_join_wrapped_text(left: str, right: str) -> str:
    left = str(left or "").rstrip()
    right = str(right or "").lstrip()
    if not left:
        return right
    if not right:
        return left
    if left.endswith("-"):
        return left[:-1] + right
    if re.search(r"[，,、；;：:。！？!?）)\]}】]$", left):
        return left + right
    if re.search(r"[\u4e00-\u9fff]$", left) or re.match(r"^[\u4e00-\u9fff，,、；;：:。！？!?）)\]}】]", right):
        return left + right
    if re.search(r"[A-Za-z]{2,}$", left) and re.match(r"^[A-Za-z]{1,8}[\u4e00-\u9fff]", right):
        return left + right
    return left + " " + right


def _pdf_join_text_block_lines(lines: list[str]) -> str:
    clean_lines = [str(line or "").strip() for line in lines if str(line or "").strip()]
    if not clean_lines:
        return ""
    if len(clean_lines) == 1:
        return clean_lines[0]

    output: list[str] = []
    for line in clean_lines:
        starts_list = bool(_PDF_TEXT_LINE_LIST_START_RE.match(line))
        if not output or starts_list or _pdf_text_is_heading_line(output[-1]) or _pdf_text_is_heading_line(line):
            output.append(line)
        else:
            output[-1] = _pdf_join_wrapped_text(output[-1], line)
    return "\n".join(output)


def _pdf_text_block_height(block: dict) -> float:
    return max(0.0, float(block.get("y1") or 0.0) - float(block.get("y0") or 0.0))


def _pdf_text_block_overlaps_row(block: dict, row: dict) -> bool:
    y0 = float(block.get("y0") or 0.0)
    y1 = float(block.get("y1") or 0.0)
    row_y0 = float(row.get("y0") or 0.0)
    row_y1 = float(row.get("y1") or 0.0)
    if _pdf_text_row_is_label_only(row):
        vertical_gap = y0 - row_y1
        if 0.0 <= vertical_gap <= 18.0:
            return True
    overlap = min(y1, row_y1) - max(y0, row_y0)
    if overlap > 0:
        min_height = max(1.0, min(y1 - y0, row_y1 - row_y0))
        return overlap / min_height >= 0.25
    center = (y0 + y1) / 2.0
    row_center = (row_y0 + row_y1) / 2.0
    return abs(center - row_center) <= PDF_TEXT_ROW_CENTER_TOLERANCE


def _pdf_text_row_is_label_only(row: dict) -> bool:
    blocks = row.get("blocks") if isinstance(row.get("blocks"), list) else []
    if not blocks:
        return False
    for block in blocks:
        text = str(block.get("text") or "").strip()
        if not text or "\n" in text or len(re.sub(r"\s+", "", text)) > PDF_TEXT_INLINE_LABEL_MAX_CHARS:
            return False
    return True


def _pdf_text_rows(blocks: list[dict]) -> list[list[dict]]:
    rows: list[dict] = []
    for block in sorted(blocks, key=lambda item: (float(item.get("y0") or 0.0), float(item.get("x0") or 0.0))):
        matched = None
        for row in rows:
            if _pdf_text_block_overlaps_row(block, row):
                matched = row
                break
        if not matched:
            matched = {
                "y0": float(block.get("y0") or 0.0),
                "y1": float(block.get("y1") or 0.0),
                "blocks": [],
            }
            rows.append(matched)
        matched["blocks"].append(block)
        matched["y0"] = min(float(matched["y0"]), float(block.get("y0") or 0.0))
        matched["y1"] = max(float(matched["y1"]), float(block.get("y1") or 0.0))
    return [
        sorted(row["blocks"], key=lambda item: (float(item.get("x0") or 0.0), float(item.get("y0") or 0.0)))
        for row in sorted(rows, key=lambda item: (float(item.get("y0") or 0.0), float(item.get("y1") or 0.0)))
    ]


def _pdf_text_is_inline_label(block: dict, next_block: dict) -> bool:
    text = re.sub(r"\s+", "", str(block.get("text") or ""))
    if not text or "\n" in str(block.get("text") or ""):
        return False
    if len(text) > PDF_TEXT_INLINE_LABEL_MAX_CHARS:
        return False
    gap = float(next_block.get("x0") or 0.0) - float(block.get("x1") or 0.0)
    if gap < -2.0 or gap > PDF_TEXT_INLINE_LABEL_GAP_MAX:
        return False
    next_text = str(next_block.get("text") or "").strip()
    if len(next_text) < 18:
        return False
    if _pdf_text_block_height(next_block) < max(1.0, _pdf_text_block_height(block)) * 1.2:
        return False
    label_center = (float(block.get("y0") or 0.0) + float(block.get("y1") or 0.0)) / 2.0
    return float(next_block.get("y0") or 0.0) - 18.0 <= label_center <= float(next_block.get("y1") or 0.0) + 18.0


def _pdf_text_row_text(blocks: list[dict]) -> str:
    units: list[str] = []
    idx = 0
    while idx < len(blocks):
        block = blocks[idx]
        text = str(block.get("text") or "").strip()
        if not text:
            idx += 1
            continue
        if idx + 1 < len(blocks) and _pdf_text_is_inline_label(block, blocks[idx + 1]):
            units.append(f"{text} {str(blocks[idx + 1].get('text') or '').strip()}")
            idx += 2
            continue
        units.append(text)
        idx += 1
    if not units:
        return ""
    row_y1 = max(float(block.get("y1") or 0.0) for block in blocks)
    if row_y1 <= 70.0 and len(units) > 1 and all("\n" not in unit and len(unit) <= 32 for unit in units):
        return " ".join(units)
    return "\n".join(units)


def _pdf_page_layout_text(page) -> str:
    try:
        payload = page.get_text("dict")
    except Exception:
        return ""
    blocks: list[dict] = []
    for block in payload.get("blocks") or []:
        if block.get("type") != 0:
            continue
        lines = [
            line_text
            for line in block.get("lines") or []
            if (line_text := _pdf_text_line_from_spans(line))
        ]
        text = _pdf_join_text_block_lines(lines)
        if not text:
            continue
        try:
            x0, y0, x1, y1 = [float(value) for value in block.get("bbox")[:4]]
        except (TypeError, ValueError):
            continue
        blocks.append({
            "text": text,
            "x0": x0,
            "y0": y0,
            "x1": x1,
            "y1": y1,
        })
    row_texts = [
        row_text
        for row in _pdf_text_rows(blocks)
        if (row_text := _pdf_text_row_text(row))
    ]
    return "\n".join(row_texts).strip()


def _pdf_page_text(page) -> str:
    layout_text = _pdf_page_layout_text(page)
    if layout_text:
        return layout_text
    try:
        return (page.get_text("text", sort=True) or "").strip()
    except TypeError:
        return (page.get_text("text") or "").strip()


def _should_add_pdf_page_references(doc) -> bool:
    page_count = len(doc)
    if page_count <= 0 or page_count > PDF_PAGE_REFERENCE_MAX_PAGES:
        return False
    checked = 0
    slide_like = 0
    for page in doc:
        checked += 1
        if _is_slide_like_pdf_page(page):
            slide_like += 1
        if checked >= min(page_count, 12):
            break
    return checked > 0 and slide_like / checked >= 0.8


def _render_pdf_page_reference(page, *, filename: str, page_num: int, text: str, asset_output_dir: str | None) -> dict:
    width = max(float(page.rect.width), 1.0)
    zoom = max(1.0, min(3.0, PDF_PAGE_REFERENCE_TARGET_WIDTH / width))
    pix = page.get_pixmap(matrix=__import__("fitz").Matrix(zoom, zoom), alpha=False)
    payload = pix.tobytes("png")
    file_path = ""
    if asset_output_dir:
        asset_filename = _safe_page_reference_filename(filename, page_num, payload)
        file_path = os.path.join(asset_output_dir, asset_filename)
        if not os.path.exists(file_path):
            with open(file_path, "wb") as f:
                f.write(payload)
    bbox = [float(page.rect.x0), float(page.rect.y0), float(page.rect.x1), float(page.rect.y1)]
    bbox_width, bbox_height, bbox_area = _bbox_metrics(bbox)
    return {
        "id": f"{filename}:p{page_num}:page",
        "source_type": "pdf",
        "source_document": filename,
        "source_page_num": page_num,
        "pdf_source_page_num": page_num,
        "chapter_id": "",
        "bbox": bbox,
        "bbox_width": bbox_width,
        "bbox_height": bbox_height,
        "bbox_area": bbox_area,
        "page_width": float(page.rect.width),
        "page_height": float(page.rect.height),
        "image_width": pix.width,
        "image_height": pix.height,
        "pixel_area": max(0, pix.width) * max(0, pix.height),
        "figure_role": "source_page",
        "content_significance": "high",
        "nearby_text": _text_preview(text),
        "file_path": file_path,
        "image_sha1": hashlib.sha1(payload).hexdigest(),
        "asset_kind": "source_page_image",
        "is_full_page_reference": True,
    }


def _image_bbox_edges(image: dict) -> tuple[float, float, float, float]:
    bbox = image.get("bbox") if isinstance(image.get("bbox"), list) else []
    try:
        x0, y0, x1, y1 = [float(value) for value in bbox[:4]]
    except (TypeError, ValueError):
        return 0.0, 0.0, 0.0, 0.0
    return x0, y0, x1, y1


def _image_page_size(image: dict) -> tuple[float, float]:
    try:
        page_width = float(image.get("page_width") or 0.0)
        page_height = float(image.get("page_height") or 0.0)
    except (TypeError, ValueError):
        return 0.0, 0.0
    return max(0.0, page_width), max(0.0, page_height)


def _image_bbox_norm(image: dict) -> tuple[float, float, float, float]:
    page_width, page_height = _image_page_size(image)
    if page_width <= 0 or page_height <= 0:
        return 0.0, 0.0, 0.0, 0.0
    x0, y0, x1, y1 = _image_bbox_edges(image)
    return (
        x0 / page_width,
        y0 / page_height,
        max(0.0, x1 - x0) / page_width,
        max(0.0, y1 - y0) / page_height,
    )


def _pdf_logo_anchor_from_image(image: dict) -> str:
    left, top, width, height = _image_bbox_norm(image)
    cx = left + width / 2
    cy = top + height / 2
    horizontal = "left" if cx < 0.5 else "right"
    vertical = "top" if cy < 0.5 else "bottom"
    return f"{vertical}-{horizontal}"


def _looks_like_repeated_pdf_logo(siblings: list[dict], pages: set[int]) -> bool:
    if len(pages) < REPEATED_PAGE_ELEMENT_MIN_PAGES:
        return False
    bounds: list[tuple[float, float, float, float]] = []
    for image in siblings:
        if image.get("asset_kind") == "source_page_image" or image.get("figure_role") == "source_page":
            return False
        if not _looks_like_repeated_page_element(image):
            return False
        left, top, width, height = _image_bbox_norm(image)
        if width <= 0 or height <= 0:
            return False
        area_ratio = width * height
        if area_ratio < PDF_REPEATED_LOGO_MIN_AREA_RATIO or area_ratio > PDF_REPEATED_LOGO_MAX_AREA_RATIO:
            return False
        aspect = width / max(height, 0.0001)
        if aspect > 12 or aspect < 0.1:
            return False
        if width > 0.32 or height > 0.16:
            return False
        cx = left + width / 2
        cy = top + height / 2
        in_brand_zone = cy <= 0.18 or cy >= 0.82 or cx <= 0.18 or cx >= 0.82
        if not in_brand_zone:
            return False
        bounds.append((left, top, width, height))
    if not bounds:
        return False
    lefts = [bound[0] for bound in bounds]
    tops = [bound[1] for bound in bounds]
    widths = [bound[2] for bound in bounds]
    heights = [bound[3] for bound in bounds]
    return (
        max(lefts) - min(lefts) <= PDF_REPEATED_LOGO_POSITION_TOLERANCE
        and max(tops) - min(tops) <= PDF_REPEATED_LOGO_POSITION_TOLERANCE
        and max(widths) - min(widths) <= PDF_REPEATED_LOGO_POSITION_TOLERANCE
        and max(heights) - min(heights) <= PDF_REPEATED_LOGO_POSITION_TOLERANCE
    )


def _looks_like_repeated_page_element(image: dict) -> bool:
    if image.get("asset_kind") == "source_page_image" or image.get("figure_role") == "source_page":
        return False
    _x0, y0, _x1, y1 = _image_bbox_edges(image)
    try:
        bbox_height = float(image.get("bbox_height") or 0)
    except (TypeError, ValueError):
        bbox_height = 0.0
    if not bbox_height:
        bbox_height = max(0.0, y1 - y0)
    return y0 <= REPEATED_PAGE_ELEMENT_TOP_BAND or bbox_height <= REPEATED_PAGE_ELEMENT_MAX_HEIGHT


def _mark_repeated_page_elements(images: list[dict]) -> None:
    by_hash: dict[str, list[dict]] = {}
    for image in images:
        image_hash = str(image.get("image_sha1") or "").strip()
        if not image_hash:
            continue
        by_hash.setdefault(image_hash, []).append(image)

    for siblings in by_hash.values():
        pages: set[int] = set()
        for image in siblings:
            try:
                page_num = int(image.get("source_page_num") or 0)
            except (TypeError, ValueError):
                page_num = 0
            if page_num > 0:
                pages.add(page_num)
        if len(pages) < REPEATED_PAGE_ELEMENT_MIN_PAGES:
            continue
        if not any(_looks_like_repeated_page_element(image) for image in siblings):
            continue
        is_logo = _looks_like_repeated_pdf_logo(siblings, pages)
        for image in siblings:
            if image.get("asset_kind") == "source_page_image":
                continue
            if is_logo:
                image["asset_kind"] = "pdf_logo"
                image["figure_role"] = "logo"
                image["content_significance"] = "high"
                image["classification"] = "pdf_repeated_logo"
                image["logo_anchor"] = _pdf_logo_anchor_from_image(image)
            else:
                image["asset_kind"] = "repeated_page_element"
                image["figure_role"] = "decorative"
                image["content_significance"] = "low"
            image["repeated_page_count"] = len(pages)


def _pdf_document_page_count(pages: list[dict]) -> int:
    page_nums: list[int] = []
    for page in pages:
        try:
            page_nums.append(int(page.get("page_num") or 0))
        except (TypeError, ValueError):
            continue
    return max(page_nums) if page_nums else 0


def _pdf_image_source_page_num(image: dict) -> int:
    try:
        return int(image.get("source_page_num") or 0)
    except (TypeError, ValueError):
        return 0


def _pdf_text_has_brand_signal(text: str) -> bool:
    cleaned = str(text or "").strip().lower()
    if not cleaned:
        return False
    return any(keyword.lower() in cleaned for keyword in PDF_DOCUMENT_LOGO_KEYWORDS)


def _is_terminal_pdf_page(page_num: int, total_pages: int) -> bool:
    if page_num <= 0 or total_pages <= 0:
        return False
    return (
        page_num <= PDF_DOCUMENT_LOGO_TERMINAL_PAGES
        or page_num >= max(1, total_pages - PDF_DOCUMENT_LOGO_TERMINAL_PAGES + 1)
    )


def _looks_like_pdf_document_logo_candidate(image: dict, total_pages: int) -> bool:
    if image.get("asset_kind") in {"source_page_image", "repeated_page_element", "pdf_logo"}:
        return False
    if image.get("figure_role") == "source_page":
        return False
    page_num = _pdf_image_source_page_num(image)
    if not _is_terminal_pdf_page(page_num, total_pages):
        return False
    if not _pdf_text_has_brand_signal(str(image.get("nearby_text") or "")):
        return False
    left, top, width, height = _image_bbox_norm(image)
    if width <= 0 or height <= 0:
        return False
    area_ratio = width * height
    if area_ratio < PDF_DOCUMENT_LOGO_MIN_AREA_RATIO or area_ratio > PDF_DOCUMENT_LOGO_MAX_AREA_RATIO:
        return False
    if width < PDF_DOCUMENT_LOGO_MIN_WIDTH_RATIO or width > PDF_DOCUMENT_LOGO_MAX_WIDTH_RATIO:
        return False
    if height > PDF_DOCUMENT_LOGO_MAX_HEIGHT_RATIO:
        return False
    aspect = width / max(height, 0.0001)
    if aspect < PDF_DOCUMENT_LOGO_MIN_ASPECT or aspect > PDF_DOCUMENT_LOGO_MAX_ASPECT:
        return False
    center_x = left + width / 2
    center_y = top + height / 2
    in_brand_zone = (
        center_y <= PDF_DOCUMENT_LOGO_EDGE_BAND
        or center_y >= 1 - PDF_DOCUMENT_LOGO_EDGE_BAND
        or center_x <= PDF_DOCUMENT_LOGO_EDGE_BAND
        or center_x >= 1 - PDF_DOCUMENT_LOGO_EDGE_BAND
    )
    return bool(in_brand_zone)


def _pdf_document_logo_score(image: dict, total_pages: int) -> float:
    left, top, width, height = _image_bbox_norm(image)
    page_num = _pdf_image_source_page_num(image)
    score = 0.0
    if page_num in {1, total_pages}:
        score += 2.0
    elif _is_terminal_pdf_page(page_num, total_pages):
        score += 1.0
    if top <= PDF_DOCUMENT_LOGO_EDGE_BAND:
        score += 1.0
    if left <= PDF_DOCUMENT_LOGO_EDGE_BAND or left + width >= 1 - PDF_DOCUMENT_LOGO_EDGE_BAND:
        score += 1.0
    score += min(width / 0.24, 1.0)
    score += min((width * height) / 0.022, 1.0)
    return score


def _mark_pdf_document_logo_candidates(images: list[dict], pages: list[dict]) -> None:
    if any(image.get("asset_kind") == "pdf_logo" or image.get("classification") == "pdf_repeated_logo" for image in images):
        return
    total_pages = _pdf_document_page_count(pages)
    candidates = [
        image for image in images
        if _looks_like_pdf_document_logo_candidate(image, total_pages)
    ]
    if not candidates:
        return
    best = max(candidates, key=lambda image: _pdf_document_logo_score(image, total_pages))
    best["asset_kind"] = "pdf_logo"
    best["figure_role"] = "logo"
    best["content_significance"] = "high"
    best["classification"] = "pdf_document_logo"
    best["logo_anchor"] = _pdf_logo_anchor_from_image(best)
    best["logo_detection"] = "terminal_page_brand_lockup"


def _looks_like_toc_line(line: str) -> bool:
    text = line.strip()
    if re.search(r"(?:\.{2,}|…{2,}|·{2,}|_{2,})\s*\d+\s*$", text):
        return True
    if re.search(r"\s+\d{1,4}\s*$", text) and len(text) > 18:
        return True
    return False


def _normalize_chapter_title(raw: str) -> str:
    title = re.sub(r"^\s{0,3}#{1,6}\s*", "", raw or "").strip()
    title = re.sub(r"\s+", " ", title)
    return title


def _heading_from_line(line: str) -> str:
    text = _normalize_chapter_title(line)
    if not text or _looks_like_toc_line(text):
        return ""
    markdown_match = re.match(r"^\s{0,3}#{1,6}\s+(.+)$", line or "")
    if markdown_match:
        return _normalize_chapter_title(markdown_match.group(1))
    if re.match(r"^第\s*[0-9一二三四五六七八九十百]+\s*[章节部篇讲课](?:\s+|[:：、.\-]).+", text):
        return text
    if re.match(r"^(?:第\s*)?[0-9一二三四五六七八九十百]+\s*部分(?:\s+|[:：、.\-]).+", text):
        return text
    if re.match(r"^[一二三四五六七八九十]+[、.]\s*.+", text):
        return text
    return ""


def _line_starts_section_term(raw: str, term: str) -> bool:
    text = re.sub(r"\s+", " ", str(raw or "").strip()).lower()
    clean_term = re.sub(r"\s+", " ", term or "").strip().lower()
    if not text or not clean_term:
        return False
    if text == clean_term:
        return True
    return bool(re.match(rf"^{re.escape(clean_term)}(?:\s+|[:：、.\-]).+", text))


def front_matter_role_for_title(title: str, roles: set[str] | None = None) -> tuple[str, str] | None:
    allowed_roles = roles or set(FRONT_MATTER_ROLE_TERMS)
    clean = _normalize_chapter_title(title)
    if not clean or _looks_like_toc_line(clean):
        return None
    for role in allowed_roles:
        for term in FRONT_MATTER_ROLE_TERMS.get(role, ()):
            if _line_starts_section_term(clean, term) and len(clean) <= 36:
                return role, term
    return None


def front_matter_role_for_text(text: str, roles: set[str] | None = None) -> tuple[str, str] | None:
    for raw_line in str(text or "").splitlines()[:8]:
        raw = str(raw_line or "").strip()
        if not raw or _looks_like_toc_line(raw):
            continue
        match = front_matter_role_for_title(raw, roles)
        if match:
            return match
    return None


def _chapters_from_pages(pages: list[dict], *, pdf_toc: list | None = None) -> list[dict]:
    candidates: list[dict] = []
    if pdf_toc:
        for item in pdf_toc:
            try:
                level, title, page_num = item[:3]
                page_num = int(page_num)
            except (TypeError, ValueError):
                continue
            if int(level or 1) > 2 or page_num <= 0:
                continue
            title = _normalize_chapter_title(str(title or ""))
            if title:
                candidates.append({"title": title, "start_page": page_num})

    if not candidates:
        for page in pages:
            for line in str(page.get("text") or "").splitlines()[:12]:
                title = _heading_from_line(line)
                if title:
                    candidates.append({"title": title, "start_page": int(page.get("page_num") or 1)})

    deduped: list[dict] = []
    seen: set[tuple[str, int]] = set()
    for item in sorted(candidates, key=lambda c: int(c.get("start_page") or 0)):
        key = (str(item.get("title") or ""), int(item.get("start_page") or 0))
        if key in seen:
            continue
        seen.add(key)
        deduped.append(item)

    max_page = max((int(page.get("page_num") or 0) for page in pages), default=0)
    chapters: list[dict] = []
    for idx, item in enumerate(deduped):
        start = int(item.get("start_page") or 1)
        next_start = int(deduped[idx + 1].get("start_page") or max_page + 1) if idx + 1 < len(deduped) else max_page + 1
        chapters.append({
            "chapter_id": f"c{idx + 1}",
            "title": str(item.get("title") or "").strip(),
            "start_page": start,
            "end_page": max(start, next_start - 1),
        })
    return chapters


def _safe_int(value, default: int = 0) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _page_num(page: dict) -> int:
    return _safe_int(page.get("page_num"))


def _source_section_record(
    *,
    section_id: str,
    section_role: str,
    title: str,
    start_page: int,
    end_page: int,
    source: str,
    chapter_id: str = "",
) -> dict:
    record = {
        "section_id": section_id,
        "section_role": section_role,
        "title": title,
        "start_page": start_page,
        "end_page": max(start_page, end_page),
        "source": source,
    }
    if chapter_id:
        record["chapter_id"] = chapter_id
    return record


def _source_structure_from_pages(pages: list[dict], chapters: list[dict]) -> list[dict]:
    sorted_pages = sorted([page for page in pages if isinstance(page, dict)], key=_page_num)
    max_page = max((_page_num(page) for page in sorted_pages), default=0)
    sections: list[dict] = []

    for idx, chapter in enumerate(chapters or [], start=1):
        if not isinstance(chapter, dict):
            continue
        start_page = _safe_int(chapter.get("start_page"))
        end_page = _safe_int(chapter.get("end_page"), start_page)
        title = _normalize_chapter_title(str(chapter.get("title") or ""))
        if start_page <= 0 or not title:
            continue
        front_role = front_matter_role_for_title(title)
        if front_role:
            role, canonical_title = front_role
            sections.append(_source_section_record(
                section_id=f"front-{role}-{start_page}",
                section_role=role,
                title=canonical_title,
                start_page=start_page,
                end_page=end_page,
                source="chapter_outline",
                chapter_id=str(chapter.get("chapter_id") or ""),
            ))
        else:
            sections.append(_source_section_record(
                section_id=str(chapter.get("chapter_id") or f"c{idx}"),
                section_role="chapter",
                title=title,
                start_page=start_page,
                end_page=end_page,
                source="chapter_outline",
                chapter_id=str(chapter.get("chapter_id") or f"c{idx}"),
            ))

    existing_starts = {
        (_safe_int(section.get("start_page")), str(section.get("section_role") or ""))
        for section in sections
    }
    page_front_starts: list[tuple[int, str, str]] = []
    for page in sorted_pages:
        page_num = _page_num(page)
        if page_num <= 0:
            continue
        match = front_matter_role_for_text(str(page.get("text") or ""))
        if not match:
            continue
        role, title = match
        if (page_num, role) in existing_starts:
            continue
        page_front_starts.append((page_num, role, title))

    body_chapter_starts = sorted(
        _safe_int(section.get("start_page"))
        for section in sections
        if str(section.get("section_role") or "") == "chapter" and _safe_int(section.get("start_page")) > 0
    )
    front_start_pages = [start for start, _role, _title in page_front_starts]
    boundaries = sorted({*body_chapter_starts, *front_start_pages, max_page + 1})
    for start_page, role, title in page_front_starts:
        next_boundaries = [page_num for page_num in boundaries if page_num > start_page]
        end_page = (min(next_boundaries) - 1) if next_boundaries else max_page
        sections.append(_source_section_record(
            section_id=f"front-{role}-{start_page}",
            section_role=role,
            title=title,
            start_page=start_page,
            end_page=end_page,
            source="page_heading",
        ))

    if not sections and sorted_pages:
        first_page = _page_num(sorted_pages[0])
        last_page = _page_num(sorted_pages[-1])
        sections.append(_source_section_record(
            section_id="document",
            section_role="document",
            title="全文",
            start_page=first_page,
            end_page=last_page,
            source="document",
        ))

    deduped: list[dict] = []
    seen: set[tuple[int, str, str]] = set()
    for section in sorted(sections, key=lambda item: (_safe_int(item.get("start_page")), str(item.get("section_id") or ""))):
        key = (
            _safe_int(section.get("start_page")),
            str(section.get("section_role") or ""),
            str(section.get("title") or ""),
        )
        if key in seen:
            continue
        seen.add(key)
        deduped.append(section)
    return deduped


def _section_for_page(sections: list[dict], page_num: int) -> dict | None:
    matches = [
        section for section in sections
        if _safe_int(section.get("start_page")) <= page_num <= _safe_int(section.get("end_page"), page_num)
    ]
    if not matches:
        return None
    return sorted(matches, key=lambda section: _safe_int(section.get("start_page")), reverse=True)[0]


def _apply_source_structure(pages: list[dict], images: list[dict], sections: list[dict]) -> None:
    for page in pages:
        if not isinstance(page, dict):
            continue
        section = _section_for_page(sections, _page_num(page))
        if not section:
            continue
        page["source_section_id"] = section.get("section_id")
        page["source_section_role"] = section.get("section_role")
        page["source_section_title"] = section.get("title")
        page["source_section_start_page"] = section.get("start_page")
        page["source_section_end_page"] = section.get("end_page")

    for image in images:
        if not isinstance(image, dict):
            continue
        page_num = _safe_int(image.get("source_page_num") or image.get("pdf_source_page_num"))
        section = _section_for_page(sections, page_num)
        if not section:
            continue
        image["source_section_id"] = section.get("section_id")
        image["source_section_role"] = section.get("section_role")
        image["source_section_title"] = section.get("title")


def _base_pack(filename: str, kind: str, pages: list[dict], chapters: list[dict], images: list[dict]) -> dict:
    source_structure = _source_structure_from_pages(pages, chapters)
    _apply_source_structure(pages, images, source_structure)
    text_chars = sum(int(page.get("text_chars") or len(str(page.get("text") or ""))) for page in pages)
    estimated = sum(int(page.get("estimated_tokens") or estimate_tokens(str(page.get("text") or ""))) for page in pages)
    return {
        "schema_version": 1,
        "document": {
            "filename": filename,
            "kind": kind,
        },
        "parse_status": {
            "status": "completed",
        },
        "stats": {
            "pages": len(pages),
            "chapters": len(chapters),
            "images": len(images),
            "text_chars": text_chars,
            "estimated_tokens": estimated,
        },
        "pages": pages,
        "chapters": chapters,
        "source_structure": source_structure,
        "images": images,
    }


def _build_text_pack(file_bytes: bytes, filename: str, kind: str) -> dict:
    text = _safe_decode(file_bytes).strip()
    page = {
        "page_num": 1,
        "text": text,
        "text_chars": len(text),
        "estimated_tokens": estimate_tokens(text),
        "ocr_status": "not_applicable",
    }
    chapters = _chapters_from_pages([page])
    return _base_pack(filename, kind, [page], chapters, [])


def _build_docx_pack(file_bytes: bytes, filename: str) -> dict:
    try:
        from docx import Document

        doc = Document(BytesIO(file_bytes))
        text = "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text and p.text.strip())
    except Exception:
        text = _safe_decode(file_bytes)
    page = {
        "page_num": 1,
        "text": text.strip(),
        "text_chars": len(text.strip()),
        "estimated_tokens": estimate_tokens(text),
        "ocr_status": "not_applicable",
    }
    chapters = _chapters_from_pages([page])
    return _base_pack(filename, "docx", [page], chapters, [])


def _append_unique_text(parts: list[str], text: str) -> None:
    cleaned = re.sub(r"[ \t]+", " ", str(text or "")).strip()
    if cleaned and cleaned not in parts:
        parts.append(cleaned)


def _extract_pptx_shape_texts(shapes, shape_type_enum) -> list[str]:
    from pptx.enum.shapes import PP_PLACEHOLDER

    texts: list[str] = []
    for shape in shapes:
        if getattr(shape, "shape_type", None) == shape_type_enum.GROUP:
            for text in _extract_pptx_shape_texts(shape.shapes, shape_type_enum):
                _append_unique_text(texts, text)
            continue

        if getattr(shape, "is_placeholder", False):
            ph_type = getattr(getattr(shape, "placeholder_format", None), "type", None)
            if ph_type in {PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.HEADER, PP_PLACEHOLDER.SLIDE_NUMBER}:
                continue

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = []
                for cell in row.cells:
                    cell_text = re.sub(r"\s+", " ", cell.text or "").strip()
                    if cell_text:
                        cells.append(cell_text)
                if cells:
                    _append_unique_text(texts, " | ".join(cells))
            continue

        if getattr(shape, "has_chart", False):
            chart_title = getattr(getattr(shape.chart, "chart_title", None), "text_frame", None)
            if chart_title is not None:
                _append_unique_text(texts, chart_title.text)

        if hasattr(shape, "text"):
            _append_unique_text(texts, shape.text)
    return texts


def _extract_pptx_notes(slide) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes = []
        for shape in slide.notes_slide.notes_text_frame.paragraphs:
            line = "".join(run.text for run in shape.runs).strip()
            if line:
                notes.append(line)
        return "\n".join(notes).strip()
    except Exception:
        return ""


def _build_pptx_pack(file_bytes: bytes, filename: str) -> dict:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(BytesIO(file_bytes))
        pages: list[dict] = []
        for page_num, slide in enumerate(prs.slides, start=1):
            slide_texts = _extract_pptx_shape_texts(slide.shapes, MSO_SHAPE_TYPE)
            notes = _extract_pptx_notes(slide)
            if notes:
                slide_texts.append("【备注】\n" + notes)
            body = "\n".join(slide_texts).strip()
            pages.append({
                "page_num": page_num,
                "text": body,
                "text_chars": len(body),
                "estimated_tokens": estimate_tokens(body),
                "ocr_status": "not_applicable",
            })
    except Exception:
        body = _safe_decode(file_bytes).strip()
        pages = [{
            "page_num": 1,
            "text": body,
            "text_chars": len(body),
            "estimated_tokens": estimate_tokens(body),
            "ocr_status": "not_applicable",
        }]
    chapters = _chapters_from_pages(pages)
    return _base_pack(filename, "pptx", pages, chapters, [])


def _build_pdf_pack(
    file_bytes: bytes,
    filename: str,
    *,
    asset_output_dir: str | None = None,
    ocr_reader: Callable[[object], str] | None = None,
) -> dict:
    import fitz

    pages: list[dict] = []
    images: list[dict] = []
    os.makedirs(asset_output_dir, exist_ok=True) if asset_output_dir else None
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        toc = doc.get_toc(simple=True)
        add_page_references = _should_add_pdf_page_references(doc)
        for index, page in enumerate(doc, start=1):
            text = _pdf_page_text(page)
            ocr_status = "not_needed"
            if len(text) < LOW_TEXT_CHAR_THRESHOLD:
                if ocr_reader:
                    try:
                        ocr_text = (ocr_reader(page) or "").strip()
                    except Exception:
                        ocr_text = ""
                    if ocr_text:
                        text = "\n\n".join(part for part in [text, ocr_text] if part)
                        ocr_status = "completed"
                    else:
                        ocr_status = "needed"
                else:
                    ocr_status = "needed"
            page_record = {
                "page_num": index,
                "text": text,
                "text_chars": len(text),
                "estimated_tokens": estimate_tokens(text),
                "ocr_status": ocr_status,
            }
            pages.append(page_record)

            if add_page_references:
                images.append(_render_pdf_page_reference(
                    page,
                    filename=filename,
                    page_num=index,
                    text=text,
                    asset_output_dir=asset_output_dir,
                ))

            for image_info in page.get_images(full=True):
                xref = int(image_info[0])
                try:
                    extracted = doc.extract_image(xref)
                except Exception:
                    extracted = {}
                payload = extracted.get("image") if isinstance(extracted, dict) else None
                if not payload:
                    continue
                ext = str(extracted.get("ext") or "png")
                try:
                    image_width = int(extracted.get("width") or 0)
                    image_height = int(extracted.get("height") or 0)
                except (TypeError, ValueError):
                    image_width = 0
                    image_height = 0
                rects = []
                try:
                    rects = page.get_image_rects(xref)
                except Exception:
                    rects = []
                if not rects:
                    rects = [page.rect]
                for rect_idx, rect in enumerate(rects, start=1):
                    file_path = ""
                    if asset_output_dir:
                        asset_filename = _safe_asset_filename(filename, index, xref, ext, payload)
                        if rect_idx > 1:
                            stem, suffix = os.path.splitext(asset_filename)
                            asset_filename = f"{stem}_{rect_idx}{suffix}"
                        file_path = os.path.join(asset_output_dir, asset_filename)
                        if not os.path.exists(file_path):
                            with open(file_path, "wb") as f:
                                f.write(payload)
                    bbox = [float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1)]
                    bbox_width, bbox_height, bbox_area = _bbox_metrics(bbox)
                    figure_role, content_significance = _pdf_figure_role(image_width, image_height, bbox)
                    images.append({
                        "id": f"{filename}:p{index}:x{xref}:{rect_idx}",
                        "source_type": "pdf",
                        "source_document": filename,
                        "source_page_num": index,
                        "pdf_source_page_num": index,
                        "chapter_id": "",
                        "bbox": bbox,
                        "bbox_width": bbox_width,
                        "bbox_height": bbox_height,
                        "bbox_area": bbox_area,
                        "page_width": float(page.rect.width),
                        "page_height": float(page.rect.height),
                        "image_width": image_width,
                        "image_height": image_height,
                        "pixel_area": max(0, image_width) * max(0, image_height),
                        "figure_role": figure_role,
                        "content_significance": content_significance,
                        "nearby_text": _text_preview(text),
                        "file_path": file_path,
                        "image_sha1": hashlib.sha1(payload).hexdigest(),
                        "asset_kind": "document_image",
                    })
        _mark_repeated_page_elements(images)
        _mark_pdf_document_logo_candidates(images, pages)
        chapters = _chapters_from_pages(pages, pdf_toc=toc)

    for image in images:
        page_num = int(image.get("source_page_num") or 0)
        chapter = next((c for c in chapters if int(c.get("start_page") or 0) <= page_num <= int(c.get("end_page") or 0)), None)
        if chapter:
            image["chapter_id"] = chapter.get("chapter_id") or ""
    return _base_pack(filename, "pdf", pages, chapters, images)


def build_source_pack(
    file_bytes: bytes,
    filename: str,
    *,
    asset_output_dir: str | None = None,
    ocr_reader: Callable[[object], str] | None = None,
) -> dict:
    kind = _document_kind(filename)
    if kind == "pdf":
        return _build_pdf_pack(file_bytes, filename, asset_output_dir=asset_output_dir, ocr_reader=ocr_reader)
    if kind == "docx":
        return _build_docx_pack(file_bytes, filename)
    if kind == "pptx":
        return _build_pptx_pack(file_bytes, filename)
    return _build_text_pack(file_bytes, filename, kind)


def source_pack_to_text(pack: dict) -> str:
    filename = str((pack.get("document") or {}).get("filename") or "")
    kind = str((pack.get("document") or {}).get("kind") or "")
    lines = [f'--- SOURCE filename="{filename}" kind="{kind}" ---']
    chapter_by_page: dict[int, dict] = {}
    for chapter in pack.get("chapters") or []:
        try:
            start = int(chapter.get("start_page") or 0)
            end = int(chapter.get("end_page") or start)
        except (TypeError, ValueError):
            continue
        for page_num in range(start, end + 1):
            chapter_by_page[page_num] = chapter
    for page in pack.get("pages") or []:
        if not isinstance(page, dict):
            continue
        page_num = int(page.get("page_num") or 1)
        chapter = chapter_by_page.get(page_num) or {}
        chapter_text = f' chapter="{chapter.get("title")}"' if chapter.get("title") else ""
        lines.append(f'--- PAGE {page_num}{chapter_text} ---')
        lines.append(str(page.get("text") or "").strip())
    return "\n".join(lines).strip()


def read_source_pack(path: str) -> dict | None:
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, dict) else None
    except (OSError, json.JSONDecodeError):
        return None


def write_source_pack(path: str, pack: dict) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(pack, f, ensure_ascii=False, indent=2)
