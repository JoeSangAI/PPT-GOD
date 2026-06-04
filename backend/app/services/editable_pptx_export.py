from __future__ import annotations

from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass, field
import json
import logging
import math
import os
import re
import time
from pathlib import Path
from typing import Any, Callable, Iterable

import numpy as np
from PIL import Image, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.core.config import settings
from app.services.editable_pptx_diagnostics import EditablePptxDiagnostics, EditablePptxPageDiagnostics
from app.services.image_analyzer import _call_vision_model

try:
    from json_repair import repair_json
except Exception:  # pragma: no cover - optional dependency guard
    repair_json = None

try:
    import cv2
except Exception:  # pragma: no cover - optional dependency guard
    cv2 = None


logger = logging.getLogger(__name__)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
SLIDE_W_EMU = Inches(SLIDE_W_IN)
SLIDE_H_EMU = Inches(SLIDE_H_IN)

OcrProvider = Callable[[str, int], list[dict[str, Any]]]
EDITABLE_ROLES = {"title", "subtitle", "body", "caption", "label", "footer"}
STRICT_NON_EDITABLE_ROLES = {"logo", "watermark", "decorative", "page_marker", "visual_only"}
MODE_OPTIONAL_TEXT_ROLES = {"image_text", "chart_internal"}
NON_EDITABLE_ROLES = STRICT_NON_EDITABLE_ROLES | MODE_OPTIONAL_TEXT_ROLES
EDITABLE_PPTX_RESTORE_MODES = {"standard", "enhanced", "aggressive"}
COMPLEX_SLIDE_HIGH_SATURATION_THRESHOLD = 0.25
COMPLEX_SLIDE_EDGE_P90_THRESHOLD = 0.15
TEXT_RESIDUAL_RETRY_THRESHOLD = 0.18
TEXT_RESIDUAL_FALLBACK_THRESHOLD = 0.26
ROLE_ALIASES = {
    "heading": "title",
    "headline": "title",
    "main_title": "title",
    "subhead": "subtitle",
    "paragraph": "body",
    "note": "caption",
    "annotation": "caption",
}

GENERIC_FONT_HINTS = {
    "arial",
    "helvetica",
    "helvetica neue",
    "sans",
    "sans-serif",
    "sans serif",
    "sf pro",
    "sf pro display",
    "source han sans cn",
    "noto sans cjk",
    "microsoft yahei",
    "microsoft yahei, sans-serif",
    "microsoft yahei, simhei",
    "微软雅黑",
    "黑体",
    "heiti sc",
    "pingfang sc",
    "pingfang sc medium",
    "pingfang sc regular",
    "sans-serif, microsoft yahei",
}


def normalize_editable_pptx_restore_mode(mode: str | None) -> str:
    normalized = str(mode or "standard").strip().lower()
    if normalized in {"default", "recommended", "stable", "标准", "標準"}:
        return "standard"
    if normalized in {"extra", "more", "advanced", "增强", "增強"}:
        return "enhanced"
    if normalized in {"max", "maximum", "full", "risky", "aggressive", "激进", "激進", "尽量全拆", "盡量全拆"}:
        return "aggressive"
    return normalized if normalized in EDITABLE_PPTX_RESTORE_MODES else "standard"


def has_cjk(text: str) -> bool:
    return bool(re.search(r"[\u4e00-\u9fff]", text or ""))


def editable_pptx_qa_worker_count(item_count: int) -> int:
    if item_count <= 1:
        return 1
    try:
        configured = int(settings.EDITABLE_PPTX_QA_MAX_WORKERS or 1)
    except (TypeError, ValueError):
        configured = 4
    cpu_count = os.cpu_count() or 2
    return max(1, min(item_count, configured, cpu_count))


@dataclass(frozen=True)
class EditablePptxResult:
    output_path: str
    slide_count: int
    text_box_count: int
    visual_asset_count: int
    ocr_failed_pages: list[int]
    qa_retry_pages: list[int] = field(default_factory=list)
    quality_fallback_pages: list[int] = field(default_factory=list)
    quality_warning_pages: list[int] = field(default_factory=list)
    diagnostics: EditablePptxDiagnostics | None = None


def clamp_box(box: dict[str, Any]) -> dict[str, float]:
    raw_x = float(box.get("x", 0) or 0)
    raw_y = float(box.get("y", 0) or 0)
    raw_w = float(box.get("width", 0) or 0)
    raw_h = float(box.get("height", 0) or 0)
    x1 = max(0.0, min(1.0, raw_x))
    y1 = max(0.0, min(1.0, raw_y))
    x2 = max(0.0, min(1.0, raw_x + raw_w))
    y2 = max(0.0, min(1.0, raw_y + raw_h))
    return {
        "x": round(x1, 6),
        "y": round(y1, 6),
        "width": round(max(0.0, x2 - x1), 6),
        "height": round(max(0.0, y2 - y1), 6),
    }


def _strip_code_fence(raw: str) -> str:
    text = str(raw or "").strip()
    text = re.sub(r"^```(?:json)?\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\s*```$", "", text)
    return text.strip()


def _extract_json_text(raw: str) -> str:
    text = _strip_code_fence(raw)
    if text.startswith("{") or text.startswith("["):
        return text
    start_obj = text.find("{")
    start_arr = text.find("[")
    starts = [idx for idx in (start_obj, start_arr) if idx >= 0]
    if not starts:
        return text
    start = min(starts)
    end = max(text.rfind("}"), text.rfind("]"))
    return text[start : end + 1] if end > start else text[start:]


def _load_json_loose(raw: str) -> Any:
    candidate = _extract_json_text(raw)
    try:
        return json.loads(candidate)
    except Exception:
        if repair_json is None:
            raise
        return json.loads(repair_json(candidate))


def _coerce_region(item: Any) -> dict[str, Any] | None:
    if not isinstance(item, dict):
        return None
    source = item.get("bbox") if isinstance(item.get("bbox"), dict) else item
    if isinstance(item.get("bbox"), list) and len(item["bbox"]) >= 4:
        source = {
            "x": item["bbox"][0],
            "y": item["bbox"][1],
            "width": item["bbox"][2],
            "height": item["bbox"][3],
        }
    text = str(item.get("text") or item.get("content") or "").strip()
    if not text:
        return None
    if not all(k in source for k in ("x", "y", "width", "height")):
        return None
    box = clamp_box(source)
    if box["width"] < 0.006 or box["height"] < 0.006:
        return None
    role = str(item.get("role") or "").strip().lower()
    role = ROLE_ALIASES.get(role, role)
    raw_editable = item.get("editable")
    editable = None
    if isinstance(raw_editable, bool):
        editable = raw_editable
    elif isinstance(raw_editable, str):
        editable = raw_editable.strip().lower() not in {"0", "false", "no", "否", "不"}
    return {
        "text": text,
        **box,
        "confidence": float(item.get("confidence", 0.8) or 0.8),
        "role": role,
        "editable": editable,
        "color_hint": str(item.get("color_hint") or item.get("color") or "").strip(),
        "weight_hint": str(item.get("weight_hint") or item.get("weight") or "").strip().lower(),
        "font_hint": str(item.get("font_hint") or item.get("font") or "").strip(),
    }


def parse_vlm_ocr_regions(raw: str) -> list[dict[str, Any]]:
    """Parse MiniMax VLM text-region JSON into top-left normalized boxes."""
    if not str(raw or "").strip():
        return []
    try:
        data = _load_json_loose(raw)
    except Exception as exc:
        logger.warning("Editable PPTX OCR JSON parse failed: %s", exc)
        return []
    if isinstance(data, dict):
        items = data.get("text_regions") or data.get("regions") or data.get("boxes") or data.get("items") or []
    else:
        items = data
    regions = []
    for item in items if isinstance(items, list) else []:
        region = _coerce_region(item)
        if region:
            regions.append(region)
    regions.sort(key=lambda box: (float(box["y"]), float(box["x"])))
    return regions


def build_minimax_ocr_prompt(page_num: int, restore_mode: str | None = None) -> str:
    mode = normalize_editable_pptx_restore_mode(restore_mode)
    if mode == "aggressive":
        mode_instruction = """当前拆分强度：激进。尽量识别所有清晰可读文字，包括卡片标签、图表主要标签、图表内部说明和图片上的说明文字。品牌 Logo、商标、水印、页码仍标为不可编辑或非主文案。"""
    elif mode == "enhanced":
        mode_instruction = """当前拆分强度：增强。除标题、正文、结论外，也要识别卡片标签、关键数字、图表主要标签和流程节点文字；复杂图表内部极小字、Logo、商标和图片内部无关文字仍不要强拆。"""
    else:
        mode_instruction = """当前拆分强度：标准。优先识别主标题、副标题、正文、底部结论和用户最可能修改的主文案；复杂图表内部小字、Logo、商标、图片内部文字尽量保留为图片。"""
    return f"""你是 PPT 分层还原 OCR 引擎。请识别这张 PPT 第 {page_num} 页截图中的可编辑文字，并输出严格 JSON。

目标：让主要文字可以在 PowerPoint 中编辑，同时尽量保留原图视觉质量。

{mode_instruction}

要求：
1. 输出每一行或一个自然文本块的文字和 bbox。
2. bbox 使用归一化坐标，左上角为原点：x, y, width, height 都在 0-1 之间。
3. 标题、正文、结论条和核心金句是最高优先级，必须尽量完整识别，不要因为它们压在图片或深色背景上就跳过。
4. 图片、产品照片、品牌 Logo、商标、设备屏幕、复杂图表内部小字，除非符合当前拆分强度并明显有编辑价值，否则不要作为可编辑文字。
5. 对不应编辑但容易被 OCR 识别的文字，可以输出 editable=false；例如 Logo、商标、装饰水印、页码、时间轴小年份、图片内部文字。
6. 尽量保持原文，不要翻译、改写、补字。
7. 对大标题、主问题、核心金句逐字复核，尤其不要漏掉单字和虚词，例如：有、是、为、的、不、吗、了、在、和、与。
8. 如果一个字在视觉上很窄或夹在两个词之间，也必须保留；不要把“已经有能力”误写成“已经能力”。
9. 只输出 JSON，不要 markdown，不要解释。

JSON 格式：
{{
  "text_regions": [
    {{
      "text": "文字",
      "x": 0.05,
      "y": 0.12,
      "width": 0.40,
      "height": 0.06,
      "role": "title|subtitle|body|caption|label|footer|logo|watermark|decorative|page_marker|image_text|chart_internal",
      "editable": true,
      "confidence": 0.95,
      "color_hint": "#000000",
      "weight_hint": "bold|regular",
      "font_hint": "可选，尽量接近原图字体气质"
    }}
  ]
}}"""


def minimax_ocr_regions(image_path: str, page_num: int, restore_mode: str | None = None) -> list[dict[str, Any]]:
    prompt = build_minimax_ocr_prompt(page_num, restore_mode=restore_mode)
    raw = _call_vision_model(
        image_path,
        prompt,
        timeout_seconds=float(settings.EDITABLE_PPTX_OCR_TIMEOUT_SECONDS or 90),
    )
    return parse_vlm_ocr_regions(raw)


def _ocr_cache_path(work_dir: Path | str, page_num: int, restore_mode: str | None = None) -> Path:
    mode = normalize_editable_pptx_restore_mode(restore_mode)
    suffix = "" if mode == "standard" else f"_{mode}"
    return Path(work_dir) / f"slide_{page_num:02d}{suffix}_ocr_regions.json"


def read_cached_ocr_regions(work_dir: Path | str, page_num: int, restore_mode: str | None = None) -> list[dict[str, Any]] | None:
    path = _ocr_cache_path(work_dir, page_num, restore_mode)
    if not path.exists():
        return None
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None
    if not isinstance(data, list):
        return None
    regions = [region for region in (_coerce_region(item) for item in data) if region]
    return regions or None


def write_cached_ocr_regions(work_dir: Path | str, page_num: int, regions: list[dict[str, Any]], restore_mode: str | None = None) -> None:
    if not regions:
        return
    path = _ocr_cache_path(work_dir, page_num, restore_mode)
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(regions, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception as exc:
        logger.warning("Editable PPTX OCR cache write failed page=%s: %s", page_num, exc)


def run_ocr_with_retries(
    provider: OcrProvider,
    image_path: str,
    page_num: int,
    *,
    retry_count: int,
    retry_delay_seconds: float,
) -> list[dict[str, Any]]:
    attempts = max(1, int(retry_count or 0) + 1)
    last_regions: list[dict[str, Any]] = []
    for attempt in range(1, attempts + 1):
        try:
            raw_regions = provider(image_path, page_num) or []
            last_regions = [region for region in (_coerce_region(r) for r in raw_regions) if region]
        except Exception as exc:
            logger.warning("Editable PPTX OCR provider raised page=%s attempt=%s/%s: %s", page_num, attempt, attempts, exc)
            last_regions = []
        if last_regions:
            return last_regions
        if attempt < attempts and retry_delay_seconds > 0:
            time.sleep(float(retry_delay_seconds) * attempt)
    return last_regions


def cached_or_run_ocr_regions(
    provider: OcrProvider,
    image_path: str,
    page_num: int,
    *,
    work_dir: Path | str,
    restore_mode: str | None = None,
    reuse_cache: bool = True,
) -> list[dict[str, Any]]:
    if reuse_cache:
        cached = read_cached_ocr_regions(work_dir, page_num, restore_mode)
        if cached:
            return cached
    try:
        retry_count = max(0, int(settings.EDITABLE_PPTX_OCR_RETRY_COUNT or 0))
    except (TypeError, ValueError):
        retry_count = 2
    try:
        retry_delay = max(0.0, float(settings.EDITABLE_PPTX_OCR_RETRY_BACKOFF_SECONDS or 0.0))
    except (TypeError, ValueError):
        retry_delay = 1.5
    regions = run_ocr_with_retries(
        provider,
        image_path,
        page_num,
        retry_count=retry_count,
        retry_delay_seconds=retry_delay,
    )
    write_cached_ocr_regions(work_dir, page_num, regions, restore_mode)
    return regions


def corner_background(rgb: np.ndarray) -> tuple[int, int, int]:
    h, w, _ = rgb.shape
    patch = max(12, min(h, w) // 14)
    samples = np.concatenate(
        [
            rgb[:patch, :patch].reshape(-1, 3),
            rgb[:patch, w - patch :].reshape(-1, 3),
            rgb[h - patch :, :patch].reshape(-1, 3),
            rgb[h - patch :, w - patch :].reshape(-1, 3),
        ],
        axis=0,
    )
    return tuple(int(x) for x in np.median(samples, axis=0))


def luminance(color: tuple[int, int, int] | np.ndarray) -> float:
    r, g, b = [float(x) for x in color[:3]]
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def rgb_distance(rgb: np.ndarray, color: tuple[int, int, int]) -> np.ndarray:
    delta = rgb.astype(np.int32) - np.array(color, dtype=np.int32)
    return np.sqrt(np.sum(delta * delta, axis=2))


def connected_components(mask: np.ndarray) -> list[tuple[int, int, int, int, int]]:
    h, w = mask.shape
    seen = np.zeros((h, w), dtype=bool)
    components: list[tuple[int, int, int, int, int]] = []
    for y0, x0 in zip(*np.where(mask & ~seen)):
        if seen[y0, x0]:
            continue
        stack = [(int(x0), int(y0))]
        seen[y0, x0] = True
        min_x = max_x = int(x0)
        min_y = max_y = int(y0)
        area = 0
        while stack:
            x, y = stack.pop()
            area += 1
            min_x = min(min_x, x)
            max_x = max(max_x, x)
            min_y = min(min_y, y)
            max_y = max(max_y, y)
            for nx, ny in ((x - 1, y), (x + 1, y), (x, y - 1), (x, y + 1)):
                if nx < 0 or ny < 0 or nx >= w or ny >= h or seen[ny, nx] or not mask[ny, nx]:
                    continue
                seen[ny, nx] = True
                stack.append((nx, ny))
        components.append((min_x, min_y, max_x + 1, max_y + 1, area))
    return components


def intersection_area(a: dict[str, Any], b: dict[str, Any]) -> float:
    ax1, ay1 = float(a["x"]), float(a["y"])
    ax2, ay2 = ax1 + float(a["width"]), ay1 + float(a["height"])
    bx1, by1 = float(b["x"]), float(b["y"])
    bx2, by2 = bx1 + float(b["width"]), by1 + float(b["height"])
    ix1, iy1 = max(ax1, bx1), max(ay1, by1)
    ix2, iy2 = min(ax2, bx2), min(ay2, by2)
    if ix2 <= ix1 or iy2 <= iy1:
        return 0.0
    return (ix2 - ix1) * (iy2 - iy1)


def box_center(box: dict[str, Any]) -> tuple[float, float]:
    return (float(box["x"]) + float(box["width"]) / 2, float(box["y"]) + float(box["height"]) / 2)


def center_inside(center: tuple[float, float], region: dict[str, Any]) -> bool:
    x, y = center
    return (
        float(region["x"]) <= x <= float(region["x"]) + float(region["width"])
        and float(region["y"]) <= y <= float(region["y"]) + float(region["height"])
    )


def detect_image_blocks(image_path: str, text_boxes: list[dict[str, Any]]) -> list[dict[str, float]]:
    img = Image.open(image_path).convert("RGB")
    rgb = np.asarray(img)
    h, w, _ = rgb.shape
    bg = corner_background(rgb)
    diff = rgb_distance(rgb, bg)
    threshold = 44 if luminance(bg) > 160 else 52
    mask = (diff > threshold).astype(np.uint8) * 255
    mask_img = Image.fromarray(mask, mode="L").filter(ImageFilter.MaxFilter(15)).filter(ImageFilter.MinFilter(7))
    closed = np.asarray(mask_img) > 0
    candidates: list[dict[str, float]] = []
    slide_area = w * h
    for x1, y1, x2, y2, _area in connected_components(closed):
        bw = x2 - x1
        bh = y2 - y1
        box_area = bw * bh
        if box_area < slide_area * 0.018 or box_area > slide_area * 0.68:
            continue
        if bw < 100 or bh < 70:
            continue
        original_density = float(np.mean(diff[y1:y2, x1:x2] > threshold))
        if original_density < 0.22:
            continue
        crop_std = float(rgb[y1:y2, x1:x2].reshape(-1, 3).std(axis=0).mean())
        if box_area < slide_area * 0.035 and crop_std < 34:
            continue
        if original_density < 0.36 and crop_std < 45:
            continue
        aspect = bw / max(1, bh)
        if original_density < 0.38 and (aspect > 5.0 or aspect < 0.22):
            continue
        block = {
            "x": max(0, x1 - 2) / w,
            "y": max(0, y1 - 2) / h,
            "width": (min(w, x2 + 2) - max(0, x1 - 2)) / w,
            "height": (min(h, y2 + 2) - max(0, y1 - 2)) / h,
        }
        block_area_norm = max(0.0001, block["width"] * block["height"])
        text_overlap = sum(intersection_area(block, text_box) for text_box in text_boxes)
        has_large_text = any(
            (
                intersection_area(block, text_box) / max(0.0001, float(text_box["width"]) * float(text_box["height"])) > 0.65
                or intersection_area(block, text_box) / block_area_norm > 0.55
            )
            and (float(text_box["height"]) > 0.045 or float(text_box["width"]) > 0.24)
            for text_box in text_boxes
        )
        if text_overlap / block_area_norm > 0.30 and has_large_text:
            continue
        candidates.append(block)
    candidates.sort(key=lambda c: c["width"] * c["height"], reverse=True)
    try:
        max_assets = max(1, int(settings.EDITABLE_PPTX_MAX_VISUAL_ASSETS_PER_SLIDE or 6))
    except (TypeError, ValueError):
        max_assets = 6
    return candidates[:max_assets]


def sample_background(rgb: np.ndarray, box: dict[str, Any]) -> tuple[int, int, int]:
    h, w, _ = rgb.shape
    x1 = max(0, int((float(box["x"]) - 0.012) * w))
    y1 = max(0, int((float(box["y"]) - 0.012) * h))
    x2 = min(w, int((float(box["x"]) + float(box["width"]) + 0.012) * w))
    y2 = min(h, int((float(box["y"]) + float(box["height"]) + 0.012) * h))
    crop = rgb[y1:y2, x1:x2].reshape(-1, 3)
    if crop.size == 0:
        return corner_background(rgb)
    bg = tuple(int(x) for x in np.median(crop, axis=0))
    if luminance(bg) > 235:
        return (255, 255, 255)
    if luminance(bg) < 25:
        return (0, 0, 0)
    return bg


def sample_background_unclamped(rgb: np.ndarray, box: dict[str, Any]) -> tuple[int, int, int]:
    h, w, _ = rgb.shape
    x1 = max(0, int((float(box["x"]) - 0.012) * w))
    y1 = max(0, int((float(box["y"]) - 0.012) * h))
    x2 = min(w, int((float(box["x"]) + float(box["width"]) + 0.012) * w))
    y2 = min(h, int((float(box["y"]) + float(box["height"]) + 0.012) * h))
    crop = rgb[y1:y2, x1:x2].reshape(-1, 3)
    if crop.size == 0:
        return corner_background(rgb)
    return tuple(int(x) for x in np.median(crop, axis=0))


def sample_text_color(rgb: np.ndarray, box: dict[str, Any], bg: tuple[int, int, int]) -> tuple[int, int, int]:
    h, w, _ = rgb.shape
    x1 = max(0, int(float(box["x"]) * w))
    y1 = max(0, int(float(box["y"]) * h))
    x2 = min(w, int((float(box["x"]) + float(box["width"])) * w))
    y2 = min(h, int((float(box["y"]) + float(box["height"])) * h))
    crop = rgb[y1:y2, x1:x2]
    if crop.size == 0:
        return (0, 0, 0) if luminance(bg) > 128 else (255, 255, 255)
    flat = crop.reshape(-1, 3)
    dist = np.sqrt(np.sum((flat.astype(np.int32) - np.array(bg, dtype=np.int32)) ** 2, axis=1))
    bg_light = luminance(bg) > 128
    lumas = np.array([luminance(px) for px in flat])
    selected = flat[(dist > 36) & (lumas < 215)] if bg_light else flat[(dist > 36) & (lumas > 45)]
    if selected.size == 0:
        return (0, 0, 0) if bg_light else (255, 255, 255)
    selected = selected.astype(np.float32)
    max_c = selected.max(axis=1)
    min_c = selected.min(axis=1)
    sat = np.zeros_like(max_c)
    np.divide(max_c - min_c, max_c, out=sat, where=max_c > 1)
    saturated = selected[sat > 0.32]
    color = np.median(saturated if len(saturated) >= 20 else selected, axis=0)
    r, g, b = [int(max(0, min(255, round(v)))) for v in color]
    if b > r + 25 and b > g + 10:
        return (0, 107, 255)
    if bg_light and luminance((r, g, b)) < 45:
        return (0, 0, 0)
    if not bg_light and luminance((r, g, b)) > 210:
        return (255, 255, 255)
    return (r, g, b)


def parse_hex_color(value: str | None, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    match = re.search(r"#?([0-9a-fA-F]{6})", str(value or ""))
    if not match:
        return fallback
    raw = match.group(1)
    return int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16)


def is_timeline_marker_text(text: str, box: dict[str, Any]) -> bool:
    token = str(text or "").strip()
    return bool(re.fullmatch(r"\d{4}", token)) and float(box["y"]) > 0.84 and float(box["height"]) < 0.055


def is_low_contrast_decorative_text(box: dict[str, Any], rgb: np.ndarray) -> bool:
    if float(box["x"]) > 0.24:
        return False
    if float(box["height"]) < 0.11 and float(box["width"]) < 0.11:
        return False
    h, w, _ = rgb.shape
    x1 = max(0, int(float(box["x"]) * w))
    y1 = max(0, int(float(box["y"]) * h))
    x2 = min(w, int((float(box["x"]) + float(box["width"])) * w))
    y2 = min(h, int((float(box["y"]) + float(box["height"])) * h))
    crop = rgb[y1:y2, x1:x2]
    if crop.size == 0:
        return False
    luma = 0.2126 * crop[:, :, 0] + 0.7152 * crop[:, :, 1] + 0.0722 * crop[:, :, 2]
    return float(np.percentile(luma, 10)) > 218 and float(np.mean(luma < 180)) < 0.02


def is_primary_editable_text(text: str, box: dict[str, Any], role: str) -> bool:
    role = ROLE_ALIASES.get(str(role or "").lower(), str(role or "").lower())
    if role in {"title", "subtitle", "body", "footer"}:
        if role in {"title", "subtitle"}:
            return float(box["height"]) >= 0.035 or estimate_text_units(text) >= 7
        if role == "footer":
            return estimate_text_units(text) >= 10 or float(box["width"]) >= 0.28
        return estimate_text_units(text) >= 12 or float(box["height"]) >= 0.055
    return False


def slide_visual_complexity(rgb: np.ndarray) -> dict[str, float]:
    if rgb.size == 0:
        return {"edge_p90": 0.0, "high_saturation_ratio": 0.0}
    h, w = rgb.shape[:2]
    step = max(1, max(h, w) // 480)
    arr = rgb[::step, ::step].astype(np.float32) / 255.0
    if arr.ndim != 3 or arr.shape[2] < 3:
        return {"edge_p90": 0.0, "high_saturation_ratio": 0.0}

    luma = 0.2126 * arr[:, :, 0] + 0.7152 * arr[:, :, 1] + 0.0722 * arr[:, :, 2]
    dx = np.abs(np.diff(luma, axis=1))
    dy = np.abs(np.diff(luma, axis=0))
    edge_p90 = max(
        float(np.percentile(dx, 90)) if dx.size else 0.0,
        float(np.percentile(dy, 90)) if dy.size else 0.0,
    )
    max_channel = arr.max(axis=2)
    min_channel = arr.min(axis=2)
    saturation = (max_channel - min_channel) / np.maximum(max_channel, 1e-6)
    high_saturation_ratio = float(np.mean(saturation > 0.45))
    return {
        "edge_p90": edge_p90,
        "high_saturation_ratio": high_saturation_ratio,
    }


def is_complex_visual_slide(visual_complexity: dict[str, float] | None) -> bool:
    metrics = visual_complexity or {}
    return (
        float(metrics.get("high_saturation_ratio") or 0.0) >= COMPLEX_SLIDE_HIGH_SATURATION_THRESHOLD
        and float(metrics.get("edge_p90") or 0.0) >= COMPLEX_SLIDE_EDGE_P90_THRESHOLD
    )


def is_auxiliary_text_on_complex_slide(
    text: str,
    box: dict[str, Any],
    role: str,
    visual_complexity: dict[str, float] | None,
) -> bool:
    if not is_complex_visual_slide(visual_complexity):
        return False
    normalized_role = ROLE_ALIASES.get(str(role or "").lower(), str(role or "").lower())
    return normalized_role == "label" and bool(re.fullmatch(r"\d{1,2}", str(text or "").strip()))


def needs_background_fill_for_complex_text(role: str, visual_complexity: dict[str, float] | None) -> bool:
    return False


def _expanded_pill_box(box: dict[str, Any]) -> dict[str, float]:
    b = clamp_box(box)
    box_h = float(b["height"])
    x_margin = min(0.028, max(0.010, box_h * 0.32))
    y_margin = min(0.012, max(0.004, box_h * 0.16))
    return clamp_box({
        "x": b["x"] - x_margin,
        "y": b["y"] - y_margin,
        "width": b["width"] + x_margin * 2,
        "height": b["height"] + y_margin * 2,
    })


def _expanded_solid_cleanup_box(group: dict[str, Any]) -> dict[str, float]:
    b = clamp_box(group.get("bbox") or {})
    box_h = float(b["height"])
    role = str(group.get("role") or "")
    if role in {"body", "caption", "footer"}:
        x_margin = min(0.020, max(0.006, box_h * 0.14))
        y_margin = min(0.030, max(0.010, box_h * 0.34))
    else:
        x_margin = min(0.024, max(0.008, box_h * 0.22))
        y_margin = min(0.020, max(0.006, box_h * 0.26))
    return clamp_box({
        "x": b["x"] - x_margin,
        "y": b["y"] - y_margin,
        "width": b["width"] + x_margin * 2,
        "height": b["height"] + y_margin * 2,
    })


def _expanded_display_cleanup_box(group: dict[str, Any]) -> dict[str, float]:
    b = clamp_box(group.get("bbox") or {})
    box_h = float(b["height"])
    role = str(group.get("role") or "")
    if role == "title":
        x_margin = min(0.048, max(0.014, box_h * 0.32))
        y_margin = min(0.054, max(0.018, box_h * 0.42))
    elif role == "subtitle":
        x_margin = min(0.038, max(0.012, box_h * 0.28))
        y_margin = min(0.044, max(0.014, box_h * 0.38))
    else:
        x_margin = min(0.030, max(0.010, box_h * 0.22))
        y_margin = min(0.034, max(0.010, box_h * 0.32))
    return clamp_box({
        "x": b["x"] - x_margin,
        "y": b["y"] - y_margin,
        "width": b["width"] + x_margin * 2,
        "height": b["height"] + y_margin * 2,
    })


def saturated_background_ratio(rgb: np.ndarray, box: dict[str, Any]) -> float:
    h, w, _ = rgb.shape
    b = clamp_box(box)
    x1 = max(0, int(float(b["x"]) * w))
    y1 = max(0, int(float(b["y"]) * h))
    x2 = min(w, int((float(b["x"]) + float(b["width"])) * w))
    y2 = min(h, int((float(b["y"]) + float(b["height"])) * h))
    crop = rgb[y1:y2, x1:x2]
    if crop.size == 0:
        return 0.0
    flat = crop.reshape(-1, 3).astype(np.float32)
    max_c = flat.max(axis=1)
    min_c = flat.min(axis=1)
    sat = np.zeros_like(max_c)
    np.divide(max_c - min_c, np.maximum(max_c, 1), out=sat, where=max_c > 1)
    luma = 0.2126 * flat[:, 0] + 0.7152 * flat[:, 1] + 0.0722 * flat[:, 2]
    return float(np.mean((sat > 0.28) & (luma > 45) & (luma < 245)))


def needs_pill_background_for_complex_text(
    role: str,
    visual_complexity: dict[str, float] | None,
    rgb: np.ndarray,
    box: dict[str, Any],
) -> bool:
    normalized_role = ROLE_ALIASES.get(str(role or "").lower(), str(role or "").lower())
    if normalized_role not in {"subtitle", "label"}:
        return False
    threshold = 0.32 if is_complex_visual_slide(visual_complexity) else 0.22
    return saturated_background_ratio(rgb, _expanded_pill_box(box)) >= threshold


def role_allowed_for_restore_mode(role: str, text: str, box: dict[str, Any], restore_mode: str) -> bool:
    if role in STRICT_NON_EDITABLE_ROLES:
        return False
    if not role or role in EDITABLE_ROLES:
        return True
    if role not in MODE_OPTIONAL_TEXT_ROLES:
        return restore_mode == "aggressive"
    if restore_mode == "aggressive":
        return float(box["height"]) >= 0.014 and estimate_text_units(text) >= 1.6
    if restore_mode == "enhanced":
        return float(box["height"]) >= 0.022 and (float(box["width"]) >= 0.055 or estimate_text_units(text) >= 4)
    return False


def should_restore_text(
    region: dict[str, Any],
    image_blocks: list[dict[str, float]],
    rgb: np.ndarray,
    restore_mode: str | None = None,
    visual_complexity: dict[str, float] | None = None,
) -> bool:
    keep, _reason = should_restore_text_with_reason(
        region,
        image_blocks,
        rgb,
        restore_mode,
        visual_complexity=visual_complexity,
    )
    return keep


def should_restore_text_with_reason(
    region: dict[str, Any],
    image_blocks: list[dict[str, float]],
    rgb: np.ndarray,
    restore_mode: str | None = None,
    visual_complexity: dict[str, float] | None = None,
) -> tuple[bool, str]:
    mode = normalize_editable_pptx_restore_mode(restore_mode)
    text = str(region.get("text") or "").strip()
    box = clamp_box(region)
    role = str(region.get("role") or "").lower()
    if not text:
        return False, "empty_text"
    if region.get("editable") is False:
        return False, "ocr_marked_non_editable"
    if not role_allowed_for_restore_mode(role, text, box, mode):
        return False, "role_not_allowed_for_mode"
    if float(region.get("confidence", 1.0) or 1.0) < 0.22:
        return False, "low_confidence"
    if box["width"] < 0.012 or box["height"] < 0.012:
        return False, "box_too_small"
    if mode == "standard" and is_auxiliary_text_on_complex_slide(text, box, role, visual_complexity):
        return False, "standard_complex_auxiliary_text"
    if mode == "standard" and role in {"body", "caption", "label"} and float(box["height"]) < 0.022:
        return False, "standard_small_auxiliary_text"
    if is_timeline_marker_text(text, box):
        return False, "timeline_marker"
    if is_low_contrast_decorative_text(box, rgb):
        return False, "low_contrast_decorative"
    inside_visual_asset = any(center_inside(box_center(box), block) for block in image_blocks)
    if inside_visual_asset and not is_primary_editable_text(text, box, role) and mode == "standard":
        return False, "standard_inside_visual_asset"
    if len(text) <= 1 and box["height"] < 0.06:
        return False, "single_character_noise"
    bg = sample_background(rgb, box)
    color = sample_text_color(rgb, box, bg)
    contrast = math.sqrt(sum((float(color[i]) - float(bg[i])) ** 2 for i in range(3)))
    if contrast < 30 and (len(text) <= 3 or box["height"] > 0.08):
        return False, "low_contrast"
    return True, "restored"


def estimate_text_units(text: str) -> float:
    lines = split_text_lines(text)
    if len(lines) > 1:
        return max(estimate_text_units(line) for line in lines)
    units = 0.0
    for ch in text:
        code = ord(ch)
        if ch.isspace():
            units += 0.28
        elif 0x4E00 <= code <= 0x9FFF:
            units += 0.95
        elif ch.isupper():
            units += 0.66
        elif ch.isdigit():
            units += 0.56
        elif ch.isalpha():
            units += 0.53
        elif ch in ".,:;!?'`，。：；！？、":
            units += 0.32
        else:
            units += 0.48
    return max(1.0, units)


def split_text_lines(text: str) -> list[str]:
    lines = [line.strip() for line in re.split(r"[\r\n\v]+", str(text or "")) if line.strip()]
    return lines or [str(text or "").strip()]


def has_explicit_line_break(text: str) -> bool:
    return bool(re.search(r"[\r\n\v]", str(text or "")))


def text_unit(ch: str) -> float:
    code = ord(ch)
    if ch.isspace():
        return 0.28
    if 0x4E00 <= code <= 0x9FFF:
        return 0.95
    if ch.isupper():
        return 0.66
    if ch.isdigit():
        return 0.56
    if ch.isalpha():
        return 0.53
    if ch in ".,:;!?'`，。：；！？、":
        return 0.32
    return 0.48


def line_capacity_units(box: dict[str, Any], font_pt: float) -> float:
    width_in = max(0.01, float(box["width"]) * SLIDE_W_IN)
    return max(1.0, width_in * 72 * 0.84 / max(1.0, font_pt))


def should_manually_wrap_text(text: str, box: dict[str, Any], role: str) -> bool:
    if has_explicit_line_break(text):
        return False
    if role not in {"body", "caption"}:
        return False
    if float(box["height"]) < 0.055:
        return False
    units = estimate_text_units(text)
    capacity_at_body_size = line_capacity_units(box, 12.0)
    return units > max(18.0, capacity_at_body_size * 1.15)


def estimated_wrapped_line_count(text: str, box: dict[str, Any], font_pt: float) -> int:
    lines = split_text_lines(text)
    if len(lines) > 1:
        return len(lines)
    return max(1, math.ceil(estimate_text_units(text) / line_capacity_units(box, font_pt)))


def wrap_text_for_box(text: str, box: dict[str, Any], font_pt: float, role: str) -> str:
    raw = str(text or "").strip()
    if not raw or not should_manually_wrap_text(raw, box, role):
        return raw
    capacity = line_capacity_units(box, font_pt) * 0.96
    tokens = re.findall(r"[A-Za-z0-9]+|\\s+|.", raw)
    lines: list[str] = []
    current = ""
    current_units = 0.0
    punctuation = set(".,:;!?'`，。：；！？、）)]}》」』")
    opening_punctuation = set("([{《「『“‘")
    for token in tokens:
        if token.isspace():
            token = " "
        token_units = sum(text_unit(ch) for ch in token)
        current_tail = current.rstrip()[-1:] if current.rstrip() else ""
        if (
            current
            and current_units + token_units > capacity
            and token not in punctuation
            and current_tail not in opening_punctuation
        ):
            lines.append(current.strip())
            current = token.lstrip()
            current_units = sum(text_unit(ch) for ch in current)
            continue
        current += token
        current_units += token_units
    if current.strip():
        lines.append(current.strip())
    return "\n".join(line for line in lines if line)


def fitted_font_size_pt(
    text: str,
    box: dict[str, Any],
    preferred_pt: float,
    *,
    min_pt: float = 6.5,
    max_pt: float = 60.0,
    wrap_long_text: bool = False,
) -> float:
    width_in = max(0.01, float(box["width"]) * SLIDE_W_IN)
    width_limited = width_in * 72 * 0.84 / estimate_text_units(text)
    lines = split_text_lines(text)
    if len(lines) > 1:
        height_in = max(0.01, float(box["height"]) * SLIDE_H_IN)
        height_limited = height_in * 72 * 0.78 / max(1, len(lines))
        preferred_pt = min(preferred_pt, height_limited)
    elif wrap_long_text:
        if float(box["width"]) < 0.22 and estimate_text_units(text) > 42:
            max_pt = min(max_pt, 11.0)
        start_pt = min(max_pt, preferred_pt)
        height_in = max(0.01, float(box["height"]) * SLIDE_H_IN)
        available_height_pt = height_in * 72 * 0.68
        steps = int(max(0, math.ceil((start_pt - min_pt) / 0.25)))
        for step in range(steps + 1):
            font_pt = round(start_pt - step * 0.25, 2)
            if font_pt < min_pt:
                break
            line_count = estimated_wrapped_line_count(text, box, font_pt)
            if font_pt * max(1, line_count) <= available_height_pt:
                return font_pt
        return round(min_pt, 2)
    return round(max(min_pt, min(max_pt, preferred_pt, width_limited)), 2)


def role_max_font_size(role: str, text: str = "") -> float:
    text_units = estimate_text_units(text)
    if len(split_text_lines(text)) > 1:
        if role == "subtitle":
            return 26.0
        if role == "label":
            return 18.0
        if role in {"body", "caption", "footer"} and text_units > 90:
            return 11.5
        if role in {"body", "caption", "footer"} and text_units > 55:
            return 12.5
        return 16.0
    if role == "label":
        return 18.0
    if role in {"body", "caption"} and text_units > 110:
        return 11.0
    if role in {"body", "caption"} and text_units > 70:
        return 12.5
    if role in {"body", "caption"} and text_units > 42:
        return 14.0
    if role in {"body", "caption"} and text_units > 18:
        return 16.0
    if role == "title":
        return 62.0 if not has_cjk(text) else 56.0
    if role == "subtitle":
        return 34.0
    return 26.0


def font_height_multiplier(role: str) -> float:
    if role == "title":
        return 0.86
    if role == "subtitle":
        return 0.82
    if role == "label":
        return 0.78
    if role == "caption":
        return 0.72
    return 0.76


def infer_role(region: dict[str, Any]) -> str:
    role = str(region.get("role") or "").lower()
    role = ROLE_ALIASES.get(role, role)
    h = float(region["height"])
    w = float(region["width"])
    text = str(region.get("text") or "").strip()
    lines = split_text_lines(text)
    if len(lines) > 1:
        if role in {"caption", "footer"}:
            return role
        if len(lines) <= 2 and estimate_text_units(text) <= 14 and h >= 0.14:
            return "subtitle"
        return "body"
    if role == "label":
        return "label"
    if role in {"body", "caption", "footer"} and estimate_text_units(text) > 18:
        return role
    if h >= 0.105:
        return "title"
    if role in {"body", "caption", "label"} and h >= 0.062 and estimate_text_units(text) <= 18:
        return "subtitle"
    if role in EDITABLE_ROLES:
        return role
    if h >= 0.09:
        return "title"
    if h >= 0.055 or (w >= 0.34 and h >= 0.045):
        return "subtitle"
    return "body"


def infer_alignment(box: dict[str, float]) -> str:
    center_x = float(box["x"]) + float(box["width"]) / 2
    left_margin = float(box["x"])
    right_margin = 1.0 - float(box["x"]) - float(box["width"])
    if 0.40 <= center_x <= 0.60 and abs(left_margin - right_margin) <= 0.065:
        return "center"
    if float(box["x"]) >= 0.12 and float(box["width"]) > 0.34 and 0.42 <= center_x <= 0.58:
        return "center"
    return "left"


def expand_line_box_for_rendering(
    box: dict[str, float],
    text: str,
    *,
    alignment: str = "left",
    role: str = "body",
) -> dict[str, float]:
    expanded = dict(box)
    if len(split_text_lines(text)) > 1 or should_manually_wrap_text(text, expanded, role):
        return clamp_box(expanded)
    preferred = min(
        role_max_font_size(role, text),
        float(expanded["height"]) * SLIDE_H_IN * 72 * font_height_multiplier(role),
    )
    desired_width_in = estimate_text_units(text) * preferred / (72 * 0.84)
    desired_width = min(0.96, max(expanded["width"], desired_width_in / SLIDE_W_IN * 1.06))
    if desired_width <= expanded["width"] + 0.004:
        return clamp_box(expanded)

    align = str(alignment or "left")
    if align == "center":
        center_x = expanded["x"] + expanded["width"] / 2
        max_width = min(0.96, center_x - 0.02, 0.98 - center_x) * 2
        expanded["width"] = min(desired_width, max_width)
        expanded["x"] = center_x - expanded["width"] / 2
    elif align == "right":
        right = expanded["x"] + expanded["width"]
        max_width = max(0.02, right - 0.02)
        expanded["width"] = min(desired_width, max_width)
        expanded["x"] = right - expanded["width"]
    else:
        max_width = max(0.02, 0.98 - expanded["x"])
        expanded["width"] = min(desired_width, max_width)
    return clamp_box(expanded)


def _natural_text_box_for_group(group: dict[str, Any]) -> dict[str, float]:
    return expand_line_box_for_rendering(
        clamp_box(group["bbox"]),
        str(group.get("text") or ""),
        alignment=str(group.get("alignment") or "left"),
        role=str(group.get("role") or "body"),
    )


def _max_font_size_by_box_height(
    box: dict[str, Any],
    role: str,
    text: str,
    original_box: dict[str, Any] | None = None,
) -> float | None:
    units = estimate_text_units(text)
    height = float(box.get("height") or 0.0)
    original = clamp_box(original_box or box)
    expansion_ratio = float(box.get("width") or 0.0) / max(0.001, float(original.get("width") or 0.0))
    is_top_expanded_title = (
        role == "title"
        and has_cjk(text)
        and float(original.get("y") or 0.0) <= 0.18
        and expansion_ratio >= 1.12
        and units >= 10
    )
    is_top_title_fragment = (
        role == "title"
        and has_cjk(text)
        and float(original.get("y") or 0.0) <= 0.18
        and units >= 10
        and str(text).rstrip().endswith(("，", "：", ",", ":"))
    )
    caps: list[float] = []
    if is_top_expanded_title or is_top_title_fragment:
        caps.append(min(28.0, max(15.0, height * 260.0)))
    if role == "title" and units >= 30:
        caps.append(min(28.0, max(13.0, height * 210.0)))
    if role in {"subtitle", "label"} and units >= 20:
        caps.append(min(20.0, max(8.0, height * 170.0)))
    if caps:
        return min(caps)
    return None


def _fitted_size_for_group(group: dict[str, Any], box: dict[str, float]) -> float:
    role = str(group.get("role") or "")
    preferred = float(box["height"]) * SLIDE_H_IN * 72 * font_height_multiplier(role)
    text = str(group.get("text") or "")
    height_cap = _max_font_size_by_box_height(box, role, text, group.get("bbox"))
    max_pt = role_max_font_size(role, text)
    if height_cap is not None:
        max_pt = min(max_pt, height_cap)
    return fitted_font_size_pt(
        text,
        box,
        preferred,
        max_pt=max_pt,
        wrap_long_text=should_manually_wrap_text(text, box, role),
    )


def normalize_same_level_text_metrics(groups: list[dict[str, Any]]) -> None:
    candidates = []
    for group in groups:
        box = clamp_box(group.get("bbox") or {})
        role = str(group.get("role") or "")
        if role not in {"title", "subtitle", "body", "caption", "label"}:
            continue
        if not (0.014 <= float(box["height"]) <= 0.12):
            continue
        natural_box = _natural_text_box_for_group(group)
        candidates.append({
            "group": group,
            "role": role,
            "alignment": str(group.get("alignment") or "left"),
            "box": natural_box,
            "size": _fitted_size_for_group(group, natural_box),
        })

    records = sorted(candidates, key=lambda item: (float(item["box"]["y"]), float(item["box"]["x"])))
    clusters = []
    i = 0
    while i < len(records):
        cluster = [records[i]]
        j = i + 1
        while j < len(records):
            first = cluster[0]["box"]
            first_record = cluster[0]
            previous = cluster[-1]["box"]
            current = records[j]["box"]
            gap = float(current["y"]) - (float(previous["y"]) + float(previous["height"]))
            height_tolerance = max(0.006, float(first["height"]) * 0.24)
            max_gap = max(0.035, float(first["height"]) * 1.7)
            alignment = first_record["alignment"]
            if alignment == "center":
                first_anchor = float(first["x"]) + float(first["width"]) / 2
                current_anchor = float(current["x"]) + float(current["width"]) / 2
                anchor_tolerance = 0.055
            elif alignment == "right":
                first_anchor = float(first["x"]) + float(first["width"])
                current_anchor = float(current["x"]) + float(current["width"])
                anchor_tolerance = 0.045
            else:
                first_anchor = float(first["x"])
                current_anchor = float(current["x"])
                anchor_tolerance = 0.045
            if (
                records[j]["role"] == first_record["role"]
                and records[j]["alignment"] == alignment
                and abs(current_anchor - first_anchor) <= anchor_tolerance
                and abs(float(current["height"]) - float(first["height"])) <= height_tolerance
                and 0 <= gap <= max_gap
            ):
                cluster.append(records[j])
                j += 1
                continue
            break
        if len(cluster) >= 2:
            clusters.append(cluster)
        i = max(j, i + 1)

    for cluster in clusters:
        natural_sizes = [float(record["size"]) for record in cluster]
        target_width = max(float(record["box"]["width"]) for record in cluster)
        if max(natural_sizes) - min(natural_sizes) < 0.15 and max(float(record["box"]["width"]) for record in cluster) - min(float(record["box"]["width"]) for record in cluster) < 0.012:
            continue
        fitted_sizes = []
        for record in cluster:
            hinted_box = dict(record["box"])
            hinted_box["width"] = target_width
            fitted_sizes.append(_fitted_size_for_group(record["group"], hinted_box))
        target_size = round(min(fitted_sizes), 2)
        for record in cluster:
            record["group"]["box_width_hint"] = round(target_width, 5)
            record["group"]["font_size_hint"] = target_size

    for role in {"title", "subtitle", "body", "caption", "label"}:
        role_records = sorted((record for record in records if record["role"] == role), key=lambda item: float(item["box"]["y"]))
        used: set[int] = set()
        for index, anchor in enumerate(role_records):
            if index in used:
                continue
            anchor_box = anchor["box"]
            anchor_center_y = float(anchor_box["y"]) + float(anchor_box["height"]) / 2
            anchor_height = float(anchor_box["height"])
            row = []
            for other_index, record in enumerate(role_records):
                if other_index in used:
                    continue
                box = record["box"]
                center_y = float(box["y"]) + float(box["height"]) / 2
                same_row = abs(center_y - anchor_center_y) <= max(0.012, anchor_height * 0.45)
                similar_height = abs(float(box["height"]) - anchor_height) <= max(0.010, anchor_height * 0.35)
                if same_row and similar_height:
                    row.append((other_index, record))
            if len(row) < 2:
                continue
            x_positions = [float(record["box"]["x"]) for _, record in row]
            if max(x_positions) - min(x_positions) < 0.16:
                continue
            row_sizes = []
            for _, record in row:
                hinted = record["group"].get("font_size_hint")
                size = float(record["size"])
                if hinted is not None:
                    size = min(size, float(hinted))
                row_sizes.append(size)
            target_size = round(min(row_sizes), 2)
            for other_index, record in row:
                current_hint = record["group"].get("font_size_hint")
                record["group"]["font_size_hint"] = target_size if current_hint is None else round(min(float(current_hint), target_size), 2)
                used.add(other_index)


def apply_complex_slide_text_safety_hints(groups: list[dict[str, Any]], visual_complexity: dict[str, float] | None) -> None:
    if not is_complex_visual_slide(visual_complexity):
        return
    for group in groups:
        role = str(group.get("role") or "")
        box = _box_for_group(group)
        text = str(group.get("text") or "")
        min_readable_hint: float | None = None
        if role in {"subtitle", "label"}:
            cap = 18.0
        elif role in {"body", "caption", "footer"}:
            if float(box["width"]) >= 0.36 and float(box["height"]) <= 0.14:
                max_width = max(0.02, 0.98 - float(box["x"]))
                group["box_width_hint"] = round(min(max_width, max(float(box["width"]), float(box["width"]) + 0.03)), 5)
                group["box_height_hint"] = round(
                    min(0.12, max(float(box["height"]) * 1.5, float(box["height"]) + 0.03, 0.095)),
                    5,
                )
                box = _box_for_group(group)
                body_role_cap = role_max_font_size(role, text)
                if body_role_cap >= 14.0:
                    cap = 14.0
                    min_readable_hint = 14.0
                else:
                    cap = body_role_cap
            elif float(box["width"]) >= 0.32:
                cap = 12.5
            else:
                cap = 11.0
        else:
            continue
        if min_readable_hint is not None:
            preferred = float(box["height"]) * SLIDE_H_IN * 72 * font_height_multiplier(role)
            fitted = fitted_font_size_pt(
                text,
                box,
                preferred,
                max_pt=cap,
                wrap_long_text=True,
            )
        else:
            fitted = _fitted_size_for_group(group, box)
        current_hint = group.get("font_size_hint")
        if current_hint is not None:
            current_size = float(current_hint)
            if (
                min_readable_hint is not None
                and role_max_font_size(role, text) >= min_readable_hint
                and fitted >= min_readable_hint - 0.5
            ):
                current_size = max(current_size, min_readable_hint)
                fitted = max(min(fitted, current_size), min_readable_hint)
            else:
                fitted = min(fitted, current_size)
        group["font_size_hint"] = round(min(fitted, cap), 2)


def choose_font_name(text: str, font_hint: str | None = None) -> str:
    hint = str(font_hint or "").strip()
    safe_hint = hint and len(hint) <= 80 and not re.search(r"[\r\n{}<>]", hint)
    normalized = re.sub(r"\s+", " ", hint).strip().lower()
    if has_cjk(text):
        return "PingFang SC"
    if safe_hint and normalized not in GENERIC_FONT_HINTS:
        if any(term in normalized for term in ("serif", "georgia", "times")):
            return hint
    return "Helvetica Neue"


def emu_from_top_box(box: dict[str, Any]) -> tuple[Emu, Emu, Emu, Emu]:
    b = clamp_box(box)
    return (
        Emu(int(b["x"] * SLIDE_W_EMU)),
        Emu(int(b["y"] * SLIDE_H_EMU)),
        Emu(int(b["width"] * SLIDE_W_EMU)),
        Emu(int(b["height"] * SLIDE_H_EMU)),
    )


def _context_fill_patch(rgb: np.ndarray, x1: int, y1: int, x2: int, y2: int) -> np.ndarray:
    h, w, _ = rgb.shape
    margin = max(8, min(w, h) // 80)
    sx1 = max(0, x1 - margin)
    sx2 = min(w, x2 + margin)
    sy1 = max(0, y1 - margin)
    sy2 = min(h, y2 + margin)
    samples_x: list[float] = []
    samples_y: list[float] = []
    samples_c: list[np.ndarray] = []
    for yy in range(sy1, sy2):
        for xx in range(sx1, sx2):
            if x1 <= xx < x2 and y1 <= yy < y2:
                continue
            near = x1 - margin <= xx < x2 + margin and y1 - margin <= yy < y2 + margin
            if not near:
                continue
            samples_x.append(xx)
            samples_y.append(yy)
            samples_c.append(rgb[yy, xx].astype(np.float64))
    if len(samples_c) < 20:
        patch = np.zeros((y2 - y1, x2 - x1, 3), dtype=np.uint8)
        patch[:, :] = np.array(corner_background(rgb), dtype=np.uint8)
        return patch
    xs = np.array(samples_x, dtype=np.float64)
    ys = np.array(samples_y, dtype=np.float64)
    design = np.stack([np.ones_like(xs), xs / max(1, w), ys / max(1, h)], axis=1)
    colors = np.stack(samples_c, axis=0)
    coeffs = [np.linalg.lstsq(design, colors[:, channel], rcond=None)[0] for channel in range(3)]
    yy, xx = np.mgrid[y1:y2, x1:x2]
    pred_design = np.stack([np.ones_like(xx, dtype=np.float64), xx / max(1, w), yy / max(1, h)], axis=2)
    patch = np.zeros((y2 - y1, x2 - x1, 3), dtype=np.float64)
    for channel, coef in enumerate(coeffs):
        patch[:, :, channel] = pred_design @ coef
    return np.clip(patch, 0, 255).astype(np.uint8)


def _fit_context_fill(img: Image.Image, box: dict[str, Any]) -> None:
    rgb = np.asarray(img).copy()
    h, w, _ = rgb.shape
    x1 = max(0, int((float(box["x"]) - 0.006) * w))
    y1 = max(0, int((float(box["y"]) - 0.008) * h))
    x2 = min(w, int((float(box["x"]) + float(box["width"]) + 0.006) * w))
    y2 = min(h, int((float(box["y"]) + float(box["height"]) + 0.008) * h))
    if x2 <= x1 or y2 <= y1:
        return
    patch = _context_fill_patch(rgb, x1, y1, x2, y2)
    mask = Image.new("L", (x2 - x1, y2 - y1), 255)
    feather = max(2, min(mask.size) // 10)
    if feather > 1:
        mask = mask.filter(ImageFilter.GaussianBlur(feather / 2))
    base = Image.fromarray(rgb)
    base.paste(Image.fromarray(patch), (x1, y1), mask)
    img.paste(base)


def _fit_solid_context_fill(img: Image.Image, box: dict[str, Any]) -> None:
    rgb = np.asarray(img).copy()
    h, w, _ = rgb.shape
    b = clamp_box(box)
    x1 = max(0, int(float(b["x"]) * w))
    y1 = max(0, int(float(b["y"]) * h))
    x2 = min(w, int((float(b["x"]) + float(b["width"])) * w))
    y2 = min(h, int((float(b["y"]) + float(b["height"])) * h))
    if x2 <= x1 or y2 <= y1:
        return
    base = Image.fromarray(rgb)
    base.paste(Image.fromarray(_solid_context_patch_array(rgb, x1, y1, x2, y2)), (x1, y1))
    img.paste(base)


def _solid_context_patch_array(rgb: np.ndarray, x1: int, y1: int, x2: int, y2: int) -> np.ndarray:
    crop = rgb[y1:y2, x1:x2].reshape(-1, 3).astype(np.float32)
    luma = 0.2126 * crop[:, 0] + 0.7152 * crop[:, 1] + 0.0722 * crop[:, 2]
    max_c = crop.max(axis=1)
    min_c = crop.min(axis=1)
    sat = np.zeros_like(max_c)
    np.divide(max_c - min_c, np.maximum(max_c, 1), out=sat, where=max_c > 1)
    median_luma = float(np.median(luma))
    if median_luma >= 172:
        selected = crop[(luma >= 205) & (sat <= 0.24)]
        if len(selected) < 20:
            selected = crop[luma >= median_luma]
    elif median_luma <= 82:
        selected = crop[(luma <= 120) & (sat <= 0.45)]
        if len(selected) < 20:
            selected = crop[luma <= median_luma]
    else:
        patch = _context_fill_patch(rgb, x1, y1, x2, y2)
        selected = patch.reshape(-1, 3).astype(np.float32)
    color = np.median(selected if len(selected) else crop, axis=0).astype(np.uint8)
    patch = np.zeros((y2 - y1, x2 - x1, 3), dtype=np.uint8)
    patch[:, :] = color
    return patch


def _fit_text_foreground_fill(img: Image.Image, box: dict[str, Any]) -> None:
    rgb = np.asarray(img).copy()
    h, w, _ = rgb.shape
    x1, y1, x2, y2 = _text_cleanup_pixel_bounds(box, w, h)
    if x2 <= x1 or y2 <= y1:
        return
    patch = _text_cleanup_patch_array(rgb, box, x1, y1, x2, y2)
    base = Image.fromarray(rgb)
    base.paste(Image.fromarray(patch), (x1, y1))
    img.paste(base)


def _text_cleanup_pixel_bounds(box: dict[str, Any], width: int, height: int) -> tuple[int, int, int, int]:
    box_height = float(box["height"])
    x_margin = max(0.006, min(0.028, box_height * 0.16))
    y_margin = max(0.008, min(0.035, box_height * 0.20))
    x1 = max(0, int((float(box["x"]) - x_margin) * width))
    y1 = max(0, int((float(box["y"]) - y_margin) * height))
    x2 = min(width, int((float(box["x"]) + float(box["width"]) + x_margin) * width))
    y2 = min(height, int((float(box["y"]) + float(box["height"]) + y_margin) * height))
    return x1, y1, x2, y2


def _cleanup_pixel_bounds_for_box(box: dict[str, Any], width: int, height: int) -> tuple[int, int, int, int]:
    b = clamp_box(box)
    if box.get("solid_fill") or box.get("full_fill"):
        x1 = max(0, int(float(b["x"]) * width))
        y1 = max(0, int(float(b["y"]) * height))
        x2 = min(width, int((float(b["x"]) + float(b["width"])) * width))
        y2 = min(height, int((float(b["y"]) + float(b["height"])) * height))
        return x1, y1, x2, y2
    return _text_cleanup_pixel_bounds(b, width, height)


def _text_foreground_mask(
    rgb: np.ndarray,
    box: dict[str, Any],
    x1: int,
    y1: int,
    x2: int,
    y2: int,
) -> tuple[np.ndarray, np.ndarray]:
    patch = _context_fill_patch(rgb, x1, y1, x2, y2)
    sampled_bg = np.array(sample_background_unclamped(rgb, box), dtype=np.int32)
    crop = rgb[y1:y2, x1:x2].astype(np.int32)
    expected = patch.astype(np.int32)
    dist = np.sqrt(np.sum((crop - expected) ** 2, axis=2))
    direct_dist = np.sqrt(np.sum((crop - sampled_bg[None, None, :]) ** 2, axis=2))
    crop_luma = 0.2126 * crop[:, :, 0] + 0.7152 * crop[:, :, 1] + 0.0722 * crop[:, :, 2]
    expected_luma = 0.2126 * expected[:, :, 0] + 0.7152 * expected[:, :, 1] + 0.0722 * expected[:, :, 2]
    bg_luma = luminance(sampled_bg)
    foreground = (dist > 22) & (np.abs(crop_luma - expected_luma) > 16)
    if bg_luma > 150:
        foreground |= (direct_dist > 32) & (crop_luma < bg_luma - 18)
    else:
        foreground |= (direct_dist > 32) & (crop_luma > bg_luma + 18)
    return foreground, patch


def _text_cleanup_patch_array(rgb: np.ndarray, box: dict[str, Any], x1: int, y1: int, x2: int, y2: int) -> np.ndarray:
    foreground, patch = _text_foreground_mask(rgb, box, x1, y1, x2, y2)
    if not bool(np.any(foreground)):
        return patch
    mask = Image.fromarray((foreground.astype(np.uint8) * 255), mode="L")
    kernel = int(max(7, min(31, round((y2 - y1) * 0.22))))
    if kernel % 2 == 0:
        kernel += 1
    mask = mask.filter(ImageFilter.MaxFilter(kernel))
    crop_image = Image.fromarray(rgb[y1:y2, x1:x2].astype(np.uint8))
    if cv2 is not None:
        crop_uint8 = rgb[y1:y2, x1:x2].astype(np.uint8)
        mask_uint8 = np.asarray(mask, dtype=np.uint8)
        radius = max(2, min(12, int(min(mask_uint8.shape) * 0.08)))
        inpainted = cv2.inpaint(crop_uint8, mask_uint8, radius, cv2.INPAINT_TELEA)
        crop_image.paste(Image.fromarray(inpainted), (0, 0), mask)
    else:
        soft_mask = mask.filter(ImageFilter.GaussianBlur(1.1))
        crop_image.paste(Image.fromarray(patch), (0, 0), soft_mask)
    return np.asarray(crop_image.convert("RGB"))


def residual_original_text_score(original_rgb: np.ndarray, underlay_rgb: np.ndarray, cleanup_box: dict[str, Any]) -> float:
    h, w, _ = original_rgb.shape
    x1, y1, x2, y2 = _cleanup_pixel_bounds_for_box(cleanup_box, w, h)
    if x2 <= x1 or y2 <= y1:
        return 0.0
    foreground, _patch = _text_foreground_mask(original_rgb, cleanup_box, x1, y1, x2, y2)
    foreground_count = int(np.count_nonzero(foreground))
    if foreground_count < 12:
        return 0.0
    original_crop = original_rgb[y1:y2, x1:x2].astype(np.int32)
    underlay_crop = underlay_rgb[y1:y2, x1:x2].astype(np.int32)
    unchanged = np.sqrt(np.sum((underlay_crop - original_crop) ** 2, axis=2)) < 10
    return float(np.count_nonzero(foreground & unchanged) / max(1, foreground_count))


def residual_text_groups(
    original_rgb: np.ndarray,
    underlay_rgb: np.ndarray,
    groups: list[dict[str, Any]],
    cleanup_boxes: list[dict[str, Any]],
    *,
    threshold: float = TEXT_RESIDUAL_RETRY_THRESHOLD,
) -> list[tuple[dict[str, Any], float]]:
    pairs = list(zip(groups, cleanup_boxes))
    if not pairs:
        return []

    def score_pair(pair: tuple[dict[str, Any], dict[str, Any]]) -> tuple[dict[str, Any], float] | None:
        group, cleanup_box = pair
        score = residual_original_text_score(original_rgb, underlay_rgb, cleanup_box)
        return (group, score) if score >= threshold else None

    workers = editable_pptx_qa_worker_count(len(pairs))
    if workers <= 1:
        scored = [score_pair(pair) for pair in pairs]
    else:
        with ThreadPoolExecutor(max_workers=workers, thread_name_prefix="editable-pptx-qa") as pool:
            scored = list(pool.map(score_pair, pairs))
    return [item for item in scored if item is not None]


def create_text_cleanup_patch(rendered: Image.Image, box: dict[str, Any], output_path: str) -> dict[str, float] | None:
    rgb = np.asarray(rendered.convert("RGB")).copy()
    h, w, _ = rgb.shape
    x1, y1, x2, y2 = _text_cleanup_pixel_bounds(box, w, h)
    if x2 <= x1 or y2 <= y1:
        return None
    patch = _text_cleanup_patch_array(rgb, box, x1, y1, x2, y2)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    Image.fromarray(patch).save(output_path)
    return {
        "x": x1 / w,
        "y": y1 / h,
        "width": (x2 - x1) / w,
        "height": (y2 - y1) / h,
    }


def create_visual_cleanup_patch(rendered: Image.Image, box: dict[str, Any], output_path: str) -> dict[str, float] | None:
    rgb = np.asarray(rendered.convert("RGB")).copy()
    h, w, _ = rgb.shape
    b = clamp_box(box)
    x1 = max(0, int((float(b["x"]) - 0.006) * w))
    y1 = max(0, int((float(b["y"]) - 0.008) * h))
    x2 = min(w, int((float(b["x"]) + float(b["width"]) + 0.006) * w))
    y2 = min(h, int((float(b["y"]) + float(b["height"]) + 0.008) * h))
    if x2 <= x1 or y2 <= y1:
        return None
    patch = _context_fill_patch(rgb, x1, y1, x2, y2)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    Image.fromarray(patch).save(output_path)
    return {
        "x": x1 / w,
        "y": y1 / h,
        "width": (x2 - x1) / w,
        "height": (y2 - y1) / h,
    }


def create_cleanup_patch(rendered: Image.Image, box: dict[str, Any], output_path: str) -> dict[str, float] | None:
    rgb = np.asarray(rendered.convert("RGB")).copy()
    h, w, _ = rgb.shape
    b = clamp_box(box)
    x1, y1, x2, y2 = _cleanup_pixel_bounds_for_box(box, w, h)
    if x2 <= x1 or y2 <= y1:
        return None
    if box.get("solid_fill"):
        patch = _solid_context_patch_array(rgb, x1, y1, x2, y2)
    elif box.get("full_fill"):
        patch = _context_fill_patch(rgb, x1, y1, x2, y2)
    else:
        patch = _text_cleanup_patch_array(rgb, b, x1, y1, x2, y2)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    Image.fromarray(patch).save(output_path)
    return {
        "x": x1 / w,
        "y": y1 / h,
        "width": (x2 - x1) / w,
        "height": (y2 - y1) / h,
    }


def apply_cleanup_box_to_image(img: Image.Image, box: dict[str, Any]) -> None:
    if box.get("solid_fill"):
        _fit_solid_context_fill(img, box)
    elif box.get("full_fill"):
        _fit_context_fill(img, box)
    else:
        _fit_text_foreground_fill(img, box)


def paste_cleanup_patch_to_image(img: Image.Image, patch_path: str, patch_box: dict[str, Any]) -> None:
    b = clamp_box(patch_box)
    width, height = img.size
    x1 = max(0, int(float(b["x"]) * width))
    y1 = max(0, int(float(b["y"]) * height))
    x2 = min(width, int((float(b["x"]) + float(b["width"])) * width))
    y2 = min(height, int((float(b["y"]) + float(b["height"])) * height))
    if x2 <= x1 or y2 <= y1:
        return
    patch = Image.open(patch_path).convert("RGB")
    if patch.size != (x2 - x1, y2 - y1):
        patch = patch.resize((x2 - x1, y2 - y1), Image.Resampling.BICUBIC)
    img.paste(patch, (x1, y1))


def paste_asset_crop_to_image(img: Image.Image, asset_path: str, region: dict[str, Any]) -> None:
    b = clamp_box(region)
    width, height = img.size
    x1 = max(0, int(float(b["x"]) * width))
    y1 = max(0, int(float(b["y"]) * height))
    x2 = min(width, int((float(b["x"]) + float(b["width"])) * width))
    y2 = min(height, int((float(b["y"]) + float(b["height"])) * height))
    if x2 <= x1 or y2 <= y1:
        return
    asset = Image.open(asset_path).convert("RGBA")
    if asset.size != (x2 - x1, y2 - y1):
        asset = asset.resize((x2 - x1, y2 - y1), Image.Resampling.BICUBIC)
    img.paste(asset.convert("RGB"), (x1, y1), asset.getchannel("A"))


def prepare_clean_background(image_path: str, text_boxes: list[dict[str, Any]], image_blocks: list[dict[str, Any]], output_path: str) -> str:
    img = Image.open(image_path).convert("RGB")
    for block in sorted(image_blocks, key=lambda b: float(b["width"]) * float(b["height"]), reverse=True):
        cleanup_box = clamp_box(block)
        cleanup_box["full_fill"] = True
        apply_cleanup_box_to_image(img, cleanup_box)
    for box in sorted(text_boxes, key=lambda b: float(b["height"])):
        apply_cleanup_box_to_image(img, box)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    img.save(output_path)
    return output_path


def crop_asset_from_image(img: Image.Image, region: dict[str, Any], output_path: str) -> str:
    img = img.convert("RGB")
    w, h = img.size
    box = clamp_box(region)
    x1 = int(box["x"] * w)
    y1 = int(box["y"] * h)
    x2 = int((box["x"] + box["width"]) * w)
    y2 = int((box["y"] + box["height"]) * h)
    crop = img.crop((x1, y1, x2, y2)).convert("RGBA")
    if box["width"] * box["height"] > 0.06 and crop.width > 8 and crop.height > 8:
        alpha = Image.new("L", crop.size, 255)
        alpha_arr = np.asarray(alpha).copy()
        fade = max(6, min(28, min(crop.size) // 9))
        if x1 > 0:
            ramp = np.linspace(0, 255, fade, dtype=np.uint8)
            alpha_arr[:, :fade] = np.minimum(alpha_arr[:, :fade], ramp[None, :])
        if y1 > 0:
            ramp = np.linspace(0, 255, fade, dtype=np.uint8)
            alpha_arr[:fade, :] = np.minimum(alpha_arr[:fade, :], ramp[:, None])
        if x2 < w:
            ramp = np.linspace(255, 0, fade, dtype=np.uint8)
            alpha_arr[:, -fade:] = np.minimum(alpha_arr[:, -fade:], ramp[None, :])
        if y2 < h:
            ramp = np.linspace(255, 0, fade, dtype=np.uint8)
            alpha_arr[-fade:, :] = np.minimum(alpha_arr[-fade:, :], ramp[:, None])
        crop.putalpha(Image.fromarray(alpha_arr))
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    crop.save(output_path)
    return output_path


def crop_asset(image_path: str, region: dict[str, Any], output_path: str) -> str:
    return crop_asset_from_image(Image.open(image_path).convert("RGB"), region, output_path)


def _group_from_region(region: dict[str, Any]) -> dict[str, Any]:
    box = clamp_box(region)
    role = infer_role(region)
    text = str(region.get("text") or "").strip()
    raw_weight = str(region.get("weight_hint") or "").strip().lower()
    if raw_weight in {"bold", "regular"}:
        weight_hint = raw_weight
    elif role == "title":
        weight_hint = "bold"
    elif role == "subtitle":
        weight_hint = "bold" if box["height"] >= 0.035 else "regular"
    elif role == "label":
        weight_hint = "bold" if box["height"] <= 0.04 and estimate_text_units(text) <= 10 else "regular"
    elif role == "body" and box["height"] >= 0.09 and estimate_text_units(text) <= 14:
        weight_hint = "bold"
    else:
        weight_hint = "regular"
    return {
        "text": text,
        "role": role,
        "bbox": box,
        "color_hint": region.get("color_hint") or "",
        "font_hint": region.get("font_hint") or "",
        "weight_hint": weight_hint,
        "alignment": infer_alignment(box),
    }


def merge_inline_prefix_body_regions(regions: list[dict[str, Any]]) -> list[dict[str, Any]]:
    merged: list[dict[str, Any]] = []
    used: set[int] = set()
    for index, region in enumerate(regions):
        if index in used:
            continue
        box = clamp_box(region)
        role = infer_role(region)
        text = str(region.get("text") or "").strip()
        prefix_like = role == "subtitle" and text.endswith((":", "："))
        if not prefix_like:
            merged.append(region)
            continue
        match_index: int | None = None
        for candidate_index, candidate in enumerate(regions):
            if candidate_index == index or candidate_index in used:
                continue
            candidate_box = clamp_box(candidate)
            if infer_role(candidate) != "body":
                continue
            same_left = abs(float(candidate_box["x"]) - float(box["x"])) <= 0.018
            starts_inline = float(candidate_box["y"]) <= float(box["y"]) + max(0.035, float(box["height"]) * 1.25)
            overlaps_vertically = float(candidate_box["y"]) + float(candidate_box["height"]) >= float(box["y"])
            if same_left and starts_inline and overlaps_vertically and float(candidate_box["width"]) > float(box["width"]):
                match_index = candidate_index
                break
        if match_index is None:
            merged.append(region)
            continue
        body = regions[match_index]
        body_box = clamp_box(body)
        x1 = min(float(box["x"]), float(body_box["x"]))
        y1 = min(float(box["y"]), float(body_box["y"]))
        x2 = max(float(box["x"]) + float(box["width"]), float(body_box["x"]) + float(body_box["width"]))
        y2 = max(float(box["y"]) + float(box["height"]), float(body_box["y"]) + float(body_box["height"]))
        merged_region = dict(body)
        merged_region.update({
            "text": f"{text}{str(body.get('text') or '').strip()}",
            "role": "body",
            "x": x1,
            "y": y1,
            "width": x2 - x1,
            "height": y2 - y1,
        })
        used.add(match_index)
        merged.append(merged_region)
    return merged


def _box_for_group(group: dict[str, Any]) -> dict[str, float]:
    box = _natural_text_box_for_group(group)
    if group.get("box_width_hint") is not None:
        box["width"] = max(float(box["width"]), float(group["box_width_hint"]))
    if group.get("box_height_hint") is not None:
        box["height"] = max(float(box["height"]), float(group["box_height_hint"]))
    return clamp_box(box)


def _pill_background_box_for_group(group: dict[str, Any]) -> dict[str, float]:
    return _expanded_pill_box(group.get("bbox") or {})


def sample_saturated_background(rgb: np.ndarray, box: dict[str, Any]) -> tuple[int, int, int]:
    h, w, _ = rgb.shape
    b = clamp_box(box)
    x1 = max(0, int(float(b["x"]) * w))
    y1 = max(0, int(float(b["y"]) * h))
    x2 = min(w, int((float(b["x"]) + float(b["width"])) * w))
    y2 = min(h, int((float(b["y"]) + float(b["height"])) * h))
    crop = rgb[y1:y2, x1:x2]
    if crop.size == 0:
        return sample_background(rgb, b)
    flat = crop.reshape(-1, 3).astype(np.float32)
    max_c = flat.max(axis=1)
    min_c = flat.min(axis=1)
    sat = np.zeros_like(max_c)
    np.divide(max_c - min_c, np.maximum(max_c, 1), out=sat, where=max_c > 1)
    luma = np.array([luminance(px) for px in flat])
    selected = flat[(sat > 0.28) & (luma > 45) & (luma < 245)]
    if len(selected) < 20:
        return sample_background(rgb, b)
    color = np.median(selected, axis=0)
    return tuple(int(max(0, min(255, round(v)))) for v in color)


def should_full_fill_text_cleanup(group: dict[str, Any], rgb: np.ndarray) -> bool:
    if group.get("background_pill"):
        return True
    role = str(group.get("role") or "")
    box = clamp_box(group.get("bbox") or {})
    if role not in {"body", "caption", "footer", "subtitle"}:
        return False
    if role == "subtitle" and float(box["height"]) > 0.06:
        return False
    h, w, _ = rgb.shape
    x1 = max(0, int(float(box["x"]) * w))
    y1 = max(0, int(float(box["y"]) * h))
    x2 = min(w, int((float(box["x"]) + float(box["width"])) * w))
    y2 = min(h, int((float(box["y"]) + float(box["height"])) * h))
    crop = rgb[y1:y2, x1:x2]
    if crop.size == 0:
        return False
    flat = crop.reshape(-1, 3).astype(np.float32)
    luma = 0.2126 * flat[:, 0] + 0.7152 * flat[:, 1] + 0.0722 * flat[:, 2]
    median_luma = float(np.median(luma))
    return median_luma >= 172 or median_luma <= 82


def should_full_fill_display_cleanup(group: dict[str, Any], visual_complexity: dict[str, float] | None) -> bool:
    role = str(group.get("role") or "")
    if role not in {"title", "subtitle", "label"}:
        return False
    if group.get("background_pill"):
        return False
    box = clamp_box(group.get("bbox") or {})
    text = str(group.get("text") or "")
    if role == "title":
        return float(box["height"]) >= 0.052 or estimate_text_units(text) >= 8
    if role == "subtitle":
        return (
            float(box["height"]) >= 0.040
            or (is_complex_visual_slide(visual_complexity) and estimate_text_units(text) >= 4)
        )
    return is_complex_visual_slide(visual_complexity) and float(box["height"]) >= 0.032


def cleanup_box_for_group(
    group: dict[str, Any],
    rgb: np.ndarray,
    visual_complexity: dict[str, float] | None,
) -> dict[str, Any]:
    if group.get("background_pill"):
        cleanup_box = _pill_background_box_for_group(group)
        cleanup_box["full_fill"] = True
        return cleanup_box
    if should_full_fill_text_cleanup(group, rgb):
        cleanup_box = _expanded_solid_cleanup_box(group)
        cleanup_box["solid_fill"] = True
        return cleanup_box
    if should_full_fill_display_cleanup(group, visual_complexity):
        cleanup_box = _expanded_display_cleanup_box(group)
        cleanup_box["full_fill"] = True
        return cleanup_box
    return clamp_box(group["bbox"])


def quality_retry_cleanup_box_for_group(
    group: dict[str, Any],
    previous_box: dict[str, Any] | None = None,
) -> dict[str, Any]:
    if previous_box:
        b = clamp_box(previous_box)
        box_h = float(b["height"])
        x_margin = min(0.008, max(0.003, box_h * 0.06))
        y_margin = min(0.010, max(0.004, box_h * 0.08))
        cleanup_box = clamp_box({
            "x": b["x"] - x_margin,
            "y": b["y"] - y_margin,
            "width": b["width"] + x_margin * 2,
            "height": b["height"] + y_margin * 2,
        })
        if previous_box.get("solid_fill"):
            cleanup_box["solid_fill"] = True
        else:
            cleanup_box["full_fill"] = True
        return cleanup_box

    role = str(group.get("role") or "")
    b = clamp_box(group.get("bbox") or {})
    box_h = float(b["height"])
    if role in {"title", "subtitle", "label"}:
        x_margin = min(0.070, max(0.020, box_h * 0.48))
        y_margin = min(0.078, max(0.024, box_h * 0.60))
    else:
        x_margin = min(0.040, max(0.014, box_h * 0.30))
        y_margin = min(0.052, max(0.018, box_h * 0.46))
    cleanup_box = clamp_box({
        "x": b["x"] - x_margin,
        "y": b["y"] - y_margin,
        "width": b["width"] + x_margin * 2,
        "height": b["height"] + y_margin * 2,
    })
    cleanup_box["full_fill"] = True
    return cleanup_box


def add_pill_background(slide, group: dict[str, Any], rendered_rgb: np.ndarray) -> tuple[int, int, int]:
    box = _pill_background_box_for_group(group)
    left, top, width, height = emu_from_top_box(box)
    color = sample_saturated_background(rendered_rgb, box)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.name = "Editable text pill background"
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(*color)
    pill.line.fill.background()
    return color


def add_textbox(slide, group: dict[str, Any], rendered_rgb: np.ndarray) -> None:
    pill_color = None
    if group.get("background_pill"):
        pill_color = add_pill_background(slide, group, rendered_rgb)
    box = _box_for_group(group)
    left, top, width, height = emu_from_top_box(box)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = f"Editable restored text - {group.get('role') or 'text'}"
    frame = textbox.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    text = str(group.get("text") or "").strip()
    role = str(group.get("role") or "body")
    frame.word_wrap = len(split_text_lines(text)) > 1 or role in {"body", "caption", "footer"}
    p = frame.paragraphs[0]
    if group.get("font_size_hint") is not None:
        font_size_pt = float(group["font_size_hint"])
    else:
        font_size_pt = _fitted_size_for_group(group, box)
    p.text = wrap_text_for_box(text, box, font_size_pt, role)
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    alignment = str(group.get("alignment") or "left")
    if alignment == "center":
        p.alignment = PP_ALIGN.CENTER
    elif alignment == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    font = p.font
    font.name = choose_font_name(text, group.get("font_hint"))
    font.bold = str(group.get("weight_hint") or "").lower() == "bold"
    font.size = Pt(font_size_pt)
    bg = sample_background(rendered_rgb, box)
    color = sample_text_color(rendered_rgb, box, bg)
    if pill_color is not None:
        color = (255, 255, 255)
    font.color.rgb = RGBColor(*parse_hex_color(group.get("color_hint"), color))


def _default_work_dir(output_path: str) -> str:
    out = Path(output_path)
    return str(out.parent / f".{out.stem}_assets")


def build_editable_pptx(
    *,
    slide_images: list[dict[str, Any]],
    output_path: str,
    ocr_provider: OcrProvider | None = None,
    progress_callback: Callable[[int, int, str], None] | None = None,
    work_dir: str | None = None,
    restore_mode: str | None = None,
    reuse_ocr_cache: bool = True,
) -> EditablePptxResult:
    mode = normalize_editable_pptx_restore_mode(restore_mode)
    provider = ocr_provider or (lambda image_path, page_num: minimax_ocr_regions(image_path, page_num, restore_mode=mode))
    work = Path(work_dir or _default_work_dir(output_path))
    work.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    text_box_count = 0
    visual_asset_count = 0
    ocr_failed_pages: list[int] = []
    qa_retry_pages: list[int] = []
    quality_fallback_pages: list[int] = []
    quality_warning_pages: list[int] = []
    sorted_slides = sorted(slide_images, key=lambda item: int(item.get("page_num") or 0))
    diagnostics = EditablePptxDiagnostics(restore_mode=mode)

    for index, slide_data in enumerate(sorted_slides, start=1):
        page_num = int(slide_data.get("page_num") or index)
        image_path = str(slide_data.get("image_path") or "")
        if progress_callback:
            progress_callback(index - 1, len(sorted_slides), f"正在解析第 {page_num} 页")
        slide = prs.slides.add_slide(blank)
        if not image_path or not os.path.exists(image_path):
            ocr_failed_pages.append(page_num)
            diagnostics.pages.append(EditablePptxPageDiagnostics(page_num=page_num, ocr_failed=True))
            continue

        rendered = Image.open(image_path).convert("RGB")
        rendered_rgb = np.asarray(rendered)
        visual_complexity = slide_visual_complexity(rendered_rgb)
        raw_regions = cached_or_run_ocr_regions(
            provider,
            image_path,
            page_num,
            work_dir=work,
            restore_mode=mode,
            reuse_cache=reuse_ocr_cache,
        )
        if not raw_regions:
            ocr_failed_pages.append(page_num)

        normalized_regions = [region for region in (_coerce_region(r) for r in raw_regions) if region]
        text_boxes_for_detection = [clamp_box(region) for region in normalized_regions]
        image_blocks = detect_image_blocks(image_path, text_boxes_for_detection)
        page_diag = EditablePptxPageDiagnostics(
            page_num=page_num,
            raw_region_count=len(raw_regions),
            normalized_region_count=len(normalized_regions),
            visual_asset_count=len(image_blocks),
            ocr_failed=not bool(raw_regions),
        )
        text_regions = []
        for region in normalized_regions:
            keep, reason = should_restore_text_with_reason(
                region,
                image_blocks,
                rendered_rgb,
                mode,
                visual_complexity=visual_complexity,
            )
            if keep:
                text_regions.append(region)
            else:
                page_diag.rejection_reasons[reason] = page_diag.rejection_reasons.get(reason, 0) + 1
        text_regions = merge_inline_prefix_body_regions(text_regions)
        groups = []
        for region in text_regions:
            group = _group_from_region(region)
            if needs_background_fill_for_complex_text(str(group.get("role") or ""), visual_complexity):
                group["background_fill"] = True
            if needs_pill_background_for_complex_text(
                str(group.get("role") or ""),
                visual_complexity,
                rendered_rgb,
                group["bbox"],
            ):
                group["background_pill"] = True
                group["alignment"] = "center"
            groups.append(group)
        page_diag.restored_text_count = len(groups)
        normalize_same_level_text_metrics(groups)
        apply_complex_slide_text_safety_hints(groups, visual_complexity)
        original_pic = slide.shapes.add_picture(image_path, 0, 0, width=SLIDE_W_EMU, height=SLIDE_H_EMU)
        original_pic.name = "Original slide image"

        text_cleanup_boxes = [cleanup_box_for_group(group, rendered_rgb, visual_complexity) for group in groups]
        asset_source = rendered.copy()
        for cleanup_box in text_cleanup_boxes:
            apply_cleanup_box_to_image(asset_source, cleanup_box)

        underlay = rendered.copy()
        cleanup_patch_count = 0

        def add_cleanup_patch(cleanup_box: dict[str, Any], patch_path: str, shape_name: str) -> bool:
            patch_box = create_cleanup_patch(underlay, cleanup_box, patch_path)
            if not patch_box:
                return False
            paste_cleanup_patch_to_image(underlay, patch_path, patch_box)
            left, top, width, height = emu_from_top_box(patch_box)
            patch_pic = slide.shapes.add_picture(patch_path, left, top, width=width, height=height)
            patch_pic.name = shape_name
            return True

        if groups or image_blocks:
            cleanup_boxes = []
            for block in image_blocks:
                cleanup_box = clamp_box(block)
                cleanup_box["full_fill"] = True
                cleanup_boxes.append(cleanup_box)
            cleanup_boxes.extend(text_cleanup_boxes)
            for cleanup_idx, cleanup_box in enumerate(cleanup_boxes, start=1):
                patch_path = str(work / f"slide_{page_num:02d}_cleanup_patch_{cleanup_idx:02d}.png")
                if add_cleanup_patch(cleanup_box, patch_path, f"Editable cleanup patch - {cleanup_idx}"):
                    cleanup_patch_count += 1

        for asset_idx, block in enumerate(image_blocks, start=1):
            asset_path = str(work / f"slide_{page_num:02d}_asset_{asset_idx:02d}.png")
            crop_asset_from_image(asset_source, block, asset_path)
            paste_asset_crop_to_image(underlay, asset_path, block)
            left, top, width, height = emu_from_top_box(block)
            pic = slide.shapes.add_picture(asset_path, left, top, width=width, height=height)
            pic.name = f"Replaceable visual asset - {asset_idx}"
            visual_asset_count += 1

        residuals = residual_text_groups(rendered_rgb, np.asarray(underlay.convert("RGB")), groups, text_cleanup_boxes)
        page_diag.qa_retry_count = len(residuals)
        if residuals:
            qa_retry_pages.append(page_num)
            cleanup_box_by_group_id = {id(group): cleanup_box for group, cleanup_box in zip(groups, text_cleanup_boxes)}
            retry_cleanup_boxes = []
            for retry_idx, (group, _score) in enumerate(residuals, start=1):
                retry_box = quality_retry_cleanup_box_for_group(group, cleanup_box_by_group_id.get(id(group)))
                retry_cleanup_boxes.append(retry_box)
                patch_path = str(work / f"slide_{page_num:02d}_qa_cleanup_patch_{retry_idx:02d}.png")
                if add_cleanup_patch(retry_box, patch_path, f"Editable QA cleanup patch - {retry_idx}"):
                    cleanup_patch_count += 1
            retry_groups = [group for group, _score in residuals]
            retry_residuals = residual_text_groups(
                rendered_rgb,
                np.asarray(underlay.convert("RGB")),
                retry_groups,
                retry_cleanup_boxes,
                threshold=TEXT_RESIDUAL_FALLBACK_THRESHOLD,
            )
            if retry_residuals:
                quality_warning_pages.append(page_num)
                page_diag.quality_warning = True
        page_diag.cleanup_patch_count = cleanup_patch_count

        for group in groups:
            add_textbox(slide, group, rendered_rgb)
            text_box_count += 1

        notes = str(slide_data.get("speaker_notes") or "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        if progress_callback:
            progress_callback(index, len(sorted_slides), f"第 {page_num} 页可编辑版已处理")
        diagnostics.pages.append(page_diag)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return EditablePptxResult(
        output_path=output_path,
        slide_count=len(sorted_slides),
        text_box_count=text_box_count,
        visual_asset_count=visual_asset_count,
        ocr_failed_pages=sorted({int(page) for page in ocr_failed_pages}),
        qa_retry_pages=sorted({int(page) for page in qa_retry_pages}),
        quality_fallback_pages=sorted({int(page) for page in quality_fallback_pages}),
        quality_warning_pages=sorted({int(page) for page in quality_warning_pages}),
        diagnostics=diagnostics,
    )


def build_project_slide_images(slides: Iterable[Any]) -> list[dict[str, Any]]:
    slide_images = []
    for slide in sorted(slides, key=lambda item: int(getattr(item, "page_num", 0) or 0)):
        image_path = str(getattr(slide, "image_path", "") or "")
        if not image_path:
            continue
        speaker_notes = ""
        content_json = getattr(slide, "content_json", None)
        if isinstance(content_json, dict):
            speaker_notes = str(content_json.get("speaker_notes") or "")
        slide_images.append({
            "page_num": int(getattr(slide, "page_num", 0) or 0),
            "image_path": image_path,
            "speaker_notes": speaker_notes,
        })
    return slide_images
