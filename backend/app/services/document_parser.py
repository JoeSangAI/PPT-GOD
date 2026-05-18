import os
from io import BytesIO
import re
from typing import Callable

from app.services.pptx_page_recovery import recover_sparse_slide_text

DocumentProgressCallback = Callable[[dict], None]


def _page_recovery_timeout_seconds() -> float:
    try:
        return max(3.0, float(os.getenv("PPTGOD_PPT_PAGE_OCR_TIMEOUT_SECONDS", "20")))
    except ValueError:
        return 20.0


def _emit_progress(progress_callback: DocumentProgressCallback | None, **payload) -> None:
    if not progress_callback:
        return
    try:
        progress_callback(payload)
    except Exception:
        pass


def parse_document(
    file_bytes: bytes,
    filename: str,
    *,
    progress_callback: DocumentProgressCallback | None = None,
    recovery_timeout_seconds: float | None = None,
) -> str:
    """根据文件扩展名解析文档，返回纯文本内容。"""
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        return _parse_pdf(file_bytes)
    elif ext in (".docx", ".doc"):
        return _parse_docx(file_bytes)
    elif ext in (".pptx", ".ppt"):
        return _parse_pptx(
            file_bytes,
            filename,
            progress_callback=progress_callback,
            recovery_timeout_seconds=recovery_timeout_seconds,
        )
    elif ext in (".md", ".txt", ".json", ".csv", ".html", ".htm"):
        return _parse_text(file_bytes)
    else:
        # 未知格式，尝试按文本读取
        try:
            return file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            raise ValueError(f"不支持的文件格式: {ext}，请上传 PDF、Word、PPT、Markdown 或文本文件。")


def _parse_pdf(file_bytes: bytes) -> str:
    import fitz  # PyMuPDF

    text_parts = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n\n".join(text_parts)


def _parse_docx(file_bytes: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _parse_pptx(
    file_bytes: bytes,
    filename: str = "",
    *,
    progress_callback: DocumentProgressCallback | None = None,
    recovery_timeout_seconds: float | None = None,
) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(BytesIO(file_bytes))
    total_pages = len(prs.slides)
    recovery_timeout = recovery_timeout_seconds if recovery_timeout_seconds is not None else _page_recovery_timeout_seconds()
    text_parts = []
    text_parts.append(f'--- PPT_SOURCE filename="{filename}" pages={total_pages} ---')
    for i, slide in enumerate(prs.slides, 1):
        _emit_progress(
            progress_callback,
            phase="pptx_page_read",
            current_page=i,
            total_pages=total_pages,
            completed_count=i - 1,
            message=f"正在识别 PPT 第 {i}/{total_pages} 页文字和截图...",
        )
        slide_texts = _extract_shape_texts(slide.shapes, MSO_SHAPE_TYPE)
        slide_texts.extend(
            text for text in recover_sparse_slide_text(
                slide,
                page_num=i,
                source_filename=filename,
                existing_lines=slide_texts,
                timeout_seconds=recovery_timeout,
            )
            if text
        )
        notes = _extract_notes_text(slide)
        if notes:
            slide_texts.append("【备注】\n" + notes)
        text_parts.append(f"--- 第{i}页 ---\n" + "\n".join(slide_texts))
        _emit_progress(
            progress_callback,
            phase="pptx_page_done",
            current_page=i,
            total_pages=total_pages,
            completed_count=i,
            message=f"已读取 PPT 第 {i}/{total_pages} 页",
        )
    return "\n\n".join(text_parts)


def _append_unique(parts: list[str], text: str) -> None:
    cleaned = re.sub(r"[ \t]+", " ", str(text or "")).strip()
    if cleaned and cleaned not in parts:
        parts.append(cleaned)


def _extract_shape_texts(shapes, shape_type_enum) -> list[str]:
    texts: list[str] = []
    for shape in shapes:
        if getattr(shape, "shape_type", None) == shape_type_enum.GROUP:
            texts.extend(text for text in _extract_shape_texts(shape.shapes, shape_type_enum) if text not in texts)
            continue

        # Skip page-level metadata placeholders (date, footer, slide number, etc.)
        # These carry no slide content value and only add noise.
        if getattr(shape, "is_placeholder", False):
            from pptx.enum.shapes import PP_PLACEHOLDER
            ph_type = getattr(getattr(shape, "placeholder_format", None), "type", None)
            if ph_type in {PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.HEADER, PP_PLACEHOLDER.SLIDE_NUMBER}:
                continue

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = []
                for cell in row.cells:
                    cell_text = re.sub(r"\s+", " ", cell.text or "").strip()
                    if cell_text:
                        cells.append(cell_text)
                if cells:
                    _append_unique(texts, " | ".join(cells))
            continue

        if getattr(shape, "has_chart", False):
            chart_title = getattr(getattr(shape.chart, "chart_title", None), "text_frame", None)
            if chart_title is not None:
                _append_unique(texts, chart_title.text)

        if hasattr(shape, "text"):
            _append_unique(texts, shape.text)
    return texts


def _extract_notes_text(slide) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes = []
        for shape in slide.notes_slide.notes_text_frame.paragraphs:
            line = "".join(run.text for run in shape.runs).strip()
            if line:
                notes.append(line)
        return "\n".join(notes).strip()
    except Exception:
        return ""


def detect_ppt_sources(documents_text: str) -> list[dict]:
    """Return PPT upload metadata embedded in extracted document text."""
    sources = []
    pattern = r'---\s*PPT_SOURCE(?:\s+filename="([^"]*)")?\s+pages=(\d+)\s*---'
    for match in re.finditer(pattern, documents_text or ""):
        sources.append({
            "filename": match.group(1) or "",
            "pages": int(match.group(2)),
        })
    return sources


def _parse_text(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8")


def extract_text_preview(text: str, max_chars: int = 300) -> str:
    """提取文本预览，截取前 max_chars 个字符。"""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "..."
