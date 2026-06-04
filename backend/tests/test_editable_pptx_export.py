from pathlib import Path

import numpy as np
from PIL import Image, ImageDraw
from pptx import Presentation

from app.services.editable_pptx_export import (
    build_editable_pptx,
    build_minimax_ocr_prompt,
    choose_font_name,
    parse_vlm_ocr_regions,
    prepare_clean_background,
    read_cached_ocr_regions,
    run_ocr_with_retries,
)
from app.services import editable_pptx_export as editable_export
from app.services.editable_pptx_diagnostics import inspect_pptx_editability


def draw_region_marks(img: Image.Image, regions: list[dict], *, color: str = "black") -> None:
    draw = ImageDraw.Draw(img)
    width, height = img.size
    for region in regions:
        x = int(region["x"] * width)
        y = int(region["y"] * height)
        w = int(region["width"] * width)
        h = int(region["height"] * height)
        draw.rectangle(
            (x + 6, y + max(3, h // 3), x + max(12, w - 6), y + max(6, h // 3 + max(3, h // 5))),
            fill=color,
        )


def pptx_text_shapes(path: Path) -> dict[str, object]:
    prs = Presentation(str(path))
    return {
        shape.text.strip(): shape
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    }


def pptx_shape_names(path: Path) -> list[str]:
    return [shape.name for shape in Presentation(str(path)).slides[0].shapes]


def test_parse_vlm_ocr_regions_accepts_fenced_json_and_clamps_boxes():
    raw = """```json
{
  "text_regions": [
    {"text": "标题", "x": -0.02, "y": 0.1, "width": 0.42, "height": 0.08, "confidence": 0.91},
    {"text": "太小", "x": 0.5, "y": 0.5, "width": 0.001, "height": 0.001},
    {"text": "", "x": 0.2, "y": 0.2, "width": 0.1, "height": 0.05}
  ]
}
```"""

    regions = parse_vlm_ocr_regions(raw)

    assert len(regions) == 1
    assert regions[0]["text"] == "标题"
    assert regions[0]["x"] == 0.0
    assert regions[0]["y"] == 0.1
    assert regions[0]["width"] == 0.4
    assert regions[0]["height"] == 0.08


def test_build_editable_pptx_normalizes_same_level_text_and_keeps_timeline_visual(tmp_path):
    slide_path = tmp_path / "slide_01.png"
    output_path = tmp_path / "editable.pptx"

    fake_regions = [
        {"text": "Steve Jobs回归Apple", "x": 0.05, "y": 0.18, "width": 0.49, "height": 0.089},
        {"text": "Jobs回来后没有先讨论CPU", "x": 0.052, "y": 0.417, "width": 0.286, "height": 0.0444},
        {"text": "没有先讨论性能", "x": 0.05, "y": 0.503, "width": 0.161, "height": 0.0444},
        {"text": "1997", "x": 0.231, "y": 0.904, "width": 0.049, "height": 0.039},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(img)
    for region in fake_regions[:3]:
        x = int(region["x"] * 1280)
        y = int(region["y"] * 720)
        w = int(region["width"] * 1280)
        h = int(region["height"] * 720)
        draw.rectangle(
            (x + 6, y + max(3, h // 3), x + max(12, w - 6), y + max(6, h // 3 + h // 5)),
            fill="black",
        )
    img.save(slide_path)

    def fake_ocr(_image_path: str, _page_num: int):
        return fake_regions

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path), "speaker_notes": ""}],
        output_path=str(output_path),
        ocr_provider=fake_ocr,
    )

    assert result.slide_count == 1
    assert result.text_box_count == 3
    text_shapes = pptx_text_shapes(output_path)
    assert "1997" not in text_shapes
    jobs_size = text_shapes["Jobs回来后没有先讨论CPU"].text_frame.paragraphs[0].font.size.pt
    performance_size = text_shapes["没有先讨论性能"].text_frame.paragraphs[0].font.size.pt
    assert jobs_size == performance_size
    assert round(text_shapes["Jobs回来后没有先讨论CPU"].width / 914400, 2) == round(text_shapes["没有先讨论性能"].width / 914400, 2)


def test_build_editable_pptx_returns_page_diagnostics(tmp_path):
    slide_path = tmp_path / "diagnostics_slide.png"
    output_path = tmp_path / "diagnostics_editable.pptx"
    fake_regions = [
        {"text": "主标题", "x": 0.1, "y": 0.1, "width": 0.3, "height": 0.08, "role": "title", "confidence": 0.95},
        {"text": "1997", "x": 0.2, "y": 0.91, "width": 0.05, "height": 0.03, "role": "page_marker", "confidence": 0.9},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
        reuse_ocr_cache=False,
    )

    assert result.diagnostics is not None
    assert result.diagnostics.page_count == 1
    page = result.diagnostics.pages[0]
    assert page.raw_region_count == 2
    assert page.normalized_region_count == 2
    assert page.restored_text_count == 1
    assert page.rejection_reasons["role_not_allowed_for_mode"] == 1
    inspection = inspect_pptx_editability(output_path)
    assert inspection.has_editable_text is True


def test_should_restore_text_reason_rejects_small_standard_label():
    image = np.zeros((720, 1280, 3), dtype=np.uint8) + 255
    region = {
        "text": "小标签",
        "x": 0.2,
        "y": 0.2,
        "width": 0.04,
        "height": 0.018,
        "role": "label",
        "editable": True,
        "confidence": 0.9,
    }

    keep, reason = editable_export.should_restore_text_with_reason(region, [], image, "standard", visual_complexity=None)

    assert keep is False
    assert reason == "standard_small_auxiliary_text"


def test_merge_cleanup_boxes_combines_nearby_text_boxes():
    boxes = [
        {"x": 0.10, "y": 0.10, "width": 0.20, "height": 0.04},
        {"x": 0.10, "y": 0.145, "width": 0.21, "height": 0.04},
        {"x": 0.70, "y": 0.70, "width": 0.10, "height": 0.03},
    ]

    merged = editable_export.merge_cleanup_boxes(boxes, gap=0.015)

    assert len(merged) == 2
    assert merged[0]["height"] > 0.08
    assert merged[1]["x"] == 0.70


def test_same_level_text_normalization_applies_outside_left_column(tmp_path):
    slide_path = tmp_path / "right_column.png"
    output_path = tmp_path / "right_column_editable.pptx"
    fake_regions = [
        {"text": "从情绪价值出发", "x": 0.56, "y": 0.34, "width": 0.16, "height": 0.042, "role": "body"},
        {"text": "让品牌成为用户表达自我的方式", "x": 0.56, "y": 0.41, "width": 0.30, "height": 0.042, "role": "body"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    assert result.text_box_count == 2
    text_shapes = pptx_text_shapes(output_path)
    first = text_shapes["从情绪价值出发"]
    second = text_shapes["让品牌成为用户表达自我的方式"]
    assert first.text_frame.paragraphs[0].font.size.pt == second.text_frame.paragraphs[0].font.size.pt
    assert round(first.width / 914400, 2) == round(second.width / 914400, 2)


def test_short_uppercase_main_title_is_not_dropped_as_logo(tmp_path):
    slide_path = tmp_path / "short_brand_title.png"
    output_path = tmp_path / "short_brand_title_editable.pptx"
    fake_regions = [
        {"text": "IBM", "x": 0.40, "y": 0.52, "width": 0.20, "height": 0.085, "role": "title"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    assert result.text_box_count == 1
    assert "IBM" in pptx_text_shapes(output_path)


def test_vlm_non_editable_regions_remain_as_image_content(tmp_path):
    slide_path = tmp_path / "non_editable_logo.png"
    output_path = tmp_path / "non_editable_logo_editable.pptx"
    fake_regions = [
        {"text": "Vibram", "x": 0.28, "y": 0.55, "width": 0.12, "height": 0.06, "role": "logo", "editable": False},
        {"text": "看得见的信任", "x": 0.52, "y": 0.18, "width": 0.32, "height": 0.08, "role": "title"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    text_shapes = pptx_text_shapes(output_path)
    assert result.text_box_count == 1
    assert "Vibram" not in text_shapes
    assert "看得见的信任" in text_shapes


def test_large_chinese_title_keeps_readable_size_and_exact_short_chars(tmp_path):
    slide_path = tmp_path / "large_chinese_title.png"
    output_path = tmp_path / "large_chinese_title_editable.pptx"
    fake_regions = [
        {
            "text": "中国已经有能力制造世界级产品吗？",
            "x": 0.033,
            "y": 0.072,
            "width": 0.74,
            "height": 0.077,
            "role": "title",
            "font_hint": "sans-serif",
        },
    ]
    img = Image.new("RGB", (1792, 1024), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 4, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    assert result.text_box_count == 1
    title = pptx_text_shapes(output_path)["中国已经有能力制造世界级产品吗？"]
    assert "有" in title.text
    assert title.text_frame.paragraphs[0].font.size.pt >= 33
    assert title.text_frame.paragraphs[0].font.name == "PingFang SC"


def test_dense_cjk_title_is_capped_by_box_height(tmp_path):
    slide_path = tmp_path / "dense_cjk_title.png"
    output_path = tmp_path / "dense_cjk_title_editable.pptx"
    title_text = "一条主线：围绕大学生手机消费场景，建立 6-12 月学生市场长周期品牌心智"
    fake_regions = [
        {
            "text": title_text,
            "x": 0.08,
            "y": 0.08,
            "width": 0.58,
            "height": 0.10,
            "role": "title",
            "weight_hint": "bold",
        },
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    assert result.text_box_count == 1
    title = pptx_text_shapes(output_path)[title_text]
    assert 13 <= title.text_frame.paragraphs[0].font.size.pt <= 22


def test_top_cjk_title_fragment_is_capped_and_normalized_with_next_line(tmp_path):
    slide_path = tmp_path / "fragmented_cjk_title.png"
    output_path = tmp_path / "fragmented_cjk_title_editable.pptx"
    first_line = "一条主线：围绕大学生手机消费场景，"
    second_line = "建立 6-12 月学生市场长周期品牌心智"
    fake_regions = [
        {"text": first_line, "x": 0.06, "y": 0.085, "width": 0.78, "height": 0.08, "role": "title"},
        {"text": second_line, "x": 0.06, "y": 0.175, "width": 0.78, "height": 0.08, "role": "title"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    assert result.text_box_count == 2
    text_shapes = pptx_text_shapes(output_path)
    first_size = text_shapes[first_line].text_frame.paragraphs[0].font.size.pt
    second_size = text_shapes[second_line].text_frame.paragraphs[0].font.size.pt
    assert first_size <= 22
    assert second_size == first_size


def test_generic_font_hints_are_normalized_to_stable_deck_fonts():
    assert choose_font_name("中国已经有能力制造世界级产品吗？", "Microsoft YaHei") == "PingFang SC"
    assert choose_font_name("品牌还有价值吗？", "sans-serif") == "PingFang SC"
    assert choose_font_name("We engineer emotions", "Arial") == "Helvetica Neue"


def test_long_body_copy_is_not_treated_as_bold_display_text(tmp_path):
    slide_path = tmp_path / "dense_body_copy.png"
    output_path = tmp_path / "dense_body_copy_editable.pptx"
    long_copy = (
        "本概要仅是为了帮助潜在购买方针对公司做出一个初步的评估，以确定是否对于该项目进行进一步考察。"
        "该报告并不能作为一个针对公司的出价基础，也不能作为任何投资决策的基础。"
    )
    fake_regions = [
        {
            "text": long_copy,
            "x": 0.10,
            "y": 0.32,
            "width": 0.76,
            "height": 0.13,
            "role": "body",
        },
    ]
    img = Image.new("RGB", (1280, 720), (18, 28, 42))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    assert result.text_box_count == 1
    shape = next(shape for shape in Presentation(str(output_path)).slides[0].shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
    font = shape.text_frame.paragraphs[0].font
    assert font.bold is False
    assert font.size.pt <= 12.5


def test_label_role_from_vlm_is_not_promoted_to_subtitle(tmp_path):
    slide_path = tmp_path / "large_label_box.png"
    output_path = tmp_path / "large_label_box_editable.pptx"
    fake_regions = [
        {"text": "多模态", "x": 0.47, "y": 0.42, "width": 0.12, "height": 0.062, "role": "label"},
    ]
    img = Image.new("RGB", (1280, 720), (10, 22, 40))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    assert result.text_box_count == 1
    shape = next(shape for shape in Presentation(str(output_path)).slides[0].shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
    assert shape.name == "Editable restored text - label"
    assert shape.text_frame.paragraphs[0].font.size.pt <= 18.5


def test_minimax_ocr_prompt_requires_exact_title_character_review():
    prompt = build_minimax_ocr_prompt(4)

    assert "逐字复核" in prompt
    assert "不要漏掉单字" in prompt
    assert "有、是、为、的、不、吗" in prompt


def test_large_display_words_are_not_capped_as_body_copy(tmp_path):
    slide_path = tmp_path / "display_words.png"
    output_path = tmp_path / "display_words_editable.pptx"
    fake_regions = [
        {"text": "梦想", "x": 0.39, "y": 0.32, "width": 0.22, "height": 0.17, "role": "body"},
        {"text": "荣耀", "x": 0.70, "y": 0.32, "width": 0.22, "height": 0.17, "role": "body"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 19, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    assert result.text_box_count == 2
    text_shapes = pptx_text_shapes(output_path)
    assert text_shapes["梦想"].text_frame.paragraphs[0].font.size.pt >= 38
    assert text_shapes["荣耀"].text_frame.paragraphs[0].font.size.pt >= 38


def test_ocr_retry_recovers_from_transient_empty_response(tmp_path):
    slide_path = tmp_path / "retry.png"
    slide_path.write_bytes(b"not-used")
    calls = []

    def flaky_provider(image_path: str, page_num: int):
        calls.append((image_path, page_num))
        if len(calls) == 1:
            return []
        return [{"text": "恢复成功", "x": 0.1, "y": 0.1, "width": 0.3, "height": 0.08}]

    regions = run_ocr_with_retries(
        flaky_provider,
        str(slide_path),
        1,
        retry_count=2,
        retry_delay_seconds=0,
    )

    assert len(calls) == 2
    assert regions[0]["text"] == "恢复成功"


def test_build_editable_pptx_writes_and_reuses_page_ocr_cache(tmp_path):
    slide_path = tmp_path / "cached_slide.png"
    output_path = tmp_path / "cached_editable.pptx"
    work_dir = tmp_path / "work"
    fake_regions = [
        {"text": "缓存文字", "x": 0.12, "y": 0.2, "width": 0.32, "height": 0.08, "role": "title"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    calls = 0

    def fake_ocr(_image_path: str, _page_num: int):
        nonlocal calls
        calls += 1
        return fake_regions

    build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=fake_ocr,
        work_dir=str(work_dir),
    )
    cached = read_cached_ocr_regions(work_dir, 1)
    assert cached and cached[0]["text"] == "缓存文字"

    second_output = tmp_path / "cached_editable_second.pptx"
    build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(second_output),
        ocr_provider=lambda *_args: (_ for _ in ()).throw(AssertionError("cache should be used")),
        work_dir=str(work_dir),
    )

    assert calls == 1
    assert "缓存文字" in pptx_text_shapes(second_output)


def test_build_editable_pptx_can_force_fresh_ocr_even_when_cache_exists(tmp_path):
    slide_path = tmp_path / "force_fresh_ocr.png"
    output_path = tmp_path / "force_fresh_ocr_editable.pptx"
    work_dir = tmp_path / "work"
    mark_region = {"x": 0.1, "y": 0.2, "width": 0.28, "height": 0.07}
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, [mark_region])
    img.save(slide_path)
    call_count = {"count": 0}

    def fake_ocr(_image_path: str, _page_num: int):
        call_count["count"] += 1
        return [
            {
                "text": f"第 {call_count['count']} 次理解",
                "x": 0.1,
                "y": 0.2,
                "width": 0.28,
                "height": 0.07,
                "role": "title",
            }
        ]

    build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=fake_ocr,
        work_dir=str(work_dir),
    )
    fresh_output = tmp_path / "force_fresh_ocr_second.pptx"
    build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(fresh_output),
        ocr_provider=fake_ocr,
        work_dir=str(work_dir),
        reuse_ocr_cache=False,
    )

    assert call_count["count"] == 2
    assert "第 2 次理解" in pptx_text_shapes(fresh_output)


def test_text_cleanup_preserves_non_text_background_inside_large_box(tmp_path):
    slide_path = tmp_path / "large_text_on_subtle_background.png"
    clean_path = tmp_path / "cleaned.png"
    img = Image.new("RGB", (400, 220), "white")
    draw = ImageDraw.Draw(img)
    text_box = {"x": 0.12, "y": 0.34, "width": 0.76, "height": 0.22}
    # A subtle background detail inside the OCR text box. Full-rectangle cleanup
    # erases this and creates visible bands on real decks.
    draw.line((70, 104, 330, 104), fill=(242, 242, 242), width=4)
    draw.rectangle((72, 84, 254, 96), fill="black")
    img.save(slide_path)

    prepare_clean_background(str(slide_path), [text_box], [], str(clean_path))

    cleaned = Image.open(clean_path).convert("RGB")
    assert cleaned.getpixel((120, 104)) == (242, 242, 242)
    assert cleaned.getpixel((120, 90)) != (0, 0, 0)


def test_multiline_copy_keeps_ocr_column_width_instead_of_title_expansion(tmp_path):
    slide_path = tmp_path / "dark_multiline_copy.png"
    output_path = tmp_path / "dark_multiline_copy_editable.pptx"
    fake_regions = [
        {
            "text": "新一代爸爸：想要和宝贝贴脸亲亲每一秒，\n不想被不净胡茬阻碍，\n想要爸爸拥有更多高质量亲子时刻",
            "x": 0.17,
            "y": 0.29,
            "width": 0.36,
            "height": 0.138,
            "role": "title",
        },
    ]
    img = Image.new("RGB", (1280, 720), (8, 18, 35))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    assert result.text_box_count == 1
    shape = next(iter(pptx_text_shapes(output_path).values()))
    assert shape.width / 914400 <= 5.2
    assert shape.text_frame.paragraphs[0].font.size.pt <= 16


def test_long_cjk_body_copy_is_manually_wrapped_in_narrow_cards(tmp_path):
    slide_path = tmp_path / "long_cjk_body.png"
    output_path = tmp_path / "long_cjk_body_editable.pptx"
    fake_regions = [
        {
            "text": "初始状态：零维时空等于纯粹振动，万物皆静，无时间、无空间、无确定性",
            "x": 0.07,
            "y": 0.55,
            "width": 0.24,
            "height": 0.15,
            "role": "body",
        },
    ]
    img = Image.new("RGB", (1280, 720), (8, 18, 35))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
    )

    assert result.text_box_count == 1
    shape = next(iter(pptx_text_shapes(output_path).values()))
    assert "\x0b" in shape.text
    assert shape.text_frame.paragraphs[0].font.size.pt >= 10


def test_restore_mode_controls_chart_internal_text_restoration(tmp_path):
    slide_path = tmp_path / "chart_internal.png"
    standard_output = tmp_path / "chart_internal_standard.pptx"
    aggressive_output = tmp_path / "chart_internal_aggressive.pptx"
    fake_regions = [
        {
            "text": "关键节点",
            "x": 0.44,
            "y": 0.42,
            "width": 0.12,
            "height": 0.04,
            "role": "chart_internal",
        }
    ]
    img = Image.new("RGB", (1280, 720), (8, 18, 35))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)

    standard_result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(standard_output),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )
    aggressive_result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(aggressive_output),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="aggressive",
    )

    assert standard_result.text_box_count == 0
    assert aggressive_result.text_box_count == 1
    assert "关键节点" in pptx_text_shapes(aggressive_output)


def test_standard_mode_skips_tiny_screenshot_body_text_but_enhanced_keeps_it(tmp_path):
    slide_path = tmp_path / "tiny_screenshot_text.png"
    standard_output = tmp_path / "tiny_screenshot_text_standard.pptx"
    enhanced_output = tmp_path / "tiny_screenshot_text_enhanced.pptx"
    fake_regions = [
        {"text": "Botlife特点", "x": 0.05, "y": 0.12, "width": 0.28, "height": 0.09, "role": "title"},
        {
            "text": "Twitter - Allow your bot to manage a Twitter account.",
            "x": 0.63,
            "y": 0.44,
            "width": 0.28,
            "height": 0.015,
            "role": "body",
        },
    ]
    img = Image.new("RGB", (1280, 720), (8, 18, 35))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)

    standard_result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(standard_output),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )
    enhanced_result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(enhanced_output),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="enhanced",
    )

    assert standard_result.text_box_count == 1
    assert "Twitter - Allow" not in "\n".join(pptx_text_shapes(standard_output))
    assert enhanced_result.text_box_count == 2
    assert any(text.startswith("Twitter - Allow") for text in pptx_text_shapes(enhanced_output))


def test_standard_mode_restores_primary_display_text_on_complex_colorful_slides(tmp_path):
    slide_path = tmp_path / "complex_colorful_display.png"
    standard_output = tmp_path / "complex_colorful_display_standard.pptx"
    enhanced_output = tmp_path / "complex_colorful_display_enhanced.pptx"
    img = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(img)
    # Dense saturated strokes approximate poster-like sports/campaign pages
    # where display text is usually stylized with outlines and shadows.
    for x in range(-720, 1280, 24):
        color = ((x * 7) % 255, (80 + x * 3) % 255, (180 + x * 5) % 255)
        draw.line((x, 720, x + 720, 0), fill=color, width=10)
    fake_regions = [
        {"text": "世界杯情绪营销", "x": 0.14, "y": 0.18, "width": 0.58, "height": 0.13, "role": "title"},
        {"text": "宽正文仍然需要可编辑", "x": 0.18, "y": 0.56, "width": 0.45, "height": 0.045, "role": "body"},
    ]
    draw_region_marks(img, fake_regions, color="black")
    img.save(slide_path)

    standard_result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(standard_output),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )
    enhanced_result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(enhanced_output),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="enhanced",
    )

    assert standard_result.text_box_count == 2
    assert "世界杯情绪营销" in pptx_text_shapes(standard_output)
    assert "宽正文仍然需要可编辑" in pptx_text_shapes(standard_output)
    assert enhanced_result.text_box_count == 2
    assert "世界杯情绪营销" in pptx_text_shapes(enhanced_output)


def test_standard_mode_preserves_complex_campaign_labels_and_short_stats(tmp_path):
    slide_path = tmp_path / "complex_campaign_labels.png"
    output_path = tmp_path / "complex_campaign_labels_standard.pptx"
    img = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(img)
    for x in range(-720, 1280, 24):
        color = ((x * 7) % 255, (80 + x * 3) % 255, (180 + x * 5) % 255)
        draw.line((x, 720, x + 720, 0), fill=color, width=10)
    fake_regions = [
        {"text": "一页总览：为什么现在必须拿下这个IP?", "x": 0.058, "y": 0.088, "width": 0.814, "height": 0.098, "role": "title"},
        {"text": "01", "x": 0.133, "y": 0.301, "width": 0.046, "height": 0.057, "role": "label"},
        {"text": "世界杯机遇", "x": 0.095, "y": 0.384, "width": 0.123, "height": 0.046, "role": "subtitle"},
        {
            "text": "世界杯机遇：小红书已拿下2026美加墨世界杯直播+短视频全版权，目标冲击2亿日活、优化男性用户结构",
            "x": 0.065,
            "y": 0.477,
            "width": 0.183,
            "height": 0.145,
            "role": "body",
        },
        {
            "text": "全球影响力：Instagram 2.33亿粉丝，2026世界杯Puma为其推出专属战靴，品牌动作印证顶级商业价值",
            "x": 0.446,
            "y": 0.448,
            "width": 0.473,
            "height": 0.072,
            "role": "body",
        },
        {
            "text": "真人AI漫剧系列：融合真实比赛画面与内马尔AI形象，演绎最后之舞等核心主题",
            "x": 0.082,
            "y": 0.268,
            "width": 0.375,
            "height": 0.092,
            "role": "body",
        },
        {
            "text": "使用已完成的9款卡通形象，设计风格统一，支持实体盲盒与数字盲盒双版本",
            "x": 0.654,
            "y": 0.305,
            "width": 0.28,
            "height": 0.105,
            "role": "body",
        },
        {"text": "Instagram", "x": 0.43, "y": 0.846, "width": 0.08, "height": 0.024, "role": "label"},
        {"text": "2.33亿粉丝", "x": 0.43, "y": 0.876, "width": 0.10, "height": 0.046, "role": "body"},
    ]
    pill = fake_regions[2]
    draw.rounded_rectangle(
        (
            int((pill["x"] - 0.012) * 1280),
            int((pill["y"] - 0.006) * 720),
            int((pill["x"] + pill["width"] + 0.012) * 1280),
            int((pill["y"] + pill["height"] + 0.006) * 720),
        ),
        radius=18,
        fill=(230, 65, 120),
    )
    draw_region_marks(img, fake_regions, color="black")
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    text_shapes = pptx_text_shapes(output_path)
    slide_shapes = Presentation(str(output_path)).slides[0].shapes
    normalized_text = "\n".join(text_shapes).replace("\x0b", "")
    assert result.text_box_count >= 6
    assert sum(1 for shape in slide_shapes if shape.name == "Editable text pill background") >= 1
    assert "全球影响力：Instagram 2.33亿粉丝，2026世界杯Puma为其推出专属战靴，品牌动作印证顶级商业价值" in normalized_text
    assert "世界杯机遇：小红书已拿下2026美加墨世界杯直播+短视频全版权，目标冲击2亿日活、优化男性用户结构" in normalized_text
    assert "真人AI漫剧系列：融合真实比赛画面与内马尔AI形象，演绎最后之舞等核心主题" in normalized_text
    assert "使用已完成的9款卡通形象，设计风格统一，支持实体盲盒与数字盲盒双版本" in normalized_text
    assert "世界杯机遇" in text_shapes
    assert text_shapes["世界杯机遇"].text_frame.paragraphs[0].font.size.pt <= 18.0
    long_body_shape = next(shape for text, shape in text_shapes.items() if text.startswith("使用已完成的9款卡通形象"))
    assert long_body_shape.text_frame.paragraphs[0].font.size.pt <= 11.0
    assert "一页总览：为什么现在必须拿下这个IP?" in text_shapes


def test_complex_subtitle_on_plain_card_does_not_get_pill_background(tmp_path):
    slide_path = tmp_path / "complex_plain_subtitle.png"
    output_path = tmp_path / "complex_plain_subtitle_standard.pptx"
    img = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(img)
    for x in range(-720, 1280, 22):
        color = ((70 + x * 5) % 255, (40 + x * 9) % 255, (160 + x * 3) % 255)
        draw.line((x, 720, x + 720, 0), fill=color, width=11)
    draw.rounded_rectangle((610, 160, 1120, 480), radius=22, fill="white")
    fake_regions = [
        {"text": "真人AI漫剧授权——", "x": 0.51, "y": 0.27, "width": 0.20, "height": 0.046, "role": "subtitle"},
        {"text": "核心权益已经就绪", "x": 0.51, "y": 0.37, "width": 0.23, "height": 0.038, "role": "body"},
    ]
    draw_region_marks(img, fake_regions, color=(210, 40, 65))
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    text_shapes = pptx_text_shapes(output_path)
    slide_shapes = Presentation(str(output_path)).slides[0].shapes
    assert result.text_box_count == 2
    assert "真人AI漫剧授权——" in text_shapes
    assert sum(1 for shape in slide_shapes if shape.name == "Editable text pill background") == 0


def test_pill_label_original_text_is_cleaned_before_overlay(tmp_path):
    slide_path = tmp_path / "pill_cleanup.png"
    output_path = tmp_path / "pill_cleanup_editable.pptx"
    work_dir = tmp_path / "work"
    region = {"text": "世界杯机遇", "x": 0.18, "y": 0.32, "width": 0.18, "height": 0.052, "role": "subtitle"}
    img = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(img)
    for x in range(-720, 1280, 24):
        color = ((x * 7) % 255, (80 + x * 3) % 255, (180 + x * 5) % 255)
        draw.line((x, 720, x + 720, 0), fill=color, width=10)
    draw.rounded_rectangle((210, 218, 490, 292), radius=24, fill=(230, 45, 78))
    draw_region_marks(img, [region], color="white")
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: [region],
        work_dir=str(work_dir),
        restore_mode="standard",
    )

    patch = Image.open(work_dir / "slide_01_cleanup_patch_01.png").convert("RGB")
    mark_x = int(region["x"] * 1280) + 12
    mark_y = int(region["y"] * 720) + int(region["height"] * 720 * 0.42)
    patch_x = mark_x - int((region["x"] - 0.010) * 1280)
    patch_y = mark_y - int((region["y"] - 0.004) * 720)
    names = pptx_shape_names(output_path)
    assert result.text_box_count == 1
    assert patch.getpixel((patch_x, patch_y)) != (255, 255, 255)
    assert "Editable text pill background" in names
    assert not any(name.startswith("Text cleanup patch") for name in names)


def test_standard_build_uses_local_cleanup_patch_not_full_clean_background(tmp_path):
    slide_path = tmp_path / "clean_background_layers.png"
    output_path = tmp_path / "clean_background_layers.pptx"
    work_dir = tmp_path / "work"
    fake_regions = [
        {"text": "可编辑标题", "x": 0.16, "y": 0.22, "width": 0.26, "height": 0.07, "role": "title"},
    ]
    img = Image.new("RGB", (1280, 720), (245, 247, 250))
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        work_dir=str(work_dir),
    )

    names = pptx_shape_names(output_path)
    assert result.text_box_count == 1
    assert names[0] == "Original slide image"
    assert names[1] == "Editable cleanup patch - 1"
    assert "Editable cleaned background" not in names
    assert not any(name.startswith("Text cleanup patch") for name in names)
    assert names.index("Editable cleanup patch - 1") < next(i for i, name in enumerate(names) if name.startswith("Editable restored text"))


def test_same_row_same_role_text_uses_consistent_font_size(tmp_path):
    slide_path = tmp_path / "same_row_subtitles.png"
    output_path = tmp_path / "same_row_subtitles_editable.pptx"
    fake_regions = [
        {"text": "IP授权就绪", "x": 0.10, "y": 0.36, "width": 0.16, "height": 0.040, "role": "subtitle"},
        {"text": "世界杯机遇", "x": 0.42, "y": 0.358, "width": 0.16, "height": 0.050, "role": "subtitle"},
        {"text": "商业闭环", "x": 0.74, "y": 0.362, "width": 0.13, "height": 0.044, "role": "subtitle"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    shapes = pptx_text_shapes(output_path)
    sizes = [shapes[item["text"]].text_frame.paragraphs[0].font.size.pt for item in fake_regions]
    assert result.text_box_count == 3
    assert max(sizes) - min(sizes) <= 0.1


def test_saturated_label_uses_pill_background_even_on_simple_slide(tmp_path):
    slide_path = tmp_path / "saturated_label.png"
    output_path = tmp_path / "saturated_label_editable.pptx"
    fake_regions = [
        {"text": "世界杯机遇", "x": 0.105, "y": 0.385, "width": 0.12, "height": 0.045, "role": "subtitle"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((120, 260, 318, 318), radius=18, fill=(238, 35, 49))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    names = pptx_shape_names(output_path)
    assert result.text_box_count == 1
    assert "Editable text pill background" in names
    assert names.index("Editable text pill background") < next(i for i, name in enumerate(names) if name.startswith("Editable restored text"))


def test_inline_prefix_and_body_merge_to_avoid_overlapping_text_boxes(tmp_path):
    slide_path = tmp_path / "inline_prefix_body.png"
    output_path = tmp_path / "inline_prefix_body_editable.pptx"
    fake_regions = [
        {"text": "全球影响力：", "x": 0.44, "y": 0.44, "width": 0.13, "height": 0.030, "role": "subtitle"},
        {
            "text": "Instagram 2.33亿粉丝，品牌动作印证顶级商业价值",
            "x": 0.44,
            "y": 0.448,
            "width": 0.48,
            "height": 0.055,
            "role": "body",
        },
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    text_shapes = pptx_text_shapes(output_path)
    assert result.text_box_count == 1
    assert "全球影响力：Instagram 2.33亿粉丝，品牌动作印证顶级商业价值" in text_shapes


def test_wide_complex_card_body_is_not_forced_to_tiny_font(tmp_path):
    slide_path = tmp_path / "wide_complex_card_body.png"
    output_path = tmp_path / "wide_complex_card_body.pptx"
    fake_regions = [
        {
            "text": "全球影响力：Instagram 2.33亿粉丝，2026世界杯Puma为其推出专属战靴，品牌动作印证顶级商业价值",
            "x": 0.44,
            "y": 0.44,
            "width": 0.48,
            "height": 0.065,
            "role": "body",
        }
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(img)
    for x in range(-720, 1280, 18):
        color = ((80 + x * 5) % 255, (20 + x * 11) % 255, (180 + x * 7) % 255)
        draw.line((x, 720, x + 720, 0), fill=color, width=9)
    draw.rounded_rectangle((540, 300, 1210, 390), radius=18, fill="white", outline=(230, 60, 60), width=4)
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    shape = next(shape for shape in Presentation(str(output_path)).slides[0].shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
    assert result.text_box_count == 1
    assert shape.text_frame.paragraphs[0].font.size.pt >= 14.0


def test_stacked_wide_complex_body_rows_do_not_inherit_tiny_normalized_font(tmp_path):
    slide_path = tmp_path / "stacked_wide_complex_body.png"
    output_path = tmp_path / "stacked_wide_complex_body.pptx"
    fake_regions = [
        {
            "text": "全球影响力：Instagram 2.33亿粉丝，2026世界杯Puma为其推出专属战靴，品牌动作印证顶级商业价值",
            "x": 0.44,
            "y": 0.44,
            "width": 0.48,
            "height": 0.065,
            "role": "body",
        },
        {
            "text": "中国情感共鸣：巴西足球兴趣用户超1亿，内马尔是情感共鸣最强的球星，无代际门槛",
            "x": 0.44,
            "y": 0.55,
            "width": 0.48,
            "height": 0.065,
            "role": "body",
        },
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(img)
    for x in range(-720, 1280, 18):
        color = ((80 + x * 5) % 255, (20 + x * 11) % 255, (180 + x * 7) % 255)
        draw.line((x, 720, x + 720, 0), fill=color, width=9)
    draw.rounded_rectangle((540, 300, 1210, 470), radius=18, fill="white", outline=(230, 60, 60), width=4)
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    shapes = [shape for shape in Presentation(str(output_path)).slides[0].shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
    sizes = [shape.text_frame.paragraphs[0].font.size.pt for shape in shapes]
    assert result.text_box_count == 2
    assert min(sizes) >= 14.0


def test_narrow_dense_body_copy_uses_smaller_fitted_text(tmp_path):
    slide_path = tmp_path / "narrow_dense_body.png"
    output_path = tmp_path / "narrow_dense_body.pptx"
    fake_regions = [
        {
            "text": "世界杯机遇：小红书已拿下2026美加墨世界杯直播+短视频全版权，目标冲击2亿日活、优化男性用户结构",
            "x": 0.065,
            "y": 0.477,
            "width": 0.183,
            "height": 0.145,
            "role": "body",
        }
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    shape = next(shape for shape in Presentation(str(output_path)).slides[0].shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
    assert result.text_box_count == 1
    assert shape.text_frame.paragraphs[0].font.size.pt <= 11.0


def test_manual_wrap_keeps_opening_quote_with_following_token():
    text = "全球影响力：Instagram 2.33亿粉丝，2026世界杯Puma为其推出「Showtime」专属战靴，品牌动作印证顶级商业价值"
    wrapped = editable_export.wrap_text_for_box(
        text,
        {"x": 0.44, "y": 0.44, "width": 0.54, "height": 0.0975},
        14.0,
        "body",
    )

    assert "「\nShowtime" not in wrapped


def test_standard_mode_keeps_large_display_text_editable_on_simple_slides(tmp_path):
    slide_path = tmp_path / "simple_display.png"
    output_path = tmp_path / "simple_display_editable.pptx"
    fake_regions = [
        {"text": "清晰大标题", "x": 0.12, "y": 0.18, "width": 0.40, "height": 0.13, "role": "title"},
    ]
    img = Image.new("RGB", (1280, 720), "white")
    draw_region_marks(img, fake_regions)
    img.save(slide_path)

    result = build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    assert result.text_box_count == 1
    assert "清晰大标题" in pptx_text_shapes(output_path)


def test_standard_mode_restores_primary_title_inside_visual_asset(tmp_path, monkeypatch):
    slide_path = tmp_path / "title_over_photo.png"
    output_path = tmp_path / "title_over_photo_editable.pptx"
    fake_regions = [
        {
            "text": "解锁武汉卡丁车馆亲子高能体验阵地",
            "x": 0.06,
            "y": 0.10,
            "width": 0.42,
            "height": 0.11,
            "role": "title",
        }
    ]
    img = Image.new("RGB", (1280, 720), (8, 18, 35))
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 760, 720), fill=(0, 35, 90))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)
    monkeypatch.setattr(
        editable_export,
        "detect_image_blocks",
        lambda *_args, **_kwargs: [{"x": 0.0, "y": 0.0, "width": 0.60, "height": 1.0}],
    )

    result = editable_export.build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        restore_mode="standard",
    )

    assert result.text_box_count == 1
    assert "解锁武汉卡丁车馆亲子高能体验阵地" in pptx_text_shapes(output_path)


def test_visual_asset_crop_uses_text_cleaned_source(tmp_path, monkeypatch):
    slide_path = tmp_path / "visual_asset_text_overlap.png"
    output_path = tmp_path / "visual_asset_text_overlap_editable.pptx"
    work_dir = tmp_path / "work"
    block = {"x": 0.12, "y": 0.22, "width": 0.58, "height": 0.42}
    region = {
        "text": "视觉块里的标题",
        "x": 0.22,
        "y": 0.36,
        "width": 0.24,
        "height": 0.085,
        "role": "title",
    }
    img = Image.new("RGB", (1280, 720), (12, 18, 30))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((154, 158, 896, 468), radius=24, fill=(0, 48, 110))
    draw_region_marks(img, [region], color="white")
    img.save(slide_path)
    monkeypatch.setattr(editable_export, "detect_image_blocks", lambda *_args, **_kwargs: [block])

    result = editable_export.build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: [region],
        work_dir=str(work_dir),
        restore_mode="standard",
    )

    asset = Image.open(work_dir / "slide_01_asset_01.png").convert("RGB")
    mark_x = int(region["x"] * 1280) + 12 - int(block["x"] * 1280)
    mark_y = int(region["y"] * 720) + int(region["height"] * 720 * 0.42) - int(block["y"] * 720)
    assert result.text_box_count == 1
    assert result.visual_asset_count == 1
    assert asset.getpixel((mark_x, mark_y)) != (255, 255, 255)
    assert "视觉块里的标题" in pptx_text_shapes(output_path)


def test_display_cleanup_patch_expands_to_cover_shadow_outside_ocr_box(tmp_path, monkeypatch):
    slide_path = tmp_path / "display_shadow_overlap.png"
    output_path = tmp_path / "display_shadow_overlap_editable.pptx"
    work_dir = tmp_path / "work"
    region = {
        "text": "金色展示标题",
        "x": 0.42,
        "y": 0.22,
        "width": 0.28,
        "height": 0.080,
        "role": "title",
    }
    shadow_y = region["y"] - 0.026
    shadow_x = region["x"] + 0.04
    gold = (230, 184, 80)
    img = Image.new("RGB", (1280, 720), (10, 14, 20))
    draw = ImageDraw.Draw(img)
    draw.rectangle(
        (
            int(shadow_x * 1280),
            int(shadow_y * 720),
            int((shadow_x + 0.12) * 1280),
            int((shadow_y + 0.014) * 720),
        ),
        fill=gold,
    )
    draw_region_marks(img, [region], color=gold)
    img.save(slide_path)
    monkeypatch.setattr(editable_export, "detect_image_blocks", lambda *_args, **_kwargs: [])

    result = editable_export.build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: [region],
        work_dir=str(work_dir),
        restore_mode="standard",
    )

    prs = Presentation(str(output_path))
    patch_shape = next(shape for shape in prs.slides[0].shapes if shape.name == "Editable cleanup patch - 1")
    patch_top = int(patch_shape.top) / int(prs.slide_height)
    patch_bottom = (int(patch_shape.top) + int(patch_shape.height)) / int(prs.slide_height)
    patch = Image.open(work_dir / "slide_01_cleanup_patch_01.png").convert("RGB")
    assert result.text_box_count == 1
    assert patch_top <= shadow_y <= patch_bottom
    assert gold not in set(patch.getdata())


def test_quality_gate_retries_when_visual_asset_reintroduces_original_text(tmp_path, monkeypatch):
    slide_path = tmp_path / "qa_asset_reintroduces_text.png"
    output_path = tmp_path / "qa_asset_reintroduces_text_editable.pptx"
    work_dir = tmp_path / "work"
    block = {"x": 0.10, "y": 0.18, "width": 0.72, "height": 0.50}
    region = {
        "text": "被裁回的原字",
        "x": 0.24,
        "y": 0.34,
        "width": 0.24,
        "height": 0.090,
        "role": "title",
    }
    img = Image.new("RGB", (1280, 720), (12, 18, 30))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((128, 130, 1050, 510), radius=24, fill=(0, 48, 110))
    draw_region_marks(img, [region], color="white")
    img.save(slide_path)
    monkeypatch.setattr(editable_export, "detect_image_blocks", lambda *_args, **_kwargs: [block])
    original_crop_asset_from_image = editable_export.crop_asset_from_image

    def buggy_crop_asset_from_image(_img, crop_region, output):
        return original_crop_asset_from_image(Image.open(slide_path).convert("RGB"), crop_region, output)

    monkeypatch.setattr(editable_export, "crop_asset_from_image", buggy_crop_asset_from_image)

    result = editable_export.build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: [region],
        work_dir=str(work_dir),
        restore_mode="standard",
    )

    names = pptx_shape_names(output_path)
    assert result.qa_retry_pages == [1]
    assert result.quality_fallback_pages == []
    assert "Editable QA cleanup patch - 1" in names
    assert "被裁回的原字" in pptx_text_shapes(output_path)


def test_quality_gate_preserves_editable_text_when_residual_persists(tmp_path, monkeypatch):
    slide_path = tmp_path / "qa_fallback.png"
    output_path = tmp_path / "qa_fallback_editable.pptx"
    region = {
        "text": "清不掉的原字",
        "x": 0.30,
        "y": 0.30,
        "width": 0.24,
        "height": 0.080,
        "role": "title",
    }
    img = Image.new("RGB", (1280, 720), (12, 18, 30))
    draw_region_marks(img, [region], color="white")
    img.save(slide_path)

    def always_residual(_original_rgb, _underlay_rgb, groups, _cleanup_boxes, **_kwargs):
        return [(groups[0], 1.0)] if groups else []

    monkeypatch.setattr(editable_export, "residual_text_groups", always_residual)

    result = editable_export.build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: [region],
        restore_mode="standard",
    )

    names = pptx_shape_names(output_path)
    assert result.qa_retry_pages == [1]
    assert result.quality_fallback_pages == []
    assert result.quality_warning_pages == [1]
    assert result.ocr_failed_pages == []
    assert "Editable QA fallback slide image" not in names
    assert "清不掉的原字" in pptx_text_shapes(output_path)


def test_quality_retry_cleanup_uses_bounded_previous_cleanup_box():
    group = {
        "text": "金色展示标题",
        "role": "title",
        "bbox": {"x": 0.42, "y": 0.22, "width": 0.28, "height": 0.08},
    }
    previous = {
        "x": 0.395,
        "y": 0.195,
        "width": 0.33,
        "height": 0.13,
        "full_fill": True,
    }

    retry = editable_export.quality_retry_cleanup_box_for_group(group, previous)

    assert retry["full_fill"] is True
    assert retry["width"] <= previous["width"] + 0.016
    assert retry["height"] <= previous["height"] + 0.020


def test_quality_gate_residual_checks_run_with_bounded_parallelism(monkeypatch):
    captured = {}

    class FakePool:
        def __init__(self, max_workers=None, thread_name_prefix=None):
            captured["max_workers"] = max_workers
            captured["thread_name_prefix"] = thread_name_prefix

        def __enter__(self):
            return self

        def __exit__(self, *_args):
            return False

        def map(self, fn, items):
            return [fn(item) for item in items]

    monkeypatch.setattr(editable_export.settings, "EDITABLE_PPTX_QA_MAX_WORKERS", 3)
    monkeypatch.setattr(editable_export.os, "cpu_count", lambda: 8)
    monkeypatch.setattr(editable_export, "ThreadPoolExecutor", FakePool)
    monkeypatch.setattr(editable_export, "residual_original_text_score", lambda *_args: 1.0)

    groups = [{"text": str(index)} for index in range(5)]
    boxes = [{"x": 0, "y": 0, "width": 0.1, "height": 0.1} for _ in groups]
    rgb = np.zeros((12, 12, 3), dtype=np.uint8)

    residuals = editable_export.residual_text_groups(rgb, rgb, groups, boxes)

    assert len(residuals) == 5
    assert captured["max_workers"] == 3
    assert captured["thread_name_prefix"] == "editable-pptx-qa"


def test_local_cleanup_patch_uses_ocr_box_not_expanded_render_box(tmp_path):
    slide_path = tmp_path / "cleanup_scope.png"
    output_path = tmp_path / "cleanup_scope_editable.pptx"
    work_dir = tmp_path / "work"
    fake_regions = [
        {
            "text": "长段正文第一行\n长段正文第二行\n长段正文第三行",
            "x": 0.10,
            "y": 0.25,
            "width": 0.25,
            "height": 0.14,
            "role": "title",
        },
    ]
    img = Image.new("RGB", (1280, 720), (12, 18, 30))
    draw = ImageDraw.Draw(img)
    draw.rectangle((850, 210, 1120, 310), fill=(0, 92, 255))
    draw_region_marks(img, fake_regions, color="white")
    img.save(slide_path)

    build_editable_pptx(
        slide_images=[{"page_num": 1, "image_path": str(slide_path)}],
        output_path=str(output_path),
        ocr_provider=lambda *_args: fake_regions,
        work_dir=str(work_dir),
    )

    patch = Image.open(work_dir / "slide_01_cleanup_patch_01.png").convert("RGB")
    assert patch.width < 500
    assert patch.height < 180
    assert (0, 92, 255) not in list(patch.getdata())
