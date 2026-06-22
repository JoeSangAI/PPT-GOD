import asyncio
import io
import json
import re
import zipfile
from xml.etree import ElementTree as ET

import pytest
from fastapi import BackgroundTasks, HTTPException
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from app.api import projects as projects_api
from app.api import documents as documents_api
from app.api import slides as slides_api
from app.services import generation_pipeline
from app.schemas.project import ProjectUpdate
from app.api.slides import _resolve_generation_page_nums
from app.models.base import Base
from app.models.models import Project, ProjectRun, ReferenceImage, Slide
from app.services.generation_pipeline import (
    _generate_one_slide,
    _load_reference_images,
)
from app.services.image_generation import _cache_key
from app.services.artifact_versions import artifact_stale
from app.services.logo_assets import prepare_logo_lockup_image, prepare_logo_overlay_image, prepare_logo_symbol_image
from app.services.logo_overlay_layout import resolve_logo_overlay_box, resolve_logo_render_policy
from app.services.logo_policy import logo_policy_for_page
from app.services.project_quality_report import build_project_quality_report
from app.services.pptx_assembler import assemble_pptx
from app.services import prompt_engine
from app.utils.text_cleaning import normalize_markdown_emphasis
from app.services.visual_plan import (
    VisualPlanGenerationError,
    _build_batch_prompt,
    _do_generate_visual_plan,
    _default_visual_asset_usage,
    _fallback_visual_plan,
    _recall_visual_assets_for_page,
    _safe_parse_json,
)
from app.services.content_plan import _annotate_ppt_source_refs, _normalize_content_markdown
from app.utils.project_docs import load_project_documents
from types import SimpleNamespace

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


STALE_PRODUCT_DOMAIN_TERMS_FOR_TEST = (
    "胡姬花", "花生油", "花生", "小油", "油瓶", "榨油", "古法小榨",
    "团购礼盒", "非遗体验官", "工业流水线",
)


def make_session():
    engine = create_engine("sqlite:///:memory:", connect_args={"check_same_thread": False})
    Base.metadata.create_all(bind=engine)
    Session = sessionmaker(bind=engine)
    return Session()


def make_shared_sessionmaker():
    engine = create_engine(
        "sqlite://",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    Base.metadata.create_all(bind=engine)
    return sessionmaker(bind=engine)


def png_upload(name="asset.png"):
    buf = io.BytesIO()
    Image.new("RGB", (10, 10), "white").save(buf, "PNG")
    buf.seek(0)
    return SimpleNamespace(filename=name, file=buf, content_type="image/png")


def txt_upload(name="brief.txt", text="开学季绿色营销推广方案"):
    return SimpleNamespace(filename=name, file=io.BytesIO(text.encode("utf-8")), content_type="text/plain")


def pdf_upload_with_picture(name="source.pdf"):
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "第1章 创新使命\n图示展示创新飞轮", fontsize=12, fontname="china-s")
    image_buf = io.BytesIO()
    Image.new("RGB", (120, 80), "white").save(image_buf, "PNG")
    page.insert_image(fitz.Rect(72, 140, 220, 240), stream=image_buf.getvalue())
    out = io.BytesIO(doc.tobytes())
    return SimpleNamespace(filename=name, file=out, content_type="application/pdf")


def pdf_upload_with_repeated_logo(name="brand-deck.pdf", count=5):
    import fitz

    logo = Image.new("RGBA", (180, 60), (255, 255, 255, 0))
    draw = ImageDraw.Draw(logo)
    draw.rounded_rectangle((4, 6, 52, 54), radius=8, fill=(36, 75, 160, 255))
    draw.rectangle((68, 16, 168, 28), fill=(36, 75, 160, 255))
    draw.rectangle((68, 36, 138, 46), fill=(36, 75, 160, 255))
    logo_buf = io.BytesIO()
    logo.save(logo_buf, "PNG")
    logo_bytes = logo_buf.getvalue()

    doc = fitz.open()
    for idx in range(count):
        page = doc.new_page()
        page.insert_image(fitz.Rect(36, 28, 180, 76), stream=logo_bytes)
        page.insert_text((72, 160), f"第{idx + 1}页 业务内容", fontsize=20, fontname="china-s")
    out = io.BytesIO(doc.tobytes())
    return SimpleNamespace(filename=name, file=out, content_type="application/pdf")


def pptx_upload_with_picture(name="source.pptx"):
    image_buf = io.BytesIO()
    img = Image.new("RGB", (240, 140), "white")
    for x in range(20, 220):
        for y in range(25, 115):
            img.putpixel((x, y), (20, 120, 210))
    img.save(image_buf, "PNG")
    image_buf.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(5.0), Inches(0.6)).text = "原 PPT 第 1 页标题"
    slide.shapes.add_picture(image_buf, Inches(1.0), Inches(1.4), width=Inches(3.0))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def pptx_upload_with_primary_product_image(name="product-source.pptx"):
    image_buf = io.BytesIO()
    img = Image.new("RGB", (260, 160), "white")
    for x in range(20, 240):
        for y in range(22, 138):
            img.putpixel((x, y), (30 + (x % 60), 105 + (y % 80), 180))
    img.save(image_buf, "PNG")
    image_buf.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(6.0), Inches(0.6)).text = "产品包装主视觉与终端陈列"
    slide.shapes.add_picture(image_buf, Inches(1.0), Inches(1.4), width=Inches(3.2))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def pptx_upload_with_content_graphic(name="graphic-source.pptx"):
    image_buf = io.BytesIO()
    img = Image.new("RGB", (300, 180), "white")
    for x in range(24, 276):
        for y in range(30, 150):
            if x % 70 < 36 or y % 54 < 24:
                img.putpixel((x, y), (45, 120, 210))
            else:
                img.putpixel((x, y), (230, 235, 245))
    img.save(image_buf, "PNG")
    image_buf.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(7.0), Inches(0.6)).text = "系统截图、数据图表与业务流程图"
    slide.shapes.add_picture(image_buf, Inches(1.0), Inches(1.4), width=Inches(3.4))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def pptx_upload_with_identity_code(name="identity-code-source.pptx"):
    image_buf = io.BytesIO()
    img = Image.new("RGB", (180, 180), "white")
    for x in range(24, 156):
        for y in range(24, 156):
            if (x // 12 + y // 12) % 2 == 0 or (x % 29 < 8 and y % 31 < 8):
                img.putpixel((x, y), (20, 20, 20))
    img.save(image_buf, "PNG")
    image_buf.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(7.0), Inches(0.6)).text = "菜鸟驿站身份码与扫码出库流程"
    slide.shapes.add_picture(image_buf, Inches(1.0), Inches(2.5), width=Inches(1.6))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def pptx_upload_with_phone_frame_chrome(name="phone-frame-source.pptx"):
    image_buf = io.BytesIO()
    img = Image.new("RGBA", (220, 460), (255, 255, 255, 0))
    for x in range(16, 204):
        for y in range(12, 448):
            on_outer = x < 26 or x > 194 or y < 24 or y > 438
            on_inner = 40 < x < 180 and 80 < y < 390 and (x + y) % 17 < 2
            if on_outer:
                img.putpixel((x, y), (30, 30, 30, 255))
            elif on_inner:
                img.putpixel((x, y), (210, 215, 225, 255))
    img.save(image_buf, "PNG")
    image_buf.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(7.0), Inches(0.6)).text = "手机边框承载小程序界面与身份码展示"
    slide.shapes.add_picture(image_buf, Inches(1.0), Inches(1.0), height=Inches(4.6))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def _patterned_png(color):
    buf = io.BytesIO()
    img = Image.new("RGB", (180, 120), "white")
    for x in range(180):
        for y in range(120):
            if (x // 12 + y // 12) % 2 == 0:
                img.putpixel((x, y), color)
            elif 30 < x < 150 and 25 < y < 95:
                img.putpixel((x, y), (30, 30, 30))
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def pptx_upload_with_many_pictures(name="many.pptx", count=5):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(6.0), Inches(0.6)).text = "渠道策略 产品矩阵 终端陈列"
    colors = [(210, 40, 40), (40, 140, 210), (40, 180, 90), (210, 150, 40), (150, 80, 210)]
    for idx in range(count):
        x = Inches(0.5 + (idx % 3) * 2.0)
        y = Inches(1.1 + (idx // 3) * 1.6)
        slide.shapes.add_picture(_patterned_png(colors[idx % len(colors)]), x, y, width=Inches(1.5))
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def pptx_upload_with_repeated_shape_fill_logos(tmp_path, name="shape-fill-logo.pptx", count=5):
    img = Image.new("RGB", (133, 57), "black")
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((5, 5, 42, 42), radius=8, fill=(255, 205, 0))
    draw.rectangle((55, 15, 124, 27), fill=(245, 245, 245))
    draw.rectangle((55, 34, 108, 44), fill=(245, 245, 245))
    logo_bytes = io.BytesIO()
    img.save(logo_bytes, "PNG")
    logo_blob = logo_bytes.getvalue()

    prs = Presentation()
    for idx in range(count):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(6.5), Inches(0.6)).text = f"业务页面 {idx + 1}"
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8.35),
            Inches(6.15),
            Inches(0.75),
            Inches(0.32),
        )
        shape.fill.solid()
        shape.line.fill.background()

    out = io.BytesIO()
    prs.save(out)

    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
        "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
    }
    for prefix, uri in ns.items():
        if prefix not in {"rel", "ct"}:
            ET.register_namespace(prefix, uri)

    patched = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(out.getvalue()), "r") as zin, zipfile.ZipFile(patched, "w") as zout:
        names = set(zin.namelist())
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                root = ET.fromstring(data)
                has_png = any(child.attrib.get("Extension") == "png" for child in root)
                if not has_png:
                    ET.SubElement(root, f"{{{ns['ct']}}}Default", {"Extension": "png", "ContentType": "image/png"})
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            elif item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                slide_num = int(item.filename.rsplit("slide", 1)[1].split(".xml", 1)[0])
                root = ET.fromstring(data)
                shapes = root.findall(".//p:sp", ns)
                target_shape = shapes[-1]
                sp_pr = target_shape.find("p:spPr", ns)
                for child in list(sp_pr):
                    if child.tag in {f"{{{ns['a']}}}solidFill", f"{{{ns['a']}}}gradFill", f"{{{ns['a']}}}pattFill", f"{{{ns['a']}}}noFill"}:
                        sp_pr.remove(child)
                blip_fill = ET.Element(f"{{{ns['a']}}}blipFill")
                ET.SubElement(blip_fill, f"{{{ns['a']}}}blip", {f"{{{ns['r']}}}embed": f"rIdLogo{slide_num}"})
                stretch = ET.SubElement(blip_fill, f"{{{ns['a']}}}stretch")
                ET.SubElement(stretch, f"{{{ns['a']}}}fillRect")
                sp_pr.append(blip_fill)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            elif item.filename.startswith("ppt/slides/_rels/slide") and item.filename.endswith(".xml.rels"):
                slide_num = int(item.filename.rsplit("slide", 1)[1].split(".xml.rels", 1)[0])
                root = ET.fromstring(data)
                ET.SubElement(root, f"{{{ns['rel']}}}Relationship", {
                    "Id": f"rIdLogo{slide_num}",
                    "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
                    "Target": "../media/shapeFillLogo.png",
                })
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            zout.writestr(item, data)
        if "ppt/media/shapeFillLogo.png" not in names:
            zout.writestr("ppt/media/shapeFillLogo.png", logo_blob)

    patched.seek(0)
    return SimpleNamespace(filename=name, file=patched, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def pptx_upload_with_two_cover_logos(name="two-logos.pptx"):
    def logo_blob(color, accent):
        buf = io.BytesIO()
        img = Image.new("RGB", (220, 72), "white")
        for x in range(12, 208):
            for y in range(18, 54):
                if x < 62 or (x // 14 + y // 12) % 2 == 0:
                    img.putpixel((x, y), color)
                else:
                    img.putpixel((x, y), accent)
        img.save(buf, "PNG")
        buf.seek(0)
        return buf

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(9.0), Inches(0.8)).text = "开学季绿色营销推广方案"
    slide.shapes.add_picture(logo_blob((35, 80, 170), (20, 20, 20)), Inches(0.65), Inches(0.35), width=Inches(1.45))
    slide.shapes.add_picture(logo_blob((30, 160, 90), (245, 165, 35)), Inches(2.35), Inches(0.35), width=Inches(1.45))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def pptx_upload_with_cropped_cover_logo(name="cropped-logo.pptx"):
    logo_buf = io.BytesIO()
    img = Image.new("RGB", (440, 72), "white")
    for x in range(16, 204):
        for y in range(18, 54):
            img.putpixel((x, y), (180, 122, 0))
    for x in range(252, 408):
        for y in range(16, 56):
            img.putpixel((x, y), (210, 20, 28))
    img.save(logo_buf, "PNG")
    logo_buf.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(9.0), Inches(0.8)).text = "校园联播网媒体介绍"
    logo = slide.shapes.add_picture(logo_buf, Inches(0.65), Inches(6.8), width=Inches(2.8))
    logo.crop_right = 0.5

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def pptx_upload_with_cropped_multicolor_cover_logo(name="cropped-multicolor-logo.pptx"):
    logo_buf = io.BytesIO()
    img = Image.new("RGB", (860, 260), "white")
    for x in range(35, 455):
        for y in range(52, 212):
            if x < 90 or y < 74 or y > 190 or (x // 11 + y // 9) % 3 == 0:
                img.putpixel((x, y), (184, 126, 8))
            elif (x + y) % 17 < 4:
                img.putpixel((x, y), (80, 40, 12))
    for x in range(500, 810):
        for y in range(50, 210):
            if (x - 500) ** 2 + (y - 130) ** 2 < 75 ** 2:
                img.putpixel((x, y), (30, 95, 210))
            elif x % 19 < 7:
                img.putpixel((x, y), (220, 30, 45))
    logo_buf.seek(0)
    img.save(logo_buf, "PNG")
    logo_buf.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(9.0), Inches(0.8)).text = "Focus Media & Cainiao Station 2026 SPRING"
    logo = slide.shapes.add_picture(logo_buf, Inches(0.65), Inches(6.75), width=Inches(1.28))
    logo.crop_left = 0.08
    logo.crop_top = 0.24
    logo.crop_right = 0.45
    logo.crop_bottom = 0.22

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return SimpleNamespace(filename=name, file=out, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def test_logo_overlay_preprocess_trims_white_background(tmp_path):
    logo_path = tmp_path / "logo.png"
    img = Image.new("RGB", (120, 80), "white")
    for x in range(42, 78):
        for y in range(18, 62):
            img.putpixel((x, y), (150, 0, 0))
    img.save(logo_path)

    overlay_path = prepare_logo_overlay_image(str(logo_path))
    overlay = Image.open(overlay_path).convert("RGBA")

    assert overlay.size[0] < 50
    assert overlay.size[1] < 60
    assert overlay.getchannel("A").getextrema()[0] == 0


def test_logo_overlay_keeps_clean_mark_without_auto_halo(tmp_path):
    logo_path = tmp_path / "mixed-logo.png"
    img = Image.new("RGB", (180, 72), "white")
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((12, 14, 52, 54), radius=8, fill=(255, 205, 0))
    draw.rectangle((70, 24, 160, 39), fill=(18, 18, 18))
    img.save(logo_path)

    overlay_path = prepare_logo_overlay_image(str(logo_path))
    overlay = Image.open(overlay_path).convert("RGBA")
    alpha = overlay.getchannel("A")

    assert overlay.width < 170
    assert overlay.height < 60
    assert alpha.getpixel((0, 0)) == 0
    assert alpha.getpixel((overlay.width - 1, overlay.height - 1)) == 0


def test_logo_symbol_preprocess_extracts_left_brand_mark(tmp_path):
    logo_path = tmp_path / "mixed-logo.png"
    img = Image.new("RGB", (260, 80), "white")
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((14, 16, 62, 64), radius=9, fill=(255, 205, 0))
    draw.rectangle((88, 28, 232, 42), fill=(18, 18, 18))
    draw.rectangle((88, 48, 198, 57), fill=(18, 18, 18))
    img.save(logo_path)

    full_path = prepare_logo_overlay_image(str(logo_path))
    symbol_path = prepare_logo_symbol_image(str(logo_path))

    assert symbol_path
    full = Image.open(full_path)
    symbol = Image.open(symbol_path)
    assert symbol.width < full.width * 0.55
    assert symbol.height >= full.height * 0.55


def test_logo_symbol_preprocess_returns_none_for_wordmark_only(tmp_path):
    logo_path = tmp_path / "wordmark.png"
    img = Image.new("RGB", (260, 80), "white")
    draw = ImageDraw.Draw(img)
    draw.rectangle((20, 26, 238, 42), fill=(18, 18, 18))
    draw.rectangle((20, 49, 190, 58), fill=(18, 18, 18))
    img.save(logo_path)

    assert prepare_logo_symbol_image(str(logo_path)) is None


def test_logo_lockup_combines_multiple_marks(tmp_path):
    paths = []
    for idx, color in enumerate([(150, 0, 0), (0, 120, 80)]):
        path = tmp_path / f"logo_{idx}.png"
        img = Image.new("RGB", (120, 72), "white")
        for x in range(18, 102):
            for y in range(20, 52):
                img.putpixel((x, y), color)
        img.save(path)
        paths.append(str(path))

    lockup_path = prepare_logo_lockup_image(paths)
    lockup = Image.open(lockup_path).convert("RGBA")

    assert lockup.width > lockup.height
    assert lockup.getchannel("A").getextrema()[1] == 255


def test_normalize_markdown_emphasis_balances_line_start_marker():
    text = "**第四部分：媒介与资本叙事\n普通正文"

    assert normalize_markdown_emphasis(text) == "**第四部分：媒介与资本叙事**\n普通正文"


def test_prompt_text_contract_strips_unbalanced_markdown():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 2,
            "type": "toc",
            "layout": "toc",
            "visual_evidence": "目录结构",
            "visual_description": "四个章节纵向排列。",
        },
        content_text={
            "headline": "战略全局",
            "body": "**第四部分：媒介与资本叙事",
        },
        style_text_override="Style: 简洁\nPalette: #FFFFFF, #111111",
    )

    assert 'Body: "第四部分：媒介与资本叙事"' in prompt
    assert 'Body: "**第四部分' not in prompt


def test_section_prompt_uses_section_title_as_structure_not_visible_copy():
    prompts = prompt_engine.generate_prompts_for_all_pages(
        visual_plan=[{
            "page_num": 2,
            "type": "section",
            "layout": "section",
            "visual_evidence": "章节过渡",
            "visual_description": "章节标题居中。",
        }],
        content_plan=[{
            "page_num": 2,
            "type": "section",
            "section_title": "第一章",
            "text_content": {
                "headline": "我们是谁？",
                "subhead": "Who We Are",
                "body": "",
            },
        }],
        style_text_override="Style: 浅色\nPalette: #FFFFFF, #6B5B7A",
    )

    prompt = prompts[0]["prompt"]
    assert "Chapter label" not in prompt
    assert "第一章" not in prompt
    assert 'Headline: "我们是谁？"' in prompt
    assert "section_title is structural metadata" in prompt


def test_section_prompt_does_not_duplicate_chapter_number_channels():
    prompts = prompt_engine.generate_prompts_for_all_pages(
        visual_plan=[{
            "page_num": 18,
            "type": "section",
            "layout": "section",
            "visual_evidence": "章节标题/编号/转场氛围",
            "visual_description": (
                "页面中央以白色无衬线粗体呈现章节编号「06」与主标题「创意概念」，"
                "下方以较小字号标注副标题「Part 6 — 《拿不准时刻》系列」。"
            ),
        }],
        content_plan=[{
            "page_num": 18,
            "type": "section",
            "section_title": "模块六",
            "text_content": {
                "headline": "创意概念",
                "subhead": "Part 6 — 《拿不准时刻》系列",
                "body": "",
            },
        }],
        style_text_override="Style: 红白活力\nPalette: #FF2442, #FFFFFF",
    )

    prompt = prompts[0]["prompt"]
    assert 'Headline: "创意概念"' in prompt
    assert 'Subhead: "Part 6 — 《拿不准时刻》系列"' in prompt
    assert 'Chapter label: "模块六"' not in prompt
    assert "模块六" not in prompt
    assert "章节编号" not in prompt
    assert "编号" not in prompt
    assert "「06」" not in prompt
    assert prompt.count("Part 6") == 1


def test_section_prompt_strips_standalone_chapter_transition_layout_language():
    prompts = prompt_engine.generate_prompts_for_all_pages(
        visual_plan=[{
            "page_num": 13,
            "type": "section",
            "layout": "section",
            "visual_evidence": "章节转场：数字章节转场与章节名主导的仪式感画面",
            "visual_description": (
                "画面中心略偏上放置一个超大细体数字章节转场（01 / 03 类），"
                "章节转场用香槟金，字距极宽；其下紧跟一行章节标题。"
            ),
        }],
        content_plan=[{
            "page_num": 13,
            "type": "section",
            "section_title": "第二章",
            "text_content": {
                "headline": "什么没变？人心仍然是终点",
                "body": "",
            },
        }],
        style_text_override="Style: 黑金极简\nPalette: #0A0A0A, #C9A961",
    )

    prompt = prompts[0]["prompt"]
    assert 'Headline: "什么没变？人心仍然是终点"' in prompt
    assert "第二章" not in prompt
    assert "数字章节" not in prompt
    assert "章节转场" not in prompt
    assert "01 / 03" not in prompt


def test_fallback_visual_plan_uses_concrete_visual_evidence():
    plan = _fallback_visual_plan(
        [
            {
                "page_num": 12,
                "type": "content",
                "text_content": {
                    "headline": "线上渠道：内容电商与社交语境统一",
                    "body": "直播间背景板与达人合作 Brief 统一古法香话术",
                },
            },
            {
                "page_num": 13,
                "type": "content",
                "text_content": {
                    "headline": "公关主线：夺取行业标准制定权",
                    "body": "发布《古法香标准白皮书》，联合权威机构举办发布会",
                },
            },
        ],
        [],
    )

    # Fallback no longer uses keyword-enumerated concrete evidence; it returns
    # open-ended descriptions driven by page_type and headline.
    assert "核心场景" in plan[0]["visual_evidence"] or "支撑" in plan[0]["visual_evidence"]
    assert "现代商务风格画面" not in plan[0]["visual_description"]
    assert "白皮书" in plan[1]["visual_evidence"] or "核心场景" in plan[1]["visual_evidence"]
    assert plan[0]["seed_family"] == "content"
    assert plan[0]["is_seed_recommended"] is True


def test_fallback_visual_plan_uses_project_neutral_business_objects():
    plan = _fallback_visual_plan(
        [
            {
                "page_num": 8,
                "type": "content",
                "text_content": {
                    "headline": "B端客户路径",
                    "body": "企业客户名单、团购合作和交付流程需要形成闭环。",
                },
            },
            {
                "page_num": 9,
                "type": "content",
                "text_content": {
                    "headline": "终端转化路径",
                    "body": "货架、导购和现场动线决定用户能否完成转化。",
                },
            },
            {
                "page_num": 10,
                "type": "content",
                "text_content": {
                    "headline": "竞品对比",
                    "body": "A方案 vs B方案，比较成本、效率和执行风险。",
                },
            },
        ],
        [],
    )

    joined = "\n".join(
        f"{item['visual_evidence']}\n{item['visual_description']}"
        for item in plan
    )
    for stale_term in STALE_PRODUCT_DOMAIN_TERMS_FOR_TEST:
        assert stale_term not in joined


def test_fallback_visual_plan_does_not_leak_old_product_domain_for_debate_pages():
    plan = _fallback_visual_plan(
        [
            {
                "page_num": 21,
                "type": "content",
                "text_content": {
                    "headline": "培养批判思维的日常工具",
                    "body": (
                        "- 工具一：苏格拉底式提问——“你怎么知道的？”“证据是什么？”\n"
                        "- 工具二：两栏日记——记录“观点”vs“事实”，每天练习区分\n"
                        "- 工具三：反方辩论——就一个观点，强制孩子站在对立面思考\n"
                        "- 研究支撑：辩论训练可提升逻辑推理能力和学业成绩（Kuhn, 1991）"
                    ),
                },
            }
        ],
        [],
    )

    evidence = plan[0]["visual_evidence"]
    description = plan[0]["visual_description"]
    # Fallback now returns open-ended descriptions rather than keyword-extracted lists.
    assert evidence and "核心场景" in evidence
    for stale_term in ("工业流水线", "古法小榨", "工坊", "花生", "油瓶", "榨油"):
        assert stale_term not in evidence
        assert stale_term not in description


def test_punchline_page_content_is_normalized_to_one_line():
    outline = _normalize_content_markdown([
        {
            "page_num": 3,
            "type": "hero",
            "text_content": {
                "headline": "管理就是把复杂留给自己",
                "subhead": "",
                "body": "- 第一条解释\n- 第二条解释",
            },
            "speaker_notes": "用于章节转折。",
            "visual_suggestion": "内容页信息图，展示三点列表",
        }
    ])

    page = outline[0]
    assert page["text_content"]["headline"] == "管理就是把复杂留给自己"
    assert page["text_content"]["subhead"] == ""
    assert page["text_content"]["body"] == ""
    assert "原金句页正文素材" in page["speaker_notes"]
    assert "金句页" in page["visual_suggestion"]
    assert "整套 PPT" in page["visual_suggestion"]


def test_single_ppt_outline_gets_source_refs_when_polishing():
    outline = [
        {"page_num": 1, "type": "cover", "text_content": {"headline": "封面", "body": ""}},
        {"page_num": 2, "type": "content", "text_content": {"headline": "策略", "body": "原内容"}},
    ]
    documents = '--- PPT_SOURCE filename="source.pptx" pages=2 ---\n\n--- 第1页 ---\n封面\n\n--- 第2页 ---\n策略'

    annotated = _annotate_ppt_source_refs(outline, documents, "请美化这个 PPT")

    assert annotated[0]["source_refs"] == [{
        "source_document": "source.pptx",
        "source_page_num": 1,
        "source_type": "pptx_slide",
        "reason": "single_ppt_page_polish",
    }]
    assert annotated[1]["source_refs"][0]["source_page_num"] == 2


def test_fallback_visual_plan_treats_hero_as_punchline_slide():
    plan = _fallback_visual_plan(
        [
            {
                "page_num": 5,
                "type": "hero",
                "text_content": {
                    "headline": "现金流就是公司的氧气",
                    "subhead": "来自 Q3 经营复盘",
                    "body": "",
                },
            }
        ],
        [],
    )

    assert plan[0]["layout"] == "hero"
    assert plan[0]["seed_family"] == "hero"
    assert "金句排版" in plan[0]["visual_evidence"]
    assert "可选署名/上下文" in plan[0]["visual_evidence"]
    assert "只保留核心短句" in plan[0]["visual_description"]
    assert plan[0]["logo_policy"]["show_logo"] is False


def test_visual_plan_json_repair_handles_multiline_llm_strings():
    raw = '{"6": {"visual_evidence": "终端货架", "visual_description": "第一行\n第二行", "visual_asset_ids": [], "visual_asset_usage": {}}}'

    parsed = _safe_parse_json(raw, 1)

    assert parsed["6"]["visual_description"] == "第一行\n第二行"


def test_visual_plan_raises_instead_of_auto_fallback_when_llm_drops_page(monkeypatch):
    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**_kwargs):
                    return SimpleNamespace(
                        choices=[SimpleNamespace(message=SimpleNamespace(content='{"2": {"visual_evidence": "别页", "visual_description": "别页描述"}}'))]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})

    with pytest.raises(VisualPlanGenerationError, match="missing page intents"):
        _do_generate_visual_plan(
            content_plan=[
                {
                    "page_num": 1,
                    "type": "content",
                    "text_content": {"headline": "当前页", "body": "本页必须得到完整视觉方案"},
                }
            ],
            style_override={"meta": {}, "body": "Style: 测试风格\nPalette: #111111, #FFFFFF\nVisual rhythm: 每页按文案生成画面证据"},
        )


def test_visual_plan_raises_instead_of_auto_fallback_when_required_fields_missing(monkeypatch):
    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**_kwargs):
                    return SimpleNamespace(
                        choices=[SimpleNamespace(message=SimpleNamespace(content='{"1": {"visual_description": "只有描述，没有画面证据"}}'))]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #FFFFFF",
    )

    with pytest.raises(VisualPlanGenerationError, match="missing visual_evidence"):
        _do_generate_visual_plan(
            content_plan=[
                {
                    "page_num": 1,
                    "type": "content",
                    "text_content": {"headline": "当前页", "body": "本页必须得到完整视觉方案"},
                }
            ],
            style_override={"meta": {}, "body": "Style: 测试风格\nPalette: #111111, #FFFFFF\nVisual rhythm: 每页按文案生成画面证据"},
        )


def test_visual_plan_retries_empty_batch_as_single_pages(monkeypatch):
    calls = []

    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**kwargs):
                    prompt = kwargs["messages"][-1]["content"]
                    calls.append(prompt)
                    if '"page_num": 6' in prompt and '"page_num": 10' in prompt:
                        return SimpleNamespace(choices=[SimpleNamespace(message=SimpleNamespace(content=""))])
                    page_nums = re.findall(r'"page_num":\s*(\d+)', prompt)
                    payload = {
                        page_num: {
                            "visual_evidence": f"第 {page_num} 页画面证据",
                            "visual_summary": f"第 {page_num} 页画面摘要",
                            "visual_description": f"第 {page_num} 页画面描述",
                            "visual_asset_ids": [],
                            "visual_asset_usage": {},
                        }
                        for page_num in page_nums
                    }
                    return SimpleNamespace(
                        choices=[SimpleNamespace(message=SimpleNamespace(content=json.dumps(payload, ensure_ascii=False)))]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setenv("PPTGOD_VISUAL_PLAN_BATCH_WORKERS", "1")
    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #FFFFFF",
    )

    plan = _do_generate_visual_plan(
        content_plan=[
            {
                "page_num": page_num,
                "type": "content",
                "text_content": {"headline": f"第 {page_num} 页", "body": "正文"},
            }
            for page_num in range(1, 11)
        ],
    )

    assert len(plan) == 10
    assert plan[5]["page_num"] == 6
    assert plan[5]["visual_description"] == "第 6 页画面描述"
    assert any('"page_num": 6' in prompt and '"page_num": 10' in prompt for prompt in calls)
    assert any('"page_num": 6' in prompt and '"page_num": 10' not in prompt for prompt in calls)


def test_visual_plan_drops_stale_numbered_visual_guidance_absent_from_visible_content(monkeypatch):
    calls = []

    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**kwargs):
                    prompt = kwargs["messages"][-1]["content"]
                    calls.append(prompt)
                    payload = {
                        "7": {
                            "visual_evidence": "章节标题转场",
                            "visual_summary": "墨黑底章节标题转场",
                            "visual_description": "纯墨黑底，只保留章节标题与大段留白。",
                            "visual_asset_ids": [],
                            "visual_asset_usage": {},
                        }
                    }
                    return SimpleNamespace(
                        choices=[SimpleNamespace(message=SimpleNamespace(content=json.dumps(payload, ensure_ascii=False)))]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setenv("PPTGOD_VISUAL_PLAN_BATCH_WORKERS", "1")
    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #FFFFFF",
    )

    _do_generate_visual_plan(
        content_plan=[
            {
                "page_num": 7,
                "type": "section",
                "text_content": {
                    "headline": "什么变了？决策不再只发生在人脑里",
                    "subhead": "",
                    "body": "",
                },
                "speaker_notes": '章节页先抛核心判断，让听众带着"被压缩过一遍"的悬念进入 1.1。',
                "visual_suggestion": "章节大字 + 1.1 / 1.2 两个小标题并列。",
                "visual_requirements": [
                    {"directive": "并排显示 1.1 和 1.2 两个章节入口", "diagram_labels": ["1.1", "1.2"]}
                ],
            }
        ],
    )

    assert calls
    assert "1.1" not in calls[0]
    assert "1.2" not in calls[0]
    assert "两个小标题并列" not in calls[0]


def test_visual_plan_sanitizes_section_numbering_returned_by_llm(monkeypatch):
    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**kwargs):
                    payload = {
                        "13": {
                            "visual_evidence": "章节转场：数字章节转场与章节名主导的仪式感画面",
                            "visual_summary": "第二章数字章节转场",
                            "visual_description": (
                                "画面中心略偏上放置一个超大细体数字章节转场（01 / 03 类），"
                                "章节转场用香槟金；其下紧跟章节标题。"
                            ),
                            "visual_asset_ids": [],
                            "visual_asset_usage": {},
                        }
                    }
                    return SimpleNamespace(
                        choices=[SimpleNamespace(message=SimpleNamespace(content=json.dumps(payload, ensure_ascii=False)))]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setenv("PPTGOD_VISUAL_PLAN_BATCH_WORKERS", "1")
    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #C9A961",
    )

    plan = _do_generate_visual_plan(
        content_plan=[
            {
                "page_num": 13,
                "type": "section",
                "section_title": "第二章",
                "text_content": {
                    "headline": "什么没变？人心仍然是终点",
                    "subhead": "",
                    "body": "",
                },
            }
        ],
    )

    joined = " ".join(
        str(plan[0].get(key) or "")
        for key in ("visual_evidence", "visual_summary", "visual_description")
    )
    assert "第二章" not in joined
    assert "数字章节" not in joined
    assert "章节转场" not in joined
    assert "01 / 03" not in joined


def test_visual_plan_drops_stale_case_cards_absent_from_visible_content(monkeypatch):
    calls = []

    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**kwargs):
                    prompt = kwargs["messages"][-1]["content"]
                    calls.append(prompt)
                    payload = {
                        "9": {
                            "visual_evidence": "私人买手式对话气泡",
                            "visual_summary": "深色对话气泡展示用户委托 AI 筛选",
                            "visual_description": "顶部用对话气泡承载提问场景，下方用两段信息层级对应正文。",
                            "visual_asset_ids": [],
                            "visual_asset_usage": {},
                        }
                    }
                    return SimpleNamespace(
                        choices=[SimpleNamespace(message=SimpleNamespace(content=json.dumps(payload, ensure_ascii=False)))]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setenv("PPTGOD_VISUAL_PLAN_BATCH_WORKERS", "1")
    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #FFFFFF",
    )

    _do_generate_visual_plan(
        content_plan=[
            {
                "page_num": 9,
                "type": "content",
                "text_content": {
                    "headline": "消费者雇了一个“全知全能”的私人买手",
                    "subhead": "",
                    "body": (
                        '- 过去搜索"防晒霜"：自己翻几百个商品、几千条评价\n'
                        '- 现在直接说："我要去三亚冲浪，预算 200 内，敏感肌，帮我推荐三款防水防晒霜"'
                    ),
                },
                "visual_suggestion": (
                    '顶部"私人买手对话框"还原场景，下方三栏讲山姆/胖东来/Costco，'
                    '底部一行大字"个人外脑"。'
                ),
            }
        ],
    )

    assert calls
    assert "山姆" not in calls[0]
    assert "胖东来" not in calls[0]
    assert "Costco" not in calls[0]
    assert "个人外脑" not in calls[0]


def test_visual_prompt_run_fails_when_any_page_prompt_generation_fails(monkeypatch):
    Session = make_shared_sessionmaker()
    db = Session()
    project = Project(
        title="Prompt partial failure",
        status="visual_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Test Style"},
    )
    db.add(project)
    db.flush()
    db.add_all(
        [
            Slide(
                project_id=project.id,
                page_num=1,
                status="pending",
                content_json={"page_num": 1, "type": "content", "text_content": {"headline": "第一页"}},
            ),
            Slide(
                project_id=project.id,
                page_num=2,
                status="pending",
                content_json={"page_num": 2, "type": "content", "text_content": {"headline": "第二页"}},
            ),
        ]
    )
    run = slides_api.create_project_run(
        db,
        project.id,
        kind="visual_prompts",
        stage="visual_planning",
        total_count=2,
    )
    db.commit()
    project_id = project.id
    run_id = run.id
    db.close()

    def fake_visual_plan(**_kwargs):
        return [
            {"page_num": 1, "visual_description": "第一页画面", "visual_evidence": "第一页证据"},
            {"page_num": 2, "visual_description": "第二页画面", "visual_evidence": "第二页证据"},
        ]

    def fake_prompt_for_page(**kwargs):
        if kwargs["page_intent"]["page_num"] == 2:
            raise RuntimeError("prompt provider returned empty output")
        return "page one prompt"

    monkeypatch.setattr(slides_api, "SessionLocal", Session)
    monkeypatch.setattr(slides_api, "generate_visual_plan", fake_visual_plan)
    monkeypatch.setattr(slides_api, "generate_prompt_for_page", fake_prompt_for_page)
    monkeypatch.setattr(slides_api, "_project_visual_assets_for_planning", lambda _project: [])
    monkeypatch.setattr(slides_api, "_derive_project_style_pack", lambda _project, _content_plan: "")

    asyncio.run(slides_api._do_generate_visual_and_prompts(project_id, None, run_id))

    verify_db = Session()
    refreshed_run = verify_db.query(ProjectRun).filter(ProjectRun.id == run_id).first()
    refreshed_project = verify_db.query(Project).filter(Project.id == project_id).first()
    refreshed_slides = (
        verify_db.query(Slide)
        .filter(Slide.project_id == project_id)
        .order_by(Slide.page_num)
        .all()
    )

    assert refreshed_run.status == "failed"
    assert refreshed_run.completed_count == 1
    assert refreshed_run.failed_count == 1
    assert "第 2 页" in (refreshed_run.error_msg or "")
    assert refreshed_project.status == "visual_ready"
    assert refreshed_slides[0].prompt_text == "page one prompt"
    assert refreshed_slides[0].status == "prompt_ready"
    assert refreshed_slides[1].prompt_text is None
    assert refreshed_slides[1].status == "visual_ready"


def test_old_visual_prompt_task_cleanup_keeps_newer_registry_task(monkeypatch):
    Session = make_shared_sessionmaker()
    db = Session()
    project = Project(title="No slides", status="planning")
    db.add(project)
    db.flush()
    run = slides_api.create_project_run(
        db,
        project.id,
        kind="visual_prompts",
        stage="visual_planning",
        total_count=0,
    )
    db.commit()
    project_id = project.id
    run_id = run.id
    db.close()

    monkeypatch.setattr(slides_api, "SessionLocal", Session)
    slides_api._running_tasks.pop(project_id, None)

    async def run_old_cleanup_with_new_task_registered():
        newer_task = asyncio.create_task(asyncio.sleep(60))
        slides_api._running_tasks[project_id] = newer_task
        try:
            await slides_api._do_generate_visual_and_prompts(project_id, None, run_id)
            assert slides_api._running_tasks.get(project_id) is newer_task
        finally:
            newer_task.cancel()
            try:
                await newer_task
            except asyncio.CancelledError:
                pass
            slides_api._running_tasks.pop(project_id, None)

    asyncio.run(run_old_cleanup_with_new_task_registered())


def test_pipeline_does_not_retry_slide_after_inner_image_retries_are_exhausted(monkeypatch, tmp_path):
    db = make_session()
    project = Project(title="No duplicate image retry", status="prompt_ready", content_plan_confirmed=True)
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        content_json={"page_num": 1, "text_content": {"headline": "图片页"}},
        visual_json={"visual_description": "图片页画面"},
        prompt_text="generate image",
    )
    db.add(slide)
    run = slides_api.create_project_run(
        db,
        project.id,
        kind="prototype_generation",
        stage="batch_generation",
        target_page_nums=[1],
        total_count=1,
    )
    db.commit()

    calls = []

    def fake_generate_one_slide(slide_arg, *_args, **_kwargs):
        calls.append(slide_arg.page_num)
        return {
            "slide": slide_arg,
            "error": "Connection aborted after image API retries",
            "image_generation_events": [
                {"status": "interrupted", "attempt": 1},
                {"status": "interrupted", "attempt": 2},
                {"status": "interrupted", "attempt": 3},
                {"status": "interrupted", "attempt": 4},
            ],
        }

    monkeypatch.setattr(generation_pipeline, "_generate_one_slide", fake_generate_one_slide)
    monkeypatch.setattr(generation_pipeline, "_pipeline_image_worker_count", lambda: 1)

    generation_pipeline.run_generation_pipeline(
        project.id,
        db,
        page_nums=[1],
        prototype=True,
        run_id=run.id,
    )

    assert calls == [1]


def test_prompt_keeps_exact_text_contract_and_visual_evidence():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 1,
            "type": "content",
            "layout": "content_split",
            "visual_evidence": "行业标准白皮书、官方印章和发布会背板",
            "visual_description": "围绕白皮书和发布会组织左右分栏。",
        },
        content_text={
            "headline": "整合营销：夺取行业标准制定权",
            "subhead": "两大事件，树立古法香正统地位",
            "body": "发布《古法香标准白皮书》\n联动权威机构举办发布会",
        },
        style_text_override="Style: 中式高端品牌\nPalette: #400000, #D4AF37\nPage type adaptation: 内容页优先可读",
    )

    assert 'Headline: "整合营销：夺取行业标准制定权"' in prompt
    assert 'Subhead: "两大事件，树立古法香正统地位"' in prompt
    assert 'Body: "发布《古法香标准白皮书》"' in prompt
    assert 'Body: "联动权威机构举办发布会"' in prompt
    assert "Visual:\n行业标准白皮书、官方印章和发布会背板" in prompt
    assert "FONT REQUIREMENT" not in prompt


def test_prompt_includes_exact_overlay_reservation():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 1,
            "type": "content",
            "layout": "content_split",
            "visual_evidence": "产品截图旁边解释三点价值",
            "visual_description": "右侧保留截图，左侧放文字。",
            "overlay_layers": [{
                "id": "ov",
                "asset_id": "asset-1",
                "enabled": True,
                "preset": "right-card",
            }],
        },
        content_text={"headline": "Exact Overlay", "body": "截图必须原样保留"},
        style_text_override="Style: clean\nPalette: #FFFFFF, #111111",
    )

    assert "Exact Overlay Reservation" in prompt
    assert "right side (approximately 34% width, center vertically)" in prompt
    assert "right-side card media slot" not in prompt
    assert "Place ALL visible text" in prompt


def test_prompt_skips_exact_overlay_reservation_when_asset_unavailable():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 1,
            "type": "content",
            "layout": "content_split",
            "visual_evidence": "产品截图旁边解释三点价值",
            "visual_description": "右侧放文字。",
            "available_overlay_asset_ids": [],
            "overlay_layers": [{
                "id": "ov",
                "asset_id": "missing-asset",
                "enabled": True,
                "preset": "right-card",
            }],
        },
        content_text={"headline": "Exact Overlay", "body": "截图不存在时不能预留空槽"},
        style_text_override="Style: clean\nPalette: #FFFFFF, #111111",
    )

    assert "Exact Overlay Reservation" not in prompt
    assert "right-side card media slot" not in prompt


def test_pptx_assembler_adds_exact_overlay_layer(tmp_path):
    bg_path = tmp_path / "slide.png"
    asset_path = tmp_path / "asset.png"
    output_path = tmp_path / "out.pptx"
    Image.new("RGB", (1792, 1024), "white").save(bg_path)
    Image.new("RGB", (320, 180), "red").save(asset_path)

    assemble_pptx(
        slide_images=[{
            "page_num": 1,
            "image_path": str(bg_path),
            "visual_json": {
                "overlay_layers": [{
                    "id": "ov",
                    "asset_id": "asset-1",
                    "enabled": True,
                    "preset": "right-card",
                    "mode": "exact_card",
                }]
            },
        }],
        output_path=str(output_path),
        overlay_assets={"asset-1": {"file_path": str(asset_path)}},
    )

    prs = Presentation(str(output_path))
    # Background image + card rectangle + exact overlay image.
    assert len(prs.slides[0].shapes) >= 3


def test_logo_overlay_layout_avoids_dense_ending_content(tmp_path):
    bg_path = tmp_path / "ending.png"
    logo_path = tmp_path / "logo.png"

    bg = Image.new("RGB", (1792, 1024), (4, 6, 18))
    draw = ImageDraw.Draw(bg)
    # Simulate a centered ending page: title, contact line, and two CTA bands.
    draw.rectangle((470, 250, 1320, 360), fill=(235, 70, 160))
    draw.rectangle((650, 450, 1140, 500), fill=(230, 230, 235))
    draw.rectangle((610, 640, 1180, 710), outline=(205, 80, 255), width=8)
    draw.rectangle((610, 750, 1180, 820), outline=(255, 80, 170), width=8)
    bg.save(bg_path)

    Image.new("RGBA", (520, 330), (255, 255, 255, 255)).save(logo_path)

    resolved = resolve_logo_overlay_box(str(bg_path), str(logo_path), "ending", "lower-center", "large")

    assert resolved is not None
    assert resolved["strategy"].startswith("smart:")
    # The old lower-center rule put the logo into the centered CTA stack.
    assert not (
        0.32 < resolved["left"] < 0.68
        and 0.55 < resolved["top"] < 0.82
    )


def test_logo_overlay_layout_avoids_light_ending_title_text(tmp_path):
    bg_path = tmp_path / "light-ending.png"
    logo_path = tmp_path / "logo.png"

    bg = Image.new("RGB", (1792, 1024), (249, 246, 241))
    draw = ImageDraw.Draw(bg)
    # Simulate the Ling Coffee ending page: a large brush headline and body
    # text centered on a light watercolor background.
    draw.rectangle((420, 210, 1490, 350), fill=(105, 85, 122))
    for y in (470, 550, 630, 710):
        draw.rectangle((560, y, 1250, y + 40), fill=(55, 48, 68))
    bg.save(bg_path)

    Image.new("RGBA", (520, 330), (90, 70, 120, 255)).save(logo_path)

    resolved = resolve_logo_overlay_box(str(bg_path), str(logo_path), "ending", "lower-center", "large")

    assert resolved is not None
    logo_box = (
        resolved["left"],
        resolved["top"],
        resolved["left"] + resolved["width"],
        resolved["top"] + resolved["height"],
    )
    title_box = (420 / 1792, 210 / 1024, 1490 / 1792, 350 / 1024)
    body_box = (560 / 1792, 470 / 1024, 1250 / 1792, 750 / 1024)

    def overlap(a, b):
        left = max(a[0], b[0])
        top = max(a[1], b[1])
        right = min(a[2], b[2])
        bottom = min(a[3], b[3])
        if right <= left or bottom <= top:
            return 0
        return (right - left) * (bottom - top)

    assert overlap(logo_box, title_box) == 0
    assert overlap(logo_box, body_box) == 0


def test_title_block_cover_logo_aligns_with_wide_title_not_densest_cluster(tmp_path):
    bg_path = tmp_path / "wide-title-cover.png"
    logo_path = tmp_path / "square-logo.png"

    bg = Image.new("RGB", (1792, 1024), (5, 8, 26))
    draw = ImageDraw.Draw(bg)
    # Simulate a cover where the Chinese title characters are denser on the
    # left, while the full title lockup spans across the page.
    for x in (330, 430, 530, 650, 760):
        draw.rectangle((x, 220, x + 70, 390), fill=(235, 228, 255))
    for x in (1210, 1370):
        draw.rectangle((x, 230, x + 95, 385), fill=(235, 228, 255))
    draw.rectangle((340, 500, 1450, 585), fill=(225, 220, 245))
    bg.save(bg_path)
    Image.new("RGBA", (600, 640), (120, 60, 255, 255)).save(logo_path)

    resolved = resolve_logo_overlay_box(str(bg_path), str(logo_path), "cover", "title-block-center", "large")

    assert resolved is not None
    assert resolved["strategy"] == "smart:above-title"
    logo_center = resolved["left"] + resolved["width"] / 2
    assert logo_center == pytest.approx(0.5, abs=0.07)


def test_large_cover_square_logo_uses_optical_height_cap(tmp_path):
    bg_path = tmp_path / "cover.png"
    logo_path = tmp_path / "square-logo.png"

    bg = Image.new("RGB", (1792, 1024), (5, 8, 26))
    draw = ImageDraw.Draw(bg)
    draw.rectangle((360, 250, 1430, 410), fill=(235, 228, 255))
    bg.save(bg_path)
    Image.new("RGBA", (600, 640), (120, 60, 255, 255)).save(logo_path)

    resolved = resolve_logo_overlay_box(str(bg_path), str(logo_path), "cover", "title-block-center", "large")

    assert resolved is not None
    assert resolved["height"] <= 0.145


def test_cover_center_logo_placement_is_physical_center(tmp_path):
    bg_path = tmp_path / "cover.png"
    logo_path = tmp_path / "logo.png"
    bg = Image.new("RGB", (1792, 1024), (10, 20, 45))
    draw = ImageDraw.Draw(bg)
    draw.rectangle((80, 260, 780, 430), fill=(245, 245, 255))
    bg.save(bg_path)
    Image.new("RGBA", (300, 120), (255, 255, 255, 255)).save(logo_path)

    resolved = resolve_logo_overlay_box(str(bg_path), str(logo_path), "cover", "center", "large")

    assert resolved is not None
    assert resolved["strategy"] == "static:center"
    assert abs((resolved["left"] + resolved["width"] / 2) - 0.5) < 0.01
    assert abs((resolved["top"] + resolved["height"] / 2) - 0.5) < 0.01


def test_small_corner_logo_is_readable_default_size(tmp_path):
    bg_path = tmp_path / "content.png"
    logo_path = tmp_path / "logo.png"
    Image.new("RGB", (1792, 1024), (20, 20, 24)).save(bg_path)
    Image.new("RGBA", (300, 120), (255, 255, 255, 255)).save(logo_path)

    resolved = resolve_logo_overlay_box(str(bg_path), str(logo_path), "content", "top-right", "small")

    assert resolved is not None
    assert resolved["width"] == pytest.approx(0.085, abs=0.002)


def test_ending_small_corner_logo_respects_small_scale(tmp_path):
    bg_path = tmp_path / "ending.png"
    logo_path = tmp_path / "logo.png"
    Image.new("RGB", (1792, 1024), (249, 246, 241)).save(bg_path)
    Image.new("RGBA", (300, 120), (80, 60, 120, 255)).save(logo_path)

    resolved = resolve_logo_overlay_box(str(bg_path), str(logo_path), "ending", "top-right", "small")

    assert resolved is not None
    assert resolved["width"] == pytest.approx(0.085, abs=0.002)
    assert resolved["strategy"] == "static:top-right"


def test_pptx_assembler_uses_resolved_logo_overlay_box(tmp_path):
    bg_path = tmp_path / "slide.png"
    logo_path = tmp_path / "logo.png"
    output_path = tmp_path / "out.pptx"
    Image.new("RGB", (1792, 1024), "black").save(bg_path)
    Image.new("RGBA", (300, 120), (255, 255, 255, 255)).save(logo_path)

    assemble_pptx(
        slide_images=[{
            "page_num": 1,
            "type": "cover",
            "image_path": str(bg_path),
            "visual_json": {
                "logo_policy": {
                    "show_logo": True,
                    "placement": "lower-center",
                    "scale": "large",
                    "resolved_overlay_box": {
                        "left": 0.12,
                        "top": 0.08,
                        "width": 0.18,
                        "height": 0.06,
                        "strategy": "smart:test",
                    },
                }
            },
        }],
        output_path=str(output_path),
        logo_config={"file_path": str(logo_path), "anchor": "top-right"},
    )

    prs = Presentation(str(output_path))
    logo_shape = prs.slides[0].shapes[1]
    assert logo_shape.left == int(prs.slide_width * 0.12)
    assert logo_shape.top == int(prs.slide_height * 0.08)
    assert logo_shape.width == int(prs.slide_width * 0.18)


def test_pptx_assembler_uses_clean_logo_without_rectangular_backplate(tmp_path):
    bg_path = tmp_path / "slide.png"
    logo_path = tmp_path / "logo.png"
    output_path = tmp_path / "out.pptx"
    Image.new("RGB", (1792, 1024), "white").save(bg_path)
    logo = Image.new("RGBA", (300, 120), (255, 255, 255, 0))
    draw = ImageDraw.Draw(logo)
    draw.rectangle((20, 35, 280, 85), fill=(18, 18, 18, 255))
    logo.save(logo_path)

    assemble_pptx(
        slide_images=[{
            "page_num": 1,
            "type": "cover",
            "image_path": str(bg_path),
            "visual_json": {},
        }],
        output_path=str(output_path),
        logo_config={"file_path": str(logo_path), "anchor": "top-right"},
    )

    prs = Presentation(str(output_path))
    shapes = prs.slides[0].shapes
    assert len(shapes) == 2
    logo_shape = shapes[1]
    assert abs((logo_shape.left + logo_shape.width / 2) - prs.slide_width / 2) < Inches(0.05)


def test_logo_render_policy_keeps_full_logo_on_required_dark_cover(tmp_path):
    bg_path = tmp_path / "dark-slide.png"
    logo_path = tmp_path / "mixed-logo.png"
    Image.new("RGB", (1792, 1024), (12, 12, 12)).save(bg_path)
    img = Image.new("RGB", (260, 80), "white")
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((14, 16, 62, 64), radius=9, fill=(255, 205, 0))
    draw.rectangle((88, 28, 232, 42), fill=(18, 18, 18))
    draw.rectangle((88, 48, 198, 57), fill=(18, 18, 18))
    img.save(logo_path)

    full_path = prepare_logo_overlay_image(str(logo_path))
    symbol_path = prepare_logo_symbol_image(str(logo_path))
    policy = resolve_logo_render_policy(
        str(bg_path),
        full_path,
        symbol_path,
        "cover",
        "lower-center",
        "large",
        {"show_logo": True, "placement": "lower-center", "scale": "large"},
    )

    assert policy["show_logo"] is True
    assert policy["render_variant"] == "full"
    assert isinstance(policy["resolved_overlay_box"], dict)


def test_logo_render_policy_uses_full_logo_on_light_quiet_page(tmp_path):
    bg_path = tmp_path / "light-slide.png"
    logo_path = tmp_path / "mixed-logo.png"
    Image.new("RGB", (1792, 1024), (246, 244, 238)).save(bg_path)
    img = Image.new("RGB", (260, 80), "white")
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((14, 16, 62, 64), radius=9, fill=(255, 205, 0))
    draw.rectangle((88, 28, 232, 42), fill=(18, 18, 18))
    draw.rectangle((88, 48, 198, 57), fill=(18, 18, 18))
    img.save(logo_path)

    full_path = prepare_logo_overlay_image(str(logo_path))
    symbol_path = prepare_logo_symbol_image(str(logo_path))
    policy = resolve_logo_render_policy(
        str(bg_path),
        full_path,
        symbol_path,
        "content",
        "top-right",
        "small",
        {"show_logo": True, "placement": "top-right", "scale": "small"},
    )

    assert policy["show_logo"] is True
    assert policy["render_variant"] == "full"
    assert policy["logo_contrast"] == "readable"


def test_logo_render_policy_keeps_required_pages_even_when_dark(tmp_path):
    bg_path = tmp_path / "dark-slide.png"
    logo_path = tmp_path / "wordmark.png"
    Image.new("RGB", (1792, 1024), (12, 12, 12)).save(bg_path)
    img = Image.new("RGB", (260, 80), "white")
    draw = ImageDraw.Draw(img)
    draw.rectangle((20, 26, 238, 42), fill=(18, 18, 18))
    draw.rectangle((20, 49, 190, 58), fill=(18, 18, 18))
    img.save(logo_path)

    full_path = prepare_logo_overlay_image(str(logo_path))
    policy = resolve_logo_render_policy(
        str(bg_path),
        full_path,
        prepare_logo_symbol_image(str(logo_path)),
        "cover",
        "lower-center",
        "large",
        {"show_logo": True, "placement": "lower-center", "scale": "large"},
    )

    assert policy["show_logo"] is True
    assert policy["render_variant"] == "full"


def test_logo_policy_ignores_stale_omit_on_required_content_page():
    policy = logo_policy_for_page(
        {
            "type": "content",
            "logo_policy": {
                "show_logo": False,
                "placement": "top-right",
                "scale": "small",
                "render_variant": "omit",
            },
        }
    )

    assert policy["show_logo"] is True
    assert policy["placement"] == "top-right"
    assert policy["scale"] == "small"
    assert "render_variant" not in policy


def test_logo_policy_ignores_auto_symbol_variant():
    policy = logo_policy_for_page(
        {
            "type": "content",
            "logo_policy": {
                "show_logo": True,
                "placement": "top-right",
                "scale": "small",
                "render_variant": "symbol",
            },
        }
    )

    assert policy["show_logo"] is True
    assert "render_variant" not in policy


def test_logo_policy_allows_omit_on_section_page():
    policy = logo_policy_for_page(
        {
            "type": "section",
            "logo_policy": {
                "show_logo": False,
                "placement": "top-right",
                "scale": "small",
                "render_variant": "omit",
            },
        }
    )

    assert policy["show_logo"] is False
    assert policy["render_variant"] == "omit"


def test_project_quality_report_flags_low_contrast_logo(tmp_path):
    logo_path = tmp_path / "logo.png"
    slide_path = tmp_path / "slide.png"
    pptx_path = tmp_path / "presentation.pptx"
    Image.new("RGBA", (120, 60), (20, 20, 20, 255)).save(logo_path)
    Image.new("RGB", (1792, 1024), (15, 15, 25)).save(slide_path)
    pptx_path.write_bytes(b"pptx")
    project = Project(id="p1", title="Deck", status="completed")
    project.reference_images = [
        ReferenceImage(
            id="logo-1",
            project_id="p1",
            role="logo",
            file_path=str(logo_path),
            asset_analysis={"review_status": "user_confirmed"},
        )
    ]
    slides = [
        Slide(
            id="s1",
            project_id="p1",
            page_num=1,
            type="content",
            status="completed",
            image_path=str(slide_path),
            content_json={"title": "增长策略", "bullets": ["统一品牌露出"]},
            visual_json={
                "type": "content",
                "logo_policy": {
                    "show_logo": True,
                    "placement": "top-right",
                    "scale": "small",
                    "logo_contrast": "low_contrast_manual_review",
                },
            },
        )
    ]

    report = build_project_quality_report(project, slides, has_pptx=True, pptx_path=str(pptx_path))

    assert report
    assert any(issue["kind"] == "logo_low_contrast" and issue["pages"] == [1] for issue in report["issues"])
    assert "Logo 对比度偏弱" in report["message"]
    assert "手动调整" in report["message"]


def test_project_quality_report_ignores_optional_logo_omission(tmp_path):
    logo_path = tmp_path / "logo.png"
    slide_path = tmp_path / "section.png"
    pptx_path = tmp_path / "presentation.pptx"
    Image.new("RGBA", (120, 60), (20, 20, 20, 255)).save(logo_path)
    Image.new("RGB", (1792, 1024), (15, 15, 25)).save(slide_path)
    pptx_path.write_bytes(b"pptx")
    project = Project(id="p1", title="Deck", status="completed")
    project.reference_images = [
        ReferenceImage(
            id="logo-1",
            project_id="p1",
            role="logo",
            file_path=str(logo_path),
            asset_analysis={"review_status": "user_confirmed"},
        )
    ]
    slides = [
        Slide(
            id="s1",
            project_id="p1",
            page_num=1,
            type="section",
            status="completed",
            image_path=str(slide_path),
            content_json={"title": "章节页"},
            visual_json={
                "type": "section",
                "logo_policy": {"show_logo": False, "render_variant": "omit"},
            },
        )
    ]

    report = build_project_quality_report(project, slides, has_pptx=True, pptx_path=str(pptx_path))

    assert report
    assert not any(issue["severity"] == "error" for issue in report["issues"])
    assert not any(issue["kind"] == "required_logo_policy_missing" for issue in report["issues"])
    assert "章节页和金句页可以不放" in report["message"]


def test_project_quality_report_waits_until_final_stage(tmp_path):
    project = Project(id="p1", title="Deck", status="prompt_ready")
    slide = Slide(id="s1", project_id="p1", page_num=1, type="content", status="pending")

    assert build_project_quality_report(project, [slide], has_pptx=False) is None


def test_ending_logo_policy_defaults_to_small_corner_signature():
    policy = logo_policy_for_page({"type": "ending"})

    assert policy["show_logo"] is True
    assert policy["placement"] == "top-right"
    assert policy["scale"] == "small"


def test_prompt_for_punchline_page_uses_punchline_treatment():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 4,
            "type": "hero",
            "layout": "hero",
            "visual_evidence": "金句排版、作者署名与纯色纹理背景",
            "visual_description": "中央大字金句，右下角小号署名，背景只有细微纸纹。",
        },
        content_text={
            "headline": "少即是多",
            "subhead": "—— 路德维希·密斯·凡德罗",
            "body": "这句名言的解释不应该出现在画面正文里",
        },
        style_text_override="Style: 冷静现代\nVisual rhythm: 内容页浅底高可读，金句页更沉浸",
    )

    assert 'Headline: "少即是多"' in prompt
    assert 'Subhead: "—— 路德维希·密斯·凡德罗"' in prompt
    assert "Body:" not in prompt
    assert "Punchline slide treatment" in prompt
    assert "same project typeface feel" in prompt
    assert "dense panels" in prompt


def test_prompt_includes_selected_global_visual_asset():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 2,
            "type": "content",
            "layout": "content_split",
            "visual_evidence": "胡姬花花生油瓶与终端货架",
            "visual_description": "右侧展示产品瓶，左侧放正文。",
            "visual_asset_usage": {"asset-1": "放在右侧产品展示区，保持瓶型和包装文字可识别"},
        },
        content_text={"headline": "终端小油瓶体验", "body": "用产品实物建立货架记忆"},
        reference_images=[
            {
                "id": "asset-1",
                "role": "visual_asset",
                "process_mode": "crop",
                "asset_name": "胡姬花花生油瓶",
                "asset_kind": "product",
                "asset_route_mode": "double_blend",
                "description": "subject=胡姬花花生油瓶; features=红色瓶盖、黄色标签",
            }
        ],
        style_text_override="Style: 品牌展示\nPalette: #FFFFFF, #B01622",
    )

    assert "Product slot: 胡姬花花生油瓶" in prompt
    assert "uploaded product image as the product source" in prompt
    assert "preserve product identity and visible details" in prompt
    assert "hidden refinement pass" not in prompt
    assert "refinement pass" not in prompt
    assert "Place the uploaded product image in the right side area" in prompt
    assert "红色瓶盖" not in prompt
    assert "黄色标签" not in prompt


def test_prompt_strips_product_details_from_layout_usage_and_style_negatives():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 1,
            "type": "cover",
            "layout": "cover",
            "visual_evidence": "品牌主视觉与胡姬花花生油产品实物",
            "visual_description": (
                "全屏深红色背景；画面中央偏下位置展示胡姬花古法花生油产品实物——"
                "5升装透明金黄桶身、金色提手盖、瓶颈深红扇形非遗吊牌尤为突出；"
                "产品后方呈现古法木榨工艺场景；不要使用科技风"
            ),
            "visual_asset_usage": {
                "asset-1": "置于画面中央偏下位置，作为品牌核心实物锚点；产品瓶颈扇形吊牌结构、红底金边标签必须完整保留"
            },
        },
        content_text={"headline": "2026胡姬花花生油品牌策略建议"},
        reference_images=[
            {
                "id": "asset-1",
                "role": "visual_asset",
                "process_mode": "crop",
                "asset_name": "胡姬花古法花生油",
                "asset_kind": "product",
                "description": "5升装透明金黄色塑料桶装花生油，瓶盖为金色带提手设计",
            }
        ],
        style_text_override=(
            "Style: 新中式古典重彩风格\n"
            "Mood: 厚重、典雅、华丽、古朴\n"
            "Visual rhythm: 内容核心更接近古法非遗、传统食品/农业品牌，应优先考虑传统质感；不要因为出现“战略”而推荐科技风"
        ),
    )

    assert prompt.index("References:") < prompt.index("Style:")
    assert "Product slot: 胡姬花古法花生油" in prompt
    assert "uploaded product image as the product source" in prompt
    assert "Place the uploaded product image in the lower center area" in prompt
    assert "产品后方呈现古法木榨工艺场景" not in prompt
    # Negative clauses are now preserved per prompt-engine policy; only strip
    # concrete product details, not style guidance containing "不要" or "科技风".
    for unwanted in ["5升", "金黄桶身", "金色提手", "非遗吊牌", "红底金边标签"]:
        assert unwanted not in prompt


def test_visual_plan_source_prompt_avoids_asset_detail_production():
    prompt = _build_batch_prompt(
        pages_summary=[
            {
                "page_num": 1,
                "type": "cover",
                "text_content": {"headline": "产品封面", "body": "展示产品与工艺场景"},
            }
        ],
        style={"meta": {"theme": "新中式", "mood": "典雅", "palette": ["#6B0B0B", "#D4AF37"]}, "body": "品牌风格"},
        global_visual_assets=[
            {
                "id": "asset-1",
                "name": "胡姬花花生油瓶",
                "kind": "product",
                "usage_note": "产品展示页使用",
            }
        ],
    )

    assert "visual_asset_usage 只能写位置和画面占比" in prompt
    assert "不要在 visual_evidence、visual_description、visual_asset_usage 里复述外观" in prompt
    assert "给用户阅读" not in prompt
    assert "下游 pipeline" not in prompt
    assert "保真要求" not in prompt


def test_visual_asset_recall_uses_structured_source_refs():
    recalled = _recall_visual_assets_for_page(
        {
            "page_num": 3,
            "text_content": {"headline": "全新叙事页", "body": "这里没有直接提产品名"},
            "source_refs": [{"source_document": "source.pptx", "source_page_num": 2}],
        },
        [
            {
                "id": "asset-1",
                "name": "不在文案里的产品图",
                "kind": "product",
                "selection_tier": "core_global",
                "importance_score": 60,
                "source_document": "source.pptx",
                "source_page_num": 2,
                "analysis_summary": "name=不在文案里的产品图",
            }
        ],
    )

    assert recalled[0]["id"] == "asset-1"
    assert "源 PPT 页" in recalled[0]["reason"]


def test_visual_asset_recall_ignores_shared_brand_terms_across_assets():
    recalled = _recall_visual_assets_for_page(
        {
            "page_num": 2,
            "type": "content",
            "text_content": {
                "headline": "vivo 校园传播链路",
                "body": "这一页讲学生换机窗口和校园媒体组合，没有指定具体机型。",
            },
        },
        [
            {
                "id": "s30",
                "name": "机型_vivo S30 KV",
                "kind": "product",
                "selection_tier": "core_global",
                "importance_score": 60,
                "analysis_summary": "keywords=vivo、S30、KV、手机、校园",
            },
            {
                "id": "s50",
                "name": "机型_vivo S50 KV",
                "kind": "product",
                "selection_tier": "core_global",
                "importance_score": 60,
                "analysis_summary": "keywords=vivo、S50、KV、手机、校园",
            },
            {
                "id": "s60",
                "name": "机型_vivo S60 KV",
                "kind": "product",
                "selection_tier": "core_global",
                "importance_score": 60,
                "analysis_summary": "keywords=vivo、S60、KV、手机、校园",
            },
        ],
    )

    assert recalled == []


def test_visual_asset_recall_keeps_specific_model_terms_with_shared_brand():
    recalled = _recall_visual_assets_for_page(
        {
            "page_num": 17,
            "type": "content",
            "text_content": {
                "headline": "女大学生内容分享型学生：主推 S60、S50、S30",
                "body": "拍照好看、自拍自然、轻薄，适合开学新生活。",
            },
        },
        [
            {
                "id": "s30",
                "name": "机型_vivo S30 KV",
                "kind": "product",
                "selection_tier": "core_global",
                "importance_score": 60,
                "analysis_summary": "keywords=vivo、S30、KV、手机、校园",
            },
            {
                "id": "z11",
                "name": "机型_iQOO Z11 KV",
                "kind": "product",
                "selection_tier": "core_global",
                "importance_score": 60,
                "analysis_summary": "keywords=iQOO、Z11、KV、游戏、手机、校园",
            },
        ],
    )

    assert [item["id"] for item in recalled] == ["s30"]


def test_visual_plan_prompt_treats_toc_as_simple_navigation():
    prompt = _build_batch_prompt(
        pages_summary=[
            {
                "page_num": 2,
                "type": "toc",
                "headline": "目录",
                "body_preview": "01 市场\n02 产品\n03 路径",
            }
        ],
        style={"meta": {"theme": "商务", "mood": "清晰", "palette": ["#111111", "#FFFFFF"]}, "body": ""},
    )

    assert "toc / 目录页：只有导航功能" in prompt
    assert "3-6 个短章节名" in prompt
    assert "不要做成花哨菜单" in prompt


def test_visual_plan_prompt_treats_section_numbers_as_visible_only_when_user_wrote_them():
    prompt = _build_batch_prompt(
        pages_summary=[
            {
                "page_num": 8,
                "type": "section",
                "section_title": "第二章",
                "text_content": {"headline": "什么没变？人心仍然是终点", "body": ""},
            }
        ],
        style={"meta": {"theme": "黑金", "mood": "克制", "palette": ["#0A0A0A", "#C9A961"]}, "body": ""},
    )

    assert "只有当章号、章节编号或序号出现在 headline、subhead 或 body 的用户可见文案中" in prompt
    assert "section_title 只作为结构元数据" in prompt


def test_visual_plan_prompt_drops_section_number_from_existing_suggestion_when_absent_from_copy():
    prompt = _build_batch_prompt(
        pages_summary=[
            {
                "page_num": 8,
                "type": "section",
                "section_title": "第六章",
                "text_content": {"headline": "行动清单", "body": ""},
                "visual_suggestion": "左上角写第六章，下面放章节主标题和一条金色细线。",
            }
        ],
        style={"meta": {"theme": "黑金", "mood": "克制", "palette": ["#0A0A0A", "#C9A961"]}, "body": ""},
    )

    assert "第六章" not in prompt
    assert "左上角写" not in prompt
    assert "章节主标题" in prompt


def test_visual_plan_prompt_guards_cover_data_and_ending_roles():
    prompt = _build_batch_prompt(
        pages_summary=[
            {"page_num": 1, "type": "cover", "headline": "年度策略"},
            {"page_num": 5, "type": "data", "headline": "关键数字", "body_preview": "转化率 18%\n客单价 320 元"},
            {"page_num": 9, "type": "ending", "headline": "谢谢"},
        ],
        style={"meta": {"theme": "商务", "mood": "清晰", "palette": ["#111111", "#FFFFFF"]}, "body": ""},
    )

    assert "cover / 封面：只负责定调和命名" in prompt
    assert "data / 数据页：只有当页面正文给出真实数字" in prompt
    assert "不要编造数值、趋势或坐标轴" in prompt
    assert "ending / 封底：只做收束、感谢、CTA 或联系方式" in prompt
    assert "不要引入新的证明素材" in prompt


def test_visual_plan_prompt_disables_logo_policy_without_project_logo():
    prompt = _build_batch_prompt(
        pages_summary=[{"page_num": 1, "type": "cover", "headline": "年度策略"}],
        style={"meta": {"theme": "商务", "mood": "清晰", "palette": ["#111111", "#FFFFFF"]}, "body": ""},
        has_project_logo=False,
    )

    assert "当前项目没有已确认的用户 Logo" in prompt
    assert '"show_logo": false' in prompt
    assert "不要为 Logo、品牌角标、标识、徽标或占位框预留空间" in prompt
    assert "title-block-center" not in prompt


def test_visual_plan_prompt_allows_logo_policy_with_project_logo():
    prompt = _build_batch_prompt(
        pages_summary=[{"page_num": 1, "type": "cover", "headline": "年度策略"}],
        style={"meta": {"theme": "商务", "mood": "清晰", "palette": ["#111111", "#FFFFFF"]}, "body": ""},
        has_project_logo=True,
    )

    assert "title-block-center" in prompt
    assert '"show_logo": true' in prompt
    assert "不要要求底图为 Logo 绘制占位框、虚线框、圆角框、底板、徽章、描边、外发光或任何容器" in prompt


def test_default_visual_asset_usage_does_not_describe_product_appearance():
    usage = _default_visual_asset_usage(
        {"name": "胡姬花花生油瓶", "kind": "product"},
        {"page_num": 1},
    )

    assert "uploaded product image" in usage
    assert "胡姬花花生油瓶" not in usage
    for unwanted in ["主体形状", "包装结构", "颜色", "品牌识别"]:
        assert unwanted not in usage


def test_prompt_does_not_reserve_logo_area_without_uploaded_logo_policy():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 1,
            "type": "cover",
            "layout": "cover",
            "visual_evidence": "品牌主张与主视觉背景",
            "visual_description": "封面突出主标题，品牌 Logo 与标题形成稳定关系。",
        },
        content_text={"headline": "胡姬花年度整合营销提案"},
        reference_images=[],
        style_text_override="Style: 品牌提案\nPalette: #FFFFFF, #B01622",
    )

    assert "Logo Placement Note" not in prompt
    assert "exact overlay" not in prompt
    assert "logo" not in prompt.lower()
    assert "LOGO" not in prompt


def test_prompt_keeps_logo_overlay_policy_out_of_image_prompt():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 1,
            "type": "cover",
            "layout": "cover",
            "visual_evidence": "品牌主张与主视觉背景",
            "visual_description": "封面突出主标题，品牌 Logo 与标题形成稳定关系。",
            "logo_policy": {"show_logo": True, "placement": "center", "scale": "large"},
        },
        content_text={"headline": "胡姬花年度整合营销提案"},
        reference_images=[],
        style_text_override="Style: 品牌提案\nPalette: #FFFFFF, #B01622",
    )

    assert "Logo Placement Note" not in prompt
    assert "Brand marks:" not in prompt
    assert "brand signature" not in prompt.lower()
    assert "safe corner" not in prompt.lower()
    assert "placeholder" not in prompt.lower()
    assert "logo" not in prompt.lower()
    assert "LOGO" not in prompt


def test_prompt_strips_generated_logo_and_watermark_instructions():
    prompt = prompt_engine.generate_prompt_for_page(
        page_intent={
            "page_num": 2,
            "type": "content",
            "layout": "content_split",
            "visual_evidence": "三列业务卡片和右上角小Logo",
            "visual_description": "浅底三列卡片；右上角小Logo；底部加入品牌抽象展翅造型装饰元素；不要出现虎课网水印。",
            "logo_policy": {"show_logo": True, "placement": "top-right", "scale": "small"},
        },
        content_text={"headline": "企业现状痛点", "body": ["账号同质化", "招聘功能缺失"]},
        reference_images=[],
        style_text_override=(
            "Style: 金黑动感\n"
            "Palette: #D3BC8E, #000000, #F0E0C0, #1A1A1A\n"
            "Visual rhythm: 继承品牌Logo金黑配色；装饰语言延续品牌抽象展翅造型；内容页右上角小Logo。"
        ),
    )

    assert "Watermarks and stray marks:" in prompt
    assert "Brand marks:" not in prompt
    assert "logo" not in prompt.lower()
    assert "LOGO" not in prompt
    assert "右上角小Logo" not in prompt
    assert "品牌抽象展翅造型" not in prompt
    assert "金黑配色" not in prompt


def test_project_logo_policy_is_disabled_when_no_confirmed_logo():
    project = SimpleNamespace(reference_images=[])

    intent = slides_api._with_project_logo_policy(
        {
            "type": "cover",
            "logo_policy": {
                "show_logo": True,
                "placement": "title-block-center",
                "scale": "large",
                "resolved_overlay_box": {"left": 0.1},
            },
        },
        project,
    )

    assert intent["logo_policy"]["show_logo"] is False
    assert intent["logo_policy"]["use_as_scene_asset"] is False
    assert "resolved_overlay_box" not in intent["logo_policy"]


def test_project_logo_policy_restores_required_content_logo(tmp_path):
    logo_path = tmp_path / "logo.png"
    Image.new("RGBA", (120, 60), (80, 60, 120, 255)).save(logo_path)
    logo = SimpleNamespace(
        id="logo-1",
        role="logo",
        slide_id=None,
        file_path=str(logo_path),
        logo_anchor="top-right",
        asset_analysis={},
    )
    project = SimpleNamespace(reference_images=[logo])

    intent = slides_api._with_project_logo_policy(
        {
            "type": "content",
            "logo_policy": {
                "show_logo": False,
                "placement": "top-right",
                "scale": "small",
                "render_variant": "omit",
            },
        },
        project,
    )

    assert intent["logo_policy"]["show_logo"] is True
    assert "render_variant" not in intent["logo_policy"]


def test_repair_project_logo_policies_updates_stale_required_content_logo(tmp_path):
    logo_path = tmp_path / "logo.png"
    slide_path = tmp_path / "slide.png"
    Image.new("RGBA", (120, 60), (80, 60, 120, 255)).save(logo_path)
    Image.new("RGB", (1792, 1024), (12, 12, 24)).save(slide_path)
    project = Project(id="p1", title="Deck", status="completed")
    project.reference_images = [
        ReferenceImage(
            id="logo-1",
            project_id="p1",
            role="logo",
            slide_id=None,
            file_path=str(logo_path),
            logo_anchor="top-right",
            asset_analysis={"review_status": "user_confirmed"},
        )
    ]
    slide = Slide(
        id="s1",
        project_id="p1",
        page_num=3,
        type="content",
        status="completed",
        image_path=str(slide_path),
        visual_json={
            "type": "content",
            "logo_policy": {
                "show_logo": False,
                "placement": "top-right",
                "scale": "small",
                "render_variant": "omit",
            },
        },
    )

    changed = slides_api._repair_project_logo_policies(project, [slide])

    assert changed is True
    assert slide.visual_json["logo_policy"]["show_logo"] is True
    assert slide.visual_json["logo_policy"]["render_variant"] == "full"
    assert "resolved_overlay_box" in slide.visual_json["logo_policy"]


def test_ending_logo_policy_uses_small_confirmed_corner_signature(tmp_path):
    logo_path = tmp_path / "logo.png"
    Image.new("RGBA", (120, 60), (80, 60, 120, 255)).save(logo_path)
    logo = SimpleNamespace(
        id="logo-1",
        role="logo",
        slide_id=None,
        file_path=str(logo_path),
        logo_anchor="top-right",
        asset_analysis={},
    )
    project = SimpleNamespace(reference_images=[logo])

    intent = slides_api._with_project_logo_policy(
        {
            "type": "ending",
            "logo_policy": {
                "show_logo": True,
                "placement": "lower-center",
                "scale": "large",
                "resolved_overlay_box": {"left": 0.45},
            },
        },
        project,
    )

    assert intent["logo_policy"]["placement"] == "top-right"
    assert intent["logo_policy"]["scale"] == "small"
    assert "resolved_overlay_box" not in intent["logo_policy"]


def test_prompt_asset_policy_filters_stale_overlay_layers(tmp_path):
    asset_path = tmp_path / "asset.png"
    Image.new("RGB", (8, 8), "white").save(asset_path)
    project_asset = SimpleNamespace(
        id="asset-ok",
        role="visual_asset",
        slide_id=None,
        file_path=str(asset_path),
    )
    project = SimpleNamespace(reference_images=[project_asset])
    slide = SimpleNamespace(id="slide-1", reference_images=[])

    intent = slides_api._with_prompt_asset_policies(
        {
            "type": "content",
            "overlay_layers": [
                {"id": "ok", "asset_id": "asset-ok", "preset": "right-card"},
                {"id": "missing", "asset_id": "asset-missing", "preset": "left-card"},
            ],
        },
        project,
        slide,
    )

    assert [layer["asset_id"] for layer in intent["overlay_layers"]] == ["asset-ok"]
    assert intent["available_overlay_asset_ids"] == ["asset-ok"]


def test_generation_loads_selected_visual_assets_without_signature_logo(tmp_path):
    page_ref_path = tmp_path / "page.png"
    logo_path = tmp_path / "logo.png"
    asset_path = tmp_path / "asset.png"
    template_path = tmp_path / "template.png"
    for path in (page_ref_path, logo_path, asset_path, template_path):
        Image.new("RGB", (8, 8), "white").save(path)

    page_ref = SimpleNamespace(
        id="page-ref",
        role="content_ref",
        file_path=str(page_ref_path),
        process_mode="blend",
    )
    logo = SimpleNamespace(
        id="logo-1",
        role="logo",
        slide_id=None,
        file_path=str(logo_path),
        process_mode="blend",
    )
    visual_asset = SimpleNamespace(
        id="asset-1",
        role="visual_asset",
        slide_id=None,
        file_path=str(asset_path),
        process_mode="crop",
        asset_kind="product",
    )
    project = SimpleNamespace(
        reference_images=[logo, visual_asset],
        selected_template_recommendations={"content": {"file_path": str(template_path)}},
    )
    slide = SimpleNamespace(
        page_num=3,
        type="content",
        visual_json={"visual_asset_ids": ["asset-1"]},
        reference_images=[page_ref],
        project=project,
    )

    refs = _load_reference_images(slide)

    assert [r["role"] for r in refs] == ["visual_asset", "content_ref", "template"]
    assert refs[0]["process_mode"] == "crop"
    assert refs[2]["image"].size == (8, 8)


def test_generation_can_load_logo_as_scene_asset_on_cover_when_blend(tmp_path):
    logo_path = tmp_path / "logo.png"
    Image.new("RGB", (8, 8), "white").save(logo_path)
    logo = SimpleNamespace(
        id="logo-1",
        role="logo",
        slide_id=None,
        file_path=str(logo_path),
        process_mode="blend",
    )
    slide = SimpleNamespace(
        page_num=1,
        type="cover",
        visual_json={
            "type": "cover",
            "layout": "cover",
            "logo_policy": {"use_as_scene_asset": True, "show_logo": False},
        },
        reference_images=[],
        project=SimpleNamespace(reference_images=[logo], selected_template_recommendations=None),
    )

    refs = _load_reference_images(slide)

    assert [r["role"] for r in refs] == ["logo"]


def test_logo_policy_skips_immersive_hero_pages(tmp_path):
    logo_path = tmp_path / "logo.png"
    template_path = tmp_path / "template.png"
    for path in (logo_path, template_path):
        Image.new("RGB", (8, 8), "white").save(path)

    logo = SimpleNamespace(
        id="logo-1",
        role="logo",
        slide_id=None,
        file_path=str(logo_path),
        process_mode="original",
    )
    project = SimpleNamespace(
        reference_images=[logo],
        selected_template_recommendations={"content": {"file_path": str(template_path)}},
    )
    slide = SimpleNamespace(
        page_num=5,
        type="hero",
        visual_json={"type": "hero", "layout": "hero"},
        reference_images=[],
        project=project,
    )

    refs = _load_reference_images(slide)

    assert "logo" not in [r["role"] for r in refs]


def test_project_refs_follow_page_logo_policy():
    logo = SimpleNamespace(id="logo-1", role="logo", slide_id=None, file_path="/tmp/logo.png", process_mode="original")
    project = SimpleNamespace(reference_images=[logo])

    content_refs = slides_api._project_refs_for_prompt(
        project,
        [],
        {"page_num": 2, "type": "content", "layout": "content_split"},
    )
    hero_refs = slides_api._project_refs_for_prompt(
        project,
        [],
        {"page_num": 3, "type": "hero", "layout": "hero"},
    )

    assert content_refs == []
    assert hero_refs == []


def test_finetune_loads_project_product_asset_when_requested(tmp_path):
    base_path = tmp_path / "base.png"
    asset_path = tmp_path / "product.png"
    for path in (base_path, asset_path):
        Image.new("RGB", (8, 8), "white").save(path)

    visual_asset = SimpleNamespace(
        id="asset-1",
        role="visual_asset",
        slide_id=None,
        file_path=str(asset_path),
        process_mode="crop",
        asset_kind="product",
        asset_name="胡姬花花生油瓶",
        usage_note=None,
        asset_analysis={"subject": "胡姬花古法花生油 5L 瓶装"},
    )
    project = SimpleNamespace(reference_images=[visual_asset], selected_template_recommendations=None)
    slide = SimpleNamespace(
        page_num=15,
        type="content",
        visual_json={
            "finetune_base_image_path": str(base_path),
            "finetune_attachment_ids": [],
            "finetune_visual_asset_ids": ["asset-1"],
        },
        reference_images=[],
        project=project,
    )

    refs = _load_reference_images(slide)

    assert [r["role"] for r in refs] == ["finetune_base", "visual_asset"]
    assert refs[1]["asset_name"] == "胡姬花花生油瓶"
    assert refs[1]["process_mode"] == "crop"
    assert refs[1]["asset_analysis"]["subject"] == "胡姬花古法花生油 5L 瓶装"


def test_direct_finetune_prompt_distinguishes_project_visual_assets():
    slide = Slide(
        page_num=15,
        content_json={
            "text_content": {
                "headline": "全年脉冲式投放",
                "subhead": "全年不断线",
                "body": "品牌宣传说明",
            }
        },
    )

    prompt = slides_api._build_direct_finetune_prompt(
        slide,
        "把油瓶换成我上传的核心资产",
        attachment_count=0,
        project_visual_asset_count=1,
    )

    assert "protected project product/material assets" in prompt
    assert "authoritative source" in prompt
    assert "conflicting brand/product already visible" in prompt


def test_direct_finetune_prompt_scopes_selected_regions():
    slide = Slide(
        page_num=3,
        content_json={
            "text_content": {
                "headline": "项目进展",
                "subhead": "重点区域需要更醒目",
                "body": "保持其他信息不变",
            }
        },
    )

    prompt = slides_api._build_direct_finetune_prompt(
        slide,
        "把框选位置的数字改成红色",
        regions=[
            {
                "id": "region-1",
                "label": "框选 1",
                "bbox": {"x": 0.125, "y": 0.25, "width": 0.5, "height": 0.2},
            }
        ],
    )

    assert "Selected edit regions" in prompt
    assert "Region 1" in prompt
    assert "x=12.5%" in prompt
    assert "width=50.0%" in prompt
    assert "Keep everything outside the selected region" in prompt


def test_direct_finetune_regions_pass_mask_and_restore_unselected_pixels(tmp_path, monkeypatch):
    calls = []
    base = Image.new("RGB", (100, 60), "black")
    base.putpixel((90, 50), (12, 34, 56))

    def fake_generate_slide_image(
        prompt,
        reference_images=None,
        resolution="4K",
        aspect_ratio="16:9",
        project_id=None,
        edit_mask=None,
    ):
        calls.append({
            "prompt": prompt,
            "reference_count": len(reference_images or []),
            "edit_mask": edit_mask.copy() if edit_mask else None,
        })
        return Image.new("RGB", base.size, "red")

    monkeypatch.setattr(generation_pipeline, "generate_slide_image", fake_generate_slide_image)

    slide = Slide(
        page_num=1,
        prompt_text="direct finetune prompt",
        visual_json={
            "finetune_base_image_path": str(tmp_path / "base.png"),
            "finetune_regions": [
                {
                    "id": "region-1",
                    "label": "框选 1",
                    "bbox": {"x": 0.1, "y": 0.2, "width": 0.2, "height": 0.2},
                }
            ],
        },
    )
    ref_data = [
        {
            "role": "finetune_base",
            "label": "Current Slide Image",
            "process_mode": "original",
            "image": base.copy(),
            "file_path": str(tmp_path / "base.png"),
        }
    ]

    result = _generate_one_slide(slide, "project-1", str(tmp_path), ref_data)

    assert result["error"] is None
    assert len(calls) == 1
    mask = calls[0]["edit_mask"]
    assert mask is not None
    assert mask.mode == "RGBA"
    assert mask.size == base.size
    assert mask.getpixel((15, 15))[3] == 0
    assert mask.getpixel((90, 50))[3] == 255

    final_img = Image.open(result["image_path"]).convert("RGB")
    assert final_img.getpixel((15, 15)) == (255, 0, 0)
    assert final_img.getpixel((90, 50)) == (12, 34, 56)


def test_direct_finetune_regions_add_visual_guide_reference(tmp_path):
    base_path = tmp_path / "base.png"
    Image.new("RGB", (100, 60), "black").save(base_path)
    slide = Slide(
        page_num=1,
        visual_json={
            "finetune_base_image_path": str(base_path),
            "finetune_regions": [
                {
                    "id": "region-1",
                    "label": "框选 1",
                    "bbox": {"x": 0.1, "y": 0.2, "width": 0.2, "height": 0.2},
                }
            ],
        },
    )

    refs = _load_reference_images(slide)

    assert [ref.get("role") for ref in refs[:2]] == ["finetune_base", "finetune_region_guide"]
    guide = refs[1]["image"].convert("RGB")
    assert guide.getpixel((10, 12)) == (255, 0, 0)
    assert guide.getpixel((90, 50)) == (0, 0, 0)


def test_generate_one_slide_uses_single_pass_by_default(tmp_path, monkeypatch):
    calls = []

    def fake_generate_slide_image(prompt, reference_images=None, resolution="4K", aspect_ratio="16:9", project_id=None):
        calls.append({
            "prompt": prompt,
            "reference_count": len(reference_images or []),
            "resolution": resolution,
            "aspect_ratio": aspect_ratio,
        })
        color = "green" if len(calls) == 2 else "blue"
        return Image.new("RGB", (16, 9), color)

    monkeypatch.setattr(generation_pipeline, "generate_slide_image", fake_generate_slide_image)

    slide = Slide(
        page_num=4,
        prompt_text="draft prompt",
        content_json={"text_content": {"headline": "场景展示"}},
        visual_json={"visual_asset_usage": {"asset-1": "放在中央场景区"}},
    )
    ref_data = [
        {
            "id": "asset-1",
            "role": "visual_asset",
            "process_mode": "blend",
            "asset_kind": "scene",
            "asset_name": "终端货架",
            "image": Image.new("RGB", (8, 8), "white"),
        }
    ]

    result = _generate_one_slide(slide, "project-1", str(tmp_path), ref_data)

    assert result["error"] is None
    assert len(calls) == 1
    assert calls[0]["prompt"] == "draft prompt"
    assert calls[0]["reference_count"] == 1

    final_img = Image.open(result["image_path"])
    assert final_img.getpixel((0, 0)) == (0, 0, 255)
    assert not (tmp_path / "project-1" / "slide_04_base.png").exists()


def test_overlay_route_material_ref_does_not_trigger_refinement(tmp_path, monkeypatch):
    calls = []

    def fake_generate_slide_image(prompt, reference_images=None, resolution="4K", aspect_ratio="16:9", project_id=None):
        calls.append({"prompt": prompt, "reference_count": len(reference_images or [])})
        return Image.new("RGB", (16, 9), "blue")

    monkeypatch.setattr(generation_pipeline, "generate_slide_image", fake_generate_slide_image)

    slide = Slide(page_num=4, prompt_text="draft prompt", visual_json={})
    ref_data = [
        {
            "id": "asset-overlay",
            "role": "visual_asset",
            "process_mode": "original",
            "asset_route_mode": "overlay",
            "asset_kind": "material",
            "asset_name": "Dashboard 截图",
            "image": Image.new("RGB", (8, 8), "white"),
        }
    ]

    result = _generate_one_slide(slide, "project-overlay-route", str(tmp_path), ref_data)

    assert result["error"] is None
    assert len(calls) == 1
    assert calls[0]["reference_count"] == 1
    assert not (tmp_path / "project-overlay-route" / "slide_04_base.png").exists()


def test_generate_one_slide_uses_reference_fidelity_product_refinement_pass(tmp_path, monkeypatch):
    calls = []

    def fake_generate_slide_image(prompt, reference_images=None, resolution="4K", aspect_ratio="16:9", project_id=None):
        calls.append({
            "prompt": prompt,
            "reference_count": len(reference_images or []),
            "resolution": resolution,
            "aspect_ratio": aspect_ratio,
        })
        color = "green" if len(calls) == 2 else "blue"
        return Image.new("RGB", (16, 9), color)

    monkeypatch.setattr(generation_pipeline, "generate_slide_image", fake_generate_slide_image)

    slide = Slide(
        page_num=4,
        prompt_text="draft prompt",
        content_json={"text_content": {"headline": "产品展示"}},
        visual_json={"visual_asset_usage": {"asset-1": "Place the uploaded product image at center-right."}},
    )
    ref_data = [
        {
            "id": "asset-1",
            "role": "visual_asset",
            "process_mode": "crop",
            "asset_kind": "product",
            "asset_name": "胡姬花花生油瓶",
            "file_path": "/tmp/uploads/huji-product.png",
            "asset_analysis": {
                "subject": "胡姬花古法花生油 5L 瓶装",
                "identity_elements": ["红色瓶盖", "黄色油液", "正面品牌标签"],
                "must_not_change": ["中文品牌名", "瓶身比例"],
            },
            "image": Image.new("RGB", (8, 8), "white"),
        }
    ]

    result = _generate_one_slide(slide, "project-1", str(tmp_path), ref_data)

    assert result["error"] is None
    assert len(calls) == 2
    assert "Reference fidelity" in calls[0]["prompt"]
    assert "FIRST PASS" not in calls[0]["prompt"]
    assert "Use the uploaded product/material image as the source" in calls[0]["prompt"]
    assert "visible text, labels, logos, and small markings" in calls[0]["prompt"]
    assert "fidelity can be approximate" not in calls[0]["prompt"]
    assert "simplify it into a generic placeholder" not in calls[0]["prompt"]
    assert calls[0]["reference_count"] == 1
    assert "第2张图是要精修融合进去的参考素材" in calls[1]["prompt"]
    assert "尽量1:1保留第2张图的所有可见细节、人物身份、文字和标识" in calls[1]["prompt"]
    assert "不要改变第1张图中的其它任何画面元素" in calls[1]["prompt"]
    assert "胡姬花古法花生油" not in calls[1]["prompt"]
    assert "/tmp/uploads/huji-product.png" not in calls[1]["prompt"]
    assert calls[1]["reference_count"] == 2

    final_img = Image.open(result["image_path"])
    assert final_img.getpixel((0, 0)) == (0, 128, 0)
    assert (tmp_path / "project-1" / "slide_04_base.png").exists()


def test_generate_one_slide_falls_back_to_base_when_refinement_fails(tmp_path, monkeypatch):
    calls = []

    def fake_generate_slide_image(prompt, reference_images=None, resolution="4K", aspect_ratio="16:9", project_id=None):
        calls.append({"prompt": prompt, "reference_count": len(reference_images or [])})
        if len(calls) == 2:
            raise RuntimeError("refinement API unavailable")
        return Image.new("RGB", (16, 9), "blue")

    monkeypatch.setattr(generation_pipeline, "generate_slide_image", fake_generate_slide_image)

    slide = Slide(
        id="slide-refine-fallback",
        page_num=5,
        prompt_text="draft prompt",
        content_json={"text_content": {"headline": "产品展示"}},
        visual_json={},
    )
    ref_data = [
        {
            "id": "asset-1",
            "role": "visual_asset",
            "asset_kind": "product",
            "asset_name": "产品图",
            "file_path": "/tmp/uploads/product.png",
            "image": Image.new("RGB", (8, 8), "white"),
        }
    ]

    result = _generate_one_slide(slide, "project-refine-fallback", str(tmp_path), ref_data)

    assert result["error"] is None
    assert len(calls) == 2
    assert (tmp_path / "project-refine-fallback" / "slide_05_base.png").exists()
    final_img = Image.open(result["image_path"])
    assert final_img.getpixel((0, 0)) == (0, 0, 255)


def test_product_refinement_pass_accepts_multiple_product_refs(tmp_path, monkeypatch):
    calls = []

    def fake_generate_slide_image(prompt, reference_images=None, resolution="4K", aspect_ratio="16:9", project_id=None):
        calls.append({"prompt": prompt, "reference_count": len(reference_images or [])})
        return Image.new("RGB", (16, 9), "green" if len(calls) == 2 else "blue")

    monkeypatch.setattr(generation_pipeline, "generate_slide_image", fake_generate_slide_image)

    slide = Slide(page_num=6, prompt_text="draft prompt", visual_json={})
    ref_data = [
        {
            "id": "asset-1",
            "role": "visual_asset",
            "asset_kind": "product",
            "asset_name": "产品 A",
            "file_path": "/tmp/uploads/product-a.png",
            "image": Image.new("RGB", (8, 8), "white"),
        },
        {
            "id": "asset-2",
            "role": "visual_asset",
            "asset_kind": "material",
            "asset_name": "产品 B",
            "file_path": "/tmp/uploads/product-b.png",
            "image": Image.new("RGB", (8, 8), "black"),
        },
    ]

    result = _generate_one_slide(slide, "project-2", str(tmp_path), ref_data)

    assert result["error"] is None
    assert len(calls) == 2
    assert "第2-3张图是要精修融合进去的参考素材" in calls[1]["prompt"]
    assert "分别把第2-3张参考图" in calls[1]["prompt"]
    assert "产品 A" not in calls[1]["prompt"]
    assert "产品 B" not in calls[1]["prompt"]
    assert "/tmp/uploads/product-a.png" not in calls[1]["prompt"]
    assert "/tmp/uploads/product-b.png" not in calls[1]["prompt"]
    assert "油瓶" not in calls[1]["prompt"]
    assert "瓶盖" not in calls[1]["prompt"]
    assert "提手" not in calls[1]["prompt"]
    assert calls[1]["reference_count"] == 3


def test_double_blend_pass_ignores_page_reference_refs(tmp_path, monkeypatch):
    calls = []

    def fake_generate_slide_image(prompt, reference_images=None, resolution="4K", aspect_ratio="16:9", project_id=None):
        calls.append({"prompt": prompt, "reference_count": len(reference_images or [])})
        return Image.new("RGB", (16, 9), "green" if len(calls) == 2 else "blue")

    monkeypatch.setattr(generation_pipeline, "generate_slide_image", fake_generate_slide_image)

    slide = Slide(page_num=8, prompt_text="draft prompt", visual_json={})
    ref_data = [
        {
            "id": "page-ref-1",
            "role": "content_ref",
            "asset_name": "原 PPT 活动海报",
            "file_path": "/tmp/uploads/page-poster.png",
            "asset_route_mode": "double_blend",
            "image": Image.new("RGB", (8, 8), "white"),
        },
    ]

    result = _generate_one_slide(slide, "project-3", str(tmp_path), ref_data)

    assert result["error"] is None
    assert len(calls) == 1
    assert calls[0]["prompt"] == "draft prompt"
    assert calls[0]["reference_count"] == 1


def test_image_cache_key_includes_reference_image_content():
    white_ref = Image.new("RGB", (8, 8), "white")
    black_ref = Image.new("RGB", (8, 8), "black")

    assert (
        _cache_key("same prompt", [white_ref], "4K", "16:9")
        != _cache_key("same prompt", [black_ref], "4K", "16:9")
    )


def test_product_visual_asset_recall_uses_generic_product_page_when_single_asset():
    recalled = _recall_visual_assets_for_page(
        {
            "page_num": 4,
            "type": "content",
            "text_content": {
                "headline": "终端货架：用产品实物建立购买记忆",
                "body": "包装、瓶身和导购体验台是本页核心画面。",
            },
        },
        [
            {
                "id": "asset-1",
                "name": "胡姬花花生油瓶",
                "kind": "product",
                "usage_note": "用于产品、包装、货架、终端展示页",
                "analysis_summary": "keywords=胡姬花、花生油、油瓶、包装、货架",
            }
        ],
    )

    assert recalled[0]["id"] == "asset-1"


def test_project_visual_assets_for_planning_filters_legacy_ppt_scenes(tmp_path):
    scene_path = tmp_path / "scene.png"
    product_path = tmp_path / "product.png"
    map_path = tmp_path / "map.png"
    for path in (scene_path, product_path, map_path):
        Image.new("RGB", (8, 8), "white").save(path)
    project = SimpleNamespace(reference_images=[
        SimpleNamespace(
            id="legacy-scene",
            role="visual_asset",
            slide_id=None,
            file_path=str(scene_path),
            asset_name="legacy scene",
            asset_kind="scene",
            process_mode="blend",
            usage_note="",
            asset_analysis={
                "source_document": "source.pptx",
                "area_ratio": 0.42,
                "source_slide_text": "校园风景与氛围背景",
            },
        ),
        SimpleNamespace(
            id="core-product",
            role="visual_asset",
            slide_id=None,
            file_path=str(product_path),
            asset_name="core product",
            asset_kind="product",
            process_mode="crop",
            usage_note="",
            asset_analysis={
                "source_document": "source.pptx",
                "selection_tier": "core_global",
                "importance_score": 35,
            },
        ),
        SimpleNamespace(
            id="legacy-map",
            role="visual_asset",
            slide_id=None,
            file_path=str(map_path),
            asset_name="legacy map",
            asset_kind="other",
            process_mode="blend",
            usage_note="",
            asset_analysis={
                "source_document": "source.pptx",
                "area_ratio": 0.12,
                "source_slide_text": "渠道地图、数据图表与业务流程",
            },
        ),
    ])

    assets = slides_api._project_visual_assets_for_planning(project)

    assert [asset["id"] for asset in assets] == ["core-product"]


def test_missing_project_assets_do_not_create_prompt_references(tmp_path):
    existing_path = tmp_path / "asset.png"
    Image.new("RGB", (8, 8), "white").save(existing_path)
    missing_path = tmp_path / "missing.png"
    existing = SimpleNamespace(
        id="asset-ok",
        role="visual_asset",
        slide_id=None,
        file_path=str(existing_path),
        asset_name="existing product",
        asset_kind="product",
        process_mode="crop",
        usage_note="",
        asset_analysis={"selection_tier": "manual"},
    )
    missing = SimpleNamespace(
        id="asset-missing",
        role="visual_asset",
        slide_id=None,
        file_path=str(missing_path),
        asset_name="missing product",
        asset_kind="product",
        process_mode="crop",
        usage_note="",
        asset_analysis={"selection_tier": "manual"},
    )
    project = SimpleNamespace(reference_images=[existing, missing])

    planning_assets = slides_api._project_visual_assets_for_planning(project)
    prompt_refs = slides_api._project_refs_for_prompt(
        project,
        ["asset-ok", "asset-missing"],
        {"visual_asset_ids": ["asset-ok", "asset-missing"]},
    )

    assert [asset["id"] for asset in planning_assets] == ["asset-ok"]
    assert [ref["id"] for ref in prompt_refs] == ["asset-ok"]


def test_missing_template_does_not_create_prompt_reference(tmp_path):
    existing_path = tmp_path / "template.png"
    Image.new("RGB", (8, 8), "white").save(existing_path)
    project = SimpleNamespace(reference_images=[
        SimpleNamespace(id="template-ok", role="template", file_path=str(existing_path), process_mode="blend"),
        SimpleNamespace(id="template-missing", role="template", file_path=str(tmp_path / "missing-template.png"), process_mode="blend"),
    ])

    refs = slides_api._project_template_refs_for_prompt(project)

    assert [ref["id"] for ref in refs] == ["template-ok"]


def test_visual_plan_auto_adds_recalled_product_asset_when_llm_misses(monkeypatch):
    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**_kwargs):
                    return SimpleNamespace(
                        choices=[
                            SimpleNamespace(
                                message=SimpleNamespace(
                                    content=(
                                        '{"7": {'
                                        '"visual_evidence": "终端货架、产品包装和体验台", '
                                        '"visual_summary": "产品货架体验画面", '
                                        '"visual_description": "以终端货架和体验台组织画面，突出产品实物展示。", '
                                        '"visual_asset_ids": [], '
                                        '"visual_asset_usage": {}'
                                        '}}'
                                    )
                                )
                            )
                        ]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #FFFFFF",
    )

    plan = _do_generate_visual_plan(
        content_plan=[
            {
                "page_num": 7,
                "type": "content",
                "text_content": {
                    "headline": "终端货架体验",
                    "body": "用产品实物、包装和小油瓶建立购买记忆",
                },
            }
        ],
        global_visual_assets=[
            {
                "id": "asset-1",
                "name": "胡姬花花生油瓶",
                "kind": "product",
                "process_mode": "crop",
                "usage_note": "用于产品、包装、货架、终端展示页",
                "analysis_summary": "keywords=胡姬花、花生油、油瓶、包装、货架",
            }
        ],
    )

    assert plan[0]["visual_asset_ids"] == ["asset-1"]
    assert "uploaded product image" in plan[0]["visual_asset_usage"]["asset-1"]
    assert "胡姬花花生油瓶" not in plan[0]["visual_asset_usage"]["asset-1"]


def test_visual_plan_forces_logo_policy_off_when_project_has_no_logo(monkeypatch):
    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**_kwargs):
                    return SimpleNamespace(
                        choices=[
                            SimpleNamespace(
                                message=SimpleNamespace(
                                    content=(
                                        '{"1": {'
                                        '"visual_evidence": "竞技场封面主视觉", '
                                        '"visual_summary": "暗红竞技场封面", '
                                        '"visual_description": "竞技场主视觉铺底；右侧预留品牌标识区域。", '
                                        '"visual_asset_ids": [], '
                                        '"visual_asset_usage": {}, '
                                        '"logo_policy": {"show_logo": true, "placement": "title-block-center", "scale": "large", "use_as_scene_asset": false}'
                                        '}}'
                                    )
                                )
                            )
                        ]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #FFFFFF",
    )

    plan = _do_generate_visual_plan(
        content_plan=[
            {
                "page_num": 1,
                "type": "cover",
                "text_content": {"headline": "古罗马角斗士", "body": ""},
            }
        ],
        has_project_logo=False,
    )

    assert plan[0]["logo_policy"]["show_logo"] is False
    assert plan[0]["logo_policy"]["placement"] == "top-right"
    assert plan[0]["logo_policy"]["use_as_scene_asset"] is False
    assert "品牌标识" not in plan[0]["visual_description"]


def test_visual_plan_locks_manual_pins_and_route_modes_when_llm_selects_other_asset(monkeypatch):
    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**_kwargs):
                    return SimpleNamespace(
                        choices=[
                            SimpleNamespace(
                                message=SimpleNamespace(
                                    content=(
                                        '{"1": {'
                                        '"visual_evidence": "校园驿站场景", '
                                        '"visual_summary": "驿站场景", '
                                        '"visual_description": "用校园驿站场景组织画面。", '
                                        '"visual_asset_ids": ["auto-1"], '
                                        '"visual_asset_usage": {"auto-1": "右侧补充"}'
                                        '}}'
                                    )
                                )
                            )
                        ]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #FFFFFF",
    )

    plan = _do_generate_visual_plan(
        content_plan=[
            {
                "page_num": 1,
                "type": "content",
                "text_content": {"headline": "校园驿站", "body": "取件机与回收场景"},
                "manual_visual_asset_ids": ["manual-1"],
                "manual_visual_asset_usage": {"manual-1": "用户指定放在左侧"},
                "asset_route_modes": {"manual-1": "blend", "auto-1": "double_blend"},
            }
        ],
        global_visual_assets=[
            {"id": "manual-1", "name": "指定驿站图", "kind": "scene", "analysis_summary": "校园驿站"},
            {"id": "auto-1", "name": "自动候选图", "kind": "scene", "analysis_summary": "取件机"},
        ],
    )

    assert plan[0]["manual_visual_asset_ids"] == ["manual-1"]
    assert plan[0]["visual_asset_ids"] == ["manual-1"]
    assert plan[0]["visual_asset_usage"]["manual-1"] == "用户指定放在左侧"
    assert plan[0]["asset_route_modes"] == {"manual-1": "blend"}


def test_prototype_without_selected_pages_samples_first_three_and_ignores_seed_flags():
    slides = [
        Slide(page_num=1, visual_json={"is_seed_recommended": False}),
        Slide(page_num=2, visual_json={"is_seed_recommended": False}),
        Slide(page_num=3, visual_json={"is_seed_recommended": False}),
        Slide(page_num=4, visual_json={"is_seed_recommended": True}),
    ]

    assert _resolve_generation_page_nums(slides, None, True) == [1, 2, 3]
    assert _resolve_generation_page_nums(slides, [2, 4], True) == [2, 4]
    assert _resolve_generation_page_nums(slides, None, False) is None


def test_confirm_prototype_autofills_missing_prompts_before_full_generation(monkeypatch):
    db = make_session()
    project = Project(
        title="Prototype confirm",
        status="prototype_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    db.add_all(
        [
            Slide(
                project_id=project.id,
                page_num=1,
                status="completed",
                content_json={"page_num": 1, "text_content": {"headline": "已打样"}},
                visual_json={"visual_description": "sample visual"},
                prompt_text="sample prompt",
                image_path="/outputs/p1.png",
            ),
            Slide(
                project_id=project.id,
                page_num=2,
                status="pending",
                content_json={"page_num": 2, "text_content": {"headline": "待生成"}},
                visual_json={"visual_description": "pending visual"},
                prompt_text=None,
            ),
            Slide(
                project_id=project.id,
                page_num=3,
                status="pending",
                content_json={"page_num": 3, "text_content": {"headline": "已有 prompt"}},
                visual_json={"visual_description": "ready visual"},
                prompt_text="ready prompt",
            ),
        ]
    )
    db.commit()

    def fake_generate_prompt_for_page(**kwargs):
        content_text = kwargs.get("content_text") or {}
        return (
            "Visible Text:\n"
            f"Headline: \"{content_text.get('headline', '')}\"\n"
            "\n\nStyle: test"
        )

    monkeypatch.setattr(slides_api, "generate_prompt_for_page", fake_generate_prompt_for_page)
    monkeypatch.setattr(slides_api, "store_current_provider_credentials", lambda _redis: "cred")
    monkeypatch.setattr(slides_api.redis_client, "set", lambda *_args, **_kwargs: True)
    captured = {}

    def fake_delay(project_id, page_nums, **kwargs):
        captured["project_id"] = project_id
        captured["page_nums"] = page_nums
        captured["kwargs"] = kwargs
        return SimpleNamespace(id="task-1")

    monkeypatch.setattr(slides_api.generate_slides_task, "delay", fake_delay)

    result = slides_api.confirm_prototype(project.id, db=db)
    refreshed = db.query(Slide).filter(Slide.project_id == project.id).order_by(Slide.page_num).all()

    assert result["page_nums"] == [2, 3]
    assert captured["page_nums"] == [2, 3]
    assert refreshed[1].prompt_text and "Headline: \"待生成\"" in refreshed[1].prompt_text
    assert refreshed[1].status == "prompt_ready"
    assert refreshed[2].prompt_text and "Headline: \"已有 prompt\"" in refreshed[2].prompt_text
    assert refreshed[2].status == "prompt_ready"


def test_enqueue_generation_task_ignores_redis_tracking_failure(monkeypatch):
    db = make_session()
    project = Project(title="Redis tracking failure", status="prompt_ready", content_plan_confirmed=True)
    db.add(project)
    db.flush()
    run = slides_api.create_project_run(
        db,
        project.id,
        kind="image_generation",
        stage="batch_generation",
        total_count=1,
        message="queued",
    )
    db.commit()

    monkeypatch.setattr(slides_api, "store_current_provider_credentials", lambda _redis: None)

    def fake_delay(project_id, page_nums, **kwargs):
        return SimpleNamespace(id="task-after-redis-failure")

    def fail_tracking_set(*_args, **_kwargs):
        raise RuntimeError("redis unavailable after dispatch")

    monkeypatch.setattr(slides_api.generate_slides_task, "delay", fake_delay)
    monkeypatch.setattr(slides_api.redis_client, "set", fail_tracking_set)

    task = slides_api._enqueue_generation_task(db, project.id, [1], run=run)
    refreshed_run = db.query(run.__class__).filter(run.__class__.id == run.id).first()

    assert task.id == "task-after-redis-failure"
    assert refreshed_run.task_id == "task-after-redis-failure"
    assert refreshed_run.status == "queued"


def test_confirm_content_plan_advances_backend_stage():
    db = make_session()
    project = Project(title="Confirm flow", status="planning", content_plan_confirmed=False)
    db.add(project)
    db.flush()
    db.add(Slide(project_id=project.id, page_num=1, status="pending", content_json={"page_num": 1}))
    db.commit()

    updated = projects_api.update_project(
        project.id,
        ProjectUpdate(content_plan_confirmed=True),
        db=db,
    )

    assert updated.content_plan_confirmed is True
    assert updated.status == "visual_ready"


def test_content_edit_preserves_confirmed_workflow_and_existing_outputs():
    db = make_session()
    project = Project(
        title="Content edit",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Old"},
        style_proposal={"proposals": [{"name": "Old"}]},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        content_json={"page_num": 1, "text_content": {"headline": "旧标题"}},
        visual_json={"visual_description": "old visual"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.commit()

    slides_api.update_slide_content(
        project.id,
        slides_api.UpdateContentRequest(
            page_num=1,
            content_json={"text_content": {"headline": "新标题"}},
        ),
        db=db,
    )
    refreshed_project = db.query(Project).filter(Project.id == project.id).first()
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).first()

    assert refreshed_project.status == "prompt_ready"
    assert refreshed_project.content_plan_confirmed is True
    assert refreshed_project.selected_style == {"name": "Old"}
    assert refreshed_project.style_proposal == {"proposals": [{"name": "Old"}]}
    assert refreshed_slide.content_json["text_content"]["headline"] == "新标题"
    assert refreshed_slide.visual_json["visual_description"] == "old visual"
    assert artifact_stale(refreshed_slide.visual_json) == {"content": True}
    assert refreshed_slide.prompt_text == "old prompt"
    assert refreshed_slide.image_path == "/tmp/old.png"
    assert refreshed_slide.status == "completed"


def test_content_edit_drops_unchanged_visual_guidance_when_text_changes():
    db = make_session()
    project = Project(
        title="Content edit visual suggestion",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    old_visual_suggestion = "标题下方并排两个金色小序号「1.1」「1.2」，下方各放一个小标题。"
    old_visual_requirements = [
        {"directive": "并排显示 1.1 和 1.2 两个章节入口", "diagram_labels": ["1.1", "1.2"]}
    ]
    slide = Slide(
        project_id=project.id,
        page_num=7,
        status="completed",
        content_json={
            "page_num": 7,
            "type": "section",
            "text_content": {
                "headline": "什么变了？决策不再只发生在人脑里",
                "body": "- 1.1 从人脑到系统\n- 1.2 从工具到代理",
            },
            "visual_suggestion": old_visual_suggestion,
            "visual_requirements": old_visual_requirements,
        },
        visual_json={"visual_description": "old visual"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.commit()

    slides_api.update_slide_content(
        project.id,
        slides_api.UpdateContentRequest(
            page_num=7,
            slide_id=slide.id,
            content_json={
                "text_content": {
                    "headline": "什么变了？决策不再只发生在人脑里",
                    "body": "",
                },
                "visual_suggestion": old_visual_suggestion,
                "visual_requirements": old_visual_requirements,
            },
        ),
        db=db,
    )
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).first()

    assert refreshed_slide.content_json["text_content"]["body"] == ""
    assert refreshed_slide.content_json.get("visual_suggestion", "") == ""
    assert refreshed_slide.content_json.get("visual_requirements", []) == []
    assert artifact_stale(refreshed_slide.visual_json) == {"content": True}


def test_content_save_drops_stale_case_visual_guidance_even_without_text_change():
    db = make_session()
    project = Project(
        title="Content save stale visual suggestion",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    stale_visual_suggestion = (
        '顶部"私人买手对话框"还原场景，下方三栏讲山姆/胖东来/Costco，'
        '底部一行大字"个人外脑"。'
    )
    body_text = (
        '- 过去搜索"防晒霜"：自己翻几百个商品、几千条评价\n'
        '- 现在直接说："我要去三亚冲浪，预算 200 内，敏感肌，帮我推荐三款防水防晒霜"'
    )
    slide = Slide(
        project_id=project.id,
        page_num=9,
        status="completed",
        content_json={
            "page_num": 9,
            "type": "content",
            "text_content": {
                "headline": "消费者雇了一个“全知全能”的私人买手",
                "body": body_text,
            },
            "visual_suggestion": stale_visual_suggestion,
        },
        visual_json={"visual_description": "old visual"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.commit()

    slides_api.update_slide_content(
        project.id,
        slides_api.UpdateContentRequest(
            page_num=9,
            slide_id=slide.id,
            content_json={
                "text_content": {
                    "headline": "消费者雇了一个“全知全能”的私人买手",
                    "body": body_text,
                },
                "visual_suggestion": stale_visual_suggestion,
            },
        ),
        db=db,
    )
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).first()

    assert refreshed_slide.content_json["text_content"]["body"] == body_text
    assert refreshed_slide.content_json.get("visual_suggestion", "") == ""
    assert artifact_stale(refreshed_slide.visual_json) == {"content": True}


def test_visual_plan_regeneration_ignores_legacy_visual_suggestion(monkeypatch):
    db = make_session()
    project = Project(
        title="Content-only visual regeneration",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=5,
        status="completed",
        content_json={
            "page_num": 5,
            "type": "content",
            "text_content": {
                "headline": "两张图看清这个时代",
                "subhead": "时间线",
                "body": "时代线——经典营销 → 移动互联网 → AI 时代",
            },
            "visual_suggestion": "左半页三段时间轴，右半页三方博弈图，下方用三条横线列张力。",
        },
        visual_json={"visual_description": "old visual"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.commit()

    captured = {}

    def fake_generate_visual_plan(**kwargs):
        captured["content_plan"] = kwargs["content_plan"]
        return [
            {
                "page_num": 5,
                "type": "content",
                "layout": "content",
                "visual_description": "只做三段时间线。",
            }
        ]

    monkeypatch.setattr(slides_api, "generate_visual_plan", fake_generate_visual_plan)

    slides_api.create_visual_plan(project.id, slides_api.PageNumsRequest(), db=db)
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).first()

    assert captured["content_plan"][0].get("visual_suggestion") == ""
    assert captured["content_plan"][0].get("visual_requirements") == []
    assert refreshed_slide.content_json.get("visual_suggestion") == ""
    assert refreshed_slide.prompt_text is None
    assert refreshed_slide.image_path is None


def test_style_reference_upload_after_confirmation_stays_in_visual_stage(tmp_path, monkeypatch):
    db = make_session()
    project = Project(
        title="Style upload",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Old"},
        style_proposal={"proposals": [{"name": "Old"}]},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        content_json={"page_num": 1},
        visual_json={"visual_description": "old visual"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.commit()
    monkeypatch.setattr(slides_api.settings, "UPLOAD_DIR", str(tmp_path))

    slides_api.upload_file(
        project.id,
        png_upload("style.png"),
        role="style_ref",
        slide_id=None,
        process_mode=None,
        asset_name=None,
        asset_kind=None,
        usage_note=None,
        db=db,
    )
    refreshed_project = db.query(Project).filter(Project.id == project.id).first()
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).first()

    assert refreshed_project.status == "prompt_ready"
    assert refreshed_project.content_plan_confirmed is True
    assert refreshed_project.selected_style == {"name": "Old"}
    assert refreshed_project.style_proposal == {"proposals": [{"name": "Old"}]}
    assert refreshed_slide.visual_json["visual_description"] == "old visual"
    assert artifact_stale(refreshed_slide.visual_json) == {"content": True}
    assert refreshed_slide.prompt_text == "old prompt"
    assert refreshed_slide.image_path == "/tmp/old.png"
    assert refreshed_slide.status == "completed"


def test_visual_edit_invalidates_prompt_and_image_only():
    db = make_session()
    project = Project(title="Visual edit", status="completed", content_plan_confirmed=True, selected_style={"name": "Brand"})
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        content_json={"page_num": 1, "text_content": {"headline": "标题"}},
        visual_json={"visual_description": "old visual"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.commit()

    slides_api.update_slide_visual(
        project.id,
        slides_api.UpdateVisualRequest(
            page_num=1,
            visual_json={"visual_description": "new visual"},
        ),
        db=db,
    )
    refreshed_project = db.query(Project).filter(Project.id == project.id).first()
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).first()

    assert refreshed_project.status == "completed"
    assert refreshed_project.content_plan_confirmed is True
    assert refreshed_project.selected_style == {"name": "Brand"}
    assert refreshed_slide.visual_json["visual_description"] == "new visual"
    assert artifact_stale(refreshed_slide.visual_json) == {"visual": True}
    assert refreshed_slide.prompt_text == "old prompt"
    assert refreshed_slide.image_path == "/tmp/old.png"
    assert refreshed_slide.status == "completed"


def test_visual_edit_can_remove_auto_visual_asset_ids():
    db = make_session()
    project = Project(title="Visual asset removal", status="completed", content_plan_confirmed=True, selected_style={"name": "Brand"})
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        visual_json={
            "visual_description": "old visual",
            "visual_asset_ids": ["asset-keep", "asset-remove"],
            "visual_asset_usage": {"asset-keep": "保留", "asset-remove": "移除"},
        },
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.commit()

    slides_api.update_slide_visual(
        project.id,
        slides_api.UpdateVisualRequest(
            page_num=1,
            visual_json={
                "visual_asset_ids": ["asset-keep"],
                "visual_asset_usage": {"asset-keep": "保留"},
            },
        ),
        db=db,
    )
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).first()

    assert refreshed_slide.visual_json["visual_asset_ids"] == ["asset-keep"]
    assert refreshed_slide.visual_json["visual_asset_usage"] == {"asset-keep": "保留"}
    assert artifact_stale(refreshed_slide.visual_json) == {"visual": True}
    assert refreshed_slide.prompt_text == "old prompt"
    assert refreshed_slide.image_path == "/tmp/old.png"


def test_visual_edit_records_removed_auto_visual_asset_ids():
    db = make_session()
    project = Project(title="Visual asset exclusion", status="completed", content_plan_confirmed=True, selected_style={"name": "Brand"})
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        visual_json={
            "visual_description": "old visual",
            "visual_asset_ids": ["asset-keep", "asset-remove"],
            "visual_asset_usage": {"asset-keep": "保留", "asset-remove": "移除"},
        },
        prompt_text="old prompt",
    )
    db.add(slide)
    db.commit()

    slides_api.update_slide_visual(
        project.id,
        slides_api.UpdateVisualRequest(
            page_num=1,
            slide_id=slide.id,
            visual_json={
                "visual_asset_ids": ["asset-keep"],
                "visual_asset_usage": {"asset-keep": "保留"},
            },
        ),
        db=db,
    )
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).first()

    assert refreshed_slide.visual_json["visual_asset_ids"] == ["asset-keep"]
    assert refreshed_slide.visual_json["excluded_visual_asset_ids"] == ["asset-remove"]


def test_visual_merge_preserves_excluded_auto_assets():
    merged = slides_api._merge_manual_pins_into_visual_json(
        {
            "visual_asset_ids": ["asset-remove", "asset-fresh"],
            "visual_asset_usage": {"asset-remove": "自动回流", "asset-fresh": "新素材"},
        },
        {"excluded_visual_asset_ids": ["asset-remove"]},
    )

    assert merged["visual_asset_ids"] == ["asset-fresh"]
    assert merged["visual_asset_usage"] == {"asset-fresh": "新素材"}
    assert merged["excluded_visual_asset_ids"] == ["asset-remove"]


def test_visual_plan_skips_user_excluded_visual_assets(monkeypatch):
    class FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**_kwargs):
                    return SimpleNamespace(
                        choices=[
                            SimpleNamespace(
                                message=SimpleNamespace(
                                    content=(
                                        '{"1": {'
                                        '"visual_evidence": "校园驿站场景", '
                                        '"visual_summary": "驿站场景", '
                                        '"visual_description": "用校园驿站场景组织画面。", '
                                        '"visual_asset_ids": ["asset-remove"], '
                                        '"visual_asset_usage": {"asset-remove": "右侧展示"}'
                                        '}}'
                                    )
                                )
                            )
                        ]
                    )

    import app.services.visual_plan as visual_plan_module

    monkeypatch.setattr(visual_plan_module, "get_llm_client", lambda: FakeClient())
    monkeypatch.setattr(visual_plan_module, "_load_style", lambda _style_id: {"meta": {}, "body": ""})
    monkeypatch.setattr(
        visual_plan_module,
        "derive_style_pack_from_content",
        lambda _content_plan: "Style: test\nPalette: #111111, #FFFFFF",
    )

    plan = _do_generate_visual_plan(
        content_plan=[
            {
                "page_num": 1,
                "type": "content",
                "text_content": {"headline": "校园驿站", "body": "取件机与回收场景"},
                "excluded_visual_asset_ids": ["asset-remove"],
            }
        ],
        global_visual_assets=[
            {"id": "asset-remove", "name": "取件机", "kind": "scene", "analysis_summary": "取件机;校园驿站"},
        ],
    )

    assert plan[0]["excluded_visual_asset_ids"] == ["asset-remove"]
    assert plan[0]["visual_asset_ids"] == []
    assert plan[0]["visual_asset_usage"] == {}


def test_visual_edit_persists_exact_overlay_layers(tmp_path):
    db = make_session()
    project = Project(title="Visual exact overlay", status="completed", content_plan_confirmed=True, selected_style={"name": "Brand"})
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        visual_json={"visual_description": "old"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    path = tmp_path / "dashboard.png"
    Image.new("RGB", (120, 70), "white").save(path)
    asset = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(path),
        process_mode="blend",
        asset_name="Dashboard 截图",
        asset_kind="material",
        asset_analysis={"subject": "dashboard"},
    )
    db.add_all([slide, asset])
    db.commit()

    slides_api.update_slide_visual(
        project.id,
        slides_api.UpdateVisualRequest(
            page_num=1,
            visual_json={
                "visual_description": "reserve background",
                "asset_route_modes": {asset.id: "overlay"},
                "overlay_layers": [
                    {
                        "asset_id": asset.id,
                        "preset": "right-card",
                        "mode": "exact_cutout",
                        "usage_note": "原样叠加",
                    }
                ],
            },
        ),
        db=db,
    )
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).one()
    refreshed_asset = db.query(ReferenceImage).filter(ReferenceImage.id == asset.id).one()

    assert refreshed_slide.visual_json["overlay_layers"][0]["asset_id"] == asset.id
    assert refreshed_slide.visual_json["overlay_layers"][0]["fit"] == "contain"
    assert refreshed_slide.visual_json["asset_route_modes"][asset.id] == "overlay"
    assert artifact_stale(refreshed_slide.visual_json) == {"visual": True}
    assert refreshed_asset.process_mode == "original"
    assert refreshed_asset.asset_analysis["exact_overlay"] is True
    assert refreshed_slide.prompt_text == "old prompt"
    assert refreshed_slide.status == "completed"


def test_visual_edit_persists_image_slots():
    db = make_session()
    project = Project(title="Visual image slots", status="completed", content_plan_confirmed=True, selected_style={"name": "Brand"})
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        visual_json={"visual_description": "old"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.commit()

    image_slots = [
        {
            "id": "A",
            "subject": "特罗卡德罗 Trocadéro 正面铁塔远景",
            "role": "primary",
            "position": "upper-center large landscape",
            "shape": "landscape postcard",
            "linked_text": ["body_1", "特罗卡德罗 Trocadéro"],
        }
    ]
    slides_api.update_slide_visual(
        project.id,
        slides_api.UpdateVisualRequest(
            page_num=1,
            visual_json={"image_slots": image_slots},
        ),
        db=db,
    )
    refreshed_slide = db.query(Slide).filter(Slide.id == slide.id).one()

    assert refreshed_slide.visual_json["image_slots"] == image_slots
    assert artifact_stale(refreshed_slide.visual_json) == {"visual": True}
    assert refreshed_slide.prompt_text == "old prompt"
    assert refreshed_slide.image_path == "/tmp/old.png"


def test_visual_asset_upload_defaults_crop_but_honors_explicit_blend(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="Asset upload", status="completed", selected_style={"name": "Brand"})
    db.add(project)
    db.flush()
    db.add(
        Slide(
            project_id=project.id,
            page_num=1,
            status="completed",
            visual_json={"visual_asset_ids": []},
            prompt_text="old prompt",
        )
    )
    db.commit()

    monkeypatch.setattr(slides_api.settings, "UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        slides_api,
        "analyze_visual_asset",
        lambda *_args, **_kwargs: {
            "detected_kind": "product",
            "subject": "测试产品瓶",
            "description": "测试产品",
            "distinctive_features": ["瓶身"],
            "suggested_keywords": ["产品"],
        },
    )

    default_result = slides_api.upload_file(
        project.id,
        png_upload("product.png"),
        role="visual_asset",
        slide_id=None,
        process_mode=None,
        asset_name=None,
        asset_kind="product",
        usage_note=None,
        db=db,
    )
    explicit_result = slides_api.upload_file(
        project.id,
        png_upload("product-hero.png"),
        role="visual_asset",
        slide_id=None,
        process_mode="blend",
        asset_name=None,
        asset_kind="product",
        usage_note=None,
        db=db,
    )
    refreshed_project = db.query(Project).filter(Project.id == project.id).first()
    refreshed_slide = db.query(Slide).filter(Slide.project_id == project.id).first()

    assert default_result["process_mode"] == "crop"
    assert explicit_result["process_mode"] == "blend"
    assert refreshed_project.selected_style == {"name": "Brand"}
    assert refreshed_project.status == "completed"
    assert refreshed_slide.prompt_text == "old prompt"
    assert artifact_stale(refreshed_slide.visual_json) == {"content": True}


def test_visual_asset_upload_returns_before_vlm_analysis_when_background_available(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="Asset upload", status="planning")
    db.add(project)
    db.commit()

    monkeypatch.setattr(slides_api.settings, "UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        slides_api,
        "analyze_visual_asset",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(AssertionError("should not analyze synchronously")),
    )
    scheduled = []
    monkeypatch.setattr(slides_api, "_submit_asset_analysis_task", lambda *args: scheduled.append(args))

    result = slides_api.upload_file(
        project.id,
        png_upload("product.png"),
        role="visual_asset",
        slide_id=None,
        process_mode=None,
        asset_name="产品图",
        asset_kind="product",
        usage_note="用于产品展示页",
        background_tasks=BackgroundTasks(),
        db=db,
    )

    assert result["asset_analysis"]["analysis_status"] == "queued"
    assert result["asset_analysis"]["subject"] == "产品图"
    assert result["process_mode"] == "crop"
    assert scheduled and scheduled[0][0] == "visual_asset"


def test_visual_asset_upload_rejects_slide_level_asset(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="Asset upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(slides_api.settings, "UPLOAD_DIR", str(tmp_path))

    with pytest.raises(HTTPException) as exc:
        slides_api.upload_file(
            project.id,
            png_upload("product.png"),
            role="visual_asset",
            slide_id=slide.id,
            process_mode=None,
            asset_name=None,
            asset_kind=None,
            usage_note=None,
            db=db,
        )

    assert exc.value.status_code == 400
    assert "project-level" in exc.value.detail


def test_page_reference_upload_attaches_directly_to_slide_and_cleans_overlay(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="Page ref upload", status="visual_ready")
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="visual_ready",
        visual_json={"visual_description": "old", "overlay_layers": []},
        prompt_text="old prompt",
    )
    db.add(slide)
    db.commit()

    monkeypatch.setattr(slides_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = slides_api.upload_file(
        project.id,
        png_upload("page-product.png"),
        role="content_ref",
        slide_id=slide.id,
        process_mode="blend",
        asset_name="本页产品图",
        asset_kind=None,
        usage_note="只给第一页使用",
        db=db,
    )

    ref = db.query(ReferenceImage).filter(ReferenceImage.id == result["id"]).first()
    ref_id = ref.id
    assert ref.role == "content_ref"
    assert ref.slide_id == slide.id
    assert ref.asset_name == "本页产品图"
    assert ref.usage_note == "只给第一页使用"
    assert result["process_mode"] == "blend"

    slides_api.update_slide_overlay_layers(
        project.id,
        slide.id,
        slides_api.OverlayLayersRequest(
            layers=[slides_api.OverlayLayerRequest(asset_id=ref_id, preset="right-card")]
        ),
        db=db,
    )
    db.expire_all()
    overlaid_slide = db.query(Slide).filter(Slide.id == slide.id).first()
    overlaid_ref = db.query(ReferenceImage).filter(ReferenceImage.id == ref_id).first()
    assert overlaid_slide.visual_json["overlay_layers"][0]["asset_id"] == ref_id
    assert overlaid_ref.process_mode == "original"

    slides_api.delete_reference_image(project.id, ref_id, db=db)
    db.expire_all()
    cleaned_slide = db.query(Slide).filter(Slide.id == slide.id).first()
    assert cleaned_slide.visual_json.get("overlay_layers") == []
    assert db.query(ReferenceImage).filter(ReferenceImage.id == ref_id).first() is None


def test_logo_upload_keeps_existing_global_logos(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="Logo upload", status="planning")
    db.add(project)
    db.commit()

    monkeypatch.setattr(slides_api.settings, "UPLOAD_DIR", str(tmp_path))

    first = slides_api.upload_file(
        project.id,
        png_upload("focus.png"),
        role="logo",
        slide_id=None,
        process_mode=None,
        asset_name=None,
        asset_kind=None,
        usage_note=None,
        logo_anchor="top-left",
        db=db,
    )
    second = slides_api.upload_file(
        project.id,
        png_upload("partner.png"),
        role="logo",
        slide_id=None,
        process_mode=None,
        asset_name=None,
        asset_kind=None,
        usage_note=None,
        logo_anchor="top-left",
        db=db,
    )

    logos = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "logo",
        ReferenceImage.slide_id.is_(None),
    ).all()
    assert len(logos) == 2
    assert {first["id"], second["id"]} == {logo.id for logo in logos}
    assert all(logo.logo_anchor == "top-left" for logo in logos)


def test_pptx_document_upload_extracts_generic_picture_as_page_ref_only(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        pptx_upload_with_picture("source.pptx"),
        db=db,
    )
    assert result["text_parse_status"] == "queued"
    assert result["asset_extraction_status"] == "queued"
    assert db.query(ReferenceImage).filter(ReferenceImage.project_id == project.id).count() == 0

    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "source.pptx"),
        "source.pptx",
        db=db,
    )

    assert stats["page_refs"] == 1
    assert stats["visual_assets"] == 0
    page_ref = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    visual_assets = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "visual_asset",
    ).all()
    assert page_ref.slide_id == slide.id
    assert page_ref.asset_analysis["source_document"] == "source.pptx"
    assert page_ref.asset_analysis["pptx_source_page_num"] == 1
    assert page_ref.asset_analysis["selection_tier"] == "page_ref"
    assert visual_assets == []
    docs = documents_api.list_documents(project.id, db=db)
    assert docs[0]["asset_extraction_status"] == "completed"


def test_pptx_asset_extraction_is_idempotent_for_page_refs(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(
        project.id,
        pptx_upload_with_picture("source.pptx"),
        db=db,
    )
    source_path = str(tmp_path / project.id / "docs" / "source.pptx")
    first = documents_api._extract_pptx_assets_for_document(project.id, source_path, "source.pptx", db=db)
    second = documents_api._extract_pptx_assets_for_document(project.id, source_path, "source.pptx", db=db)

    page_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).all()

    assert first["page_refs"] == 1
    assert second["page_refs"] == 0
    assert len(page_refs) == 1
    assert page_refs[0].slide_id == slide.id


def test_pdf_document_upload_extracts_source_pack_and_page_ref_images(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PDF upload", status="planning")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        pdf_upload_with_picture("book.pdf"),
        db=db,
    )
    assert result["text_parse_status"] == "queued"
    assert result["source_parse_status"] == "queued"
    assert result["asset_extraction_status"] == "queued"

    stats = documents_api._extract_pdf_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "book.pdf"),
        "book.pdf",
        db=db,
    )

    assert stats["page_refs"] == 1
    page_ref = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    assert page_ref.asset_analysis["source_type"] == "pdf"
    assert page_ref.asset_analysis["source_document"] == "book.pdf"
    assert page_ref.asset_analysis["source_page_num"] == 1
    assert page_ref.asset_analysis["pdf_source_page_num"] == 1
    assert page_ref.asset_analysis["bbox"]
    assert page_ref.asset_analysis["nearby_text"]

    docs = documents_api.list_documents(project.id, db=db)
    assert docs[0]["source_parse_status"] == "completed"
    assert docs[0]["source_stats"]["pages"] == 1
    assert docs[0]["source_stats"]["images"] == 1
    assert docs[0]["asset_extraction_status"] == "completed"


def test_pdf_document_upload_promotes_repeated_brand_mark_to_logo(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PDF repeated logo", status="planning")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(
        project.id,
        pdf_upload_with_repeated_logo("brand-deck.pdf"),
        db=db,
    )

    stats = documents_api._extract_pdf_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "brand-deck.pdf"),
        "brand-deck.pdf",
        db=db,
    )

    logos = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "logo",
        ReferenceImage.slide_id.is_(None),
    ).all()
    page_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).all()

    assert stats["logos"] == 1
    assert stats["page_refs"] == 0
    assert len(logos) == 1
    assert page_refs == []
    assert logos[0].process_mode == "original"
    assert logos[0].asset_analysis["source_type"] == "pdf"
    assert logos[0].asset_analysis["classification"] == "pdf_repeated_logo"
    assert logos[0].asset_analysis["review_status"] == "auto_confirmed"
    assert logos[0].logo_anchor == "top-left"


def test_pptx_shape_fill_repeated_logo_goes_to_logo_library(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(
        project.id,
        pptx_upload_with_repeated_shape_fill_logos(tmp_path, "shape-fill-logo.pptx"),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "shape-fill-logo.pptx"),
        "shape-fill-logo.pptx",
        db=db,
    )

    logos = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "logo",
        ReferenceImage.slide_id.is_(None),
    ).all()
    page_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).all()

    assert stats["logos"] == 1
    assert stats["page_refs"] == 0
    assert len(logos) == 1
    assert page_refs == []
    assert logos[0].asset_analysis["classification"] == "logo"
    assert logos[0].asset_analysis["asset_scope"] == "project_logo"
    assert logos[0].asset_analysis["shape_bounds"]["width"] < 0.1
    assert logos[0].process_mode == "original"


def test_template_logo_candidates_only_auto_confirm_one_lockup(tmp_path):
    db = make_session()
    project = Project(title="Template logos", status="planning")
    db.add(project)
    db.commit()

    paths = []
    for index, color in enumerate(("black", "navy", "gold"), start=1):
        path = tmp_path / f"template-logo-{index}.png"
        Image.new("RGBA", (240, 80), color).save(path)
        paths.append(path)

    assets = [
        SimpleNamespace(
            file_path=str(path),
            asset_name=f"候选 Logo {index}",
            usage_note=None,
            classification="logo",
            metadata={
                "pptx_image_sha1": f"candidate-{index}",
                "shape_bounds": {"left": 0.86, "top": 0.03, "width": 0.1, "height": 0.04},
            },
        )
        for index, path in enumerate(paths, start=1)
    ]

    attached = slides_api._attach_template_logo_assets(project, assets, "template.pptx", db)
    db.commit()
    logos = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "logo",
        ReferenceImage.slide_id.is_(None),
    ).order_by(ReferenceImage.asset_name).all()
    db.refresh(project)

    assert attached == 3
    assert len(logos) == 3
    assert [slides_api.logo_review_status(logo) for logo in logos].count("auto_confirmed") == 1
    assert [logo.id for logo in slides_api._project_logo_refs(project)] == [logos[0].id]
    assert "额外 Logo 候选" in logos[1].asset_analysis["review_reason"]


def test_logo_review_after_generation_reassembles_without_clearing_images(tmp_path, monkeypatch):
    db = make_session()
    output_dir = tmp_path / "outputs"
    monkeypatch.setattr(slides_api.settings, "OUTPUT_DIR", str(output_dir))

    project = Project(
        id="generated-logo-review",
        title="Generated logo review",
        status="prototype_ready",
        content_plan_confirmed=True,
        selected_style={"name": "沿用模板"},
    )
    db.add(project)
    image_paths = []
    for page_num in (1, 2):
        image_path = output_dir / project.id / f"slide_{page_num:02d}.png"
        image_path.parent.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (1792, 1024), (255, 255, 255)).save(image_path)
        image_paths.append(str(image_path))
        db.add(Slide(
            project_id=project.id,
            page_num=page_num,
            type="content",
            status="completed",
            image_path=str(image_path),
            content_json={"speaker_notes": f"notes {page_num}"},
            visual_json={"logo_policy": {"show_logo": True, "placement": "top-right", "scale": "small"}},
            prompt_text="prompt",
        ))

    for index, color in enumerate(((255, 203, 34, 255), (20, 24, 30, 255)), start=1):
        path = tmp_path / f"logo-{index}.png"
        Image.new("RGBA", (220, 80), color).save(path)
        db.add(ReferenceImage(
            id=f"logo-{index}",
            project_id=project.id,
            role="logo",
            file_path=str(path),
            process_mode="original",
            asset_name=f"Logo {index}",
            logo_anchor="top-right",
            asset_analysis={"review_status": "auto_confirmed"},
        ))
    db.commit()

    result = slides_api.update_reference_image(
        project.id,
        "logo-2",
        {"review_status": "needs_review"},
        db=db,
    )
    db.refresh(project)
    slides = db.query(Slide).filter(Slide.project_id == project.id).order_by(Slide.page_num).all()

    assert result["review_status"] == "needs_review"
    assert project.status == "prototype_ready"
    assert [slide.status for slide in slides] == ["completed", "completed"]
    assert [slide.image_path for slide in slides] == image_paths
    assert all(artifact_stale(slide.visual_json) == {} for slide in slides)
    assert (output_dir / project.id / "prototype.pptx").exists()


def test_pptx_document_upload_queues_asset_extraction_with_background_tasks(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))
    background_tasks = BackgroundTasks()
    scheduled = []
    monkeypatch.setattr(
        documents_api,
        "_dispatch_document_processing",
        lambda *args: scheduled.append(args),
    )

    result = documents_api.upload_document(
        project.id,
        pptx_upload_with_primary_product_image("product-source.pptx"),
        background_tasks=background_tasks,
        db=db,
    )

    assert result["asset_extraction_status"] == "queued"
    assert result["text_parse_status"] == "queued"
    assert result["extracted_assets"]["total"] == 0
    assert len(background_tasks.tasks) == 0
    assert scheduled == [(
        project.id,
        str(tmp_path / project.id / "docs" / "product-source.pptx"),
        "product-source.pptx",
        ".pptx",
    )]
    assert db.query(ReferenceImage).filter(ReferenceImage.project_id == project.id).count() == 0
    docs = documents_api.list_documents(project.id, db=db)
    assert docs[0]["asset_extraction_status"] == "queued"
    assert docs[0]["text_parse_status"] == "queued"


def test_uploaded_document_text_is_parsed_on_generate_when_background_not_done(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="Text upload", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        txt_upload("brief.txt", "预算有限\n目标是绿色开学季营销"),
        db=db,
    )

    assert result["text_parse_status"] == "queued"
    docs_before = documents_api.list_documents(project.id, db=db)
    assert docs_before[0]["filename"] == "brief.txt"
    assert docs_before[0]["text_parse_status"] == "queued"

    documents_text = load_project_documents(project.id, parse_missing=True)
    assert "预算有限" in documents_text
    docs_after = documents_api.list_documents(project.id, db=db)
    assert docs_after[0]["text_parse_status"] == "completed"
    assert docs_after[0]["char_count"] > 0


def test_pptx_upload_keeps_parallel_picture_group_complete_and_out_of_global_library(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        pptx_upload_with_many_pictures("many.pptx", count=5),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "many.pptx"),
        "many.pptx",
        db=db,
    )

    page_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).all()
    library_assets = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "visual_asset",
        ReferenceImage.slide_id.is_(None),
    ).all()

    assert result["extracted_assets"]["total"] == 0
    assert stats["total"] == 5
    assert len(page_refs) == 5
    assert len(library_assets) == 0
    assert {ref.asset_analysis["asset_group_size"] for ref in page_refs} == {5}
    assert sorted(ref.asset_analysis["asset_group_index"] for ref in page_refs) == [1, 2, 3, 4, 5]


def test_pptx_document_upload_promotes_primary_core_asset(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        pptx_upload_with_primary_product_image("product-source.pptx"),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "product-source.pptx"),
        "product-source.pptx",
        db=db,
    )

    assert result["asset_extraction_status"] == "queued"
    assert stats["page_refs"] == 1
    assert stats["visual_assets"] == 1
    visual_asset = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "visual_asset",
    ).one()
    assert visual_asset.slide_id is None
    assert visual_asset.asset_kind == "product"
    assert visual_asset.asset_analysis["selection_tier"] == "core_global"
    assert visual_asset.asset_analysis["library_role"] == "core_global_asset"


def test_pptx_document_upload_keeps_content_graphic_page_level(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        pptx_upload_with_content_graphic("graphic-source.pptx"),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "graphic-source.pptx"),
        "graphic-source.pptx",
        db=db,
    )

    assert result["asset_extraction_status"] == "queued"
    assert stats["page_refs"] == 1
    assert stats["visual_assets"] == 0
    page_ref = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    assert page_ref.asset_analysis["selection_tier"] == "page_ref"
    assert page_ref.asset_analysis["selection_reason"] == "content graphic kept as page-level evidence"


def test_pptx_document_upload_keeps_identity_code_out_of_global_library(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        pptx_upload_with_identity_code("identity-code-source.pptx"),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "identity-code-source.pptx"),
        "identity-code-source.pptx",
        db=db,
    )

    assert result["asset_extraction_status"] == "queued"
    assert stats["page_refs"] == 1
    assert stats["visual_assets"] == 0
    page_ref = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    assert page_ref.asset_analysis["selection_tier"] == "page_ref"
    assert page_ref.asset_analysis["selection_reason"] == "QR/identity code kept as page-specific evidence"


def test_pptx_document_upload_ignores_phone_frame_layout_chrome(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        pptx_upload_with_phone_frame_chrome("phone-frame-source.pptx"),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "phone-frame-source.pptx"),
        "phone-frame-source.pptx",
        db=db,
    )

    assert result["asset_extraction_status"] == "queued"
    assert stats["page_refs"] == 0
    assert stats["visual_assets"] == 0
    assert db.query(ReferenceImage).filter(ReferenceImage.project_id == project.id).count() == 0


def test_pptx_document_upload_promotes_multiple_cover_logos(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    result = documents_api.upload_document(
        project.id,
        pptx_upload_with_two_cover_logos("two-logos.pptx"),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "two-logos.pptx"),
        "two-logos.pptx",
        db=db,
    )

    logos = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "logo",
        ReferenceImage.slide_id.is_(None),
    ).all()

    assert result["asset_extraction_status"] == "queued"
    assert stats["logos"] == 2
    assert len(logos) == 2
    assert all(logo.process_mode == "original" for logo in logos)
    assert all(logo.asset_analysis["classification"] == "logo_candidate" for logo in logos)


def test_pptx_document_upload_respects_picture_crop_for_logo(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(
        project.id,
        pptx_upload_with_cropped_cover_logo("cropped-logo.pptx"),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "cropped-logo.pptx"),
        "cropped-logo.pptx",
        db=db,
    )

    logo = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "logo",
        ReferenceImage.slide_id.is_(None),
    ).one()
    extracted = Image.open(logo.file_path).convert("RGB")
    right_half_pixels = [
        extracted.getpixel((x, y))
        for x in range(extracted.width // 2, extracted.width)
        for y in range(extracted.height)
    ]

    assert stats["logos"] == 1
    assert extracted.width <= 225
    assert logo.asset_analysis["pptx_crop"] == [0.0, 0.0, 0.5, 0.0]
    assert not any(r > 180 and g < 60 and b < 70 for r, g, b in right_half_pixels)


def test_pptx_document_upload_promotes_cropped_multicolor_cover_logo(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(
        project.id,
        pptx_upload_with_cropped_multicolor_cover_logo("cropped-multicolor-logo.pptx"),
        db=db,
    )
    stats = documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "cropped-multicolor-logo.pptx"),
        "cropped-multicolor-logo.pptx",
        db=db,
    )

    logos = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "logo",
        ReferenceImage.slide_id.is_(None),
    ).all()
    page_refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).all()

    assert stats["logos"] == 1
    assert stats["page_refs"] == 0
    assert len(logos) == 1
    assert page_refs == []
    assert logos[0].asset_analysis["classification"] == "logo_candidate"
    assert logos[0].asset_analysis["review_status"] == "needs_review"
    assert max(logos[0].asset_analysis["pptx_crop"]) > 0


def test_unconfirmed_logos_do_not_enter_project_logo_refs(tmp_path):
    db = make_session()
    project = Project(title="Logo review", status="planning")
    db.add(project)
    db.flush()
    needs_review_path = tmp_path / "maybe_logo.png"
    confirmed_path = tmp_path / "confirmed_logo.png"
    Image.new("RGB", (80, 40), "white").save(needs_review_path)
    Image.new("RGB", (80, 40), "white").save(confirmed_path)
    needs_review = ReferenceImage(
        project_id=project.id,
        file_path=str(needs_review_path),
        role="logo",
        process_mode="original",
        asset_analysis={"classification": "logo_candidate", "review_status": "needs_review"},
    )
    confirmed = ReferenceImage(
        project_id=project.id,
        file_path=str(confirmed_path),
        role="logo",
        process_mode="original",
        asset_analysis={"review_status": "user_confirmed"},
    )
    db.add_all([needs_review, confirmed])
    db.commit()
    db.refresh(project)

    refs = slides_api._project_logo_refs(project)

    assert [ref.id for ref in refs] == [confirmed.id]


def test_missing_confirmed_logo_does_not_enable_logo_policy(tmp_path):
    db = make_session()
    project = Project(title="Missing logo", status="planning")
    db.add(project)
    db.flush()
    ref = ReferenceImage(
        project_id=project.id,
        file_path=str(tmp_path / "missing-logo.png"),
        role="logo",
        process_mode="original",
        asset_analysis={"review_status": "user_confirmed"},
    )
    db.add(ref)
    db.commit()
    db.refresh(project)

    intent = slides_api._with_project_logo_policy(
        {"type": "cover", "logo_policy": {"show_logo": True, "placement": "center", "scale": "large"}},
        project,
    )

    assert slides_api._project_logo_refs(project) == []
    assert intent["logo_policy"]["show_logo"] is False


def test_logo_review_status_can_promote_candidate_back_to_logo(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="Logo review", status="planning")
    db.add(project)
    db.flush()
    logo_path = tmp_path / "maybe_logo.png"
    Image.new("RGB", (80, 40), "white").save(logo_path)
    ref = ReferenceImage(
        project_id=project.id,
        file_path=str(logo_path),
        role="logo",
        process_mode="original",
        asset_analysis={"classification": "logo_candidate", "review_status": "needs_review"},
    )
    db.add(ref)
    db.commit()

    result = slides_api.update_reference_image(
        project.id,
        ref.id,
        {"review_status": "user_confirmed"},
        db=db,
    )
    db.refresh(project)

    assert result["review_status"] == "user_confirmed"
    assert [logo.id for logo in slides_api._project_logo_refs(project)] == [ref.id]


def test_pending_pptx_page_refs_link_after_content_plan_exists(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(
        project.id,
        pptx_upload_with_picture("source.pptx"),
        db=db,
    )
    documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "source.pptx"),
        "source.pptx",
        db=db,
    )
    pending = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    assert pending.slide_id is None

    slide = Slide(project_id=project.id, page_num=1, status="pending")
    db.add(slide)
    db.commit()

    linked = slides_api._link_pending_pptx_page_refs(project.id, db)
    db.commit()

    assert linked == 1
    db.refresh(pending)
    assert pending.slide_id == slide.id


def test_pending_pptx_page_refs_link_by_source_refs_when_pages_are_reordered(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PPT upload", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(
        project.id,
        pptx_upload_with_picture("source.pptx"),
        db=db,
    )
    documents_api._extract_pptx_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "source.pptx"),
        "source.pptx",
        db=db,
    )
    pending = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    assert pending.slide_id is None

    slide = Slide(
        project_id=project.id,
        page_num=3,
        status="pending",
        content_json={
            "page_num": 3,
            "source_refs": [{"source_document": "source.pptx", "source_page_num": 1}],
        },
    )
    db.add(slide)
    db.commit()

    linked = slides_api._link_pending_pptx_page_refs(project.id, db)
    db.commit()

    assert linked == 1
    db.refresh(pending)
    assert pending.slide_id == slide.id


def test_pending_pdf_page_refs_require_explicit_figure_refs(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PDF source refs", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(project.id, pdf_upload_with_picture("book.pdf"), db=db)
    documents_api._extract_pdf_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "book.pdf"),
        "book.pdf",
        db=db,
    )
    pending = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    assert pending.slide_id is None

    slide = Slide(
        project_id=project.id,
        page_num=5,
        status="pending",
        content_json={
            "page_num": 5,
            "source_refs": [{"source_document": "book.pdf", "source_page_num": 1, "source_type": "pdf"}],
        },
    )
    db.add(slide)
    db.commit()

    linked = slides_api._link_pending_pptx_page_refs(project.id, db)
    db.commit()

    assert linked == 0
    db.refresh(pending)
    assert pending.slide_id is None


def test_pending_pdf_page_refs_link_by_figure_refs(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="PDF figure refs", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(project.id, pdf_upload_with_picture("book.pdf"), db=db)
    documents_api._extract_pdf_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "book.pdf"),
        "book.pdf",
        db=db,
    )
    pending = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    figure_id = pending.asset_analysis["id"]

    slide = Slide(
        project_id=project.id,
        page_num=9,
        status="pending",
        content_json={
            "page_num": 9,
            "figure_refs": [{
                "source_document": "book.pdf",
                "source_page_num": 1,
                "source_type": "pdf",
                "figure_id": figure_id,
            }],
        },
    )
    db.add(slide)
    db.commit()

    linked = slides_api._link_pending_pptx_page_refs(project.id, db)
    db.commit()

    assert linked == 1
    db.refresh(pending)
    assert pending.slide_id == slide.id


def test_pending_pdf_page_refs_clone_when_same_content_figure_is_used_on_multiple_slides(tmp_path, monkeypatch):
    db = make_session()
    project = Project(title="Repeated PDF figure refs", status="draft")
    db.add(project)
    db.commit()

    monkeypatch.setattr(documents_api.settings, "UPLOAD_DIR", str(tmp_path))

    documents_api.upload_document(project.id, pdf_upload_with_picture("book.pdf"), db=db)
    documents_api._extract_pdf_assets_for_document(
        project.id,
        str(tmp_path / project.id / "docs" / "book.pdf"),
        "book.pdf",
        db=db,
    )
    pending = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).one()
    figure_id = pending.asset_analysis["id"]

    slides = [
        Slide(
            project_id=project.id,
            page_num=page_num,
            status="pending",
            content_json={
                "page_num": page_num,
                "figure_refs": [{
                    "source_document": "book.pdf",
                    "source_page_num": 1,
                    "source_type": "pdf",
                    "figure_id": figure_id,
                }],
            },
        )
        for page_num in (5, 6)
    ]
    db.add_all(slides)
    db.commit()

    linked = slides_api._link_pending_pptx_page_refs(project.id, db)
    db.commit()

    refs = db.query(ReferenceImage).filter(
        ReferenceImage.project_id == project.id,
        ReferenceImage.role == "content_ref",
    ).order_by(ReferenceImage.slide_id).all()
    assert linked == 2
    assert len(refs) == 2
    assert {ref.slide_id for ref in refs} == {slide.id for slide in slides}
    assert {ref.file_path for ref in refs} == {pending.file_path}
    assert {ref.asset_analysis["id"] for ref in refs} == {figure_id}


def test_delete_visual_asset_cleans_slide_selection_and_invalidates_outputs(tmp_path):
    db = make_session()
    asset_path = tmp_path / "asset.png"
    Image.new("RGB", (10, 10), "white").save(asset_path)
    project = Project(title="Asset delete", status="completed", selected_style={"name": "Brand"})
    db.add(project)
    db.flush()
    asset = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(asset_path),
        process_mode="crop",
        asset_name="测试产品瓶",
        asset_kind="product",
    )
    db.add(asset)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        visual_json={
            "manual_visual_asset_ids": [asset.id],
            "manual_visual_asset_usage": {asset.id: "锁定到本页"},
            "visual_asset_ids": [asset.id],
            "visual_asset_usage": {asset.id: "右侧展示"},
            "overlay_layers": [{
                "id": "ov_asset",
                "asset_id": asset.id,
                "enabled": True,
                "preset": "right-card",
                "fit": "contain",
                "mode": "exact_card",
                "usage_note": "原样展示",
            }],
        },
        prompt_text="old prompt",
    )
    db.add(slide)
    db.commit()

    result = slides_api.delete_reference_image(project.id, asset.id, db=db)
    refreshed_project = db.query(Project).filter(Project.id == project.id).first()
    refreshed_slide = db.query(Slide).filter(Slide.project_id == project.id).first()

    assert result["message"] == "Deleted"
    assert refreshed_slide.visual_json["visual_asset_ids"] == []
    assert refreshed_slide.visual_json["visual_asset_usage"] == {}
    assert refreshed_slide.visual_json["manual_visual_asset_ids"] == []
    assert refreshed_slide.visual_json["manual_visual_asset_usage"] == {}
    assert refreshed_slide.visual_json["overlay_layers"] == []
    assert artifact_stale(refreshed_slide.visual_json) == {"content": True}
    assert refreshed_slide.prompt_text == "old prompt"
    assert refreshed_slide.status == "completed"
    assert refreshed_project.status == "completed"


def test_asset_pins_replace_reorder_and_survive_visual_merge(tmp_path):
    db = make_session()
    project = Project(title="Pins", status="completed")
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        visual_json={"visual_asset_ids": ["old-auto"], "visual_asset_usage": {"old-auto": "旧自动"}},
        prompt_text="old",
    )
    db.add(slide)
    assets = []
    for idx in range(2):
        path = tmp_path / f"asset_{idx}.png"
        Image.new("RGB", (10, 10), "white").save(path)
        asset = ReferenceImage(
            project_id=project.id,
            role="visual_asset",
            file_path=str(path),
            process_mode="crop",
            asset_name=f"素材{idx}",
            asset_kind="scene",
        )
        db.add(asset)
        assets.append(asset)
    db.commit()

    result = slides_api.update_slide_asset_pins(
        project.id,
        slide.id,
        slides_api.AssetPinsRequest(
            asset_ids=[assets[1].id, assets[0].id],
            usage={assets[1].id: "放在左侧"},
        ),
        db=db,
    )
    db.refresh(slide)

    assert result["manual_visual_asset_ids"] == [assets[1].id, assets[0].id]
    assert slide.visual_json["visual_asset_ids"][:2] == [assets[1].id, assets[0].id]
    assert slide.visual_json["manual_visual_asset_usage"][assets[1].id] == "放在左侧"
    assert artifact_stale(slide.visual_json) == {"visual": True}
    assert slide.prompt_text == "old"
    assert slide.status == "completed"

    slide.visual_json["asset_route_modes"] = {
        assets[1].id: "blend",
        assets[0].id: "double_blend",
        "old-auto": "overlay",
    }
    merged = slides_api._merge_manual_pins_into_visual_json(
        {"visual_asset_ids": ["auto-1"], "visual_asset_usage": {"auto-1": "自动"}},
        slide.visual_json,
    )
    assert merged["visual_asset_ids"][:3] == [assets[1].id, assets[0].id, "auto-1"]
    assert merged["asset_route_modes"] == {
        assets[1].id: "blend",
        assets[0].id: "double_blend",
    }


def test_overlay_layers_endpoint_and_visual_merge(tmp_path):
    db = make_session()
    project = Project(title="Overlay", status="completed")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="completed", prompt_text="old", visual_json={})
    db.add(slide)
    path = tmp_path / "screenshot.png"
    Image.new("RGB", (120, 70), "white").save(path)
    asset = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(path),
        process_mode="blend",
        asset_name="后台截图",
        asset_kind="material",
        asset_analysis={"subject": "后台截图"},
    )
    db.add(asset)
    db.commit()

    result = slides_api.update_slide_overlay_layers(
        project.id,
        slide.id,
        slides_api.OverlayLayersRequest(layers=[
            slides_api.OverlayLayerRequest(
                asset_id=asset.id,
                preset="right-card",
                mode="exact_card",
                usage_note="原样叠加",
            )
        ]),
        db=db,
    )
    db.refresh(slide)

    assert result["overlay_layers"][0]["asset_id"] == asset.id
    assert result["overlay_layers"][0]["fit"] == "contain"
    assert slide.visual_json["overlay_layers"][0]["preset"] == "right-card"
    assert artifact_stale(slide.visual_json) == {"visual": True}
    assert slide.prompt_text == "old"
    assert slide.status == "completed"
    db.refresh(asset)
    assert asset.process_mode == "original"
    assert asset.asset_analysis["exact_overlay"] is True

    merged = slides_api._merge_manual_pins_into_visual_json({"visual_asset_ids": []}, slide.visual_json)
    assert merged["overlay_layers"][0]["asset_id"] == asset.id


def test_overlay_layers_accept_page_reference_assets(tmp_path):
    db = make_session()
    project = Project(title="Page reference overlay", status="completed")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=2, status="completed", prompt_text="old", visual_json={})
    db.add(slide)
    db.flush()
    path = tmp_path / "page-ref.png"
    Image.new("RGB", (120, 70), "white").save(path)
    page_ref = ReferenceImage(
        project_id=project.id,
        slide_id=slide.id,
        role="content_ref",
        file_path=str(path),
        process_mode="blend",
        asset_name="原 PPT 第 2 页参考图",
        asset_analysis={"subject": "页面图片"},
    )
    db.add(page_ref)
    db.commit()

    result = slides_api.update_slide_overlay_layers(
        project.id,
        slide.id,
        slides_api.OverlayLayersRequest(layers=[
            slides_api.OverlayLayerRequest(
                asset_id=page_ref.id,
                preset="center-card",
                mode="exact_card",
                usage_note="原样保留",
            )
        ]),
        db=db,
    )
    db.refresh(slide)
    db.refresh(page_ref)

    assert result["overlay_layers"][0]["asset_id"] == page_ref.id
    assert slide.visual_json["overlay_layers"][0]["preset"] == "center-card"
    assert page_ref.process_mode == "original"
    assert page_ref.asset_analysis["exact_overlay"] is True

    serialized = slides_api.list_slides(project.id, db=db)
    serialized_slide = next(item for item in serialized if item["id"] == slide.id)
    assert serialized_slide["visual_json"]["overlay_layers"][0]["asset_id"] == page_ref.id
    assert serialized_slide["reference_images"][0]["id"] == page_ref.id
    assert serialized_slide["reference_images"][0]["process_mode"] == "original"


def test_page_reference_original_mode_creates_exact_overlay_layer(tmp_path):
    db = make_session()
    project = Project(title="Page ref original sync", status="completed", content_plan_confirmed=True)
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=6,
        status="completed",
        prompt_text="old prompt",
        image_path="/tmp/old.png",
        visual_json={"visual_description": "old"},
    )
    db.add(slide)
    db.flush()
    path = tmp_path / "page-ref.png"
    Image.new("RGB", (160, 90), "white").save(path)
    page_ref = ReferenceImage(
        project_id=project.id,
        slide_id=slide.id,
        role="content_ref",
        file_path=str(path),
        process_mode="blend",
        asset_name="原 PDF 第 6 页参考图",
        asset_analysis={"source_page_num": 6},
    )
    db.add(page_ref)
    db.commit()

    slides_api.update_reference_image(
        project.id,
        page_ref.id,
        {"process_mode": "original"},
        db=db,
    )
    db.refresh(slide)
    db.refresh(page_ref)

    layers = slide.visual_json["overlay_layers"]
    assert len(layers) == 1
    assert layers[0]["asset_id"] == page_ref.id
    assert layers[0]["mode"] == "exact_cutout"
    assert page_ref.asset_analysis["exact_overlay"] is True
    assert artifact_stale(slide.visual_json) == {"visual": True}
    assert slide.prompt_text == "old prompt"


def test_original_page_reference_is_synced_before_reference_context(tmp_path):
    db = make_session()
    project = Project(title="Page ref context sync", status="visual_ready")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=8, status="visual_ready", visual_json={})
    db.add(slide)
    db.flush()
    path = tmp_path / "chart.png"
    Image.new("RGB", (160, 90), "white").save(path)
    page_ref = ReferenceImage(
        project_id=project.id,
        slide_id=slide.id,
        role="content_ref",
        file_path=str(path),
        process_mode="original",
        asset_name="原 PDF 第 8 页截图",
    )
    db.add(page_ref)
    db.commit()

    contexts, hints = slides_api._build_slide_reference_contexts([slide])

    assert contexts == {}
    assert hints == {}
    assert slide.visual_json["overlay_layers"][0]["asset_id"] == page_ref.id


def test_exact_page_reference_sync_is_noop_without_original_refs(tmp_path):
    db = make_session()
    project = Project(title="No exact sync", status="visual_ready")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=3, status="visual_ready", visual_json={})
    db.add(slide)
    db.flush()
    path = tmp_path / "blend.png"
    Image.new("RGB", (160, 90), "white").save(path)
    page_ref = ReferenceImage(
        project_id=project.id,
        slide_id=slide.id,
        role="content_ref",
        file_path=str(path),
        process_mode="blend",
        asset_name="智能融合参考图",
    )
    db.add(page_ref)
    db.commit()

    changed = slides_api._sync_exact_page_reference_overlay_layers(slide, [page_ref])

    assert changed is False
    assert slide.visual_json == {}


def test_asset_pins_unpin_removes_overlay_layer(tmp_path):
    db = make_session()
    project = Project(title="Unpin overlay", status="completed")
    db.add(project)
    db.flush()
    path = tmp_path / "asset.png"
    Image.new("RGB", (10, 10), "white").save(path)
    asset = ReferenceImage(project_id=project.id, role="visual_asset", file_path=str(path), asset_kind="other")
    db.add(asset)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        visual_json={
            "manual_visual_asset_ids": [asset.id],
            "visual_asset_ids": [asset.id],
            "overlay_layers": [{"id": "ov", "asset_id": asset.id, "enabled": True, "preset": "right-card"}],
        },
    )
    db.add(slide)
    db.commit()

    slides_api.update_slide_asset_pins(
        project.id,
        slide.id,
        slides_api.AssetPinsRequest(asset_ids=[]),
        db=db,
    )
    db.refresh(slide)

    assert slide.visual_json["manual_visual_asset_ids"] == []
    assert slide.visual_json["overlay_layers"] == []


def test_pipeline_loads_manual_pins_before_auto_assets(tmp_path):
    db = make_session()
    project = Project(title="Pipeline pins", status="prompt_ready")
    db.add(project)
    db.flush()
    paths = []
    for idx, color in enumerate(["red", "blue"]):
        path = tmp_path / f"asset_{idx}.png"
        Image.new("RGB", (20, 20), color).save(path)
        paths.append(path)
    manual = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(paths[0]),
        process_mode="blend",
        asset_name="手动锁定",
        asset_kind="other",
    )
    auto = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(paths[1]),
        process_mode="blend",
        asset_name="自动候选",
        asset_kind="product",
    )
    db.add_all([manual, auto])
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        prompt_text="prompt",
        visual_json={
            "manual_visual_asset_ids": [manual.id],
            "visual_asset_ids": [auto.id, manual.id],
        },
    )
    db.add(slide)
    db.commit()

    loaded_slide = db.query(Slide).filter(Slide.id == slide.id).one()
    refs = _load_reference_images(loaded_slide)

    assert refs[0]["id"] == manual.id
    assert refs[0]["manual_pin"] is True
    assert [ref["id"] for ref in refs[:2]] == [manual.id, auto.id]


def test_pipeline_skips_exact_overlay_assets_as_image_inputs(tmp_path):
    db = make_session()
    project = Project(title="Pipeline overlay", status="prompt_ready")
    db.add(project)
    db.flush()
    overlay_path = tmp_path / "overlay.png"
    auto_path = tmp_path / "auto.png"
    Image.new("RGB", (20, 20), "red").save(overlay_path)
    Image.new("RGB", (20, 20), "blue").save(auto_path)
    overlay_asset = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(overlay_path),
        process_mode="original",
        asset_name="Exact 截图",
        asset_kind="material",
    )
    auto_asset = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(auto_path),
        process_mode="blend",
        asset_name="场景参考",
        asset_kind="scene",
    )
    db.add_all([overlay_asset, auto_asset])
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        prompt_text="prompt",
        visual_json={
            "visual_asset_ids": [overlay_asset.id, auto_asset.id],
            "overlay_layers": [{
                "id": "ov_exact",
                "asset_id": overlay_asset.id,
                "enabled": True,
                "preset": "right-card",
            }],
        },
    )
    db.add(slide)
    db.commit()

    refs = _load_reference_images(db.query(Slide).filter(Slide.id == slide.id).one())

    assert overlay_asset.id not in [ref.get("id") for ref in refs]
    assert auto_asset.id in [ref.get("id") for ref in refs]


def test_pipeline_skips_route_overlay_assets_even_without_layers(tmp_path):
    db = make_session()
    project = Project(title="Route overlay", status="prompt_ready")
    db.add(project)
    db.flush()
    overlay_path = tmp_path / "overlay.png"
    Image.new("RGB", (20, 20), "red").save(overlay_path)
    overlay_asset = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(overlay_path),
        process_mode="original",
        asset_name="Exact 截图",
        asset_kind="material",
    )
    db.add(overlay_asset)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        prompt_text="prompt",
        visual_json={
            "visual_asset_ids": [overlay_asset.id],
            "asset_route_modes": {overlay_asset.id: "overlay"},
        },
    )
    db.add(slide)
    db.commit()

    loaded_slide = db.query(Slide).filter(Slide.id == slide.id).one()
    refs = _load_reference_images(loaded_slide)
    prompt_refs = slides_api._project_refs_for_prompt(
        db.query(Project).filter(Project.id == project.id).one(),
        [overlay_asset.id],
        loaded_slide.visual_json,
    )

    assert overlay_asset.id not in [ref.get("id") for ref in refs]
    assert prompt_refs == []


def test_project_asset_route_overrides_stale_process_mode(tmp_path):
    db = make_session()
    project = Project(title="Route override", status="prompt_ready")
    db.add(project)
    db.flush()
    path = tmp_path / "asset.png"
    Image.new("RGB", (20, 20), "red").save(path)
    asset = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(path),
        process_mode="original",
        asset_name="已从精确粘贴切回的素材",
        asset_kind="material",
    )
    db.add(asset)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        prompt_text="prompt",
        visual_json={
            "visual_asset_ids": [asset.id],
            "asset_route_modes": {asset.id: "blend"},
        },
    )
    db.add(slide)
    db.commit()

    refs = _load_reference_images(db.query(Slide).filter(Slide.id == slide.id).one())
    prompt_refs = slides_api._project_refs_for_prompt(
        db.query(Project).filter(Project.id == project.id).one(),
        [asset.id],
        slide.visual_json,
    )

    assert refs[0]["id"] == asset.id
    assert refs[0]["asset_route_mode"] == "blend"
    assert refs[0]["process_mode"] == "blend"
    assert prompt_refs[0]["asset_route_mode"] == "blend"
    assert prompt_refs[0]["process_mode"] == "blend"


def test_precision_page_references_use_refinement_pass():
    refs = [
        {"id": "person-1", "role": "content_ref", "process_mode": "crop", "image": object()},
        {"id": "person-2", "role": "content_ref", "process_mode": "crop", "image": object()},
        {"id": "blend-ref", "role": "content_ref", "process_mode": "blend", "image": object()},
        {"id": "overlay-ref", "role": "content_ref", "process_mode": "original", "image": object()},
    ]

    refined = generation_pipeline._product_refinement_refs(refs)

    assert [ref["id"] for ref in refined] == ["person-1", "person-2"]


def test_pipeline_skips_page_reference_overlay_as_image_input(tmp_path):
    db = make_session()
    project = Project(title="Pipeline page reference overlay", status="prompt_ready")
    db.add(project)
    db.flush()
    overlay_path = tmp_path / "page-overlay.png"
    blend_path = tmp_path / "page-blend.png"
    Image.new("RGB", (20, 20), "red").save(overlay_path)
    Image.new("RGB", (20, 20), "blue").save(blend_path)
    slide = Slide(
        project_id=project.id,
        page_num=3,
        status="prompt_ready",
        prompt_text="prompt",
        visual_json={
            "overlay_layers": [{
                "id": "ov_page_ref",
                "asset_id": "pending",
                "enabled": True,
                "preset": "right-card",
            }],
        },
    )
    db.add(slide)
    db.flush()
    overlay_ref = ReferenceImage(
        project_id=project.id,
        slide_id=slide.id,
        role="content_ref",
        file_path=str(overlay_path),
        process_mode="original",
        asset_name="原样保留素材",
    )
    blend_ref = ReferenceImage(
        project_id=project.id,
        slide_id=slide.id,
        role="content_ref",
        file_path=str(blend_path),
        process_mode="blend",
        asset_name="融合素材",
    )
    db.add_all([overlay_ref, blend_ref])
    db.flush()
    slide.visual_json = {
        "overlay_layers": [{
            "id": "ov_page_ref",
            "asset_id": overlay_ref.id,
            "enabled": True,
            "preset": "right-card",
        }],
    }
    db.commit()

    refs = _load_reference_images(db.query(Slide).filter(Slide.id == slide.id).one())

    assert overlay_ref.id not in [ref.get("id") for ref in refs]
    assert blend_ref.id in [ref.get("id") for ref in refs]


def test_pipeline_skips_chart_reference_overlay_as_image_input(tmp_path):
    db = make_session()
    project = Project(title="Pipeline chart reference overlay", status="prompt_ready")
    db.add(project)
    db.flush()
    chart_path = tmp_path / "chart-overlay.png"
    Image.new("RGB", (20, 20), "green").save(chart_path)
    slide = Slide(project_id=project.id, page_num=4, status="prompt_ready", prompt_text="prompt", visual_json={})
    db.add(slide)
    db.flush()
    chart_ref = ReferenceImage(
        project_id=project.id,
        slide_id=slide.id,
        role="chart_ref",
        file_path=str(chart_path),
        process_mode="original",
        asset_name="原样保留图表",
    )
    db.add(chart_ref)
    db.flush()
    slide.visual_json = {
        "overlay_layers": [{
            "id": "ov_chart_ref",
            "asset_id": chart_ref.id,
            "enabled": True,
            "preset": "center-card",
        }],
    }
    db.commit()

    refs = _load_reference_images(db.query(Slide).filter(Slide.id == slide.id).one())

    assert chart_ref.id not in [ref.get("id") for ref in refs]


def test_reference_image_library_filters_and_facets(tmp_path):
    db = make_session()
    project = Project(title="Library", status="visual_ready")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=3, status="visual_ready")
    db.add(slide)
    asset_path = tmp_path / "station.png"
    Image.new("RGB", (10, 10), "white").save(asset_path)
    asset = ReferenceImage(
        project_id=project.id,
        role="visual_asset",
        file_path=str(asset_path),
        process_mode="blend",
        asset_name="菜鸟驿站取件机",
        asset_kind="scene",
        asset_analysis={
            "source_document": "媒体介绍.pptx",
            "pptx_source_page_num": 6,
            "asset_tags": ["驿站", "取件机", "低碳"],
            "source_slide_text": "校园驿站低碳行为",
        },
    )
    db.add(asset)
    db.commit()

    library = slides_api.list_reference_images(
        project.id,
        q="取件机",
        source_document="媒体介绍.pptx",
        source_page_num=6,
        recommend_slide_id=slide.id,
        db=db,
    )

    assert library["total"] == 1
    assert library["items"][0]["id"] == asset.id
    assert library["items"][0]["source_document"] == "媒体介绍.pptx"
    assert "媒体介绍.pptx" in library["facets"]["source_documents"]


def test_default_reference_image_list_hides_legacy_low_value_visual_assets(tmp_path):
    db = make_session()
    project = Project(title="Library", status="visual_ready")
    db.add(project)
    db.flush()
    scene_path = tmp_path / "scene.png"
    core_path = tmp_path / "core.png"
    phone_path = tmp_path / "phone.png"
    Image.new("RGB", (10, 10), "white").save(scene_path)
    Image.new("RGB", (10, 10), "white").save(core_path)
    Image.new("RGB", (10, 10), "white").save(phone_path)
    db.add_all([
        ReferenceImage(
            project_id=project.id,
            role="visual_asset",
            file_path=str(scene_path),
            process_mode="blend",
            asset_name="校园风景背景",
            asset_kind="scene",
            asset_analysis={
                "source_document": "source.pptx",
                "area_ratio": 0.5,
                "source_slide_text": "校园风景与氛围背景",
            },
        ),
        ReferenceImage(
            project_id=project.id,
            role="visual_asset",
            file_path=str(core_path),
            process_mode="crop",
            asset_name="产品主视觉",
            asset_kind="product",
            asset_analysis={
                "source_document": "source.pptx",
                "selection_tier": "core_global",
                "importance_score": 35,
            },
        ),
        ReferenceImage(
            project_id=project.id,
            role="visual_asset",
            file_path=str(phone_path),
            process_mode="crop",
            asset_name="手机边框",
            asset_kind="product",
            asset_analysis={
                "source_document": "source.pptx",
                "selection_tier": "core_global",
                "importance_score": 35,
                "source_slide_text": "手机边框承载小程序界面与身份码展示",
            },
        ),
    ])
    db.commit()

    items = slides_api.list_reference_images(project.id, db=db)

    assert [item["asset_name"] for item in items] == ["产品主视觉"]
