import io

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

from app.services.document_parser import parse_document


def test_pptx_parser_extracts_grouped_text_and_table_cells():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.5)).text = "顶层标题"

    group = slide.shapes.add_group_shape()
    group.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(3), Inches(0.5)).text = "组合内关键数字 3000+"

    table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(2.0), Inches(4), Inches(1.0))
    table_shape.table.cell(0, 0).text = "指标"
    table_shape.table.cell(0, 1).text = "数据"
    table_shape.table.cell(1, 0).text = "覆盖城市"
    table_shape.table.cell(1, 1).text = "300城"

    out = io.BytesIO()
    prs.save(out)

    text = parse_document(out.getvalue(), "source.pptx")

    assert "顶层标题" in text
    assert "组合内关键数字 3000+" in text
    assert "指标 | 数据" in text
    assert "覆盖城市 | 300城" in text


def _pptx_with_full_slide_picture(label: str = "截图里的标题") -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(image)
    draw.text((80, 80), label, fill="black")
    image_buffer = io.BytesIO()
    image.save(image_buffer, format="PNG")
    image_buffer.seek(0)
    slide.shapes.add_picture(
        image_buffer,
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def test_parse_pptx_recovers_text_from_image_only_slide(monkeypatch):
    from app.services import pptx_page_recovery

    def fake_reader(image_path: str, *, page_num: int, source_filename: str, timeout_seconds: float | None = None) -> dict:
        return {
            "ocr_text": "疯火轮 AI 营销平台\n核心能力：智能投放、素材管理、数据复盘",
            "page_intent": "介绍平台核心能力",
            "key_facts": ["智能投放", "素材管理", "数据复盘"],
            "confidence": 0.91,
        }

    monkeypatch.setattr(pptx_page_recovery, "read_ppt_page_image", fake_reader)

    text = parse_document(_pptx_with_full_slide_picture(), "平台介绍.pptx")

    assert '--- PPT_SOURCE filename="平台介绍.pptx" pages=1 ---' in text
    assert "--- 第1页 ---" in text
    assert "【截图识别文字】" in text
    assert "疯火轮 AI 营销平台" in text
    assert "【页面意图】" not in text
    assert "【识别置信度】" not in text


def test_parse_pptx_recovers_only_ocr_section_from_unstructured_vision_response(monkeypatch):
    from app.services import pptx_page_recovery

    def fake_describe(*args, **kwargs) -> str:
        return """
你好，我是 PPT Agent 的读图助手。以下是截图内容的详细解读：

### 1. OCR文字
* **左侧大标题：** 团队一起用
* **右侧卡片标题：** 个人版是工具，团队版是环境

### 2. 图像内容
这是一张产品介绍页截图。

### 3. 可用于 PPT 的信息
团队协作、知识库、工作台。

### 4. 视觉参考
* **动感线条：** 闪电图形和底部运动场线条赋予了静态页面一种“速度感”。

【识别置信度】0.60
"""

    monkeypatch.setattr(pptx_page_recovery, "describe_context_image", fake_describe)

    text = parse_document(_pptx_with_full_slide_picture(), "平台介绍.pptx")

    assert "团队一起用" in text
    assert "个人版是工具，团队版是环境" in text
    assert "视觉参考" not in text
    assert "识别置信度" not in text
    assert "动感线条" not in text
    assert "PPT Agent 的读图助手" not in text


def test_parse_pptx_does_not_call_reader_for_dense_editable_text(monkeypatch):
    from app.services import pptx_page_recovery

    called = False

    def fake_reader(*args, **kwargs):
        nonlocal called
        called = True
        return {}

    monkeypatch.setattr(pptx_page_recovery, "read_ppt_page_image", fake_reader)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text = "这是足够长的可编辑正文。" * 10
    output = io.BytesIO()
    prs.save(output)

    text = parse_document(output.getvalue(), "editable.pptx")

    assert "这是足够长的可编辑正文" in text
    assert "【截图识别文字】" not in text
    assert called is False


def test_parse_pptx_reports_page_progress_and_bounds_ocr_timeout(monkeypatch):
    from app.services import pptx_page_recovery

    seen = {"timeout": None}

    def fake_reader(image_path: str, *, page_num: int, source_filename: str, timeout_seconds: float | None = None) -> dict:
        seen["timeout"] = timeout_seconds
        return {
            "ocr_text": "截图文字",
            "page_intent": "恢复原页内容",
            "key_facts": [],
            "confidence": 0.8,
        }

    monkeypatch.setattr(pptx_page_recovery, "read_ppt_page_image", fake_reader)
    events = []

    text = parse_document(
        _pptx_with_full_slide_picture(),
        "平台介绍.pptx",
        progress_callback=events.append,
        recovery_timeout_seconds=7,
    )

    assert "截图文字" in text
    assert seen["timeout"] == 7
    assert any(event.get("phase") == "pptx_page_read" and event.get("current_page") == 1 for event in events)
    assert any(event.get("phase") == "pptx_page_done" and event.get("completed_count") == 1 for event in events)
