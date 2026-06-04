from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.services.editable_pptx_diagnostics import inspect_pptx_editability


def test_inspect_pptx_editability_counts_text_shapes(tmp_path: Path):
    path = tmp_path / "editable.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, 1000000, 500000)
    box.text_frame.text = "可编辑标题"
    prs.save(path)

    result = inspect_pptx_editability(path)

    assert result.slide_count == 1
    assert result.text_shape_count == 1
    assert result.text_run_count == 1
    assert result.picture_shape_count == 0
    assert result.has_editable_text is True


def test_inspect_pptx_editability_rejects_image_only_deck(tmp_path: Path):
    image_path = tmp_path / "slide.png"
    Image.new("RGB", (1280, 720), "white").save(image_path)
    path = tmp_path / "image_only.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(13.333), height=Inches(7.5))
    prs.save(path)

    result = inspect_pptx_editability(path)

    assert result.slide_count == 1
    assert result.text_shape_count == 0
    assert result.text_run_count == 0
    assert result.picture_shape_count == 1
    assert result.has_editable_text is False
