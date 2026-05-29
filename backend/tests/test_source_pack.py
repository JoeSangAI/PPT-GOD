import io

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation

from app.services.source_pack import build_source_pack


def _sample_pdf() -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "目录\n第1章 创新使命 ........ 2\n第2章 客户痴迷 ........ 4", fontsize=12, fontname="china-s")

    page = doc.new_page()
    page.insert_text((72, 72), "第1章 创新使命\n微软被卡住了，需要重新找到存在主义使命。", fontsize=12, fontname="china-s")

    page = doc.new_page()
    page.insert_text((72, 72), "让组织愿景与个人北极星相一致。", fontsize=12, fontname="china-s")
    image = Image.new("RGB", (120, 60), "white")
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    page.insert_image(fitz.Rect(72, 140, 220, 220), stream=buf.getvalue())

    page = doc.new_page()
    page.insert_text((72, 72), "第2章 客户痴迷\n客户永远不会满足。", fontsize=12, fontname="china-s")
    return doc.tobytes()


def test_pdf_source_pack_preserves_pages_chapters_and_images(tmp_path):
    pack = build_source_pack(
        _sample_pdf(),
        "创新.pdf",
        asset_output_dir=str(tmp_path),
    )

    assert pack["document"]["filename"] == "创新.pdf"
    assert pack["document"]["kind"] == "pdf"
    assert pack["stats"]["pages"] == 4
    assert pack["pages"][1]["page_num"] == 2
    assert "微软被卡住了" in pack["pages"][1]["text"]

    chapter_titles = [chapter["title"] for chapter in pack["chapters"]]
    assert "第1章 创新使命" in chapter_titles
    assert "第2章 客户痴迷" in chapter_titles
    chapter1 = next(chapter for chapter in pack["chapters"] if chapter["title"] == "第1章 创新使命")
    assert chapter1["start_page"] == 2
    assert chapter1["end_page"] == 3

    assert len(pack["images"]) == 1
    image = pack["images"][0]
    assert image["source_page_num"] == 3
    assert image["source_document"] == "创新.pdf"
    assert image["bbox"]
    assert image["nearby_text"]


def test_slide_like_pdf_source_pack_adds_full_page_reference_image(tmp_path):
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    page.insert_text((72, 72), "Slide title")
    page.insert_text((72, 120), "Slide body")

    pack = build_source_pack(
        doc.tobytes(),
        "deck.pdf",
        asset_output_dir=str(tmp_path),
    )

    page_refs = [
        image for image in pack["images"]
        if image.get("asset_kind") == "source_page_image"
    ]
    assert len(page_refs) == 1
    ref = page_refs[0]
    assert ref["id"] == "deck.pdf:p1:page"
    assert ref["source_page_num"] == 1
    assert ref["figure_role"] == "source_page"
    assert ref["content_significance"] == "high"
    assert ref["bbox"] == [0.0, 0.0, 960.0, 540.0]
    assert ref["image_width"] >= 1280
    assert ref["image_height"] >= 720
    assert ref["nearby_text"] == "Slide title Slide body"
    from pathlib import Path
    assert Path(ref["file_path"]).exists()


def test_pdf_source_pack_orders_text_by_page_position():
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    page.insert_textbox(
        fitz.Rect(72, 160, 880, 260),
        "本融资概要仅是非凡资本为该项目而准备。\n本概要仅是为了帮助潜在购买方做出初步评估。",
        fontsize=16,
        fontname="china-s",
    )
    page.insert_text((72, 72), "免责申明", fontsize=32, fontname="china-s")

    pack = build_source_pack(doc.tobytes(), "deck.pdf")

    assert pack["pages"][0]["text"].splitlines()[0] == "免责申明"


def test_slide_like_pdf_source_pack_keeps_multicolumn_cards_together():
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    page.insert_text((72, 40), "Botlife.ai", fontsize=18)
    page.insert_text((790, 40), "新一代AI社交平台", fontsize=18, fontname="china-s")
    page.insert_text((425, 120), "核心团队", fontsize=32, fontname="china-s")
    page.insert_text((72, 190), "朱德栋", fontsize=18, fontname="china-s")
    page.insert_textbox(
        fitz.Rect(123, 185, 462, 300),
        "电子科技大学计算机系硕士。前Intel高级系统架构师；18年大型软件研发和管理经验；系统架构、智能算法领域专家；Botlife项目创始人。",
        fontsize=16,
        fontname="china-s",
    )
    page.insert_text((502, 190), "王颖", fontsize=18, fontname="china-s")
    page.insert_textbox(
        fitz.Rect(550, 185, 890, 300),
        "西南财经大学金融学士。先后在马士基、英特尔工作。10年产品设计及运营经验。Botlife项目联合创始人。",
        fontsize=16,
        fontname="china-s",
    )

    pack = build_source_pack(doc.tobytes(), "deck.pdf")
    text = pack["pages"][0]["text"]

    assert "朱德栋 电子科技大学计算机系硕士" in text
    assert "王颖 西南财经大学金融学士" in text
    assert "18年大型软件研发" in text
    assert "10年产品设计" in text
    assert text.index("朱德栋") < text.index("王颖")
    assert "朱德栋 西南财经大学" not in text
    assert "王颖 电子科技大学" not in text


def test_pdf_source_pack_uses_layout_for_inline_labels_not_domain_terms():
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    page.insert_text((72, 190), "甲乙丙", fontsize=18, fontname="china-s")
    page.insert_textbox(
        fitz.Rect(123, 185, 462, 260),
        "负责跨部门交付、客户沟通和上线验收，带领项目从需求梳理推进到最终发布。",
        fontsize=16,
        fontname="china-s",
    )

    pack = build_source_pack(doc.tobytes(), "deck.pdf")
    text = pack["pages"][0]["text"]

    assert "甲乙丙 负责跨部门交付" in text
    assert "甲乙丙负责跨部门交付" not in text


def test_pdf_source_pack_marks_repeated_header_images_as_decorative(tmp_path):
    doc = fitz.open()
    header = Image.new("RGB", (320, 24), "white")
    header_buf = io.BytesIO()
    header.save(header_buf, format="PNG")
    header_bytes = header_buf.getvalue()
    for page_num in range(1, 5):
        page = doc.new_page(width=960, height=540)
        page.insert_text((72, 120), f"第{page_num}页正文")
        page.insert_image(fitz.Rect(40, 8, 920, 42), stream=header_bytes)

    pack = build_source_pack(
        doc.tobytes(),
        "deck.pdf",
        asset_output_dir=str(tmp_path),
    )

    repeated = [
        image for image in pack["images"]
        if image.get("asset_kind") == "repeated_page_element"
    ]
    assert len(repeated) == 4
    assert {image["figure_role"] for image in repeated} == {"decorative"}
    assert {image["content_significance"] for image in repeated} == {"low"}


def test_pdf_source_pack_marks_single_top_header_strip_as_auxiliary(tmp_path):
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    page.insert_text((72, 140), "正文内容")
    header = Image.new("RGB", (320, 80), "white")
    header_buf = io.BytesIO()
    header.save(header_buf, format="PNG")
    page.insert_image(fitz.Rect(730, 18, 930, 82), stream=header_buf.getvalue())

    pack = build_source_pack(
        doc.tobytes(),
        "deck.pdf",
        asset_output_dir=str(tmp_path),
    )

    header_images = [
        image for image in pack["images"]
        if image.get("asset_kind") == "document_image"
    ]
    assert len(header_images) == 1
    assert header_images[0]["figure_role"] == "auxiliary"
    assert header_images[0]["content_significance"] == "low"


def test_pdf_source_pack_marks_terminal_brand_lockup_as_logo(tmp_path):
    doc = fitz.open()
    logo = Image.new("RGB", (320, 88), "white")
    draw = ImageDraw.Draw(logo)
    draw.rectangle((14, 18, 78, 74), fill=(220, 0, 180))
    draw.text((96, 22), "UNIQUE", fill=(20, 35, 80))
    draw.text((96, 50), "CAPITAL", fill=(20, 35, 80))
    logo_buf = io.BytesIO()
    logo.save(logo_buf, format="PNG")
    logo_bytes = logo_buf.getvalue()

    page = doc.new_page(width=960, height=540)
    page.insert_text((72, 140), "本融资概要仅是非凡资本为该项目而准备。", fontsize=16, fontname="china-s")
    page.insert_image(fitz.Rect(730, 18, 930, 82), stream=logo_bytes)

    page = doc.new_page(width=960, height=540)
    page.insert_text((72, 140), "普通正文，不应把其他页眉当作 Logo。", fontsize=16, fontname="china-s")

    page = doc.new_page(width=960, height=540)
    page.insert_text((72, 140), "本轮融资由非凡资本担任财务顾问。", fontsize=16, fontname="china-s")
    page.insert_image(fitz.Rect(725, 22, 940, 78), stream=logo_bytes)

    pack = build_source_pack(
        doc.tobytes(),
        "brand-deck.pdf",
        asset_output_dir=str(tmp_path),
    )

    logos = [
        image for image in pack["images"]
        if image.get("asset_kind") == "pdf_logo"
    ]
    assert len(logos) == 1
    assert logos[0]["classification"] == "pdf_document_logo"
    assert logos[0]["figure_role"] == "logo"
    assert logos[0]["content_significance"] == "high"
    assert logos[0]["logo_anchor"] == "top-right"


def test_pdf_source_pack_profiles_front_matter_sections_before_body_chapters():
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "前言\n作者为什么写这本书", fontsize=12, fontname="china-s")
    page = doc.new_page()
    page.insert_text((72, 72), "绪论\n大企业为什么会失去创新能力", fontsize=12, fontname="china-s")
    page = doc.new_page()
    page.insert_text((72, 72), "阅读指南\n保持开放和乐观的心态", fontsize=12, fontname="china-s")
    page = doc.new_page()
    page.insert_text((72, 72), "第1章 创新使命\n微软需要重新发现灵魂。", fontsize=12, fontname="china-s")

    pack = build_source_pack(doc.tobytes(), "创新.pdf")

    sections = pack["source_structure"]
    assert [
        (section["section_role"], section["title"], section["start_page"], section["end_page"])
        for section in sections
    ] == [
        ("preface", "前言", 1, 1),
        ("intro", "绪论", 2, 2),
        ("guide", "阅读指南", 3, 3),
        ("chapter", "第1章 创新使命", 4, 4),
    ]
    assert pack["pages"][1]["source_section_role"] == "intro"


def test_pdf_source_pack_distinguishes_content_figures_from_auxiliary_images():
    doc = fitz.open()
    large = Image.new("RGB", (640, 480), "white")
    tiny = Image.new("RGB", (48, 24), "black")
    large_buf = io.BytesIO()
    tiny_buf = io.BytesIO()
    large.save(large_buf, format="PNG")
    tiny.save(tiny_buf, format="PNG")

    page = doc.new_page()
    page.insert_text((72, 72), "框架图展示企业使命与个人北极星如何对齐", fontsize=12, fontname="china-s")
    page.insert_image(fitz.Rect(72, 120, 500, 520), stream=large_buf.getvalue())

    page = doc.new_page()
    page.insert_text((72, 72), "正文旁边有一个小装饰符号", fontsize=12, fontname="china-s")
    page.insert_image(fitz.Rect(72, 120, 110, 138), stream=tiny_buf.getvalue())

    pack = build_source_pack(doc.tobytes(), "创新.pdf")

    figures_by_page = {image["source_page_num"]: image for image in pack["images"]}
    assert figures_by_page[1]["figure_role"] == "content"
    assert figures_by_page[1]["content_significance"] == "high"
    assert figures_by_page[1]["image_width"] == 640
    assert figures_by_page[1]["image_height"] == 480
    assert figures_by_page[2]["figure_role"] == "auxiliary"
    assert figures_by_page[2]["content_significance"] == "low"


def test_markdown_source_pack_uses_headings_as_chapters():
    pack = build_source_pack(
        "# 总论\n\n开场。\n\n## 第一章 增长逻辑\n\n正文一。\n\n## 第二章 组织逻辑\n\n正文二。".encode("utf-8"),
        "讲稿.md",
    )

    assert pack["document"]["kind"] == "markdown"
    assert pack["stats"]["pages"] == 1
    assert [chapter["title"] for chapter in pack["chapters"]] == [
        "总论",
        "第一章 增长逻辑",
        "第二章 组织逻辑",
    ]
    assert "正文一" in pack["pages"][0]["text"]


def test_pptx_source_pack_uses_lightweight_text_extraction_without_sparse_recovery(monkeypatch):
    from app.services import pptx_page_recovery

    def fail_recovery(*_args, **_kwargs):
        raise AssertionError("SourcePack PPTX parsing must not call VLM sparse page recovery")

    monkeypatch.setattr(pptx_page_recovery, "recover_sparse_slide_text", fail_recovery)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "流量红利 VS 人心红利"
    slide.placeholders[1].text = "品牌增长需要从流量套利转向心智占领"
    buf = io.BytesIO()
    prs.save(buf)

    pack = build_source_pack(buf.getvalue(), "source.pptx")

    assert pack["document"]["kind"] == "pptx"
    assert pack["stats"]["pages"] == 1
    assert "流量红利 VS 人心红利" in pack["pages"][0]["text"]
    assert "心智占领" in pack["pages"][0]["text"]
