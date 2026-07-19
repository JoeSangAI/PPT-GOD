from datetime import datetime, timedelta

import pytest

from fastapi import HTTPException
from PIL import Image
from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from app.core.config import settings
from app.models.base import Base
from app.models.models import Project, ProjectRun, ReferenceImage, Slide
from app.services.artifact_versions import artifact_stale, dependency_signature, with_artifact_meta, with_stale_flags
from app.services.run_state import (
    apply_project_rollback,
    create_project_run,
    finish_run,
    image_generation_run_stage,
    normalize_confirmed_project_stage,
    reconcile_project_state,
    serialize_run,
    serialize_workflow_status,
    stale_inactive_run_if_needed,
    stale_queued_run_if_needed,
    utc_now,
    target_counts,
    update_run_progress,
)
from app.api import slides as slides_api


def make_session():
    engine = create_engine("sqlite:///:memory:", connect_args={"check_same_thread": False})
    Base.metadata.create_all(bind=engine)
    Session = sessionmaker(bind=engine)
    return Session()


def test_target_progress_is_scoped_to_run_pages():
    db = make_session()
    project = Project(title="Scoped progress", status="prompt_ready")
    db.add(project)
    db.flush()
    for page_num in range(1, 5):
        db.add(
            Slide(
                project_id=project.id,
                page_num=page_num,
                status="completed" if page_num in {1, 2, 3} else "prompt_ready",
                prompt_text="prompt",
            )
        )
    db.flush()

    run = create_project_run(
        db,
        project.id,
        kind="page_generation",
        stage="batch_generation",
        target_page_nums=[2, 3],
        total_count=2,
    )
    slides = db.query(Slide).filter(Slide.project_id == project.id).order_by(Slide.page_num).all()

    target_total, target_completed, target_failed = target_counts(run, slides)
    payload = serialize_run(run, slides)

    assert target_total == 2
    assert target_completed == 2
    assert target_failed == 0
    assert payload["completed_count"] == 2
    assert payload["completed_count"] <= payload["total_count"]


def test_image_generation_run_stage_matches_run_kind():
    assert image_generation_run_stage(kind="prototype_generation") == "prototype_generation"
    assert image_generation_run_stage(kind="page_generation") == "page_generation"
    assert image_generation_run_stage(kind="retry_failed") == "retry_failed"
    assert image_generation_run_stage(kind="finetune") == "finetune"
    assert image_generation_run_stage(kind="batch_generation") == "batch_generation"
    assert image_generation_run_stage(prototype=True) == "prototype_generation"
    assert image_generation_run_stage(page_nums=[3]) == "page_generation"


def test_active_run_guard_rejects_second_run():
    db = make_session()
    project = Project(title="One active run", status="planning")
    db.add(project)
    db.flush()

    create_project_run(db, project.id, kind="content_plan", stage="content_plan", total_count=10)

    try:
        create_project_run(db, project.id, kind="style_proposal", stage="style_proposal", total_count=3)
    except RuntimeError as exc:
        assert "已有任务正在运行" in str(exc)
    else:
        raise AssertionError("expected active run guard to reject a second run")


def test_stale_queued_run_unblocks_next_run():
    db = make_session()
    project = Project(title="Stale queued run", status="visual_ready")
    db.add(project)
    db.flush()

    stale_run = create_project_run(db, project.id, kind="style_proposal", stage="style_proposal", total_count=1)
    stale_run.started_at = utc_now() - timedelta(seconds=300)
    db.flush()

    stale_queued_run_if_needed(db, project.id, timeout_seconds=120)
    next_run = create_project_run(db, project.id, kind="style_proposal", stage="style_proposal", total_count=1)

    assert stale_run.status == "stale"
    assert "后台生成服务未接收" in stale_run.error_msg
    assert next_run.id != stale_run.id
    assert next_run.status == "queued"


def test_celery_queued_run_uses_longer_queue_timeout():
    db = make_session()
    project = Project(title="Queued Celery run", status="prompt_ready")
    db.add(project)
    db.flush()

    run = create_project_run(
        db,
        project.id,
        kind="prototype_generation",
        stage="batch_generation",
        total_count=1,
    )
    run.task_id = "queued-task"
    run.started_at = utc_now() - timedelta(seconds=300)
    db.flush()

    stale_inactive_run_if_needed(db, project.id)

    assert run.status == "queued"


def test_stale_running_run_without_heartbeat():
    db = make_session()
    project = Project(title="Dead background task", status="visual_ready")
    db.add(project)
    db.flush()

    run = create_project_run(db, project.id, kind="visual_prompts", stage="visual_planning", total_count=10)
    update_run_progress(db, run.id, completed_count=5, message="正在生成视觉方案")
    run.updated_at = utc_now() - timedelta(seconds=600)
    db.flush()

    stale_inactive_run_if_needed(db, project.id, heartbeat_timeout_seconds=300)

    assert run.status == "stale"
    assert "没有进度更新" in run.error_msg


def test_content_plan_running_run_uses_longer_default_heartbeat(monkeypatch):
    db = make_session()
    project = Project(title="Long manuscript", status="planning")
    db.add(project)
    db.flush()
    monkeypatch.setattr(settings, "RUN_HEARTBEAT_TIMEOUT_SECONDS", 300)
    monkeypatch.setattr(settings, "CONTENT_PLAN_HEARTBEAT_TIMEOUT_SECONDS", 1800)

    run = create_project_run(db, project.id, kind="content_plan", stage="analyzing", total_count=50)
    update_run_progress(db, run.id, message="正在读取材料")
    run.updated_at = utc_now() - timedelta(seconds=600)
    db.flush()

    stale_inactive_run_if_needed(db, project.id)

    assert run.status == "running"


def test_non_content_plan_running_run_keeps_default_heartbeat(monkeypatch):
    db = make_session()
    project = Project(title="Visual planning", status="visual_ready")
    db.add(project)
    db.flush()
    monkeypatch.setattr(settings, "RUN_HEARTBEAT_TIMEOUT_SECONDS", 300)
    monkeypatch.setattr(settings, "CONTENT_PLAN_HEARTBEAT_TIMEOUT_SECONDS", 1800)

    run = create_project_run(db, project.id, kind="visual_prompts", stage="visual_planning", total_count=10)
    update_run_progress(db, run.id, completed_count=1, message="正在生成视觉方案")
    run.updated_at = utc_now() - timedelta(seconds=600)
    db.flush()

    stale_inactive_run_if_needed(db, project.id)

    assert run.status == "stale"


def test_content_plan_writeback_guard_rejects_cancelled_run():
    db = make_session()
    project = Project(title="Cancelled content plan", status="planning")
    db.add(project)
    db.flush()

    run = create_project_run(db, project.id, kind="content_plan", stage="content_plan", total_count=10)
    finish_run(db, run.id, status="cancelled", message="任务被取消", error_msg="任务被取消")
    db.flush()

    class Logger:
        def info(self, *_args):
            return None

    assert not slides_api._content_plan_run_is_active_for_writeback(
        db,
        run.id,
        Logger(),
        project.id,
        "before_slide_replace",
    )


def test_content_plan_writeback_guard_rejects_superseded_active_run():
    db = make_session()
    project = Project(title="Superseded content plan", status="planning")
    db.add(project)
    db.flush()

    older = create_project_run(db, project.id, kind="content_plan", stage="content_plan", total_count=50)
    older.status = "running"
    older.started_at = utc_now() - timedelta(seconds=30)
    newer = ProjectRun(
        project_id=project.id,
        kind="content_plan",
        status="running",
        stage="content_plan",
        total_count=50,
        started_at=utc_now(),
    )
    db.add(newer)
    db.flush()

    class Logger:
        def info(self, *_args):
            return None

    assert not slides_api._content_plan_run_is_active_for_writeback(
        db,
        older.id,
        Logger(),
        project.id,
        "before_slide_replace",
    )
    assert slides_api._content_plan_run_is_active_for_writeback(
        db,
        newer.id,
        Logger(),
        project.id,
        "before_slide_replace",
    )


def test_cancelled_content_plan_exception_does_not_reset_project_status(monkeypatch):
    db = make_session()
    project = Project(title="Cancelled late failure", status="planning")
    db.add(project)
    db.flush()
    run = create_project_run(db, project.id, kind="content_plan", stage="content_plan", total_count=50)
    finish_run(db, run.id, status="cancelled", message="任务被取消", error_msg="任务被取消")
    db.commit()

    def fail_load_source_packs(*_args, **_kwargs):
        raise RuntimeError("late background failure")

    monkeypatch.setattr(slides_api, "SessionLocal", lambda: db)
    monkeypatch.setattr(db, "close", lambda: None)
    monkeypatch.setattr(slides_api, "load_project_source_packs", fail_load_source_packs)

    slides_api._generate_content_plan_bg(
        project.id,
        "取消后旧后台任务失败",
        page_count=50,
        run_id=run.id,
    )

    db.refresh(project)
    db.refresh(run)
    assert project.status == "planning"
    assert run.status == "cancelled"


def test_serialize_run_marks_naive_datetimes_as_utc_for_browser_math():
    db = make_session()
    project = Project(title="UTC progress", status="planning")
    db.add(project)
    db.flush()

    run = create_project_run(db, project.id, kind="content_plan", stage="content_plan", total_count=10)
    run.started_at = datetime(2026, 5, 14, 13, 47, 0, 123456)
    run.updated_at = datetime(2026, 5, 14, 13, 47, 5)
    run.finished_at = datetime(2026, 5, 14, 13, 48, 0)
    db.flush()

    payload = serialize_run(run)

    assert payload["started_at"] == "2026-05-14T13:47:00.123456Z"
    assert payload["updated_at"] == "2026-05-14T13:47:05Z"
    assert payload["finished_at"] == "2026-05-14T13:48:00Z"


def test_confirmed_planning_project_normalizes_to_visual_ready():
    db = make_session()
    project = Project(title="Legacy confirmed", status="planning", content_plan_confirmed=True)
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, status="pending", content_json={"page_num": 1})
    db.add(slide)
    db.flush()

    normalize_confirmed_project_stage(project, [slide])

    assert project.status == "visual_ready"


def test_reconcile_clears_downstream_outputs_without_selected_style():
    db = make_session()
    project = Project(title="Broken nonlinear state", status="prompt_ready", content_plan_confirmed=True)
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        content_json={"page_num": 1},
        visual_json={"layout": "old"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.flush()

    reconcile_project_state(project, [slide])

    assert project.status == "visual_ready"
    assert slide.visual_json == {}
    assert slide.prompt_text is None
    assert slide.image_path is None
    assert slide.status == "pending"


def test_reconcile_preserves_signed_visual_outputs_when_content_changes():
    db = make_session()
    project = Project(
        title="Signed outputs",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        content_json={"page_num": 1, "title": "旧标题"},
        visual_json={"layout": "old"},
        prompt_text="old prompt",
    )
    db.add(slide)
    db.flush()
    slide.visual_json = with_artifact_meta(
        slide.visual_json,
        kind="visual_plan",
        dependencies=dependency_signature(project, [slide]),
    )
    slide.content_json = {"page_num": 1, "title": "新标题"}
    db.flush()

    reconcile_project_state(project, [slide])

    assert project.status == "prompt_ready"
    assert slide.prompt_text == "old prompt"
    assert slide.status == "prompt_ready"
    assert slide.visual_json["layout"] == "old"


def test_reconcile_keeps_completed_prototype_when_dependency_metadata_lags():
    db = make_session()
    project = Project(
        title="Completed prototype",
        status="prototype_ready",
        content_plan_confirmed=True,
    )
    db.add(project)
    db.flush()
    sample = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        content_json={"page_num": 1},
        visual_json={"layout": "sample"},
        prompt_text="sample prompt",
        image_path="/tmp/sample.png",
    )
    remaining = Slide(
        project_id=project.id,
        page_num=2,
        status="prompt_ready",
        content_json={"page_num": 2},
        visual_json={"layout": "remaining"},
        prompt_text="remaining prompt",
    )
    db.add_all([sample, remaining])
    db.flush()
    old_deps = dependency_signature(project, [sample, remaining])
    sample.visual_json = with_artifact_meta(sample.visual_json, kind="visual_plan", dependencies=old_deps)
    remaining.visual_json = with_artifact_meta(remaining.visual_json, kind="visual_plan", dependencies=old_deps)
    project.selected_style = {"name": "Brand"}
    run = create_project_run(
        db,
        project.id,
        kind="prototype_generation",
        stage="batch_generation",
        target_page_nums=[1],
        total_count=1,
    )
    run.status = "succeeded"
    db.flush()

    reconcile_project_state(project, [sample, remaining], run)

    assert project.status == "prototype_ready"
    assert sample.status == "completed"
    assert sample.prompt_text == "sample prompt"
    assert sample.image_path == "/tmp/sample.png"
    assert remaining.status == "prompt_ready"
    assert remaining.prompt_text == "remaining prompt"


def test_cancelled_prototype_with_preserved_old_image_returns_to_prompt_ready():
    db = make_session()
    project = Project(
        title="Cancelled prototype",
        status="prototype_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        content_json={"page_num": 1},
        visual_json={"layout": "new"},
        prompt_text="new prompt",
        image_path="/tmp/old-sample.png",
        error_msg="用户手动停止",
    )
    db.add(slide)
    db.flush()
    run = create_project_run(
        db,
        project.id,
        kind="prototype_generation",
        stage="batch_generation",
        target_page_nums=[1],
        total_count=1,
    )
    run.status = "cancelled"
    run.message = "用户手动停止"
    run.error_msg = "用户手动停止"
    db.flush()

    reconcile_project_state(project, [slide], run)

    assert project.status == "prompt_ready"
    assert slide.image_path == "/tmp/old-sample.png"
    assert slide.status == "prompt_ready"


def test_update_slide_type_accepts_canonical_toc_and_syncs_content_json():
    db = make_session()
    project = Project(title="Manual type", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=2,
        type="content",
        content_json={
            "page_num": 2,
            "type": "content",
            "section_title": "内容地图",
            "text_content": {"headline": "课程地图", "body": "- 模块一\n- 模块二\n- 模块三"},
        },
    )
    db.add(slide)
    db.flush()

    result = slides_api.update_slide_type(
        project.id,
        slides_api.UpdateTypeRequest(page_num=2, slide_id=slide.id, type="toc"),
        db=db,
    )
    db.refresh(slide)

    assert result["type"] == "toc"
    assert slide.type == "toc"
    assert slide.content_json["type"] == "toc"
    assert slide.type_locked is True


def test_update_slide_type_rejects_legacy_or_unknown_values():
    db = make_session()
    project = Project(title="Strict manual type", status="planning")
    db.add(project)
    db.flush()
    slide = Slide(project_id=project.id, page_num=1, type="content", content_json={"type": "content"})
    db.add(slide)
    db.flush()

    for invalid_type in ("agenda", "content_split", "content_magic"):
        with pytest.raises(HTTPException) as exc_info:
            slides_api.update_slide_type(
                project.id,
                slides_api.UpdateTypeRequest(page_num=1, slide_id=slide.id, type=invalid_type),
                db=db,
            )
        assert exc_info.value.status_code == 400

    db.refresh(slide)
    assert slide.type == "content"
    assert slide.content_json["type"] == "content"


def test_rollback_to_prompt_requires_selected_style():
    db = make_session()
    project = Project(title="Prompt rollback without style", status="completed", content_plan_confirmed=True)
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        content_json={"page_num": 1},
        visual_json={"layout": "old"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.flush()

    apply_project_rollback(project, [slide], "prompt_ready")

    assert project.status == "visual_ready"
    assert project.selected_style is None
    assert slide.visual_json == {}
    assert slide.prompt_text is None
    assert slide.image_path is None


def test_workflow_status_exposes_unified_progress_view_model():
    db = make_session()
    project = Project(title="Unified progress", status="prompt_ready")
    db.add(project)
    db.flush()
    for page_num in range(1, 4):
        db.add(
            Slide(
                project_id=project.id,
                page_num=page_num,
                status="completed" if page_num == 1 else "generating",
                prompt_text="prompt",
            )
        )
    db.flush()
    run = create_project_run(
        db,
        project.id,
        kind="batch_generation",
        stage="batch_generation",
        target_page_nums=[1, 2, 3],
        total_count=3,
        message="正在生成图片",
    )
    update_run_progress(db, run.id, completed_count=1, failed_count=0)
    slides = db.query(Slide).filter(Slide.project_id == project.id).order_by(Slide.page_num).all()

    payload = serialize_workflow_status(project, slides, active_run=run, latest_run=run)

    assert payload["project_phase"] == "prompt_ready"
    assert payload["project_status"] == "prompt_ready"
    assert payload["active_run"]["id"] == run.id
    assert payload["active_run"]["updated_at"]
    assert payload["progress"]["label"] == "批量生成进度"
    assert payload["progress"]["current"] == 1
    assert payload["progress"]["total"] == 3
    assert payload["progress"]["unit"] == "页"
    assert payload["progress"]["updated_at"] == payload["active_run"]["updated_at"]
    assert payload["progress"]["active_page_nums"] == [2, 3]
    assert payload["progress"]["running_count"] == 2
    assert payload["progress"]["can_cancel"] is True


def test_content_dependent_invalidation_preserves_generated_outputs():
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Preserve generated deck",
        status="completed",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        content_json={"page_num": 1},
        visual_json={"layout": "old"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.flush()

    slides_api._invalidate_content_dependent_outputs(project)

    assert project.status == "completed"
    assert project.content_plan_confirmed is True
    assert slide.prompt_text == "old prompt"
    assert slide.image_path == "/tmp/old.png"
    assert slide.visual_json["layout"] == "old"
    assert artifact_stale(slide.visual_json) == {"content": True}


def test_visual_invalidation_preserves_image_and_marks_prompt_stale():
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Preserve image",
        status="prototype_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="completed",
        content_json={"page_num": 1},
        visual_json={"layout": "old"},
        prompt_text="old prompt",
        image_path="/tmp/old.png",
    )
    db.add(slide)
    db.flush()

    slides_api._invalidate_visual_plan_dependent_outputs(project, [slide])

    assert project.status == "prototype_ready"
    assert slide.prompt_text == "old prompt"
    assert slide.image_path == "/tmp/old.png"
    assert slide.status == "completed"
    assert artifact_stale(slide.visual_json) == {"visual": True}


def test_start_generation_blocks_pages_with_unapplied_plan_changes():
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Block stale generation",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    db.add(
        Slide(
            project_id=project.id,
            page_num=1,
            status="completed",
            content_json={"page_num": 1},
            visual_json={"visual_description": "old view", "_artifact": {"stale": {"content": True}}},
            prompt_text="old prompt",
            image_path="/tmp/old.png",
        )
    )
    db.commit()

    try:
        slides_api.start_generation(project.id, slides_api.PageNumsRequest(page_nums=[1]), db)
    except HTTPException as exc:
        assert exc.status_code == 400
        assert "未应用的内容或画面修改" in exc.detail
    else:
        raise AssertionError("expected stale page to block direct generation")


def test_generation_preflight_allows_legacy_prompt_when_only_another_slide_changed():
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Legacy deck dependency",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    first = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        content_json={
            "page_num": 1,
            "type": "content",
            "text_content": {"headline": "第一页标题", "body": "第一页正文"},
        },
        visual_json={"page_num": 1, "visual_description": "第一页画面"},
        prompt_text=(
            "Visible Text:\n"
            "Headline: \"第一页标题\"\n"
            "Body: \"第一页正文\"\n\n"
            "Layout:\n"
            "第一页画面"
        ),
    )
    second = Slide(
        project_id=project.id,
        page_num=2,
        status="prompt_ready",
        content_json={
            "page_num": 2,
            "type": "content",
            "text_content": {"headline": "第二页旧标题", "body": "第二页正文"},
        },
        visual_json={"page_num": 2, "visual_description": "第二页画面"},
        prompt_text=(
            "Visible Text:\n"
            "Headline: \"第二页旧标题\"\n"
            "Body: \"第二页正文\"\n\n"
            "Layout:\n"
            "第二页画面"
        ),
    )
    db.add_all([first, second])
    db.flush()
    legacy_deck_deps = dependency_signature(project, [first, second])
    first.visual_json = with_artifact_meta(
        first.visual_json,
        kind="visual_plan",
        dependencies=legacy_deck_deps,
        prompt_dependencies=legacy_deck_deps,
        prompt_visual_signature=slides_api._prompt_visual_plan_signature(first.visual_json),
    )
    second.visual_json = with_artifact_meta(
        second.visual_json,
        kind="visual_plan",
        dependencies=legacy_deck_deps,
        prompt_dependencies=legacy_deck_deps,
        prompt_visual_signature=slides_api._prompt_visual_plan_signature(second.visual_json),
    )
    db.flush()

    second.content_json = {
        "page_num": 2,
        "type": "content",
        "text_content": {"headline": "第二页新标题", "body": "第二页正文"},
    }
    db.flush()
    slides = db.query(Slide).filter(Slide.project_id == project.id).order_by(Slide.page_num).all()
    first = slides[0]

    slides_api._raise_if_generation_inputs_not_current(project, slides, [first])


def test_start_generation_refreshes_prompt_when_exact_overlay_syncs(monkeypatch, tmp_path):
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Exact overlay generation",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=6,
        status="prompt_ready",
        content_json={
            "page_num": 6,
            "type": "content",
            "text_content": {
                "headline": "创意样张",
                "body": "把公众号事件声量转成电梯里的即时提醒",
            },
        },
        visual_json={"page_num": 6, "visual_description": "三张电梯提醒海报并排展示"},
        prompt_text=(
            "Visible Text:\n"
            "Headline: \"创意样张\"\n"
            "Body: \"把公众号事件声量转成电梯里的即时提醒\"\n\n"
            "Layout:\n"
            "旧 Prompt 没有为原样保留素材预留位置。"
        ),
    )
    db.add(slide)
    db.flush()
    image_path = tmp_path / "page-ref.png"
    Image.new("RGB", (160, 90), "white").save(image_path)
    page_ref = ReferenceImage(
        project_id=project.id,
        slide_id=slide.id,
        role="content_ref",
        file_path=str(image_path),
        process_mode="original",
        asset_name="原样保留素材",
    )
    db.add(page_ref)
    db.commit()

    refreshed_prompts = []

    def fake_generate_prompt_for_page(**kwargs):
        page_intent = kwargs["page_intent"]
        refreshed_prompts.append(page_intent)
        assert page_intent["overlay_layers"][0]["asset_id"] == page_ref.id
        return (
            "Visible Text:\n"
            "Headline: \"创意样张\"\n"
            "Body: \"把公众号事件声量转成电梯里的即时提醒\"\n\n"
            "Exact Overlay Reservation:\n"
            "Reserve a clean slot for the original page asset."
        )

    class FakeTask:
        id = "task-1"

    monkeypatch.setattr(slides_api, "generate_prompt_for_page", fake_generate_prompt_for_page)
    monkeypatch.setattr(slides_api, "ensure_generation_worker_ready", lambda: True)
    monkeypatch.setattr(slides_api, "_enqueue_generation_task", lambda *args, **kwargs: FakeTask())

    result = slides_api.start_generation(
        project.id,
        body=slides_api.PageNumsRequest(page_nums=[6], prototype=True),
        db=db,
    )

    refreshed = db.query(Slide).filter(Slide.id == slide.id).one()
    assert result["message"] == "Generation started"
    assert refreshed_prompts
    assert refreshed.visual_json["overlay_layers"][0]["asset_id"] == page_ref.id
    assert artifact_stale(refreshed.visual_json) == {}
    assert "Exact Overlay Reservation" in refreshed.prompt_text


def test_generation_preflight_repairs_residual_visual_stale_when_inputs_match(tmp_path):
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Residual visual stale",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()

    image_path = tmp_path / "page-ref.png"
    Image.new("RGB", (160, 90), "white").save(image_path)
    slide = Slide(
        project_id=project.id,
        page_num=6,
        status="prompt_ready",
        content_json={
            "page_num": 6,
            "type": "content",
            "text_content": {
                "headline": "创意样张",
                "body": "把公众号事件声量转成电梯里的即时提醒",
            },
        },
        visual_json={
            "page_num": 6,
            "visual_description": "新画面，预留本页素材卡位",
            "overlay_layers": [
                {
                    "id": "ov_ref",
                    "asset_id": "ref-6",
                    "enabled": True,
                    "preset": "right-card",
                    "fit": "contain",
                    "mode": "exact_cutout",
                    "valign": "bottom",
                    "usage_note": "本页素材：原图保留",
                    "z_index": 0,
                }
            ],
        },
        prompt_text=(
            "Visible Text:\n"
            "Headline: \"创意样张\"\n"
            "Body: \"把公众号事件声量转成电梯里的即时提醒\"\n\n"
            "References:\n"
            "- Page reference: use this uploaded image as the page visual source. Context: page-ref.png\n\n"
            "Layout:\n"
            "预留本页素材卡位。"
        ),
    )
    ref = ReferenceImage(
        id="ref-6",
        project_id=project.id,
        slide_id=slide.id,
        role="content_ref",
        file_path=str(image_path),
        process_mode="original",
        asset_name="本页素材",
    )
    db.add_all([slide, ref])
    db.flush()
    project_slides = db.query(Slide).filter(Slide.project_id == project.id).order_by(Slide.page_num).all()
    deps = dependency_signature(project, project_slides)
    slide.visual_json = with_artifact_meta(
        slide.visual_json,
        kind="visual_plan",
        dependencies=deps,
        prompt_dependencies=deps,
    )
    slide.visual_json = with_stale_flags(slide.visual_json, visual=True)
    db.commit()

    refreshed = db.query(Slide).filter(Slide.id == slide.id).one()
    assert artifact_stale(refreshed.visual_json) == {"visual": True}

    slides_api._raise_if_generation_inputs_not_current(project, [refreshed], [refreshed])

    assert artifact_stale(refreshed.visual_json) == {}


def test_generation_preflight_keeps_visual_stale_when_visual_signature_changed():
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Changed visual stays stale",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    slide = Slide(
        project_id=project.id,
        page_num=1,
        status="prompt_ready",
        content_json={
            "page_num": 1,
            "type": "content",
            "text_content": {"headline": "当前标题"},
        },
        visual_json={
            "page_num": 1,
            "visual_description": "旧画面",
        },
        prompt_text=(
            "Visible Text:\n"
            "Headline: \"当前标题\"\n\n"
            "Layout:\n"
            "旧画面"
        ),
    )
    db.add(slide)
    db.flush()
    deps = dependency_signature(project, [slide])
    old_signature = slides_api._prompt_visual_plan_signature(slide.visual_json)
    slide.visual_json = with_artifact_meta(
        slide.visual_json,
        kind="visual_plan",
        dependencies=deps,
        prompt_dependencies=deps,
        prompt_visual_signature=old_signature,
    )
    slide.visual_json["visual_description"] = "新的手工画面描述"
    slide.visual_json = with_stale_flags(slide.visual_json, visual=True)
    db.commit()

    refreshed = db.query(Slide).filter(Slide.id == slide.id).one()
    try:
        slides_api._raise_if_generation_inputs_not_current(project, [refreshed], [refreshed])
    except HTTPException as exc:
        assert exc.status_code == 400
        assert "未应用的内容或画面修改" in exc.detail
    else:
        raise AssertionError("expected changed visual signature to keep blocking generation")


def test_start_generation_blocks_visual_guidance_with_old_visible_labels():
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Old visual labels guard",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    db.add(
        Slide(
            project_id=project.id,
            page_num=9,
            status="prompt_ready",
            content_json={
                "page_num": 9,
                "type": "content",
                "visual_suggestion": "下方三栏讲山姆/胖东来/Costco，底部大字写个人外脑。",
                "text_content": {
                    "headline": "消费者雇了一个“全知全能”的私人买手",
                    "body": "过去搜索防晒霜：自己翻几百个商品、几千条评价",
                },
            },
            visual_json={
                "visual_description": "根据当前正文做极简信息页。",
            },
            prompt_text=(
                "Visible Text:\n"
                "Headline: \"消费者雇了一个“全知全能”的私人买手\"\n"
                "Body: \"过去搜索防晒霜：自己翻几百个商品、几千条评价\"\n"
                "\n\nStyle: test"
            ),
        )
    )
    db.commit()

    try:
        slides_api.start_generation(project.id, slides_api.PageNumsRequest(page_nums=[9]), db)
    except HTTPException as exc:
        assert exc.status_code == 400
        assert "当前内容不一致" in exc.detail
    else:
        raise AssertionError("expected old visual labels to block direct generation")


def test_start_generation_blocks_prompt_visible_text_from_old_content():
    from app.api import slides as slides_api

    db = make_session()
    project = Project(
        title="Old prompt visible text guard",
        status="prompt_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()
    db.add(
        Slide(
            project_id=project.id,
            page_num=1,
            status="prompt_ready",
            content_json={
                "page_num": 1,
                "type": "content",
                "text_content": {
                    "headline": "当前标题",
                    "body": "这是最终确认的正文",
                },
            },
            visual_json={"visual_description": "根据当前正文做极简信息页。"},
            prompt_text=(
                "Visible Text:\n"
                "Headline: \"当前标题\"\n"
                "Body: \"这是上一版已经删除的正文\"\n"
                "\n\nStyle: test"
            ),
        )
    )
    db.commit()

    try:
        slides_api.start_generation(project.id, slides_api.PageNumsRequest(page_nums=[1]), db)
    except HTTPException as exc:
        assert exc.status_code == 400
        assert "生图 Prompt 不是基于最新内容" in exc.detail
    else:
        raise AssertionError("expected stale prompt visible text to block direct generation")


def test_download_exports_partial_deck_with_blank_pages(tmp_path, monkeypatch):
    from app.api import slides as slides_api

    db = make_session()
    monkeypatch.setattr(slides_api.settings, "OUTPUT_DIR", str(tmp_path))
    project = Project(
        title="Partial export",
        status="prototype_ready",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()

    generated_path = tmp_path / project.id / "slide_01.png"
    generated_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (1792, 1024), "white").save(generated_path)

    db.add_all([
        Slide(
            project_id=project.id,
            page_num=1,
            status="completed",
            content_json={"page_num": 1},
            visual_json={"visual_description": "old view", "_artifact": {"stale": {"image": True}}},
            prompt_text="new prompt",
            image_path=str(generated_path),
        ),
        Slide(
            project_id=project.id,
            page_num=2,
            status="prompt_ready",
            content_json={"page_num": 2},
            visual_json={"visual_description": "planned"},
            prompt_text="prompt",
        ),
        Slide(
            project_id=project.id,
            page_num=3,
            status="prompt_ready",
            content_json={"page_num": 3},
            visual_json={"visual_description": "planned"},
            prompt_text="prompt",
        ),
    ])
    db.commit()

    response = slides_api.download_pptx(project.id, db=db)
    prs = Presentation(response.path)

    assert len(prs.slides) == 3
    assert len(prs.slides[0].shapes) >= 1
    assert len(prs.slides[1].shapes) == 0
    assert len(prs.slides[2].shapes) == 0


def test_download_completed_project_serves_final_presentation(tmp_path, monkeypatch):
    from app.api import slides as slides_api

    db = make_session()
    monkeypatch.setattr(slides_api.settings, "OUTPUT_DIR", str(tmp_path))
    project = Project(
        title="Completed export",
        status="completed",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()

    original_slide_path = tmp_path / project.id / "slide_01.png"
    original_slide_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (1792, 1024), "red").save(original_slide_path)

    final_pptx_path = tmp_path / project.id / "presentation.pptx"
    assemble_input = [{"page_num": 1, "image_path": str(original_slide_path), "speaker_notes": ""}]
    slides_api.assemble_pptx(assemble_input, str(final_pptx_path))

    db.add(
        Slide(
            project_id=project.id,
            page_num=1,
            status="completed",
            content_json={"page_num": 1},
            visual_json={"visual_description": "final"},
            prompt_text="prompt",
            image_path=str(original_slide_path),
        )
    )
    db.commit()

    def fail_partial_assemble(*_args, **_kwargs):
        raise AssertionError("completed downloads should serve the final presentation")

    monkeypatch.setattr(slides_api, "_assemble_partial_project_pptx", fail_partial_assemble)

    response = slides_api.download_pptx(project.id, db=db)

    assert response.path == str(final_pptx_path)


def test_download_completed_project_rebuilds_stale_final_presentation_with_missing_pages(tmp_path, monkeypatch):
    from app.api import slides as slides_api

    db = make_session()
    monkeypatch.setattr(slides_api.settings, "OUTPUT_DIR", str(tmp_path))
    project = Project(
        title="Completed export",
        status="completed",
        content_plan_confirmed=True,
        selected_style={"name": "Brand"},
    )
    db.add(project)
    db.flush()

    project_dir = tmp_path / project.id
    project_dir.mkdir(parents=True, exist_ok=True)
    slide_paths = []
    for page_num, color in [(1, "red"), (2, "blue")]:
        slide_path = project_dir / f"slide_{page_num:02d}.png"
        Image.new("RGB", (1792, 1024), color).save(slide_path)
        slide_paths.append(slide_path)
        db.add(
            Slide(
                project_id=project.id,
                page_num=page_num,
                status="completed",
                content_json={"page_num": page_num},
                visual_json={"visual_description": f"page {page_num}"},
                prompt_text="prompt",
                image_path=str(slide_path),
            )
        )

    final_pptx_path = project_dir / "presentation.pptx"
    slides_api.assemble_pptx(
        [{"page_num": 1, "image_path": str(slide_paths[0]), "speaker_notes": ""}],
        str(final_pptx_path),
    )
    db.commit()

    response = slides_api.download_pptx(project.id, db=db)
    rebuilt = Presentation(response.path)

    assert response.path == str(final_pptx_path)
    assert len(rebuilt.slides) == 2
