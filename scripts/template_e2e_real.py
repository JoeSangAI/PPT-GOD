from __future__ import annotations

import os
import sys
import tempfile
import time
from pathlib import Path

import fitz
import requests
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


API = os.environ.get("PPTGOD_E2E_API", "http://localhost:8000")
KEEP_PROJECTS = os.environ.get("KEEP_E2E_PROJECTS") == "1"
TIMEOUT = 120
TEMPLATE_KEYS = ("cover", "toc", "section", "content", "data", "quote", "ending")

session = requests.Session()


def request_json(method: str, path: str, *, expected=(200,), **kwargs):
    response = session.request(method, API + path, timeout=TIMEOUT, **kwargs)
    if response.status_code not in expected:
        raise AssertionError(f"{method} {path} -> {response.status_code}: {response.text[:500]}")
    if response.content and "application/json" in response.headers.get("content-type", ""):
        return response.json()
    return response


def wait_template(project_id: str, job_id: str | None = None, timeout: int = 180) -> dict:
    start = time.time()
    last = None
    while time.time() - start < timeout:
        status = request_json("GET", f"/projects/{project_id}/template-status")
        if job_id and status.get("job_id") and status.get("job_id") != job_id:
            time.sleep(0.5)
            continue
        last = status
        if status.get("status") == "completed":
            return status
        if status.get("status") == "failed":
            raise AssertionError(f"template extraction failed: {status}")
        time.sleep(0.75)
    raise AssertionError(f"template extraction timeout: {last}")


def wait_template_failed(project_id: str, job_id: str, timeout: int = 120) -> dict:
    start = time.time()
    last = None
    while time.time() - start < timeout:
        status = request_json("GET", f"/projects/{project_id}/template-status")
        if status.get("job_id") != job_id:
            time.sleep(0.5)
            continue
        last = status
        if status.get("status") == "failed":
            return status
        if status.get("status") == "completed":
            raise AssertionError(f"corrupt template unexpectedly completed: {status}")
        time.sleep(0.75)
    raise AssertionError(f"corrupt template did not fail in time: {last}")


def add_text(slide, x, y, w, h, text, *, size=24, bold=False, color=(20, 24, 36)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def make_logo_png(path: Path, text="ACME") -> Path:
    image = Image.new("RGBA", (300, 90), (0, 0, 0, 0))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((4, 8, 286, 78), radius=18, fill=(18, 24, 38, 245))
    draw.ellipse((22, 22, 76, 76), fill=(72, 211, 153, 255))
    draw.text((92, 27), text, fill=(255, 255, 255, 255))
    image.save(path)
    return path


def make_manual_logo(path: Path) -> Path:
    image = Image.new("RGBA", (240, 80), (0, 0, 0, 0))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((6, 10, 232, 70), radius=14, fill=(220, 38, 38, 255))
    draw.text((28, 30), "MANUAL LOGO", fill=(255, 255, 255, 255))
    image.save(path)
    return path


def make_template_pptx(path: Path, logo_path: Path, *, pages=7) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    labels = [
        ("封面标题", "cover"),
        ("目录 Agenda\n01 背景\n02 方案\n03 预算", "toc"),
        ("章节 01\n增长机会", "section"),
        ("内容页标题\n这里是正文占位，展示左右分栏和信息层级。", "content"),
        ("数据指标\n同比 +42%\nROI 3.8x\nKPI 完成率 91%", "data"),
        ("核心观点 quote\n一句强有力的结论放在这里", "quote"),
        ("谢谢观看", "ending"),
    ]
    while len(labels) < pages:
        labels.insert(-1, (f"内容页标题 {len(labels)}\n正文占位信息，不包含真实业务正文。", "content"))
    for index, (text, kind) in enumerate(labels[:pages], start=1):
        slide = presentation.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(248, 249, 251)
        slide.shapes.add_picture(str(logo_path), Inches(11.55), Inches(0.28), width=Inches(1.25))
        bar = slide.shapes.add_shape(1, Inches(0.65), Inches(0.6), Inches(0.08), Inches(6.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(72, 211, 153)
        add_text(slide, 0.95, 0.9, 7.4, 1.4, text, size=34 if kind in {"cover", "section"} else 24, bold=True)
        if kind == "content":
            add_text(slide, 0.95, 2.55, 4.5, 2.4, "正文占位\n- 要点一\n- 要点二\n- 要点三", size=18)
            block = slide.shapes.add_shape(1, Inches(7.2), Inches(2.1), Inches(4.6), Inches(3.3))
            block.fill.solid()
            block.fill.fore_color.rgb = RGBColor(223, 232, 239)
        if kind == "data":
            for offset, value in enumerate(("42%", "3.8x", "91%")):
                add_text(slide, 1.0 + offset * 3.6, 2.8, 2.5, 1.2, value, size=42, bold=True, color=(15, 118, 110))
        if kind == "quote":
            block = slide.shapes.add_shape(1, Inches(1.0), Inches(2.2), Inches(10.5), Inches(2.6))
            block.fill.solid()
            block.fill.fore_color.rgb = RGBColor(232, 240, 254)
        add_text(slide, 0.95, 6.95, 2, 0.3, f"{index:02d}", size=10, color=(120, 130, 145))
    presentation.save(path)
    return path


def make_finished_pptx(path: Path, logo_path: Path) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    dense = " ".join(["这是成品 PPT 正文内容，包含真实段落、业务描述、事实、结论和行动项。"] * 12)
    for index in range(1, 11):
        slide = presentation.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        slide.shapes.add_picture(str(logo_path), Inches(11.55), Inches(0.28), width=Inches(1.25))
        add_text(slide, 0.8, 0.65, 9.8, 0.7, f"第 {index} 页：真实业务章节标题", size=26, bold=True)
        add_text(slide, 0.8, 1.55, 11.2, 4.8, dense, size=13)
        add_text(slide, 0.8, 6.6, 2, 0.3, f"{index:02d}", size=10)
    presentation.save(path)
    return path


def make_pdf_template(path: Path, logo_path: Path) -> Path:
    doc = fitz.open()
    for index, title in enumerate(("封面标题", "目录", "章节 01", "数据指标", "核心观点", "谢谢观看"), start=1):
        page = doc.new_page(width=1280, height=720)
        page.draw_rect(fitz.Rect(0, 0, 1280, 720), color=(0.96, 0.97, 0.98), fill=(0.96, 0.97, 0.98))
        page.draw_rect(fitz.Rect(72, 80, 82, 640), color=(0.20, 0.75, 0.60), fill=(0.20, 0.75, 0.60))
        page.insert_image(fitz.Rect(1095, 34, 1235, 76), filename=str(logo_path), keep_proportion=True)
        page.insert_text(fitz.Point(110, 150), title, fontsize=36, color=(0.08, 0.10, 0.14))
        page.insert_text(fitz.Point(110, 240), "示例占位内容，不作为当前项目正文读取。", fontsize=20, color=(0.30, 0.34, 0.40))
        page.insert_text(fitz.Point(110, 680), f"{index:02d}", fontsize=12, color=(0.45, 0.49, 0.55))
    doc.save(path)
    doc.close()
    return path


def create_project(title: str) -> dict:
    return request_json("POST", "/projects", json={"title": title})


def delete_project(project_id: str) -> None:
    try:
        request_json("DELETE", f"/projects/{project_id}")
    except Exception as exc:
        print(f"cleanup failed for {project_id}: {exc}", file=sys.stderr)


def create_slide(project_id: str, page_num: int, slide_type: str, headline: str):
    return request_json(
        "POST",
        f"/projects/{project_id}/slides",
        json={
            "page_num": page_num,
            "content_json": {
                "page_num": page_num,
                "type": slide_type,
                "text_content": {"headline": headline, "subhead": "", "body": "保留的原始正文"},
            },
        },
    )


def upload_template(project_id: str, file_path: Path) -> dict:
    with file_path.open("rb") as handle:
        result = request_json("POST", f"/projects/{project_id}/extract-template", files={"file": (file_path.name, handle, "application/octet-stream")})
    if result.get("status") == "completed":
        return result
    if result.get("status") != "processing":
        raise AssertionError(f"unexpected template upload response: {result}")
    return wait_template(project_id, result.get("job_id"))


def list_template_pages(project_id: str) -> list[dict]:
    pages = request_json("GET", f"/projects/{project_id}/template-pages")
    if not isinstance(pages, list):
        raise AssertionError(f"template-pages should return a list: {pages}")
    return pages


def list_reference_images(project_id: str) -> list[dict]:
    refs = request_json("GET", f"/projects/{project_id}/reference-images")
    if not isinstance(refs, list):
        raise AssertionError(f"reference-images should return a list: {refs}")
    return refs


def assert_static_url(url: str) -> None:
    response = session.get(API + url, timeout=TIMEOUT)
    if response.status_code != 200 or len(response.content) < 100:
        raise AssertionError(f"static url failed {url}: {response.status_code}, size={len(response.content)}")


def confirm_recommendations(project_id: str, pages: list[dict]) -> dict:
    by_category = {}
    for page in pages:
        by_category.setdefault(page.get("category") or "content", page)
    selection = {
        key: {
            "page_num": (by_category.get(key) or pages[0])["page_num"],
            "application_strength": "strong" if key == "cover" else "standard",
        }
        for key in TEMPLATE_KEYS
    }
    project = request_json("PATCH", f"/projects/{project_id}/template-recommendations", json={"recommendations": selection})
    recommendations = project.get("selected_template_recommendations") or {}
    missing = [key for key in TEMPLATE_KEYS if not recommendations.get(key) or not recommendations[key].get("file_path")]
    if missing:
        raise AssertionError(f"hydrated recommendations missing file_path for {missing}: {recommendations}")
    if recommendations["cover"].get("application_strength") != "strong":
        raise AssertionError(f"application strength was not preserved: {recommendations['cover']}")
    return recommendations


def scenario_true_template(root: Path, logo_path: Path, created: list[str]) -> tuple[str, dict]:
    project = create_project("E2E true PPTX template")
    created.append(project["id"])
    pptx = make_template_pptx(root / "true-template.pptx", logo_path, pages=7)
    status = upload_template(project["id"], pptx)
    pages = list_template_pages(project["id"])
    if status.get("document_kind") != "template":
        raise AssertionError(status)
    if len(pages) != 7:
        raise AssertionError(pages)
    categories = {page.get("category") for page in pages}
    if not set(TEMPLATE_KEYS).issubset(categories):
        raise AssertionError(f"missing categories: {categories}")
    if not any(page.get("logo_removed") for page in pages):
        raise AssertionError(f"expected logo-removed template refs: {pages}")
    for page in pages:
        if not page.get("file_exists"):
            raise AssertionError(page)
        assert_static_url(page["url"])
        assert_static_url(page["layout_url"])
    recommendations = confirm_recommendations(project["id"], pages)
    refs = list_reference_images(project["id"])
    logos = [ref for ref in refs if ref["role"] == "logo"]
    templates = [ref for ref in refs if ref["role"] == "template"]
    if len(templates) != 7 or not logos:
        raise AssertionError(refs)
    if not any(ref.get("review_status") == "auto_confirmed" for ref in logos):
        raise AssertionError(logos)
    return "true_template_pptx", {"pages": len(pages), "logos": len(logos), "cover_path": bool(recommendations["cover"]["file_path"])}


def scenario_manual_logo(root: Path, logo_path: Path, manual_logo: Path, created: list[str]) -> tuple[str, dict]:
    project = create_project("E2E manual logo preserved")
    created.append(project["id"])
    with manual_logo.open("rb") as handle:
        request_json(
            "POST",
            f"/projects/{project['id']}/upload",
            files={"file": ("manual-logo.png", handle, "image/png")},
            data={"role": "logo", "process_mode": "original", "logo_anchor": "top-left"},
        )
    upload_template(project["id"], make_template_pptx(root / "manual-logo-template.pptx", logo_path, pages=7))
    logos = [ref for ref in list_reference_images(project["id"]) if ref["role"] == "logo"]
    manual = [ref for ref in logos if not (ref.get("asset_analysis") or {}).get("template_logo_source")]
    extracted = [ref for ref in logos if (ref.get("asset_analysis") or {}).get("template_logo_source") == "layout_template"]
    if not manual or manual[0].get("review_status") != "auto_confirmed":
        raise AssertionError(logos)
    if not extracted or not all(ref.get("review_status") == "needs_review" for ref in extracted):
        raise AssertionError(logos)
    return "manual_logo_preserved", {"manual": len(manual), "template_candidates": len(extracted)}


def scenario_finished_ppt(root: Path, logo_path: Path, created: list[str]) -> tuple[str, dict]:
    project = create_project("E2E finished PPT as layout only")
    created.append(project["id"])
    create_slide(project["id"], 1, "cover", "Original cover must stay")
    create_slide(project["id"], 2, "content", "Original body must stay")
    before = request_json("GET", f"/projects/{project['id']}/slides")
    status = upload_template(project["id"], make_finished_pptx(root / "finished-deck.pptx", logo_path))
    after = request_json("GET", f"/projects/{project['id']}/slides")
    if status.get("document_kind") != "finished_ppt":
        raise AssertionError(status)
    if [slide["content_json"] for slide in before] != [slide["content_json"] for slide in after]:
        raise AssertionError(after)
    if len(list_template_pages(project["id"])) != 10:
        raise AssertionError("finished PPT should still extract all template page refs")
    return "finished_ppt_layout_only", {"document_kind": status["document_kind"], "slides_preserved": len(after)}


def scenario_pdf_template(root: Path, logo_path: Path, created: list[str]) -> tuple[str, dict]:
    project = create_project("E2E PDF repeated logo")
    created.append(project["id"])
    status = upload_template(project["id"], make_pdf_template(root / "pdf-template.pdf", logo_path))
    pages = list_template_pages(project["id"])
    logos = [ref for ref in list_reference_images(project["id"]) if ref["role"] == "logo"]
    if len(pages) != 6 or not any(page.get("logo_removed") for page in pages):
        raise AssertionError(pages)
    if not any((ref.get("asset_analysis") or {}).get("template_logo_detection") == "rendered_corner_repeat" for ref in logos):
        raise AssertionError(logos)
    return "pdf_repeated_logo", {"pages": len(pages), "logos": len(logos), "extracted_logos": status.get("extracted_logos")}


def scenario_corrupt_stability(root: Path, logo_path: Path, created: list[str]) -> tuple[str, dict]:
    project = create_project("E2E corrupt template keeps old refs")
    created.append(project["id"])
    upload_template(project["id"], make_template_pptx(root / "stable-template.pptx", logo_path, pages=7))
    before = list_template_pages(project["id"])
    broken = root / "broken-template.pptx"
    broken.write_bytes(b"not a real pptx file")
    with broken.open("rb") as handle:
        result = request_json("POST", f"/projects/{project['id']}/extract-template", files={"file": (broken.name, handle, "application/octet-stream")})
    failed = wait_template_failed(project["id"], result["job_id"])
    after = list_template_pages(project["id"])
    if len(before) != len(after) or not all(page.get("file_exists") for page in after):
        raise AssertionError(after)
    return "corrupt_template_stability", {"pages_kept": len(after), "error": (failed.get("error") or "")[:80]}


def scenario_consecutive_uploads(root: Path, logo_path: Path, created: list[str]) -> tuple[str, dict]:
    project = create_project("E2E consecutive template uploads")
    created.append(project["id"])
    old_pptx = make_template_pptx(root / "old-template-12-pages.pptx", logo_path, pages=12)
    new_pptx = make_template_pptx(root / "new-template-5-pages.pptx", logo_path, pages=5)
    with old_pptx.open("rb") as handle:
        first = request_json("POST", f"/projects/{project['id']}/extract-template", files={"file": (old_pptx.name, handle, "application/octet-stream")})
    with new_pptx.open("rb") as handle:
        second = request_json("POST", f"/projects/{project['id']}/extract-template", files={"file": (new_pptx.name, handle, "application/octet-stream")})
    status = wait_template(project["id"], second.get("job_id"), timeout=220)
    pages = list_template_pages(project["id"])
    if status.get("filename") != new_pptx.name or len(pages) != 5:
        raise AssertionError({"status": status, "pages": pages})
    return "consecutive_upload_latest_wins", {"first_job": first.get("job_id"), "second_job": second.get("job_id"), "pages": len(pages)}


def main() -> int:
    print(f"API={API}")
    print("health", request_json("GET", "/health"))
    login = request_json("POST", "/auth/tester-login", json={"display_name": f"E2E template tester {int(time.time())}", "passcode": ""})
    session.headers.update({"x-pptgod-tester-id": login["tester_id"]})
    print("tester", login["tester_id"])
    created: list[str] = []
    results: list[tuple[str, dict]] = []
    try:
        with tempfile.TemporaryDirectory(prefix="pptgod-template-e2e-") as tmp:
            root = Path(tmp)
            logo = make_logo_png(root / "acme-logo.png")
            manual_logo = make_manual_logo(root / "manual-logo.png")
            scenarios = (
                lambda: scenario_true_template(root, logo, created),
                lambda: scenario_manual_logo(root, logo, manual_logo, created),
                lambda: scenario_finished_ppt(root, logo, created),
                lambda: scenario_pdf_template(root, logo, created),
                lambda: scenario_corrupt_stability(root, logo, created),
                lambda: scenario_consecutive_uploads(root, logo, created),
            )
            for scenario in scenarios:
                name, detail = scenario()
                results.append((name, detail))
                print(f"PASS {name}: {detail}")
    finally:
        if KEEP_PROJECTS:
            print(f"kept_projects={created}")
        else:
            for project_id in created:
                delete_project(project_id)
            print(f"cleaned_projects={len(created)}")

    print("\nE2E_RESULTS")
    for name, detail in results:
        print(f"- {name}: ok {detail}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
